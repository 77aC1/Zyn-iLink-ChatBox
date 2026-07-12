#衍生/开发 请 标注 原仓库"https://github.com/zynsync/Zyn-iLink-ChatBox" 与原作者  。仓库受到开源证书保护!请合规使用!
#Derivative/Development Please attribute the original repository: "https://github.com/zynsync/Zyn-iLink-ChatBox" and the original author. This repository is protected by an open source license. Please comply with the license terms!
#我是屎山代码的塔尖
import threading
import time
import json
import uuid
import base64
import random
import sys
import subprocess
import os
import hashlib
import struct
import shutil

try:
    from Cryptodome.Cipher import AES as _CryptoAES
    _HAS_PYCRYPTODOME = True
except ImportError:
    _HAS_PYCRYPTODOME = False
from pathlib import Path
from typing import Optional, Dict, List, Union
from datetime import datetime, timedelta
import io
import socket
import socketserver
from http.server import SimpleHTTPRequestHandler
import select

_original_print = print
_SUPPRESSED_PREFIXES = ('[AI]', '[POLL]', '[USER]', '[QR]', '[添加用户]', '[ADD-USER]', '正在获取连接二维码', '获取二维码失败', '二维码已过期', '消息发送与二维码扫描请去本地网页操作', '[WEB] 开始添加用户', '[WEB] 添加用户异常')

def _filtered_print(*args, **kwargs):
    try:
        msg = ' '.join(str(a) for a in args)
        if any(msg.startswith(p) for p in _SUPPRESSED_PREFIXES):
            return
    except Exception:
        pass
    _original_print(*args, **kwargs)

print = _filtered_print

import urllib.request
import urllib.error
import urllib.parse

def is_termux():
    if sys.platform != "linux":
        return False
    
    checks = [
        "termux" in sys.prefix.lower(),
        "com.termux" in sys.prefix.lower(),
        "termux" in sys.executable.lower(),
        "com.termux" in sys.executable.lower(),
    ]
    
    if os.environ.get("TERMUX") or os.environ.get("PREFIX", "").startswith("/data/data/com.termux"):
        return True
    
    termux_paths = [
        "/data/data/com.termux",
        "/data/data/com.termux/files/usr/bin/python",
    ]
    for path in termux_paths:
        try:
            if os.path.exists(path):
                return True
        except Exception:
            pass
    
    return any(checks)

def setup_termux_compat():
    if not is_termux():
        return
    
    print("=" * 60)
    print("[Zyn] 检测到 Termux 环境")
    print("[Zyn] 正在启用兼容模式...")
    print("=" * 60)
    
    env_vars = {
        'TMPDIR': '/data/data/com.termux/files/usr/tmp',
        'TEMP': '/data/data/com.termux/files/usr/tmp',
        'TMP': '/data/data/com.termux/files/usr/tmp',
        'TERMUX': '1',
        'LD_LIBRARY_PATH': '/data/data/com.termux/files/usr/lib',
        'PATH': '/data/data/com.termux/files/usr/bin:' + os.environ.get('PATH', ''),
    }
    
    for key, value in env_vars.items():
        os.environ.setdefault(key, value)
        print(f"[TERMUX]   ✓ 设置 {key}")
    
    tmp_dir = Path("/data/data/com.termux/files/usr/tmp")
    try:
        tmp_dir.mkdir(parents=True, exist_ok=True)
        print(f"[TERMUX]   ✓ 确保临时目录存在: {tmp_dir}")
    except Exception as e:
        print(f"[TERMUX]   ⚠ 无法创建临时目录: {e}")
    
    tools = ["pkg", "python", "pip"]
    for tool in tools:
        if shutil.which(tool):
            print(f"[TERMUX]   ✓ {tool} 可用")
        else:
            print(f"[TERMUX]   ⚠ {tool} 未找到")
    
    print("[TERMUX] 兼容性设置完成")
    print("=" * 60)

setup_termux_compat()

def ensure_pip_available():
    if is_termux():
        print("[Zyn] 检测到 Termux 环境，尝试使用 pkg 安装 pip...")
        try:
            subprocess.check_call(["pkg", "install", "-y", "python-pip"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            print("[TERMUX] Termux pip 安装成功")
            return True
        except (subprocess.CalledProcessError, FileNotFoundError):
            print("[TERMUX] pkg 安装失败，尝试其他方式...")
    
    try:
        import pip
        return True
    except ImportError:
        pass
    
    try:
        result = subprocess.run([sys.executable, "-m", "pip", "--version"], capture_output=True, text=True)
        if result.returncode == 0:
            return True
    except Exception:
        pass
    
    try:
        import ensurepip
        print("正在通过 ensurepip 安装 pip...")
        
        try:
            ensurepip.bootstrap(upgrade=True)
        except Exception as bootstrap_err:
            print(f"  bootstrap 升级失败，尝试基础安装...")
            ensurepip.bootstrap()
        
        import site
        site.main()
        
        result = subprocess.run([sys.executable, "-m", "pip", "--version"], capture_output=True, text=True)
        if result.returncode == 0:
            print("pip 安装成功")
            return True
        else:
            print(f"  验证失败: {result.stderr}")
            
            import importlib
            try:
                if 'pip' in sys.modules:
                    del sys.modules['pip']
                importlib.import_module('pip')
                print("pip 安装成功（通过重新导入）")
                return True
            except ImportError:
                pass
                
    except Exception as e:
        print(f"ensurepip 安装失败: {e}")
    
    try:
        print("正在下载 get-pip.py...")
        import tempfile
        get_pip_url = "https://bootstrap.pypa.io/get-pip.py"
        
        if is_termux():
            temp_dir = Path("/data/data/com.termux/files/usr/tmp") / "temp_pip_install"
        else:
            temp_dir = Path(tempfile.gettempdir()) / "temp_pip_install"
        
        temp_dir.mkdir(exist_ok=True, parents=True)
        get_pip_path = temp_dir / "get-pip.py"
        
        urllib.request.urlretrieve(get_pip_url, str(get_pip_path))
        
        print("正在运行 get-pip.py 安装 pip...")
        result = subprocess.run(
            [sys.executable, str(get_pip_path), "--user", "--no-warn-script-location"],
            capture_output=True,
            text=True
        )
        
        if result.returncode != 0:
            print(f"用户目录安装失败，尝试使用临时目录...")
            target_dir = Path(temp_dir) / "pip_target"
            target_dir.mkdir(exist_ok=True)
            
            result = subprocess.run(
                [sys.executable, str(get_pip_path), "--target", str(target_dir)],
                capture_output=True,
                text=True
            )
            
            if result.returncode != 0:
                raise Exception(result.stderr)
                
            sys.path.insert(0, str(target_dir))
        
        get_pip_path.unlink(missing_ok=True)
        print("pip 安装成功")
        return True
    except Exception as e:
        print(f"get-pip.py 安装失败: {e}")
    
    return False

_PIP_MIRRORS = [
    ("https://pypi.tuna.tsinghua.edu.cn/simple", "pypi.tuna.tsinghua.edu.cn"),
    ("https://mirrors.aliyun.com/pypi/simple", "mirrors.aliyun.com"),
    ("https://mirrors.cloud.tencent.com/pypi/simple", "mirrors.cloud.tencent.com"),
    ("https://pypi.mirrors.ustc.edu.cn/simple", "pypi.mirrors.ustc.edu.cn"),
    ("https://mirrors.huaweicloud.com/repository/pypi/simple", "mirrors.huaweicloud.com"),
]

def _get_pip_index_args():
    import urllib.request
    for url, host in _PIP_MIRRORS:
        try:
            urllib.request.urlopen(url + "/", timeout=5)
            return ["-i", url, "--trusted-host", host]
        except Exception:
            continue
    return []

def install_package(package):
    index_args = _get_pip_index_args()
    if index_args:
        print(f"  使用镜像源: {index_args[1]}")

    install_commands = []
    if index_args:
        install_commands.append([sys.executable, "-m", "pip", "install", package] + index_args)
        install_commands.append([sys.executable, "-m", "pip", "install", "--user", package] + index_args)
    install_commands.append([sys.executable, "-m", "pip", "install", package])
    install_commands.append([sys.executable, "-m", "pip", "install", "--user", package])

    pip_exe = shutil.which("pip") or shutil.which("pip3")
    if pip_exe:
        if index_args:
            install_commands.insert(0, [pip_exe, "install", package] + index_args)
        install_commands.append([pip_exe, "install", package])
    
    for cmd in install_commands:
        try:
            print(f"  尝试: {' '.join(cmd[:6])}{'...' if len(cmd) > 6 else ''}")
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=180
            )
            
            if result.returncode == 0:
                return True
            else:
                err_msg = result.stderr.strip()
                if len(err_msg) > 200:
                    err_msg = err_msg[-200:]
                print(f"  失败: {err_msg}")
        except subprocess.TimeoutExpired:
            print(f"  超时(180s)")
        except FileNotFoundError:
            continue
        except Exception as e:
            print(f"  错误: {e}")
    
    if is_termux():
        try:
            print("  [TERMUX] 尝试 Termux 方式...")
            termux_cmd = ["pip", "install", package]
            if index_args:
                termux_cmd += index_args
            subprocess.check_call(termux_cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            return True
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass
    
    return False

def check_and_install_dependencies():
    required_packages = {
        "qrcode": "qrcode"
    }
    
    missing_packages = []
    for pip_name, import_name in required_packages.items():
        try:
            __import__(import_name)
        except ImportError:
            missing_packages.append(pip_name)
    
    if missing_packages:
        print(f"需要安装的库: {', '.join(missing_packages)}")
        
        has_pip = True
        try:
            import pip
        except ImportError:
            print("未检测到 pip，正在尝试自动安装...")
            has_pip = ensure_pip_available()
        
        if not has_pip:
            print("错误: 无法安装 pip，请手动安装后重试")
            print("  - Windows: python -m ensurepip --upgrade")
            print("  - Linux/Mac: python3 -m ensurepip --upgrade")
            print("  - Termux: pkg install python-pip")
            sys.exit(1)
        
        for package in missing_packages:
            print(f"正在安装 {package}...")
            if install_package(package):
                print(f"{package} 安装完成")
            else:
                print(f"{package} 安装失败，请手动安装: pip install {package}")
                sys.exit(1)

try:
    from Cryptodome.Cipher import AES
except ImportError:
    print()
    print("=" * 56)
    print("  未安装 pycryptodomex 库，媒体解密将极慢（纯Python实现）")
    print("  正在尝试自动安装...（一到两分钟）")
    print("=" * 56)
    print()
    
    if install_package("pycryptodomex"):
        print("[依赖] ✓ pycryptodomex 安装成功！")
        try:
            from Cryptodome.Cipher import AES
            print("[依赖] ✓ 导入验证通过，媒体解密加速已生效")
        except ImportError:
            print("[依赖] ⚠ 安装成功但导入失败，请手动检查")
    else:
        print()
        print("=" * 56)
        print("  ⚠ 自动安装失败，请手动运行：")
        print()
        print("    pip install pycryptodomex")
        print()
        print("  安装后重启程序即可生效")
        print("=" * 56)
        print()

    try:
        import pilk
    except ImportError:
        print()
        print("=" * 56)
        print("  未安装 pilk 库，SILK 语音将依赖 ffmpeg 转码")
        print("  建议安装 pilk 以获得更快的语音解码速度，请运行：")
        print()
        print("    pip install pilk -i https://pypi.tuna.tsinghua.edu.cn/simple")
        print("    (阿里云: -i https://mirrors.aliyun.com/pypi/simple)")
        print("    (腾讯云: -i https://mirrors.cloud.tencent.com/pypi/simple)")
        print()
        print("  安装后重启程序即可生效")
        print("=" * 56)
        print()

check_and_install_dependencies()

import qrcode

CONFIG_FILE = "wechat_bot_config.json"
MESSAGES_FILE = "wechat_messages.json"
AI_CONFIG_FILE = "ai_config.json"
USER_PROMPTS_FILE = "user_prompts.json"
WEB_PASSWORD_FILE = "web_password.json"
MEDIA_CACHE_DIR = "media_cache"
USER_DATA_DIR = "user_data"
ACCOUNTS_FILE = "accounts/accounts.json"

class BotAccount:
    def __init__(self, username: str, password_hash: str, salt: str, created_at: str, email: str = ""):
        self.username = username
        self.password_hash = password_hash
        self.salt = salt
        self.created_at = created_at
        self.email = email
        self.data_dir = Path(f"accounts/{username}")
        self.data_dir.mkdir(parents=True, exist_ok=True)
        self.config_file = self.data_dir / "wechat_bot_config.json"
        self.ai_config_file = self.data_dir / "ai_config.json"
        self.user_prompts_file = self.data_dir / "user_prompts.json"
        self.media_cache_dir = self.data_dir / "media_cache"
        self.media_cache_dir.mkdir(parents=True, exist_ok=True)
        self.user_data_dir = self.data_dir / "user_data"
        self.user_data_dir.mkdir(parents=True, exist_ok=True)
        self.token = None
        self.bot_id = None
        self.user_id = None
        self._cursor = ""
        self._context_tokens = {}
        self._current_user = None
        self._login_done = False
        self._qrcode_key = None
        self._qrcode_matrix = None
        self._messages = []
        self._last_msg_id = 0
        self._max_messages_per_user = 500
        self._total_max_messages = 2000
        self._media_memory = {}
        self._media_memory_max = 20
        self._active_timers = {}
        self._daily_timers = {}
        self._bot_accounts = {}
        self._user_token_map = {}
        self._poll_threads = []
        self._pending_qrcode = None
        self._media_downloading = {}
        self._media_download_lock = threading.Lock()
        self._add_user_lock = threading.Lock()
        self._msg_lock = threading.Lock()
        self._ai_reply_lock = threading.Lock()
        self.is_admin = False
        self.last_ip = ""
        self._last_ai_reply_time = {}
        self.ai_config = self._load_ai_config()
        self.user_prompts = self._load_user_prompts()
        self._load_messages()
        self._load_config()

    def _load_ai_config(self):
        default_config = {
            "auto_reply": False, "scheduled_reply": False, "api_url": "", "api_key": "",
            "model": "gpt-3.5-turbo", "active_interval": 60, "min_words": 10, "max_words": 200,
            "system_prompt": "你是一个微信聊天助手，请用自然的中文回复。",
            "vision_api_url": "", "vision_api_key": "", "vision_model": "gpt-4o", "vision_enabled": False,
            "image_gen_api_url": "", "image_gen_api_key": "", "image_gen_model": "dall-e-3", "image_gen_enabled": False,
            "file_recognize_enabled": False, "file_recognize_api_url": "", "file_recognize_api_key": "",
            "file_recognize_model": "gpt-4o", "file_recognize_max_size": 512, "file_recognize_compat_mode": False,
            "ai_cooldown": 5,
            "daily_reply": False, "daily_time": "09:00"
        }
        try:
            if self.ai_config_file.exists():
                with open(self.ai_config_file, "r", encoding="utf-8") as f:
                    saved = json.load(f)
                    if "enabled" in saved and "auto_reply" not in saved:
                        saved["auto_reply"] = saved.pop("enabled")
                    default_config.update(saved)
        except Exception:
            pass
        return default_config

    def _save_ai_config(self):
        try:
            with open(self.ai_config_file, "w", encoding="utf-8") as f:
                json.dump(self.ai_config, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def _load_user_prompts(self):
        try:
            if self.user_prompts_file.exists():
                with open(self.user_prompts_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    if isinstance(data, dict) and "prompts" in data:
                        return data
                    return {"prompts": data, "ai_enabled": {}}
        except Exception:
            pass
        return {"prompts": {}, "ai_enabled": {}}

    def _save_user_prompts(self):
        try:
            with open(self.user_prompts_file, "w", encoding="utf-8") as f:
                json.dump(self.user_prompts, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def _load_messages(self):
        self._messages = []
        self._last_msg_id = 0

    def _save_messages(self):
        try:
            with self._msg_lock:
                if len(self._messages) > self._total_max_messages:
                    self._messages = self._messages[-self._total_max_messages:]
            self._save_all_messages()
        except Exception:
            pass

    def _add_message_to_history(self, msg):
        with self._msg_lock:
            self._last_msg_id += 1
            msg['id'] = self._last_msg_id
            if 'time' not in msg:
                msg['time'] = datetime.now().strftime('%H:%M:%S')
            self._messages.append(msg)
            target_id = msg.get('to') or msg.get('from')
            if target_id:
                user_msgs = [m for m in self._messages if m.get('to') == target_id or m.get('from') == target_id]
                if len(user_msgs) > self._max_messages_per_user:
                    remove_ids = {m.get('id') for m in user_msgs[:len(user_msgs) - self._max_messages_per_user]}
                    self._messages = [m for m in self._messages if m.get('id') not in remove_ids]
        threading.Thread(target=self._save_messages, daemon=True).start()

    def get_user_messages(self, user_id, limit=50):
        with self._msg_lock:
            if not user_id:
                return list(self._messages[-limit:]) if self._messages else []
            user_msgs = [m for m in self._messages if m.get('from') == user_id or m.get('to') == user_id]
            return list(user_msgs[-limit:]) if limit > 0 else list(user_msgs)

    def _get_user_dir(self, user_id):
        safe_id = hashlib.md5(user_id.encode('utf-8')).hexdigest()[:16]
        user_dir = self.user_data_dir / safe_id
        user_dir.mkdir(parents=True, exist_ok=True)
        return user_dir

    def _get_user_dir_path(self, user_id):
        safe_id = hashlib.md5(user_id.encode('utf-8')).hexdigest()[:16]
        return self.user_data_dir / safe_id

    def _get_user_media_dir(self, user_id):
        media_dir = self._get_user_dir(user_id) / "media"
        media_dir.mkdir(parents=True, exist_ok=True)
        return media_dir

    def _get_user_token_file(self, user_id):
        return self._get_user_dir(user_id) / "token.json"

    def _get_user_messages_file(self, user_id):
        return self._get_user_dir(user_id) / "messages.json"

    def _save_user_token(self, user_id, context_token):
        try:
            data = {"user_id": user_id, "context_token": context_token, "saved_at": datetime.now().isoformat()}
            with open(self._get_user_token_file(user_id), "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def _load_user_token(self, user_id):
        token_file = self._get_user_token_file(user_id)
        try:
            if token_file.exists():
                with open(token_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    return data.get("context_token")
        except Exception:
            pass
        return None

    def _save_user_messages(self, user_id):
        if not user_id:
            return
        try:
            user_msgs = [m for m in self._messages if m.get('from') == user_id or m.get('to') == user_id]
            if len(user_msgs) > self._max_messages_per_user * 2:
                user_msgs = user_msgs[-self._max_messages_per_user:]
            data = {"user_id": user_id, "messages": user_msgs, "count": len(user_msgs), "saved_at": datetime.now().isoformat()}
            with open(self._get_user_messages_file(user_id), "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def _load_user_messages(self, user_id):
        msg_file = self._get_user_messages_file(user_id)
        try:
            if msg_file.exists():
                with open(msg_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    return data.get("messages", [])
        except Exception:
            pass
        return []

    def _load_all_user_messages(self):
        all_msgs = []
        loaded_ids = set()
        for user_id in self._context_tokens.keys():
            user_msgs = self._load_user_messages(user_id)
            for msg in user_msgs:
                msg_id = msg.get('id')
                if msg_id and msg_id not in loaded_ids:
                    all_msgs.append(msg)
                    loaded_ids.add(msg_id)
        all_msgs.sort(key=lambda m: m.get('id', 0))
        self._messages = all_msgs
        if self._messages:
            self._last_msg_id = max(msg.get('id', 0) for msg in self._messages)
        else:
            self._last_msg_id = 0

    def _save_all_messages(self):
        for user_id in self._context_tokens.keys():
            self._save_user_messages(user_id)

    def _save_config(self):
        config = {
            "token": self.token, "bot_id": self.bot_id, "user_id": self.user_id,
            "cursor": self._cursor, "context_tokens": self._context_tokens,
            "current_user": self._current_user,
            "bot_accounts": {k: v for k, v in self._bot_accounts.items()},
            "user_token_map": dict(self._user_token_map),
        }
        with open(self.config_file, "w", encoding="utf-8") as f:
            json.dump(config, f, indent=2, ensure_ascii=False)
        try:
            self.config_file.chmod(0o600)
        except (OSError, AttributeError, NotImplementedError):
            pass
        for user_id, ctx_token in self._context_tokens.items():
            self._save_user_token(user_id, ctx_token)

    def _load_config(self):
        try:
            if not self.config_file.exists():
                return False
            with open(self.config_file, "r", encoding="utf-8") as f:
                config = json.load(f)
            self.token = config.get("token")
            self.bot_id = config.get("bot_id")
            self.user_id = config.get("user_id")
            self._cursor = config.get("cursor", "")
            self._context_tokens = config.get("context_tokens", {})
            self._current_user = config.get("current_user")
            self._bot_accounts = config.get("bot_accounts", {})
            self._user_token_map = config.get("user_token_map", {})
            if self.token and self.token not in self._bot_accounts:
                self._bot_accounts[self.token] = {
                    "bot_id": self.bot_id or "", "user_id": self.user_id or "",
                    "cursor": self._cursor, "context_tokens": dict(self._context_tokens)
                }
            for user_id in list(self._context_tokens.keys()):
                if user_id not in self._user_token_map:
                    self._user_token_map[user_id] = self.token
            stale_users = [uid for uid in self._user_token_map if uid not in self._context_tokens]
            for uid in stale_users:
                self._user_token_map.pop(uid, None)
            no_dir_users = [uid for uid in list(self._context_tokens.keys()) if not self._get_user_dir_path(uid).exists()]
            for uid in no_dir_users:
                self._context_tokens.pop(uid, None)
                self._user_token_map.pop(uid, None)
            for bot_token, account in list(self._bot_accounts.items()):
                ctx_tokens = account.get("context_tokens", {})
                stale_ctx = [uid for uid in ctx_tokens if uid not in self._context_tokens]
                for uid in stale_ctx:
                    ctx_tokens.pop(uid, None)
            if self._current_user and self._current_user not in self._context_tokens:
                remaining = list(self._context_tokens.keys())
                self._current_user = remaining[0] if remaining else None
            self._load_all_user_messages()
            if self.token:
                print(f"[账号 {self.username}] 加载配置成功，{len(self._context_tokens)} 个会话，{len(self._bot_accounts)} 个 bot 账号，{len(self._messages)} 条消息")
                return True
            return False
        except FileNotFoundError:
            return False
        except Exception:
            return False

    def _get_token_for_user(self, user_id):
        return self._user_token_map.get(user_id) or self.token

    def _register_user_to_account(self, user_id, ctx_token, bot_token):
        self._context_tokens[user_id] = ctx_token
        self._user_token_map[user_id] = bot_token
        if bot_token not in self._bot_accounts:
            self._bot_accounts[bot_token] = {
                "bot_id": self.bot_id or "", "user_id": self.user_id or "",
                "cursor": "", "context_tokens": {}
            }
        self._bot_accounts[bot_token]["context_tokens"][user_id] = ctx_token
        self._save_user_token(user_id, ctx_token)
        if not self._current_user:
            self._current_user = user_id

    def _on_new_user(self, user_id):
        pass

    def list_users(self):
        return list(self._context_tokens.keys())

    def get_current_user(self):
        return self._current_user

    def set_current_user(self, user_id):
        if user_id in self._context_tokens:
            self._current_user = user_id
            self._save_config()

    def remove_user(self, user_id):
        if not user_id or user_id not in self._context_tokens:
            return False
        self._context_tokens.pop(user_id, None)
        bot_token = self._user_token_map.pop(user_id, None)
        if bot_token and bot_token in self._bot_accounts:
                self._bot_accounts[bot_token].get("context_tokens", {}).pop(user_id, None)
        if user_id in self._active_timers:
            timer = self._active_timers.pop(user_id)
            if timer:
                timer.cancel()
        if user_id in self._daily_timers:
            timer = self._daily_timers.pop(user_id)
            if timer:
                timer.cancel()
        self._messages = [m for m in self._messages if m.get('from') != user_id and m.get('to') != user_id]
        try:
            user_dir = self._get_user_dir_path(user_id)
            if user_dir.exists():
                shutil.rmtree(str(user_dir))
        except Exception:
            pass
        if self._current_user == user_id:
            remaining = list(self._context_tokens.keys())
            self._current_user = remaining[0] if remaining else None
        self._save_config()
        return True

    def get_effective_system_prompt(self, user_id=""):
        prompts = self.user_prompts.get("prompts", {})
        if user_id and user_id in prompts and prompts[user_id].strip():
            return prompts[user_id]
        system_prompt = self.ai_config.get("system_prompt", "")
        if not system_prompt:
            system_prompt = "你是一个微信聊天助手，请用自然的中文回复。"
        return system_prompt

    def is_ai_enabled_for_user(self, user_id):
        ai_enabled = self.user_prompts.get("ai_enabled", {})
        if user_id in ai_enabled:
            return ai_enabled[user_id]
        return self.ai_config.get("auto_reply", False)

    def is_scheduled_enabled_for_user(self, user_id):
        scheduled_enabled = self.user_prompts.get("scheduled_enabled", {})
        if user_id in scheduled_enabled:
            return scheduled_enabled[user_id]
        return self.ai_config.get("scheduled_reply", False)

    def is_daily_enabled_for_user(self, user_id):
        daily_enabled = self.user_prompts.get("daily_enabled", {})
        if user_id in daily_enabled:
            return daily_enabled[user_id]
        return self.ai_config.get("daily_reply", False)

    def _save_media_memory(self, user_id, media_type, description, filename="", cdn_info=None):
        if user_id not in self._media_memory:
            self._media_memory[user_id] = []
        entry = {"type": media_type, "description": description, "filename": filename,
                 "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'), "cdn_info": cdn_info}
        self._media_memory[user_id].append(entry)
        if len(self._media_memory[user_id]) > self._media_memory_max:
            self._media_memory[user_id] = self._media_memory[user_id][-self._media_memory_max:]

    def _get_media_memory(self, user_id):
        return self._media_memory.get(user_id, [])

    def _get_user_media_cache_path(self, user_id, cache_key):
        return self._get_user_media_dir(user_id) / cache_key

    def _get_user_media_meta_path(self, user_id, cache_key):
        return self._get_user_media_dir(user_id) / (cache_key + ".meta")

    def _save_user_media_cache(self, user_id, cache_key, media_data, mime, filename=""):
        try:
            self._get_user_media_cache_path(user_id, cache_key).write_bytes(media_data)
            meta = {'mime': mime, 'filename': filename, 'size': len(media_data)}
            self._get_user_media_meta_path(user_id, cache_key).write_text(json.dumps(meta, ensure_ascii=False), 'utf-8')
        except Exception:
            pass

    def _get_user_cached_media(self, user_id, cache_key):
        data_path = self._get_user_media_cache_path(user_id, cache_key)
        meta_path = self._get_user_media_meta_path(user_id, cache_key)
        if data_path.exists() and meta_path.exists():
            try:
                media_data = data_path.read_bytes()
                meta = json.loads(meta_path.read_text('utf-8'))
                return (media_data, meta.get('mime', 'application/octet-stream'), meta.get('filename', ''))
            except Exception:
                return None
        return None

    def _media_cache_path(self, cache_key):
        return self.media_cache_dir / cache_key

    def _media_meta_path(self, cache_key):
        return self.media_cache_dir / (cache_key + ".meta")

    def _get_cached_media(self, cache_key):
        data_path = self._media_cache_path(cache_key)
        meta_path = self._media_meta_path(cache_key)
        if data_path.exists() and meta_path.exists():
            try:
                media_data = data_path.read_bytes()
                meta = json.loads(meta_path.read_text('utf-8'))
                return (media_data, meta.get('mime', 'application/octet-stream'), meta.get('filename', ''))
            except Exception:
                return None
        return None

    def _save_media_cache(self, cache_key, media_data, mime, filename=""):
        try:
            self._media_cache_path(cache_key).write_bytes(media_data)
            meta = {'mime': mime, 'filename': filename, 'size': len(media_data)}
            self._media_meta_path(cache_key).write_text(json.dumps(meta, ensure_ascii=False), 'utf-8')
        except Exception:
            pass

    def _enrich_msg_with_cache_id(self, msg):
        if msg.get('media_cdn') and msg.get('media_type'):
            try:
                cdn_info = json.loads(msg['media_cdn']) if isinstance(msg['media_cdn'], str) else msg['media_cdn']
                cache_key = hashlib.md5((cdn_info.get("encrypt_query_param") or cdn_info.get("encrypted_query_param") or "").encode('utf-8')).hexdigest()
                user_id = msg.get('from') if msg.get('type') == 'in' else msg.get('to')
                cached = None
                if user_id:
                    cached = self._get_user_cached_media(user_id, cache_key)
                if not cached:
                    cached = self._get_cached_media(cache_key)
                if cached:
                    msg['media_cache_id'] = cache_key
                    msg['media_cache_user'] = user_id
            except Exception:
                pass
        return msg

class WeChatiLinkBot:
    ILINK_BASE_URL = "https://ilinkai.weixin.qq.com"
    MEDIA_TYPE_MAP = {"image": 2, "voice": 3, "file": 4, "video": 5}
    MEDIA_TYPE_NAMES = {2: "图片", 3: "语音", 4: "文件", 5: "视频"}
    MEDIA_TYPE_PREFIXES = {"image": "[图片]", "video": "[视频]", "file": "[文件]", "voice": "[语音]"}
    EXPIRED_CODES = {-14, 40014, 1002}
    SCRIPT_VERSION = "3.1.9-patched"
    AUTHOR_NAME = "ZynSync"
    
    def __init__(self):
        self.token: Optional[str] = None
        self.bot_id: Optional[str] = None
        self.user_id: Optional[str] = None
        self._cursor: str = ""
        self._context_tokens: Dict[str, str] = {}
        self._current_user: Optional[str] = None
        self._timeout = 35
        self._running = True 
        self._qrcode_matrix: Optional[List[List[str]]] = None
        self._http_server = None
        self._server_thread = None
        self._qrcode_key = None
        self._login_done = False
        self._web_port = 1145
        self._messages: List[dict] = []
        self._message_callback = None
        self._max_messages_per_user = 500
        self._total_max_messages = 2000
        self.ai_config = self._load_ai_config()
        self.user_prompts = self._load_user_prompts()
        self._media_memory: Dict[str, List[dict]] = {}
        self._media_memory_max = 20
        self._active_timers: Dict[str, threading.Timer] = {}
        self._daily_timers: Dict[str, threading.Timer] = {}
        self._session_tokens: Dict[str, float] = {}
        self._verified_sessions: Dict[str, float] = {}
        self._web_password_config: dict = self._load_web_password_config()
        self._verification_codes: Dict[str, dict] = {}
        self._captcha_store: Dict[str, dict] = {}
        self._login_attempts: Dict[str, dict] = {}
        self._send_code_attempts: Dict[str, dict] = {}
        self._account_reset_codes: Dict[str, dict] = {}
        self._email_code_cooldown: Dict[str, float] = {}
        self._email_verification_codes: Dict[str, dict] = {}
        self._email_bind_codes: Dict[str, dict] = {}
        self._media_cache_dir = Path(MEDIA_CACHE_DIR)
        self._media_cache_dir.mkdir(parents=True, exist_ok=True)
        self._user_data_dir = Path(USER_DATA_DIR)
        self._user_data_dir.mkdir(parents=True, exist_ok=True)
        self._media_downloading: Dict[str, threading.Event] = {}
        self._media_download_lock = threading.Lock()
        self._add_user_lock = threading.Lock()
        self._pending_qrcode: Optional[dict] = None
        self._msg_lock = threading.Lock()
        self._ai_reply_lock = threading.Lock()
        self._last_ai_reply_time: Dict[str, float] = {}
        
        self._bot_accounts: Dict[str, dict] = {}
        self._user_token_map: Dict[str, str] = {}
        self._poll_threads: List[threading.Thread] = []
        self._cf_process = None
        self._cf_url = ""
        
        self._accounts: Dict[str, BotAccount] = {}
        self._account_sessions: Dict[str, str] = {}
        self._fingerprint_sessions: Dict[str, str] = {}
        self._sse_connections: Dict[str, list] = {}
        self._sse_lock = threading.Lock()
        self._load_accounts()

        self._system_announcements = self._load_system_announcements()
        self._banned_ips = self._load_banned_ips()

        self._load_messages()
        for account in self._accounts.values():
            for user_id in account._context_tokens.keys():
                self._schedule_active_message_for_account(account, user_id)
                self._schedule_daily_message_for_account(account, user_id)
    
    def _save_media_memory(self, user_id: str, media_type: str, description: str, filename: str = "", cdn_info: dict = None):
        if user_id not in self._media_memory:
            self._media_memory[user_id] = []
        entry = {
            "type": media_type,
            "description": description,
            "filename": filename,
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "cdn_info": cdn_info
        }
        self._media_memory[user_id].append(entry)
        if len(self._media_memory[user_id]) > self._media_memory_max:
            self._media_memory[user_id] = self._media_memory[user_id][-self._media_memory_max:]
    
    def _get_media_memory(self, user_id: str) -> List[dict]:
        return self._media_memory.get(user_id, [])
    
    def _format_media_memory_for_prompt(self, user_id: str) -> str:
        memories = self._get_media_memory(user_id)
        if not memories:
            return ""
        lines = ["[历史媒体记录]"]
        for i, m in enumerate(memories):
            mt = "图片" if m.get('type') == 'image' else "文件"
            fn = m.get('filename', '')
            desc = m.get('description', '')
            ts = m.get('timestamp', '')
            lines.append(f"  {i+1}. [{mt}] {fn} - {desc} ({ts})")
        lines.append("[/历史媒体记录]")
        return "\n".join(lines)
    
    def _load_ai_config(self) -> dict:
        default_config = {
            "auto_reply": False,
            "scheduled_reply": False,
            "api_url": "",
            "api_key": "",
            "model": "gpt-3.5-turbo",
            "active_interval": 60,
            "min_words": 10,
            "max_words": 200,
            "system_prompt": "你是一个微信聊天助手，请用自然的中文回复。",
            "vision_api_url": "",
            "vision_api_key": "",
            "vision_model": "gpt-4o",
            "vision_enabled": False,
            "image_gen_api_url": "",
            "image_gen_api_key": "",
            "image_gen_model": "dall-e-3",
            "image_gen_enabled": False,
            "file_recognize_enabled": False,
            "file_recognize_api_url": "",
            "file_recognize_api_key": "",
            "file_recognize_model": "gpt-4o",
            "file_recognize_max_size": 512,
            "file_recognize_compat_mode": False,
            "ai_cooldown": 5,
            "daily_reply": False, "daily_time": "09:00"
        }
        try:
            if Path(AI_CONFIG_FILE).exists():
                with open(AI_CONFIG_FILE, "r", encoding="utf-8") as f:
                    saved = json.load(f)
                    if "enabled" in saved and "auto_reply" not in saved:
                        saved["auto_reply"] = saved.pop("enabled")
                    default_config.update(saved)
                    print(f"[AI] 已加载 AI 配置: auto_reply={default_config.get('auto_reply')}, scheduled_reply={default_config.get('scheduled_reply')}, api_url={default_config.get('api_url', '')[:50]}, api_key={'已设置' if default_config.get('api_key') else '未设置'}")
            else:
                print("[AI] 未找到 AI 配置文件，使用默认配置")
        except Exception as e:
            print(f"[AI] 加载 AI 配置失败: {e}")
        return default_config
    
    def _save_ai_config(self):
        try:
            with open(AI_CONFIG_FILE, "w", encoding="utf-8") as f:
                json.dump(self.ai_config, f, ensure_ascii=False, indent=2)
            print(f"[AI] 配置已保存: auto_reply={self.ai_config.get('auto_reply')}, scheduled_reply={self.ai_config.get('scheduled_reply')}, api_url={self.ai_config.get('api_url', '')[:50]}, api_key={'已设置' if self.ai_config.get('api_key') else '未设置'}")
        except Exception as e:
            print(f"[AI] 保存 AI 配置失败: {e}")
    
    def _load_user_prompts(self) -> dict:
        try:
            if Path(USER_PROMPTS_FILE).exists():
                with open(USER_PROMPTS_FILE, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    if isinstance(data, dict) and "prompts" in data:
                        return data
                    return {"prompts": data, "ai_enabled": {}}
        except Exception as e:
            print(f"[AI] 加载用户提示词失败: {e}")
        return {"prompts": {}, "ai_enabled": {}}
    
    def _save_user_prompts(self):
        try:
            with open(USER_PROMPTS_FILE, "w", encoding="utf-8") as f:
                json.dump(self.user_prompts, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"[AI] 保存用户提示词失败: {e}")

    def _load_web_password_config(self) -> dict:
        try:
            if Path(WEB_PASSWORD_FILE).exists():
                with open(WEB_PASSWORD_FILE, "r", encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            pass
        return {}

    def _save_web_password_config(self):
        try:
            with open(WEB_PASSWORD_FILE, "w", encoding="utf-8") as f:
                json.dump(self._web_password_config, f, ensure_ascii=False, indent=2)
            try:
                Path(WEB_PASSWORD_FILE).chmod(0o600)
            except (OSError, AttributeError, NotImplementedError):
                pass
        except Exception:
            pass

    def _load_system_announcements(self):
        try:
            path = Path("system_announcements.json")
            if path.exists():
                with open(path, "r", encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            pass
        return []

    def _save_system_announcements(self):
        try:
            with open("system_announcements.json", "w", encoding="utf-8") as f:
                json.dump(self._system_announcements, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def _load_banned_ips(self):
        try:
            path = Path("banned_ips.json")
            if path.exists():
                with open(path, "r", encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            pass
        return []

    def _save_banned_ips(self):
        try:
            with open("banned_ips.json", "w", encoding="utf-8") as f:
                json.dump(self._banned_ips, f, ensure_ascii=False, indent=2)
        except Exception:
            pass

    def _is_ip_banned(self, ip):
        return ip in self._banned_ips

    def _load_accounts(self):
        try:
            Path("accounts").mkdir(parents=True, exist_ok=True)
            if Path(ACCOUNTS_FILE).exists():
                with open(ACCOUNTS_FILE, "r", encoding="utf-8") as f:
                    data = json.load(f)
                for username, info in data.get("accounts", {}).items():
                    account = BotAccount(
                        username=username,
                        password_hash=info.get("password_hash", ""),
                        salt=info.get("salt", ""),
                        created_at=info.get("created_at", ""),
                        email=info.get("email", "")
                    )
                    account.is_admin = info.get("is_admin", False)
                    self._accounts[username] = account
                print(f"[多账户] 已加载 {len(self._accounts)} 个账户")
            else:
                print("[多账户] 未找到账户文件，首次运行")
        except Exception as e:
            print(f"[多账户] 加载账户失败: {e}")

    def _save_accounts(self):
        try:
            Path("accounts").mkdir(parents=True, exist_ok=True)
            data = {"accounts": {}}
            for username, account in self._accounts.items():
                data["accounts"][username] = {
                    "password_hash": account.password_hash,
                    "salt": account.salt,
                    "created_at": account.created_at,
                    "is_admin": account.is_admin,
                    "email": account.email
                }
            with open(ACCOUNTS_FILE, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            try:
                Path(ACCOUNTS_FILE).chmod(0o600)
            except (OSError, AttributeError, NotImplementedError):
                pass
        except Exception as e:
            print(f"[多账户] 保存账户失败: {e}")

    def _register_account(self, username: str, password: str, client_ip: str = "", email: str = "") -> dict:
        if not username or len(username) < 2 or len(username) > 32:
            return {"success": False, "error": "用户名长度需在2-32之间"}
        if not password or len(password) < 6:
            return {"success": False, "error": "密码长度不能少于6位"}
        if not all(c.isalnum() or c in '_-.' for c in username):
            return {"success": False, "error": "用户名只能包含字母、数字、_-."}
        if username in self._accounts:
            return {"success": False, "error": "用户名已存在"}
        salt = uuid.uuid4().hex
        dk = hashlib.pbkdf2_hmac('sha256', (salt + password).encode("utf-8"), salt.encode("utf-8"), 100000)
        password_hash = dk.hex()
        created_at = datetime.now().isoformat()
        account = BotAccount(username=username, password_hash=password_hash, salt=salt, created_at=created_at, email=email)
        account.last_ip = client_ip
        self._accounts[username] = account
        self._save_accounts()
        session_token = self._generate_session_token()
        self._account_sessions[session_token] = username
        self._verified_sessions[session_token] = time.time() + 86400 * 30
        print(f"[多账户] 新账户注册: {username} IP: {client_ip}")
        return {"success": True, "session_token": session_token, "username": username}

    def _login_account(self, username: str, password: str, fingerprint: str = "", client_ip: str = "") -> dict:
        if not username or not password:
            return {"success": False, "error": "用户名和密码不能为空"}
        account = self._accounts.get(username)
        if not account:
            return {"success": False, "error": "用户名或密码错误"}
        dk = hashlib.pbkdf2_hmac('sha256', (account.salt + password).encode("utf-8"), account.salt.encode("utf-8"), 100000)
        if dk.hex() != account.password_hash:
            return {"success": False, "error": "用户名或密码错误"}
        old_sessions = [tok for tok, uname in self._account_sessions.items() if uname == username]
        for tok in old_sessions:
            self._account_sessions.pop(tok, None)
            self._verified_sessions.pop(tok, None)
            self._session_tokens.pop(tok, None)
        session_token = self._generate_session_token()
        self._account_sessions[session_token] = username
        self._verified_sessions[session_token] = time.time() + 86400 * 30
        if fingerprint:
            self._fingerprint_sessions[fingerprint] = username
        account.last_ip = client_ip
        print(f"[多账户] 账户登录: {username} IP: {client_ip}")
        return {"success": True, "session_token": session_token, "username": username}

    def _fingerprint_login(self, fingerprint: str, client_ip: str = "") -> dict:
        if not fingerprint:
            return {"success": False, "error": "指纹为空"}
        username = self._fingerprint_sessions.get(fingerprint)
        if not username:
            return {"success": False, "error": "未找到关联账户"}
        if username not in self._accounts:
            self._fingerprint_sessions.pop(fingerprint, None)
            return {"success": False, "error": "账户不存在"}
        old_sessions = [tok for tok, uname in self._account_sessions.items() if uname == username]
        for tok in old_sessions:
            self._account_sessions.pop(tok, None)
            self._verified_sessions.pop(tok, None)
            self._session_tokens.pop(tok, None)
        session_token = self._generate_session_token()
        self._account_sessions[session_token] = username
        self._verified_sessions[session_token] = time.time() + 86400 * 30
        self._accounts[username].last_ip = client_ip
        print(f"[多账户] 指纹登录: {username} IP: {client_ip}")
        return {"success": True, "session_token": session_token, "username": username}

    def _change_account_password(self, username: str, old_password: str, new_password: str) -> dict:
        account = self._accounts.get(username)
        if not account:
            return {"success": False, "error": "账户不存在"}
        dk = hashlib.pbkdf2_hmac('sha256', (account.salt + old_password).encode("utf-8"), account.salt.encode("utf-8"), 100000)
        if dk.hex() != account.password_hash:
            return {"success": False, "error": "原密码错误"}
        if not new_password or len(new_password) < 6:
            return {"success": False, "error": "新密码长度不能少于6位"}
        salt = uuid.uuid4().hex
        dk = hashlib.pbkdf2_hmac('sha256', (salt + new_password).encode("utf-8"), salt.encode("utf-8"), 100000)
        account.password_hash = dk.hex()
        account.salt = salt
        self._save_accounts()
        old_sessions = [tok for tok, uname in self._account_sessions.items() if uname == username]
        for tok in old_sessions:
            self._account_sessions.pop(tok, None)
            self._verified_sessions.pop(tok, None)
        for fp, uname in list(self._fingerprint_sessions.items()):
            if uname == username:
                del self._fingerprint_sessions[fp]
        print(f"[多账户] 账户修改密码: {username}")
        return {"success": True}

    def _delete_account(self, username: str, password: str) -> dict:
        account = self._accounts.get(username)
        if not account:
            return {"success": False, "error": "账户不存在"}
        dk = hashlib.pbkdf2_hmac('sha256', (account.salt + password).encode("utf-8"), account.salt.encode("utf-8"), 100000)
        if dk.hex() != account.password_hash:
            return {"success": False, "error": "密码错误"}
        old_sessions = [tok for tok, uname in self._account_sessions.items() if uname == username]
        for tok in old_sessions:
            self._account_sessions.pop(tok, None)
            self._verified_sessions.pop(tok, None)
        for fp, uname in list(self._fingerprint_sessions.items()):
            if uname == username:
                del self._fingerprint_sessions[fp]
        del self._accounts[username]
        self._save_accounts()
        print(f"[多账户] 账户已注销: {username}")
        return {"success": True}

    def _reset_account_password(self, username: str, code: str, new_password: str, client_ip: str = '') -> dict:
        account = self._accounts.get(username)
        if not account:
            return {"success": False, "error": "账户不存在"}
        code_info = self._account_reset_codes.get(code)
        if not code_info:
            return {"success": False, "error": "验证码无效"}
        if code_info.get('used'):
            return {"success": False, "error": "验证码已使用"}
        if time.time() > code_info.get('expiry', 0):
            self._account_reset_codes.pop(code, None)
            return {"success": False, "error": "验证码已过期"}
        if code_info.get('username') != username:
            return {"success": False, "error": "验证码无效"}
        if code_info.get('ip') and code_info['ip'] != client_ip:
            return {"success": False, "error": "验证码无效"}
        if not new_password or len(new_password) < 6:
            return {"success": False, "error": "新密码长度不能少于6位"}
        if len(new_password) > 128:
            return {"success": False, "error": "新密码长度不能超过128位"}
        code_info['used'] = True
        self._account_reset_codes.pop(code, None)
        salt = uuid.uuid4().hex
        dk = hashlib.pbkdf2_hmac('sha256', (salt + new_password).encode("utf-8"), salt.encode("utf-8"), 100000)
        account.password_hash = dk.hex()
        account.salt = salt
        self._save_accounts()
        old_sessions = [tok for tok, uname in self._account_sessions.items() if uname == username]
        for tok in old_sessions:
            self._account_sessions.pop(tok, None)
            self._verified_sessions.pop(tok, None)
        for fp, uname in list(self._fingerprint_sessions.items()):
            if uname == username:
                del self._fingerprint_sessions[fp]
        print(f"[多账户] 账户重置密码: {username}")
        return {"success": True}

    def _get_account_from_session(self, session_token: str) -> Optional[BotAccount]:
        if not session_token:
            return None
        username = self._account_sessions.get(session_token)
        if not username:
            return None
        return self._accounts.get(username)

    def _start_account_polling(self, account: BotAccount):
        if account.token and account.token not in account._bot_accounts:
            account._bot_accounts[account.token] = {
                "bot_id": account.bot_id or "",
                "user_id": account.user_id or "",
                "cursor": account._cursor,
                "context_tokens": dict(account._context_tokens)
            }
        seen_tokens = set()
        for user_id, bot_token in account._user_token_map.items():
            if bot_token and bot_token not in seen_tokens:
                seen_tokens.add(bot_token)
                bot_account = account._bot_accounts.get(bot_token)
                if bot_account:
                    self._start_account_poll_thread(account, bot_token, bot_account)

    def _start_account_poll_thread(self, account: BotAccount, bot_token: str, bot_account: dict):
        for t in account._poll_threads:
            if t.is_alive() and getattr(t, '_bot_token', None) == bot_token:
                return
        def poll():
            cursor = bot_account.get("cursor", "")
            while self._running:
                try:
                    body = {"get_updates_buf": cursor}
                    result = self._post("getupdates", body, timeout=25, token=bot_token)
                    if result.get("get_updates_buf"):
                        cursor = result["get_updates_buf"]
                        bot_account["cursor"] = cursor
                        account._save_config()
                    messages = result.get("msgs", [])
                    for msg in messages:
                        from_user = msg.get("from_user_id")
                        ctx_token = msg.get("context_token")
                        text, media_info = self._process_message_items(msg.get("item_list", []))
                        msg_text = text
                        msg_type = 'in'
                        msg_metadata = {}
                        if media_info:
                            media_type_int = self.MEDIA_TYPE_MAP.get(media_info["type"], 0)
                            media_prefix = self.MEDIA_TYPE_PREFIXES.get(media_info["type"], f"[{media_info['type']}]")
                            if text:
                                msg_text = f"{media_prefix} {text}"
                            else:
                                msg_text = f"{media_prefix} {media_info.get('filename', '')}"
                            msg_metadata = {
                                'media_type': media_type_int,
                                'media_filename': media_info.get('filename', ''),
                                'media_duration': media_info.get('duration', 0),
                                'has_media': True
                            }
                            media_item = media_info.get("item", {})
                            cdn_media = self._extract_cdn_media(media_item)
                            if cdn_media:
                                msg_metadata['media_cdn'] = json.dumps(cdn_media)
                                _prefetch_fn = media_info.get('filename', '')
                                threading.Thread(target=self._prefetch_media_for_account, args=(account, cdn_media, _prefetch_fn, from_user), daemon=True).start()
                            print(f"\n[收到{media_info['type']}] [{account.username}] {from_user}: {media_info.get('filename', '')}")
                        elif text:
                            print(f"\n[收到消息] [{account.username}] {from_user}: {text}")
                        if msg_text:
                            new_msg = {
                                'from': from_user, 'to': 'me', 'text': msg_text,
                                'time': datetime.now().strftime('%H:%M:%S'), 'type': msg_type, **msg_metadata
                            }
                            account._add_message_to_history(new_msg)
                            if media_info and media_info.get("type") == "image" and cdn_media:
                                def _vision_reply(fn_from_user=from_user, fn_cdn=cdn_media, fn_text=text, fn_filename=media_info.get('filename', '')):
                                    try:
                                        downloaded = self.download_media_for_account(account, fn_cdn, filename=fn_filename, user_id=fn_from_user)
                                        if downloaded:
                                            img_b64 = base64.b64encode(downloaded).decode('utf-8')
                                            self._auto_ai_reply_with_vision_for_account(account, fn_from_user, img_b64, original_text=fn_text, cdn_info=fn_cdn)
                                        else:
                                            if fn_text and account.ai_config.get("auto_reply"):
                                                self._auto_ai_reply_for_account(account, fn_from_user, fn_text)
                                    except Exception:
                                        if fn_text and account.ai_config.get("auto_reply"):
                                            self._auto_ai_reply_for_account(account, fn_from_user, fn_text)
                                threading.Thread(target=_vision_reply, daemon=True).start()
                            elif media_info and media_info.get("type") == "file" and cdn_media:
                                def _file_reply(fn_from_user=from_user, fn_cdn=cdn_media, fn_text=text, fn_filename=media_info.get('filename', '')):
                                    try:
                                        downloaded = self.download_media_for_account(account, fn_cdn, filename=fn_filename, user_id=fn_from_user)
                                        if downloaded:
                                            self._auto_ai_reply_with_file_for_account(account, fn_from_user, downloaded, fn_filename, original_text=fn_text)
                                        else:
                                            if fn_text and account.ai_config.get("auto_reply"):
                                                self._auto_ai_reply_for_account(account, fn_from_user, fn_text)
                                    except Exception:
                                        if fn_text and account.ai_config.get("auto_reply"):
                                            self._auto_ai_reply_for_account(account, fn_from_user, fn_text)
                                threading.Thread(target=_file_reply, daemon=True).start()
                            elif text:
                                threading.Thread(target=self._auto_ai_reply_for_account, args=(account, from_user, text), daemon=True).start()
                        if from_user and ctx_token:
                            is_new = from_user not in account._context_tokens
                            account._register_user_to_account(from_user, ctx_token, bot_token)
                            account._save_config()
                            if is_new:
                                print(f"[USER] 新用户 {from_user} (账号 {bot_token[:8]}..., 账户 {account.username})")
                            self._schedule_active_message_for_account(account, from_user)
                            self._schedule_daily_message_for_account(account, from_user)

                except Exception:
                    time.sleep(0.5)
        thread = threading.Thread(target=poll, daemon=True)
        thread._bot_token = bot_token
        thread.start()
        account._poll_threads.append(thread)
        print(f"[POLL] 已启动轮询线程: {account.username}/{bot_token[:8]}...")

    def _schedule_active_message_for_account(self, account: BotAccount, user_id: str):
        if not account.is_scheduled_enabled_for_user(user_id):
            old_timer = account._active_timers.pop(user_id, None)
            if old_timer:
                old_timer.cancel()
            return
        interval = account.ai_config.get("active_interval", 60)
        if interval <= 0:
            return
        
        if user_id in account._active_timers:
            old_timer = account._active_timers[user_id]
            if old_timer:
                old_timer.cancel()
        
        print(f"[AI] 为 {user_id} 安排主动发送，间隔 {interval} 秒")
        timer = threading.Timer(interval, self._send_active_message_for_account, args=[account, user_id])
        timer.daemon = True
        timer.start()
        account._active_timers[user_id] = timer

    def _send_active_message_for_account(self, account: BotAccount, user_id: str):
        if not account.is_scheduled_enabled_for_user(user_id):
            return
        if not self._running:
            return
        if user_id not in account._context_tokens:
            print(f"[AI] 用户 {user_id} 已不存在，取消主动发送")
            if user_id in account._active_timers:
                del account._active_timers[user_id]
            return

        print(f"[AI] 主动发送定时器触发，准备向 {user_id} 发送消息...")
        self.send_typing_for_account(account, user_id)
        history = account.get_user_messages(user_id, 200)
        media_memory_text = self._format_media_memory_for_account(account, user_id)

        response = self._call_ai_api_for_account(account, "", history, is_active=True, media_memory_text=media_memory_text, user_id=user_id)

        if response:
            self.send_text_for_account(account, user_id, response)
        else:
            print(f"[AI] 主动发送未能获取有效回复")

        if account.is_scheduled_enabled_for_user(user_id) and self._running and user_id in account._context_tokens:
            self._schedule_active_message_for_account(account, user_id)

    def _schedule_daily_message_for_account(self, account: BotAccount, user_id: str):
        if not account.is_daily_enabled_for_user(user_id):
            old_timer = account._daily_timers.pop(user_id, None)
            if old_timer:
                old_timer.cancel()
            return
        daily_time = account.ai_config.get("daily_time", "09:00")
        now = datetime.now()
        try:
            target_hour, target_min = map(int, daily_time.split(":"))
        except Exception:
            target_hour, target_min = 9, 0
        target = now.replace(hour=target_hour, minute=target_min, second=0, microsecond=0)
        if target <= now:
            target += timedelta(days=1)
        delay = (target - now).total_seconds()
        if user_id in account._daily_timers:
            old_timer = account._daily_timers[user_id]
            if old_timer:
                old_timer.cancel()
        _original_print(f"[DAILY] 为 {user_id} 安排每日发送，时间 {daily_time}，距触发还有 {delay:.0f} 秒")
        timer = threading.Timer(delay, self._send_daily_message_for_account, args=[account, user_id])
        timer.daemon = True
        timer.start()
        account._daily_timers[user_id] = timer

    def _send_daily_message_for_account(self, account: BotAccount, user_id: str):
        if not account.is_daily_enabled_for_user(user_id):
            return
        if not self._running:
            return
        if user_id not in account._context_tokens:
            if user_id in account._daily_timers:
                del account._daily_timers[user_id]
            return
        _original_print(f"[DAILY] 每日定时触发，准备向 {user_id} 发送消息...")
        self.send_typing_for_account(account, user_id)
        history = account.get_user_messages(user_id, 200)
        media_memory_text = self._format_media_memory_for_account(account, user_id)
        response = self._call_ai_api_for_account(account, "", history, is_active=True, media_memory_text=media_memory_text, user_id=user_id)
        if response:
            self.send_text_for_account(account, user_id, response)
        else:
            _original_print(f"[DAILY] 每日发送未能获取有效回复")
        if account.is_daily_enabled_for_user(user_id) and self._running and user_id in account._context_tokens:
            self._schedule_daily_message_for_account(account, user_id)

    def _auto_ai_reply_for_account(self, account: BotAccount, from_user: str, text: str):
        with account._ai_reply_lock:
            last_time = account._last_ai_reply_time.get(from_user, 0)
            cooldown = account.ai_config.get("ai_cooldown", 5)
            if time.time() - last_time < cooldown:
                return
            account._last_ai_reply_time[from_user] = time.time()
            if not account.is_ai_enabled_for_user(from_user):
                return
            try:
                self.send_typing_for_account(account, from_user)
                history = account.get_user_messages(from_user, limit=20)
                media_memory_text = self._format_media_memory_for_account(account, from_user)
                reply = self._call_ai_api_for_account(account, text, history, media_memory_text=media_memory_text, user_id=from_user)
                if reply:
                    reply_text, s = self._parse_ai_s(reply)
                    if s:
                        for tc in s:
                            self._handle_ai__for_account(account, from_user, tc)
                    if reply_text:
                        self.send_text_for_account(account, from_user, reply_text)
                    self._schedule_active_message_for_account(account, from_user)
            except Exception as e:
                _original_print(f"[AI-ERR] 自动回复异常 [{account.username}]: {e}")

    def _auto_ai_reply_with_vision_for_account(self, account: BotAccount, from_user: str, img_b64: str, original_text: str = "", cdn_info: dict = None):
        if not account.ai_config.get("vision_enabled") or not account.ai_config.get("vision_api_url"):
            if original_text and account.ai_config.get("auto_reply"):
                self._auto_ai_reply_for_account(account, from_user, original_text)
            return
        try:
            account._save_media_memory(from_user, "image", original_text or "图片", "", cdn_info)
            vision_url = account.ai_config.get("vision_api_url")
            vision_key = account.ai_config.get("vision_api_key")
            vision_model = account.ai_config.get("vision_model", "gpt-4o")
            system_prompt = account.get_effective_system_prompt(from_user)
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": [
                    {"type": "text", "text": original_text or "请描述这张图片"},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}}
                ]}
            ]
            payload = {"model": vision_model, "messages": messages, "max_tokens": 500}
            headers = {"Content-Type": "application/json"}
            if vision_key:
                headers["Authorization"] = f"Bearer {vision_key}"
            req = urllib.request.Request(vision_url, data=json.dumps(payload).encode("utf-8"), headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                reply = result.get("choices", [{}])[0].get("message", {}).get("content", "")
                if reply:
                    self.send_text_for_account(account, from_user, reply)
        except Exception as e:
            print(f"[VISION] 识图异常 [{account.username}]: {e}")
            if original_text and account.ai_config.get("auto_reply"):
                self._auto_ai_reply_for_account(account, from_user, original_text)

    def _auto_ai_reply_with_file_for_account(self, account: BotAccount, from_user: str, file_data: bytes, filename: str, original_text: str = ""):
        if not account.ai_config.get("file_recognize_enabled") or not account.ai_config.get("file_recognize_api_url"):
            if original_text and account.ai_config.get("auto_reply"):
                self._auto_ai_reply_for_account(account, from_user, original_text)
            return
        try:
            account._save_media_memory(from_user, "file", original_text or filename, filename)
            max_size = account.ai_config.get("file_recognize_max_size", 512) * 1024
            if len(file_data) > max_size:
                if original_text and account.ai_config.get("auto_reply"):
                    self._auto_ai_reply_for_account(account, from_user, original_text)
                return
            file_b64 = base64.b64encode(file_data).decode('utf-8')
            api_url = account.ai_config.get("file_recognize_api_url")
            api_key = account.ai_config.get("file_recognize_api_key")
            model = account.ai_config.get("file_recognize_model", "gpt-4o")
            system_prompt = account.get_effective_system_prompt(from_user)
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": [
                    {"type": "text", "text": f"文件名: {filename}\n{original_text or '请分析这个文件'}"},
                    {"type": "image_url", "image_url": {"url": f"data:application/octet-stream;base64,{file_b64}"}}
                ]}
            ]
            payload = {"model": model, "messages": messages, "max_tokens": 500}
            headers = {"Content-Type": "application/json"}
            if api_key:
                headers["Authorization"] = f"Bearer {api_key}"
            req = urllib.request.Request(api_url, data=json.dumps(payload).encode("utf-8"), headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                reply = result.get("choices", [{}])[0].get("message", {}).get("content", "")
                if reply:
                    self.send_text_for_account(account, from_user, reply)
        except Exception as e:
            print(f"[FILE_RECOGNIZE] 文件识别异常 [{account.username}]: {e}")
            if original_text and account.ai_config.get("auto_reply"):
                self._auto_ai_reply_for_account(account, from_user, original_text)

    def _call_ai_api_for_account(self, account: BotAccount, user_message: str, history: list, is_active: bool = False, custom_instruction: str = "", media_memory_text: str = "", user_id: str = "") -> Optional[str]:
        if not account.ai_config.get("api_url"):
            return None
        api_url = account.ai_config.get("api_url")
        if not api_url or not api_url.startswith(('http://', 'https://')):
            return None
        system_prompt = account.get_effective_system_prompt(user_id)
        tool_prompt = self._build_tool_prompt_for_account(account)
        full_system = system_prompt
        if tool_prompt:
            full_system += "\n" + tool_prompt
        if media_memory_text:
            full_system += "\n" + media_memory_text
        messages = []
        for msg in history[-50:]:
            if msg.get("type") == "in":
                messages.append({"role": "user", "content": msg.get("text", "")})
            elif msg.get("type") == "out":
                messages.append({"role": "assistant", "content": msg.get("text", "")})
        if is_active:
            final_prompt = "现在没有用户的新消息，你需要主动发起一个话题。请严格按照你的性格要求来回复。"
        else:
            if custom_instruction:
                final_prompt = f"用户说：{user_message}\n\n额外要求：{custom_instruction}\n\n请严格按照你的性格要求和额外要求回复。"
            else:
                final_prompt = f"用户说：{user_message}\n\n请严格按照你的性格要求回复。"
        if messages and messages[0]["role"] == "user":
            messages[0]["content"] = f"[系统指令]\n{full_system}\n\n[用户消息]\n{messages[0]['content']}"
            messages.append({"role": "user", "content": final_prompt})
        else:
            messages.append({"role": "user", "content": f"[系统指令]\n{full_system}\n\n[用户消息]\n{final_prompt}"})
        payload = {"model": account.ai_config.get("model", "gpt-3.5-turbo"), "messages": messages, "temperature": 1.2, "max_tokens": 500}
        headers = {"Content-Type": "application/json"}
        if account.ai_config.get("api_key"):
            headers["Authorization"] = f"Bearer {account.ai_config.get('api_key')}"
        req = urllib.request.Request(api_url, data=json.dumps(payload).encode("utf-8"), headers=headers, method="POST")
        try:
            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                return result.get("choices", [{}])[0].get("message", {}).get("content", "")
        except Exception as e:
            _original_print(f"[AI-ERR] API 调用失败: {e}")
            return None

    def _build_tool_prompt_for_account(self, account: BotAccount) -> str:
        tools = []
        if account.ai_config.get("vision_enabled") and account.ai_config.get("vision_api_url"):
            tools.append("识图AI")
        if account.ai_config.get("file_recognize_enabled") and account.ai_config.get("file_recognize_api_url"):
            tools.append("文件AI")
        if account.ai_config.get("image_gen_enabled") and account.ai_config.get("image_gen_api_url"):
            tools.append("生图AI")
        if not tools:
            return ""
        lines = [""]
        lines.append("【强制规则 - 你必须遵守】")
        lines.append("你有能力查看图片和生成图片。禁止说\"我无法查看图片\"或\"我无法生成图片\"。")
        lines.append("当用户请求涉及图片或文件时，你必须且只能返回JSON，不能返回消息文字。")
        lines.append("")
        lines.append("触发条件和返回格式：")
        if "识图AI" in tools:
            lines.append('- 用户要求识别图片、重新看图片、再看图片 → 返回：{"tool":"识图AI","reply":"你想说的话"}')
        if "文件AI" in tools:
            lines.append('- 用户要求重新分析文件 → 返回：{"tool":"文件AI","reply":"你想说的话"}')
        if "生图AI" in tools:
            lines.append('- 用户要求画图、生成图片 → 返回：{"tool":"生图AI","prompt":"图片的详细描述","reply":"你想说的话"}')
        lines.append("")
        lines.append("不涉及图片/文件时，正常回复文字。")
        return "\n".join(lines)

    def _handle_ai_tool_call_for_account(self, account: BotAccount, from_user: str, tool_call: dict) -> bool:
        tool_type = tool_call.get("type")
        if tool_type == "vision":
            return self._tool_call_vision_for_account(account, from_user)
        elif tool_type == "file":
            return self._tool_call_file_for_account(account, from_user)
        elif tool_type == "image_gen":
            return self._tool_call_image_gen_for_account(account, from_user, tool_call.get("prompt", ""))
        return False

    def _tool_call_vision_for_account(self, account: BotAccount, from_user: str) -> bool:
        memories = account._get_media_memory(from_user)
        last_image = None
        for m in reversed(memories):
            if m.get("type") == "image":
                last_image = m
                break
        if not last_image or not last_image.get("cdn_info"):
            self.send_text_for_account(account, from_user, "[识图AI调用失败：没有找到可识别的图片]")
            return False
        try:
            downloaded = self.download_media_for_account(account, last_image["cdn_info"], user_id=from_user)
            if downloaded:
                img_b64 = base64.b64encode(downloaded).decode('utf-8')
                self._auto_ai_reply_with_vision_for_account(account, from_user, img_b64, original_text=last_image.get("description", ""))
                return True
            else:
                self.send_text_for_account(account, from_user, "[识图AI调用失败：图片下载失败]")
                return False
        except Exception:
            self.send_text_for_account(account, from_user, "[识图AI调用失败：图片处理异常]")
            return False

    def _tool_call_file_for_account(self, account: BotAccount, from_user: str) -> bool:
        memories = account._get_media_memory(from_user)
        last_file = None
        for m in reversed(memories):
            if m.get("type") == "file":
                last_file = m
                break
        if not last_file:
            self.send_text_for_account(account, from_user, "[文件AI调用失败：没有找到可分析的文件]")
            return False
        self.send_text_for_account(account, from_user, "[文件AI] 正在重新分析文件...")
        return True

    def _tool_call_image_gen_for_account(self, account: BotAccount, from_user: str, prompt: str) -> bool:
        if not account.ai_config.get("image_gen_enabled") or not account.ai_config.get("image_gen_api_url"):
            self.send_text_for_account(account, from_user, "[生图AI未启用]")
            return False
        try:
            api_url = account.ai_config.get("image_gen_api_url")
            api_key = account.ai_config.get("image_gen_api_key")
            model = account.ai_config.get("image_gen_model", "dall-e-3")
            payload = {"model": model, "prompt": prompt, "n": 1, "size": "1024x1024"}
            headers = {"Content-Type": "application/json"}
            if api_key:
                headers["Authorization"] = f"Bearer {api_key}"
            req = urllib.request.Request(api_url, data=json.dumps(payload).encode("utf-8"), headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=120) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                image_url = result.get("data", [{}])[0].get("url", "")
                if image_url:
                    img_req = urllib.request.Request(image_url)
                    with urllib.request.urlopen(img_req, timeout=60) as img_resp:
                        image_bytes = img_resp.read()
                    self.send_image_for_account(account, from_user, image_bytes, filename="ai_generated.png", description=prompt[:50])
                    return True
                b64 = result.get("data", [{}])[0].get("b64_json", "")
                if b64:
                    image_bytes = base64.b64decode(b64)
                    self.send_image_for_account(account, from_user, image_bytes, filename="ai_generated.png", description=prompt[:50])
                    return True
            self.send_text_for_account(account, from_user, "[生图AI调用失败]")
            return False
        except Exception as e:
            print(f"[生图AI] 异常: {e}")
            self.send_text_for_account(account, from_user, "[生图AI调用失败]")
            return False

    def _format_media_memory_for_account(self, account: BotAccount, user_id: str) -> str:
        memories = account._get_media_memory(user_id)
        if not memories:
            return ""
        lines = ["[历史媒体记录]"]
        for i, m in enumerate(memories):
            mt = "图片" if m.get('type') == 'image' else "文件"
            fn = m.get('filename', '')
            desc = m.get('description', '')
            ts = m.get('timestamp', '')
            lines.append(f"  {i+1}. [{mt}] {fn} - {desc} ({ts})")
        lines.append("[/历史媒体记录]")
        return "\n".join(lines)

    def _prefetch_media_for_account(self, account: BotAccount, cdn_media_info: dict, filename: str = "", user_id: str = ""):
        try:
            cache_key = hashlib.md5((cdn_media_info.get("encrypt_query_param") or cdn_media_info.get("encrypted_query_param") or "").encode('utf-8')).hexdigest()
            if user_id and account._get_user_cached_media(user_id, cache_key):
                return
            if account._get_cached_media(cache_key):
                return
            result = self.download_media_for_account(account, cdn_media_info, filename=filename, user_id=user_id)
            if result:
                print(f"[媒体预取] 完成: {cache_key[:12]}..., {len(result)} bytes")
        except Exception:
            pass

    def send_text_for_account(self, account: BotAccount, to_user_id: str, text: str) -> bool:
        context_token = account._context_tokens.get(to_user_id)
        if not context_token:
            return False
        use_token = account._get_token_for_user(to_user_id)
        client_id = f"msg-{uuid.uuid4().hex[:16]}"
        body = {
            "msg": {
                "from_user_id": "", "to_user_id": to_user_id, "client_id": client_id,
                "message_type": 2, "message_state": 2, "context_token": context_token,
                "item_list": [{"type": 1, "text_item": {"text": text}}]
            }
        }
        result = self._post("sendmessage", body, token=use_token)
        errcode = result.get("errcode")
        ret = result.get("ret")
        success = (ret is None or ret == 0) and (errcode is None or errcode == 0)
        if success:
            print(f"[发送成功] [{account.username}] 给 {to_user_id}: {text[:50]}...")
            out_msg = {'from': 'me', 'to': to_user_id, 'text': text, 'time': datetime.now().strftime('%H:%M:%S'), 'type': 'out'}
            account._add_message_to_history(out_msg)
            return True
        if errcode in self.EXPIRED_CODES:
            account._context_tokens.pop(to_user_id, None)
            bot_token = account._user_token_map.pop(to_user_id, None)
            if bot_token and bot_token in account._bot_accounts:
                account._bot_accounts[bot_token].get("context_tokens", {}).pop(to_user_id, None)
            account._save_config()
        print(f"[发送失败] [{account.username}] ret={ret}, errcode={errcode}")
        return False

    def send_typing_for_account(self, account: BotAccount, to_user_id: str) -> bool:
        context_token = account._context_tokens.get(to_user_id)
        if not context_token:
            return False
        use_token = account._get_token_for_user(to_user_id)
        try:
            config_body = {"ilink_user_id": to_user_id, "context_token": context_token}
            config_result = self._post("getconfig", config_body, token=use_token)
            typing_ticket = config_result.get("typing_ticket")
            if not typing_ticket:
                return False
            typing_body = {"ilink_user_id": to_user_id, "typing_ticket": typing_ticket, "status": 1}
            self._post("sendtyping", typing_body, token=use_token)
            return True
        except Exception:
            return False

    def send_image_for_account(self, account: BotAccount, to_user_id: str, image_bytes: bytes, filename: str = "image.jpg", description: str = "", media_data: str = "") -> bool:
        uploaded = self._upload_media_for_account(account, image_bytes, filename, media_type=1, to_user_id=to_user_id)
        if not uploaded:
            return False
        image_item = {"media": uploaded["media"], "aeskey": uploaded["aes_key_hex"], "mid_size": uploaded["encrypted_size"]}
        media_item = {"type": 2, "image_item": image_item}
        return self._send_media_message_for_account(account, to_user_id, media_item, description, media_data=media_data, media_filename=filename)

    def _upload_media_for_account(self, account: BotAccount, file_bytes: bytes, filename: str, media_type: int, to_user_id: str) -> Optional[dict]:
        try:
            use_token = account._get_token_for_user(to_user_id)
            aes_key_hex = self._random_hex(16)
            aes_key_bytes = bytes.fromhex(aes_key_hex)
            encrypted = self._aes_ecb_encrypt(file_bytes, aes_key_bytes)
            filekey = self._random_hex(16)
            raw_md5 = self._md5_hex(file_bytes)
            body = {
                "filekey": filekey, "media_type": media_type, "to_user_id": to_user_id,
                "rawsize": len(file_bytes), "rawfilemd5": raw_md5, "filesize": len(encrypted),
                "no_need_thumb": True, "aeskey": aes_key_hex
            }
            result = self._post("getuploadurl", body, token=use_token)
            ret = result.get("ret")
            errcode = result.get("errcode")
            if (ret is not None and ret != 0) or (errcode is not None and errcode != 0):
                return None
            upload_param = result.get("upload_param")
            if not upload_param:
                return None
            cdn_url = self.CDN_BASE + "/upload?encrypted_query_param=" + urllib.parse.quote(upload_param, safe='') + "&filekey=" + urllib.parse.quote(filekey, safe='')
            req = urllib.request.Request(cdn_url, data=encrypted, method='POST', headers={'Content-Type': 'application/octet-stream'})
            with urllib.request.urlopen(req, timeout=120) as resp:
                encrypted_param = resp.headers.get('x-encrypted-param', '')
                if not encrypted_param:
                    return None
                aes_key_b64 = base64.b64encode(aes_key_hex.encode('utf-8')).decode('utf-8')
                cdn_media = {"encrypt_query_param": encrypted_param, "aes_key": aes_key_b64, "encrypt_type": 1}
                return {
                    "filekey": filekey, "media": cdn_media, "aes_key_hex": aes_key_hex,
                    "raw_size": len(file_bytes), "encrypted_size": len(encrypted), "md5": raw_md5, "filename": filename
                }
        except Exception:
            return None

    def _send_media_message_for_account(self, account: BotAccount, to_user_id: str, media_item: dict, description: str = "", media_data: str = "", media_filename: str = "", media_duration: int = 0) -> bool:
        context_token = account._context_tokens.get(to_user_id)
        if not context_token:
            return False
        use_token = account._get_token_for_user(to_user_id)
        client_id = f"ilink-sdk:{int(time.time() * 1000)}-{uuid.uuid4().hex[:8]}"
        if description:
            text_item = {"type": 1, "text_item": {"text": description}}
            media_item_list = [media_item, text_item]
        else:
            media_item_list = [media_item]
        body = {
            "msg": {
                "from_user_id": "", "to_user_id": to_user_id, "client_id": client_id,
                "message_type": 2, "message_state": 2, "context_token": context_token,
                "item_list": media_item_list
            }
        }
        result = self._post("sendmessage", body, token=use_token)
        errcode = result.get("errcode")
        ret = result.get("ret")
        success = (ret is None or ret == 0) and (errcode is None or errcode == 0)
        if success:
            type_name = self.MEDIA_TYPE_NAMES.get(media_item.get("type", 0), "媒体")
            out_msg = {
                'from': 'me', 'to': to_user_id,
                'text': f"[{type_name}]" + (f" {description}" if description else ""),
                'time': datetime.now().strftime('%H:%M:%S'), 'type': 'out',
                'media_type': media_item.get("type"), 'media_data': media_data,
                'media_filename': media_filename or description, 'media_duration': media_duration
            }
            cdn_media = self._extract_cdn_media(media_item)
            if cdn_media:
                out_msg['media_cdn'] = json.dumps(cdn_media)
            account._add_message_to_history(out_msg)
            return True
        return False

    def download_media_for_account(self, account: BotAccount, cdn_media_info: dict, filename: str = "", user_id: str = "") -> Optional[bytes]:
        cache_key = hashlib.md5((cdn_media_info.get("encrypt_query_param") or cdn_media_info.get("encrypted_query_param") or "").encode('utf-8')).hexdigest()
        if user_id:
            cached = account._get_user_cached_media(user_id, cache_key)
            if cached:
                return cached[0]
        cached = account._get_cached_media(cache_key)
        if cached:
            return cached[0]
        with account._media_download_lock:
            if cache_key in account._media_downloading:
                wait_event = account._media_downloading[cache_key]
            else:
                wait_event = None
        if wait_event:
            wait_event.wait(timeout=60)
            if user_id:
                cached = account._get_user_cached_media(user_id, cache_key)
                if cached:
                    return cached[0]
            cached = account._get_cached_media(cache_key)
            if cached:
                return cached[0]
            return None
        event = threading.Event()
        with account._media_download_lock:
            account._media_downloading[cache_key] = event
        try:
            encrypt_query_param = cdn_media_info.get("encrypt_query_param") or cdn_media_info.get("encrypted_query_param")
            aes_key_b64 = cdn_media_info.get("aes_key")
            if not encrypt_query_param or not aes_key_b64:
                return None
            download_url = self.CDN_BASE + "/download?encrypted_query_param=" + urllib.parse.quote(encrypt_query_param, safe='')
            req = urllib.request.Request(download_url)
            with urllib.request.urlopen(req, timeout=60) as resp:
                data = resp.read()
                decoded_key = base64.b64decode(aes_key_b64)
                if len(decoded_key) == 16:
                    aes_key_bytes = decoded_key
                else:
                    aes_key_hex = decoded_key.decode('utf-8')
                    aes_key_bytes = bytes.fromhex(aes_key_hex)
                decrypted = self._aes_ecb_decrypt(data, aes_key_bytes)
                mime = self._detect_mime(decrypted)
                if mime == 'audio/silk':
                    wav_data = self._silk_to_wav_for_account(account, decrypted)
                    if wav_data:
                        decrypted = wav_data
                        mime = 'audio/wav'
                        filename = filename.replace('.silk', '.wav') if filename else 'voice.wav'
                elif mime == 'audio/amr':
                    wav_data = self._ffmpeg_to_wav_for_account(account, decrypted)
                    if wav_data:
                        decrypted = wav_data
                        mime = 'audio/wav'
                        filename = filename.replace('.amr', '.wav') if filename else 'voice.wav'
                account._save_media_cache(cache_key, decrypted, mime, filename)
                if user_id:
                    account._save_user_media_cache(user_id, cache_key, decrypted, mime, filename)
                return decrypted
        except Exception:
            return None
        finally:
            with account._media_download_lock:
                account._media_downloading.pop(cache_key, None)
            event.set()

    def _silk_to_wav_for_account(self, account: BotAccount, silk_data: bytes) -> Optional[bytes]:
        try:
            import pilk
        except ImportError:
            return self._ffmpeg_to_wav_for_account(account, silk_data)
        if silk_data[:1] == b'\x02' and len(silk_data) > 10 and silk_data[1:10] == b'#!SILK_V3':
            silk_data = silk_data[1:]
        if silk_data[:9] != b'#!SILK_V3':
            return self._ffmpeg_to_wav_for_account(account, silk_data)
        try:
            tmp_in = account.media_cache_dir / ('_silk_tmp_in_' + uuid.uuid4().hex[:12] + '.silk')
            tmp_out = account.media_cache_dir / ('_silk_tmp_out_' + uuid.uuid4().hex[:12] + '.pcm')
            tmp_in.write_bytes(silk_data)
            pilk.decode(str(tmp_in), str(tmp_out), pcm_rate=24000)
            if not tmp_out.exists() or tmp_out.stat().st_size == 0:
                return self._ffmpeg_to_wav_for_account(account, silk_data)
            pcm_data = tmp_out.read_bytes()
            sample_rate = 24000
            num_channels = 1
            bits_per_sample = 16
            byte_rate = sample_rate * num_channels * bits_per_sample // 8
            block_align = num_channels * bits_per_sample // 8
            data_size = len(pcm_data)
            wav_buf = io.BytesIO()
            wav_buf.write(b'RIFF')
            wav_buf.write(struct.pack('<I', 36 + data_size))
            wav_buf.write(b'WAVE')
            wav_buf.write(b'fmt ')
            wav_buf.write(struct.pack('<I', 16))
            wav_buf.write(struct.pack('<H', 1))
            wav_buf.write(struct.pack('<H', num_channels))
            wav_buf.write(struct.pack('<I', sample_rate))
            wav_buf.write(struct.pack('<I', byte_rate))
            wav_buf.write(struct.pack('<H', block_align))
            wav_buf.write(struct.pack('<H', bits_per_sample))
            wav_buf.write(b'data')
            wav_buf.write(struct.pack('<I', data_size))
            wav_buf.write(pcm_data)
            return wav_buf.getvalue()
        except Exception:
            return self._ffmpeg_to_wav_for_account(account, silk_data)
        finally:
            try:
                tmp_in.unlink(missing_ok=True)
            except Exception:
                pass
            try:
                tmp_out.unlink(missing_ok=True)
            except Exception:
                pass

    def _ffmpeg_to_wav_for_account(self, account: BotAccount, audio_data: bytes) -> Optional[bytes]:
        tmp_in = None
        tmp_out = None
        try:
            tmp_in = account.media_cache_dir / ('_ffmpeg_tmp_in_' + uuid.uuid4().hex[:12])
            tmp_out = account.media_cache_dir / ('_ffmpeg_tmp_out_' + uuid.uuid4().hex[:12] + '.wav')
            tmp_in.write_bytes(audio_data)
            result = subprocess.run(
                ['ffmpeg', '-y', '-i', str(tmp_in), '-f', 'wav', '-ar', '24000', '-ac', '1', str(tmp_out)],
                capture_output=True, timeout=30
            )
            if tmp_out.exists() and tmp_out.stat().st_size > 44:
                return tmp_out.read_bytes()
            return None
        except Exception:
            return None
        finally:
            for tmp in (tmp_in, tmp_out):
                if tmp:
                    try:
                        if tmp.exists():
                            tmp.unlink()
                    except Exception:
                        pass

    def _hash_password(self, password: str) -> str:
        salt = self._web_password_config.get("salt", "")
        if not salt:
            salt = uuid.uuid4().hex
            self._web_password_config["salt"] = salt
        iterations = self._web_password_config.get("iterations", 100000)
        dk = hashlib.pbkdf2_hmac('sha256', (salt + password).encode("utf-8"), salt.encode("utf-8"), iterations)
        return dk.hex()

    def _verify_web_password(self, password: str) -> bool:
        stored_hash = self._web_password_config.get("password_hash", "")
        if not stored_hash:
            return False
        if self._hash_password(password) == stored_hash:
            return True
        salt = self._web_password_config.get("salt", "")
        if salt and hashlib.sha256((salt + password).encode("utf-8")).hexdigest() == stored_hash:
            self._web_password_config["password_hash"] = self._hash_password(password)
            self._web_password_config["iterations"] = 100000
            self._save_web_password_config()
            return True
        return False

    def _is_web_password_set(self) -> bool:
        return bool(self._web_password_config.get("password_hash"))

    def _generate_verification_code(self, client_ip: str = '') -> str:
        import secrets
        code = str(secrets.randbelow(900000) + 100000)
        self._verification_codes[code] = {"expiry": time.time() + 300, "used": False, "ip": client_ip}
        expired_codes = [c for c, v in self._verification_codes.items() if v["expiry"] < time.time() or v["used"]]
        for c in expired_codes:
            del self._verification_codes[c]
        return code

    def _send_email_verification(self, to_email: str, code: str):
        def _do_send():
            try:
                import smtplib
                from email.mime.text import MIMEText
                from email.mime.multipart import MIMEMultipart
                from email.header import Header
                sender = 'ZynSync@local-server.local'
                current_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                html_body = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Zyn iLink</title>
<style>
    * {{ margin: 0; padding: 0; box-sizing: border-box; }}
    body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'PingFang SC', 'Microsoft YaHei', 'Helvetica Neue', sans-serif; background: #fafafc; color: #1a1a2e; -webkit-font-smoothing: antialiased; -webkit-text-size-adjust: 100%; -ms-text-size-adjust: 100%; }}
    .container {{ max-width: 520px; margin: 0 auto; padding: 40px 16px; }}
    .card {{ background: #ffffff; border-radius: 20px; border: 1px solid #e8e8f0; overflow: hidden; box-shadow: 0 1px 3px rgba(0,0,0,0.04), 0 8px 40px rgba(0,0,0,0.06); }}
    .card-inner {{ padding: 44px 36px; }}
    .eyebrow {{ display: flex; align-items: center; gap: 10px; margin-bottom: 36px; }}
    .brand-dot {{ width: 10px; height: 10px; border-radius: 3px; background: #7c3aed; animation: dotMorph 3s ease-in-out infinite; box-shadow: 0 0 12px rgba(124, 58, 237, 0.4); }}
    @keyframes dotMorph {{ 0%, 100% {{ border-radius: 3px; transform: rotate(0deg); }} 50% {{ border-radius: 50%; transform: rotate(45deg); }} }}
    .eyebrow-text {{ font-size: 13px; font-weight: 600; letter-spacing: 1.8px; color: #9494a8; text-transform: uppercase; }}
    .hero-block {{ text-align: center; margin-bottom: 36px; }}
    .wave-hand {{ display: inline-block; font-size: 52px; margin-bottom: 18px; animation: gentleWave 2.5s ease-in-out infinite; transform-origin: 70% 70%; line-height: 1; }}
    @keyframes gentleWave {{ 0%, 100% {{ transform: rotate(0deg) scale(1); }} 15% {{ transform: rotate(14deg) scale(1.08); }} 30% {{ transform: rotate(-8deg) scale(0.96); }} 45% {{ transform: rotate(6deg) scale(1.04); }} 60% {{ transform: rotate(-4deg) scale(0.98); }} 75% {{ transform: rotate(2deg) scale(1.01); }} }}
    .hero-title {{ font-size: 26px; font-weight: 700; letter-spacing: -0.6px; margin-bottom: 6px; color: #1a1a2e; }}
    .hero-sub {{ font-size: 15px; color: #5c5c78; font-weight: 400; }}
    .code-block {{ background: #f4f4f8; border: 1px solid #e8e8f0; border-radius: 10px; padding: 28px 20px 20px; text-align: center; margin-bottom: 28px; }}
    .code-block-label {{ font-size: 11px; font-weight: 600; letter-spacing: 2px; color: #9494a8; text-transform: uppercase; margin-bottom: 16px; }}
    .code-digits {{ display: flex; justify-content: center; gap: 10px; margin-bottom: 20px; }}
    .digit {{ width: 48px; height: 58px; background: #ffffff; border: 1.5px solid #e8e8f0; border-radius: 6px; display: flex; align-items: center; justify-content: center; font-family: 'SF Mono', 'Menlo', 'Consolas', 'Courier New', monospace; font-size: 26px; font-weight: 700; color: #1a1a2e; animation: digitPopIn 0.5s cubic-bezier(0.34, 1.56, 0.64, 1) both; box-shadow: 0 2px 8px rgba(0,0,0,0.04); }}
    .digit:nth-child(1) {{ animation-delay: 0.05s; }} .digit:nth-child(2) {{ animation-delay: 0.1s; }} .digit:nth-child(3) {{ animation-delay: 0.15s; }} .digit:nth-child(4) {{ animation-delay: 0.2s; }} .digit:nth-child(5) {{ animation-delay: 0.25s; }} .digit:nth-child(6) {{ animation-delay: 0.3s; }}
    @keyframes digitPopIn {{ from {{ opacity: 0; transform: translateY(12px) scale(0.7); }} to {{ opacity: 1; transform: translateY(0) scale(1); }} }}
    .timer-badge {{ display: inline-flex; align-items: center; gap: 7px; font-size: 12px; font-weight: 500; color: #d97706; background: #fffbeb; padding: 7px 14px; border-radius: 20px; letter-spacing: 0.3px; }}
    .timer-dot {{ width: 7px; height: 7px; border-radius: 50%; background: #d97706; animation: softPulse 1.8s ease-in-out infinite; }}
    @keyframes softPulse {{ 0%, 100% {{ opacity: 0.4; transform: scale(0.8); }} 50% {{ opacity: 1; transform: scale(1.2); }} }}
    .meta-stack {{ display: flex; flex-direction: column; gap: 1px; margin-bottom: 24px; }}
    .meta-row {{ display: flex; align-items: center; gap: 12px; padding: 12px 16px; background: #f4f4f8; border-radius: 6px; font-size: 13px; }}
    .meta-icon {{ font-size: 17px; flex-shrink: 0; width: 22px; text-align: center; }}
    .meta-label {{ color: #9494a8; font-weight: 500; font-size: 11px; letter-spacing: 0.5px; text-transform: uppercase; min-width: 56px; }}
    .meta-value {{ color: #1a1a2e; font-weight: 600; margin-left: auto; text-align: right; }}
    .divider-line {{ height: 1px; background: #e8e8f0; margin: 24px 0; }}
    .security-notice {{ display: flex; align-items: flex-start; gap: 12px; padding: 16px; background: #fff1f2; border-radius: 10px; border: 1px solid rgba(225, 29, 72, 0.12); }}
    .security-notice .meta-icon {{ font-size: 20px; flex-shrink: 0; }}
    .security-text {{ font-size: 13px; color: #e11d48; font-weight: 500; line-height: 1.6; }}
    .security-text strong {{ font-weight: 700; }}
    .footer {{ text-align: center; padding: 28px 36px 20px; border-top: 1px solid #e8e8f0; background: #f4f4f8; }}
    .footer-brand {{ font-size: 14px; font-weight: 700; color: #1a1a2e; letter-spacing: -0.3px; margin-bottom: 6px; }}
    .footer-tag {{ font-size: 11px; color: #9494a8; letter-spacing: 1.5px; text-transform: uppercase; }}
    .footer-copy {{ font-size: 11px; color: #9494a8; margin-top: 14px; line-height: 1.7; }}
    @media (max-width: 480px) {{ .card-inner {{ padding: 32px 20px; }} .digit {{ width: 40px; height: 50px; font-size: 22px; }} .code-digits {{ gap: 6px; }} .hero-title {{ font-size: 22px; }} }}
</style>
</head>
<body>
<div class="container">
    <div class="card">
        <div class="card-inner">
            <div class="eyebrow"><div class="brand-dot"></div><span class="eyebrow-text">Zyn iLink · ChatBox</span></div>
            <div class="hero-block"><div class="wave-hand">\U0001f44b</div><div class="hero-title">验证你的账号</div><div class="hero-sub">输入下方验证码以完成注册</div></div>
            <div class="code-block"><div class="code-block-label">安全验证码</div><div class="code-digits"><div class="digit">{code[0]}</div><div class="digit">{code[1]}</div><div class="digit">{code[2]}</div><div class="digit">{code[3]}</div><div class="digit">{code[4]}</div><div class="digit">{code[5]}</div></div><div class="timer-badge"><span class="timer-dot"></span>5 分钟内有效</div></div>
            <div class="meta-stack"><div class="meta-row"><span class="meta-icon">&#x1F550;</span><span class="meta-label">时间</span><span class="meta-value">{current_time}</span></div><div class="meta-row"><span class="meta-icon">&#x1F4CD;</span><span class="meta-label">来源</span><span class="meta-value">ChatBox 客户端</span></div><div class="meta-row"><span class="meta-icon">&#x1F194;</span><span class="meta-label">请求</span><span class="meta-value">账号注册</span></div></div>
            <div class="divider-line"></div>
            <div class="security-notice"><span class="meta-icon">&#x1F510;</span><span class="security-text"><strong>请勿分享</strong> — 我们绝不会通过任何渠道索要此验证码。如非本人操作，请忽略此邮件。</span></div>
        </div>
        <div class="footer"><div class="footer-brand">ZynSync</div><div class="footer-tag">Intelligent ChatBox</div><div class="footer-copy">&copy; {datetime.now().year} Zyn iLink &middot; All rights reserved.<br>此邮件由系统自动发送</div></div>
    </div>
</div>
</body>
</html>"""
                text_body = f"Zyn iLink ChatBox - 注册验证\n\n验证码: {code}\n发送时间: {current_time}\n\n请勿将验证码透露给他人。如非本人操作，请忽略此邮件。\n\n© {datetime.now().year} Zyn iLink\n"
                msg = MIMEMultipart('alternative')
                msg['From'] = sender
                msg['To'] = to_email
                msg['Subject'] = Header('Zyn iLink ChatBox - 注册验证', 'utf-8')
                msg.attach(MIMEText(text_body, 'plain', 'utf-8'))
                msg.attach(MIMEText(html_body, 'html', 'utf-8'))
                mx_servers = ['mx3.qq.com', 'mx1.qq.com', 'mx2.qq.com']
                for mx_server in mx_servers:
                    try:
                        server = smtplib.SMTP(mx_server, 25, timeout=10)
                        server.set_debuglevel(0)
                        server.ehlo('local-server')
                        server.sendmail(sender, [to_email], msg.as_string())
                        server.quit()
                        print(f"[邮件] 验证码已发送至 {to_email}")
                        return
                    except Exception:
                        continue
                print(f"[邮件] 发送失败: 所有MX服务器均不可达")
            except Exception as e:
                print(f"[邮件] 发送失败: {e}")
        threading.Thread(target=_do_send, daemon=True).start()

    def _send_forgot_password_email(self, to_email: str, code: str, username: str):
        def _do_send():
            try:
                import smtplib
                from email.mime.text import MIMEText
                from email.mime.multipart import MIMEMultipart
                from email.header import Header
                sender = 'ZynSync@local-server.local'
                current_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                html_body = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Zyn iLink</title>
<style>
    * {{ margin: 0; padding: 0; box-sizing: border-box; }}
    body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'PingFang SC', 'Microsoft YaHei', 'Helvetica Neue', sans-serif; background: #fafafc; color: #1a1a2e; -webkit-font-smoothing: antialiased; -webkit-text-size-adjust: 100%; -ms-text-size-adjust: 100%; }}
    .container {{ max-width: 520px; margin: 0 auto; padding: 40px 16px; }}
    .card {{ background: #ffffff; border-radius: 20px; border: 1px solid #e8e8f0; overflow: hidden; box-shadow: 0 1px 3px rgba(0,0,0,0.04), 0 8px 40px rgba(0,0,0,0.06); }}
    .card-inner {{ padding: 44px 36px; }}
    .eyebrow {{ display: flex; align-items: center; gap: 10px; margin-bottom: 36px; }}
    .brand-dot {{ width: 10px; height: 10px; border-radius: 3px; background: #7c3aed; animation: dotMorph 3s ease-in-out infinite; box-shadow: 0 0 12px rgba(124, 58, 237, 0.4); }}
    @keyframes dotMorph {{ 0%, 100% {{ border-radius: 3px; transform: rotate(0deg); }} 50% {{ border-radius: 50%; transform: rotate(45deg); }} }}
    .eyebrow-text {{ font-size: 13px; font-weight: 600; letter-spacing: 1.8px; color: #9494a8; text-transform: uppercase; }}
    .hero-block {{ text-align: center; margin-bottom: 36px; }}
    .hero-icon {{ display: inline-block; font-size: 52px; margin-bottom: 18px; line-height: 1; }}
    .hero-title {{ font-size: 26px; font-weight: 700; letter-spacing: -0.6px; margin-bottom: 6px; color: #1a1a2e; }}
    .hero-sub {{ font-size: 15px; color: #5c5c78; font-weight: 400; }}
    .code-block {{ background: #f4f4f8; border: 1px solid #e8e8f0; border-radius: 10px; padding: 28px 20px 20px; text-align: center; margin-bottom: 28px; }}
    .code-block-label {{ font-size: 11px; font-weight: 600; letter-spacing: 2px; color: #9494a8; text-transform: uppercase; margin-bottom: 16px; }}
    .code-digits {{ display: flex; justify-content: center; gap: 10px; margin-bottom: 20px; }}
    .digit {{ width: 48px; height: 58px; background: #ffffff; border: 1.5px solid #e8e8f0; border-radius: 6px; display: flex; align-items: center; justify-content: center; font-family: 'SF Mono', 'Menlo', 'Consolas', 'Courier New', monospace; font-size: 26px; font-weight: 700; color: #1a1a2e; animation: digitPopIn 0.5s cubic-bezier(0.34, 1.56, 0.64, 1) both; box-shadow: 0 2px 8px rgba(0,0,0,0.04); }}
    .digit:nth-child(1) {{ animation-delay: 0.05s; }} .digit:nth-child(2) {{ animation-delay: 0.1s; }} .digit:nth-child(3) {{ animation-delay: 0.15s; }} .digit:nth-child(4) {{ animation-delay: 0.2s; }} .digit:nth-child(5) {{ animation-delay: 0.25s; }} .digit:nth-child(6) {{ animation-delay: 0.3s; }}
    @keyframes digitPopIn {{ from {{ opacity: 0; transform: translateY(12px) scale(0.7); }} to {{ opacity: 1; transform: translateY(0) scale(1); }} }}
    .timer-badge {{ display: inline-flex; align-items: center; gap: 7px; font-size: 12px; font-weight: 500; color: #d97706; background: #fffbeb; padding: 7px 14px; border-radius: 20px; letter-spacing: 0.3px; }}
    .timer-dot {{ width: 7px; height: 7px; border-radius: 50%; background: #d97706; animation: softPulse 1.8s ease-in-out infinite; }}
    @keyframes softPulse {{ 0%, 100% {{ opacity: 0.4; transform: scale(0.8); }} 50% {{ opacity: 1; transform: scale(1.2); }} }}
    .meta-stack {{ display: flex; flex-direction: column; gap: 1px; margin-bottom: 24px; }}
    .meta-row {{ display: flex; align-items: center; gap: 12px; padding: 12px 16px; background: #f4f4f8; border-radius: 6px; font-size: 13px; }}
    .meta-icon {{ font-size: 17px; flex-shrink: 0; width: 22px; text-align: center; }}
    .meta-label {{ color: #9494a8; font-weight: 500; font-size: 11px; letter-spacing: 0.5px; text-transform: uppercase; min-width: 56px; }}
    .meta-value {{ color: #1a1a2e; font-weight: 600; margin-left: auto; text-align: right; }}
    .divider-line {{ height: 1px; background: #e8e8f0; margin: 24px 0; }}
    .security-notice {{ display: flex; align-items: flex-start; gap: 12px; padding: 16px; background: #fff1f2; border-radius: 10px; border: 1px solid rgba(225, 29, 72, 0.12); }}
    .security-notice .meta-icon {{ font-size: 20px; flex-shrink: 0; }}
    .security-text {{ font-size: 13px; color: #e11d48; font-weight: 500; line-height: 1.6; }}
    .security-text strong {{ font-weight: 700; }}
    .footer {{ text-align: center; padding: 28px 36px 20px; border-top: 1px solid #e8e8f0; background: #f4f4f8; }}
    .footer-brand {{ font-size: 14px; font-weight: 700; color: #1a1a2e; letter-spacing: -0.3px; margin-bottom: 6px; }}
    .footer-tag {{ font-size: 11px; color: #9494a8; letter-spacing: 1.5px; text-transform: uppercase; }}
    .footer-copy {{ font-size: 11px; color: #9494a8; margin-top: 14px; line-height: 1.7; }}
    @media (max-width: 480px) {{ .card-inner {{ padding: 32px 20px; }} .digit {{ width: 40px; height: 50px; font-size: 22px; }} .code-digits {{ gap: 6px; }} .hero-title {{ font-size: 22px; }} }}
</style>
</head>
<body>
<div class="container">
    <div class="card">
        <div class="card-inner">
            <div class="eyebrow"><div class="brand-dot"></div><span class="eyebrow-text">Zyn iLink · ChatBox</span></div>
            <div class="hero-block"><div class="hero-icon">&#x1F512;</div><div class="hero-title">重置你的密码</div><div class="hero-sub">输入下方验证码以重置账户密码</div></div>
            <div class="code-block"><div class="code-block-label">安全验证码</div><div class="code-digits"><div class="digit">{code[0]}</div><div class="digit">{code[1]}</div><div class="digit">{code[2]}</div><div class="digit">{code[3]}</div><div class="digit">{code[4]}</div><div class="digit">{code[5]}</div></div><div class="timer-badge"><span class="timer-dot"></span>5 分钟内有效</div></div>
            <div class="meta-stack"><div class="meta-row"><span class="meta-icon">&#x1F550;</span><span class="meta-label">时间</span><span class="meta-value">{current_time}</span></div><div class="meta-row"><span class="meta-icon">&#x1F4CD;</span><span class="meta-label">来源</span><span class="meta-value">ChatBox 客户端</span></div><div class="meta-row"><span class="meta-icon">&#x1F194;</span><span class="meta-label">请求</span><span class="meta-value">账户 {username} 密码重置</span></div></div>
            <div class="divider-line"></div>
            <div class="security-notice"><span class="meta-icon">&#x1F510;</span><span class="security-text"><strong>请勿分享</strong> — 我们绝不会通过任何渠道索要此验证码。如非本人操作，请忽略此邮件。</span></div>
        </div>
        <div class="footer"><div class="footer-brand">ZynSync</div><div class="footer-tag">Intelligent ChatBox</div><div class="footer-copy">&copy; {datetime.now().year} Zyn iLink &middot; All rights reserved.<br>此邮件由系统自动发送</div></div>
    </div>
</div>
</body>
</html>"""
                text_body = f"Zyn iLink ChatBox - 密码重置\n\n账户: {username}\n验证码: {code}\n发送时间: {current_time}\n\n请勿将验证码透露给他人。如非本人操作，请忽略此邮件。\n\n© {datetime.now().year} Zyn iLink\n"
                msg = MIMEMultipart('alternative')
                msg['From'] = sender
                msg['To'] = to_email
                msg['Subject'] = Header('Zyn iLink ChatBox - 密码重置', 'utf-8')
                msg.attach(MIMEText(text_body, 'plain', 'utf-8'))
                msg.attach(MIMEText(html_body, 'html', 'utf-8'))
                mx_servers = ['mx3.qq.com', 'mx1.qq.com', 'mx2.qq.com']
                for mx_server in mx_servers:
                    try:
                        server = smtplib.SMTP(mx_server, 25, timeout=10)
                        server.set_debuglevel(0)
                        server.ehlo('local-server')
                        server.sendmail(sender, [to_email], msg.as_string())
                        server.quit()
                        print(f"[邮件] 密码重置验证码已发送至 {to_email}")
                        return
                    except Exception:
                        continue
                print(f"[邮件] 发送失败: 所有MX服务器均不可达")
            except Exception as e:
                print(f"[邮件] 发送失败: {e}")
        threading.Thread(target=_do_send, daemon=True).start()

    def _send_bind_email_verification(self, to_email: str, code: str):
        def _do_send():
            try:
                import smtplib
                from email.mime.text import MIMEText
                from email.mime.multipart import MIMEMultipart
                from email.header import Header
                sender = 'ZynSync@local-server.local'
                current_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                html_body = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Zyn iLink</title>
<style>
    * {{ margin: 0; padding: 0; box-sizing: border-box; }}
    body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', 'PingFang SC', 'Microsoft YaHei', 'Helvetica Neue', sans-serif; background: #fafafc; color: #1a1a2e; -webkit-font-smoothing: antialiased; -webkit-text-size-adjust: 100%; -ms-text-size-adjust: 100%; }}
    .container {{ max-width: 520px; margin: 0 auto; padding: 40px 16px; }}
    .card {{ background: #ffffff; border-radius: 20px; border: 1px solid #e8e8f0; overflow: hidden; box-shadow: 0 1px 3px rgba(0,0,0,0.04), 0 8px 40px rgba(0,0,0,0.06); }}
    .card-inner {{ padding: 44px 36px; }}
    .eyebrow {{ display: flex; align-items: center; gap: 10px; margin-bottom: 36px; }}
    .brand-dot {{ width: 10px; height: 10px; border-radius: 3px; background: #7c3aed; animation: dotMorph 3s ease-in-out infinite; box-shadow: 0 0 12px rgba(124, 58, 237, 0.4); }}
    @keyframes dotMorph {{ 0%, 100% {{ border-radius: 3px; transform: rotate(0deg); }} 50% {{ border-radius: 50%; transform: rotate(45deg); }} }}
    .eyebrow-text {{ font-size: 13px; font-weight: 600; letter-spacing: 1.8px; color: #9494a8; text-transform: uppercase; }}
    .hero-block {{ text-align: center; margin-bottom: 36px; }}
    .hero-icon {{ display: inline-block; font-size: 52px; margin-bottom: 18px; line-height: 1; }}
    .hero-title {{ font-size: 26px; font-weight: 700; letter-spacing: -0.6px; margin-bottom: 6px; color: #1a1a2e; }}
    .hero-sub {{ font-size: 15px; color: #5c5c78; font-weight: 400; }}
    .code-block {{ background: #f4f4f8; border: 1px solid #e8e8f0; border-radius: 10px; padding: 28px 20px 20px; text-align: center; margin-bottom: 28px; }}
    .code-block-label {{ font-size: 11px; font-weight: 600; letter-spacing: 2px; color: #9494a8; text-transform: uppercase; margin-bottom: 16px; }}
    .code-digits {{ display: flex; justify-content: center; gap: 10px; margin-bottom: 20px; }}
    .digit {{ width: 48px; height: 58px; background: #ffffff; border: 1.5px solid #e8e8f0; border-radius: 6px; display: flex; align-items: center; justify-content: center; font-family: 'SF Mono', 'Menlo', 'Consolas', 'Courier New', monospace; font-size: 26px; font-weight: 700; color: #1a1a2e; animation: digitPopIn 0.5s cubic-bezier(0.34, 1.56, 0.64, 1) both; box-shadow: 0 2px 8px rgba(0,0,0,0.04); }}
    .digit:nth-child(1) {{ animation-delay: 0.05s; }} .digit:nth-child(2) {{ animation-delay: 0.1s; }} .digit:nth-child(3) {{ animation-delay: 0.15s; }} .digit:nth-child(4) {{ animation-delay: 0.2s; }} .digit:nth-child(5) {{ animation-delay: 0.25s; }} .digit:nth-child(6) {{ animation-delay: 0.3s; }}
    @keyframes digitPopIn {{ from {{ opacity: 0; transform: translateY(12px) scale(0.7); }} to {{ opacity: 1; transform: translateY(0) scale(1); }} }}
    .timer-badge {{ display: inline-flex; align-items: center; gap: 7px; font-size: 12px; font-weight: 500; color: #d97706; background: #fffbeb; padding: 7px 14px; border-radius: 20px; letter-spacing: 0.3px; }}
    .timer-dot {{ width: 7px; height: 7px; border-radius: 50%; background: #d97706; animation: softPulse 1.8s ease-in-out infinite; }}
    @keyframes softPulse {{ 0%, 100% {{ opacity: 0.4; transform: scale(0.8); }} 50% {{ opacity: 1; transform: scale(1.2); }} }}
    .meta-stack {{ display: flex; flex-direction: column; gap: 1px; margin-bottom: 24px; }}
    .meta-row {{ display: flex; align-items: center; gap: 12px; padding: 12px 16px; background: #f4f4f8; border-radius: 6px; font-size: 13px; }}
    .meta-icon {{ font-size: 17px; flex-shrink: 0; width: 22px; text-align: center; }}
    .meta-label {{ color: #9494a8; font-weight: 500; font-size: 11px; letter-spacing: 0.5px; text-transform: uppercase; min-width: 56px; }}
    .meta-value {{ color: #1a1a2e; font-weight: 600; margin-left: auto; text-align: right; }}
    .divider-line {{ height: 1px; background: #e8e8f0; margin: 24px 0; }}
    .security-notice {{ display: flex; align-items: flex-start; gap: 12px; padding: 16px; background: #fff1f2; border-radius: 10px; border: 1px solid rgba(225, 29, 72, 0.12); }}
    .security-notice .meta-icon {{ font-size: 20px; flex-shrink: 0; }}
    .security-text {{ font-size: 13px; color: #e11d48; font-weight: 500; line-height: 1.6; }}
    .security-text strong {{ font-weight: 700; }}
    .footer {{ text-align: center; padding: 28px 36px 20px; border-top: 1px solid #e8e8f0; background: #f4f4f8; }}
    .footer-brand {{ font-size: 14px; font-weight: 700; color: #1a1a2e; letter-spacing: -0.3px; margin-bottom: 6px; }}
    .footer-tag {{ font-size: 11px; color: #9494a8; letter-spacing: 1.5px; text-transform: uppercase; }}
    .footer-copy {{ font-size: 11px; color: #9494a8; margin-top: 14px; line-height: 1.7; }}
    @media (max-width: 480px) {{ .card-inner {{ padding: 32px 20px; }} .digit {{ width: 40px; height: 50px; font-size: 22px; }} .code-digits {{ gap: 6px; }} .hero-title {{ font-size: 22px; }} }}
</style>
</head>
<body>
<div class="container">
    <div class="card">
        <div class="card-inner">
            <div class="eyebrow"><div class="brand-dot"></div><span class="eyebrow-text">Zyn iLink · ChatBox</span></div>
            <div class="hero-block"><div class="hero-icon">&#x1F4E7;</div><div class="hero-title">绑定你的邮箱</div><div class="hero-sub">输入下方验证码以完成邮箱绑定</div></div>
            <div class="code-block"><div class="code-block-label">安全验证码</div><div class="code-digits"><div class="digit">{code[0]}</div><div class="digit">{code[1]}</div><div class="digit">{code[2]}</div><div class="digit">{code[3]}</div><div class="digit">{code[4]}</div><div class="digit">{code[5]}</div></div><div class="timer-badge"><span class="timer-dot"></span>5 分钟内有效</div></div>
            <div class="meta-stack"><div class="meta-row"><span class="meta-icon">&#x1F550;</span><span class="meta-label">时间</span><span class="meta-value">{current_time}</span></div><div class="meta-row"><span class="meta-icon">&#x1F4CD;</span><span class="meta-label">来源</span><span class="meta-value">ChatBox 客户端</span></div><div class="meta-row"><span class="meta-icon">&#x1F194;</span><span class="meta-label">请求</span><span class="meta-value">邮箱绑定</span></div></div>
            <div class="divider-line"></div>
            <div class="security-notice"><span class="meta-icon">&#x1F510;</span><span class="security-text"><strong>请勿分享</strong> — 我们绝不会通过任何渠道索要此验证码。如非本人操作，请忽略此邮件。</span></div>
        </div>
        <div class="footer"><div class="footer-brand">ZynSync</div><div class="footer-tag">Intelligent ChatBox</div><div class="footer-copy">&copy; {datetime.now().year} Zyn iLink &middot; All rights reserved.<br>此邮件由系统自动发送</div></div>
    </div>
</div>
</body>
</html>"""
                text_body = f"Zyn iLink ChatBox - 邮箱绑定\n\n验证码: {code}\n发送时间: {current_time}\n\n请勿将验证码透露给他人。如非本人操作，请忽略此邮件。\n\n© {datetime.now().year} Zyn iLink\n"
                msg = MIMEMultipart('alternative')
                msg['From'] = sender
                msg['To'] = to_email
                msg['Subject'] = Header('Zyn iLink ChatBox - 邮箱绑定', 'utf-8')
                msg.attach(MIMEText(text_body, 'plain', 'utf-8'))
                msg.attach(MIMEText(html_body, 'html', 'utf-8'))
                mx_servers = ['mx3.qq.com', 'mx1.qq.com', 'mx2.qq.com']
                for mx_server in mx_servers:
                    try:
                        server = smtplib.SMTP(mx_server, 25, timeout=10)
                        server.set_debuglevel(0)
                        server.ehlo('local-server')
                        server.sendmail(sender, [to_email], msg.as_string())
                        server.quit()
                        print(f"[邮件] 邮箱绑定验证码已发送至 {to_email}")
                        return
                    except Exception:
                        continue
                print(f"[邮件] 发送失败: 所有MX服务器均不可达")
            except Exception as e:
                print(f"[邮件] 发送失败: {e}")
        threading.Thread(target=_do_send, daemon=True).start()

    def _generate_captcha(self, session_token: str) -> tuple:
        import string
        import secrets
        chars = string.ascii_uppercase + string.digits
        captcha_text = ''.join(secrets.choice(chars) for _ in range(4))
        self._captcha_store[session_token] = {"answer": captcha_text, "expiry": time.time() + 300}
        expired = [k for k, v in self._captcha_store.items() if v["expiry"] < time.time()]
        for k in expired:
            del self._captcha_store[k]
        try:
            from PIL import Image, ImageDraw, ImageFont
            img = Image.new('RGB', (120, 40), (245, 245, 245))
            draw = ImageDraw.Draw(img)
            for _ in range(5):
                x1 = random.randint(0, 120)
                y1 = random.randint(0, 40)
                x2 = random.randint(0, 120)
                y2 = random.randint(0, 40)
                draw.line([(x1, y1), (x2, y2)], fill=(200, 200, 200), width=1)
            for _ in range(30):
                x = random.randint(0, 120)
                y = random.randint(0, 40)
                draw.point((x, y), fill=(180, 180, 180))
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 28)
            except Exception:
                try:
                    font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 28)
                except Exception:
                    font = ImageFont.load_default()
            for i, ch in enumerate(captcha_text):
                x = 10 + i * 26 + random.randint(-3, 3)
                y = 4 + random.randint(-3, 3)
                r = random.randint(0, 100)
                g = random.randint(0, 100)
                b = random.randint(0, 100)
                draw.text((x, y), ch, fill=(r, g, b), font=font)
            buf = io.BytesIO()
            img.save(buf, format='PNG')
            return ('image/png', buf.getvalue())
        except ImportError:
            svg_parts = []
            for i, ch in enumerate(captcha_text):
                x = 15 + i * 26
                y = 28 + random.randint(-3, 3)
                r = random.randint(0, 100)
                g = random.randint(0, 100)
                b = random.randint(0, 100)
                rot = random.randint(-15, 15)
                safe_ch = ch.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;').replace("'", '&#x27;').replace('`', '&#x60;')
                svg_parts.append(f'<text x="{x}" y="{y}" fill="rgb({r},{g},{b})" font-size="28" font-weight="bold" font-family="monospace" transform="rotate({rot},{x},{y})">{safe_ch}</text>')
            for _ in range(5):
                x1 = random.randint(0, 120)
                y1 = random.randint(0, 40)
                x2 = random.randint(0, 120)
                y2 = random.randint(0, 40)
                svg_parts.append(f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="#ccc" stroke-width="1"/>')
            svg = f'<svg xmlns="http://www.w3.org/2000/svg" width="120" height="40" viewBox="0 0 120 40"><rect width="120" height="40" fill="#f5f5f5"/>{"".join(svg_parts)}</svg>'
            return ('image/svg+xml', svg.encode('utf-8'))
    
    def get_effective_system_prompt(self, user_id: str = "") -> str:
        prompts = self.user_prompts.get("prompts", {})
        if user_id and user_id in prompts and prompts[user_id].strip():
            return prompts[user_id]
        system_prompt = self.ai_config.get("system_prompt", "")
        if not system_prompt:
            system_prompt = "你是一个微信聊天助手，请用自然的中文回复。"
        return system_prompt
    
    def is_ai_enabled_for_user(self, user_id: str) -> bool:
        ai_enabled = self.user_prompts.get("ai_enabled", {})
        if user_id in ai_enabled:
            return ai_enabled[user_id]
        return self.ai_config.get("auto_reply", False)

    def is_scheduled_enabled_for_user(self, user_id: str) -> bool:
        scheduled_enabled = self.user_prompts.get("scheduled_enabled", {})
        if user_id in scheduled_enabled:
            return scheduled_enabled[user_id]
        return self.ai_config.get("scheduled_reply", False)

    def is_daily_enabled_for_user(self, user_id: str) -> bool:
        daily_enabled = self.user_prompts.get("daily_enabled", {})
        if user_id in daily_enabled:
            return daily_enabled[user_id]
        return self.ai_config.get("daily_reply", False)

    def _generate_session_token(self) -> str:
        import secrets
        token = secrets.token_hex(32)
        self._session_tokens[token] = time.time() + 3600
        if len(self._session_tokens) > 100:
            self._cleanup_expired_sessions()
        return token
    
    def _verify_session_token(self, token: str) -> bool:
        if not token:
            return False
        if token in self._session_tokens:
            if time.time() < self._session_tokens[token]:
                return True
            else:
                del self._session_tokens[token]
        if token in self._account_sessions:
            return True
        if token in self._verified_sessions and time.time() <= self._verified_sessions[token]:
            return True
        return False
    
    def _cleanup_expired_sessions(self):
        now = time.time()
        expired = [t for t, exp in self._session_tokens.items() if exp <= now]
        for t in expired:
            del self._session_tokens[t]
            self._verified_sessions.pop(t, None)

    def _sse_add(self, username, wfile):
        with self._sse_lock:
            if username not in self._sse_connections:
                self._sse_connections[username] = []
            self._sse_connections[username].append(wfile)

    def _sse_remove(self, username, wfile):
        with self._sse_lock:
            if username in self._sse_connections:
                try:
                    self._sse_connections[username].remove(wfile)
                except ValueError:
                    pass
                if not self._sse_connections[username]:
                    del self._sse_connections[username]

    def _sse_send(self, username, event, data):
        with self._sse_lock:
            connections = list(self._sse_connections.get(username, []))
        msg = f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"
        encoded = msg.encode("utf-8")
        dead = []
        for wfile in connections:
            try:
                wfile.write(encoded)
                wfile.flush()
            except Exception:
                dead.append(wfile)
        for wfile in dead:
            self._sse_remove(username, wfile)
    
    def _get_token_for_user(self, user_id: str) -> Optional[str]:
        return self._user_token_map.get(user_id) or self.token
    
    def _register_user_to_account(self, user_id: str, ctx_token: str, bot_token: str):
        self._context_tokens[user_id] = ctx_token
        self._user_token_map[user_id] = bot_token
        
        if bot_token not in self._bot_accounts:
            self._bot_accounts[bot_token] = {
                "bot_id": self.bot_id or "",
                "user_id": self.user_id or "",
                "cursor": "",
                "context_tokens": {}
            }
        self._bot_accounts[bot_token]["context_tokens"][user_id] = ctx_token
        
        self._save_user_token(user_id, ctx_token)
        if not self._current_user:
            self._current_user = user_id
    
    def _get_user_dir(self, user_id: str) -> Path:
        safe_id = hashlib.md5(user_id.encode('utf-8')).hexdigest()[:16]
        user_dir = self._user_data_dir / safe_id
        user_dir.mkdir(parents=True, exist_ok=True)
        return user_dir
    
    def _get_user_dir_path(self, user_id: str) -> Path:
        safe_id = hashlib.md5(user_id.encode('utf-8')).hexdigest()[:16]
        return self._user_data_dir / safe_id
    
    def _get_user_media_dir(self, user_id: str) -> Path:
        media_dir = self._get_user_dir(user_id) / "media"
        media_dir.mkdir(parents=True, exist_ok=True)
        return media_dir
    
    def _get_user_token_file(self, user_id: str) -> Path:
        return self._get_user_dir(user_id) / "token.json"
    
    def _get_user_messages_file(self, user_id: str) -> Path:
        return self._get_user_dir(user_id) / "messages.json"
    
    def _save_user_token(self, user_id: str, context_token: str):
        try:
            data = {
                "user_id": user_id,
                "context_token": context_token,
                "saved_at": datetime.now().isoformat()
            }
            with open(self._get_user_token_file(user_id), "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"[USER] 保存用户 token 失败 ({user_id}): {e}")
    
    def _load_user_token(self, user_id: str) -> Optional[str]:
        token_file = self._get_user_token_file(user_id)
        try:
            if token_file.exists():
                with open(token_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    return data.get("context_token")
        except Exception as e:
            print(f"[USER] 加载用户 token 失败 ({user_id}): {e}")
        return None
    
    def _save_user_messages(self, user_id: str):
        if not user_id:
            return
        try:
            user_msgs = [m for m in self._messages 
                        if m.get('from') == user_id or m.get('to') == user_id]
            
            if len(user_msgs) > self._max_messages_per_user * 2:
                user_msgs = user_msgs[-self._max_messages_per_user:]
            
            data = {
                "user_id": user_id,
                "messages": user_msgs,
                "count": len(user_msgs),
                "saved_at": datetime.now().isoformat()
            }
            with open(self._get_user_messages_file(user_id), "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"[USER] 保存用户消息失败 ({user_id}): {e}")
    
    def _load_user_messages(self, user_id: str) -> List[dict]:
        msg_file = self._get_user_messages_file(user_id)
        try:
            if msg_file.exists():
                with open(msg_file, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    return data.get("messages", [])
        except Exception as e:
            print(f"[USER] 加载用户消息失败 ({user_id}): {e}")
        return []
    
    def _load_all_user_messages(self):
        all_msgs = []
        loaded_ids = set()
        
        for user_id in self._context_tokens.keys():
            user_msgs = self._load_user_messages(user_id)
            for msg in user_msgs:
                msg_id = msg.get('id')
                if msg_id and msg_id not in loaded_ids:
                    all_msgs.append(msg)
                    loaded_ids.add(msg_id)
        
        all_msgs.sort(key=lambda m: m.get('id', 0))
        self._messages = all_msgs
        
        if self._messages:
            self._last_msg_id = max(msg.get('id', 0) for msg in self._messages)
        else:
            self._last_msg_id = 0
        
        print(f"[MSG] 已从用户文件夹加载 {len(self._messages)} 条消息 ({len(self._context_tokens)} 个用户)")
    
    def _save_all_messages(self):
        for user_id in self._context_tokens.keys():
            self._save_user_messages(user_id)
    
    def _get_user_media_cache_path(self, user_id: str, cache_key: str) -> Path:
        return self._get_user_media_dir(user_id) / cache_key
    
    def _get_user_media_meta_path(self, user_id: str, cache_key: str) -> Path:
        return self._get_user_media_dir(user_id) / (cache_key + ".meta")
    
    def _save_user_media_cache(self, user_id: str, cache_key: str, media_data: bytes, mime: str, filename: str = ""):
        try:
            self._get_user_media_cache_path(user_id, cache_key).write_bytes(media_data)
            meta = {'mime': mime, 'filename': filename, 'size': len(media_data)}
            self._get_user_media_meta_path(user_id, cache_key).write_text(json.dumps(meta, ensure_ascii=False), 'utf-8')
        except Exception as e:
            print(f"[媒体缓存] 保存失败: {e}")
    
    def _get_user_cached_media(self, user_id: str, cache_key: str) -> Optional[tuple]:
        data_path = self._get_user_media_cache_path(user_id, cache_key)
        meta_path = self._get_user_media_meta_path(user_id, cache_key)
        if data_path.exists() and meta_path.exists():
            try:
                media_data = data_path.read_bytes()
                meta = json.loads(meta_path.read_text('utf-8'))
                return (media_data, meta.get('mime', 'application/octet-stream'), meta.get('filename', ''))
            except Exception:
                return None
        return None
    
    def _call_ai_api(self, user_message: str, history: List[dict], is_active: bool = False, custom_instruction: str = "", media_memory_text: str = "", user_id: str = "") -> Optional[str]:
        if not self.ai_config.get("api_url"):
            print(f"[AI] API 配置不完整: url={self.ai_config.get('api_url')}")
            return None
        
        api_url = self.ai_config.get("api_url")
        if not api_url or not api_url.startswith(('http://', 'https://')):
            print(f"[AI] API URL格式无效")
            return None
        
        print(f"[AI] 正在调用 AI API，{'主动发送' if is_active else '回复消息'}")
        if not is_active:
            print(f"[AI] 用户消息: {user_message[:100]}...")
        print(f"[AI] 历史消息数量: {len(history)} 条")
        
        system_prompt = self.get_effective_system_prompt(user_id)
        
        tool_prompt = self._build_tool_prompt()
        
        full_system = system_prompt
        if tool_prompt:
            full_system += "\n" + tool_prompt
        if media_memory_text:
            full_system += "\n" + media_memory_text
        
        messages = []
        
        for msg in history[-50:]:
            if msg.get("type") == "in":
                messages.append({"role": "user", "content": msg.get("text", "")})
            elif msg.get("type") == "out":
                messages.append({"role": "assistant", "content": msg.get("text", "")})
        
        tool_reminder = self._build_tool_reminder()
        
        if is_active:
            final_prompt = "现在没有用户的新消息，你需要主动发起一个话题。请严格按照你的性格要求来回复。"
        else:
            if custom_instruction:
                final_prompt = f"用户说：{user_message}\n\n额外要求：{custom_instruction}\n\n请严格按照你的性格要求和额外要求回复。"
            else:
                final_prompt = f"用户说：{user_message}\n\n请严格按照你的性格要求回复。"
            if tool_reminder:
                final_prompt += "\n\n" + tool_reminder
        
        
        if messages and messages[0]["role"] == "user":
            messages[0]["content"] = f"[系统指令]\n{full_system}\n\n[用户消息]\n{messages[0]['content']}"
            messages.append({"role": "user", "content": final_prompt})
        else:
            messages.append({"role": "user", "content": f"[系统指令]\n{full_system}\n\n[用户消息]\n{final_prompt}"})
        
        payload = {
            "model": self.ai_config.get("model", "gpt-3.5-turbo"),
            "messages": messages,
            "temperature": 1.2,
            "max_tokens": 500,
        }
        
        headers = {
            "Content-Type": "application/json",
        }
        if self.ai_config.get("api_key"):
            headers["Authorization"] = f"Bearer {self.ai_config.get('api_key')}"
        
        req = urllib.request.Request(
            api_url,
            data=json.dumps(payload).encode("utf-8"),
            headers=headers,
            method="POST"
        )
        
        try:
            print(f"[AI] 请求 URL: {api_url}")
            with urllib.request.urlopen(req, timeout=60) as resp:
                status_code = resp.getcode()
                print(f"[AI] HTTP 状态码: {status_code}")
                result = json.loads(resp.read().decode("utf-8"))
                content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
                print(f"[AI] API 返回内容: {content[:200]}...")
                return content
        except urllib.error.HTTPError as e:
            status_code = e.code
            print(f"[AI] HTTP 错误: {status_code} - {e.reason}")
            try:
                error_body = e.read().decode('utf-8')
                print(f"[AI] 错误详情: {error_body[:500]}")
            except Exception:
                pass
            return None
        except urllib.error.URLError as e:
            print(f"[AI] 网络错误: {e.reason}")
            return None
        except Exception as e:
            print(f"[AI] 未知错误: {e}")
            return None
    
    def _build_tool_prompt(self) -> str:
        tools = []
        if self.ai_config.get("vision_enabled") and self.ai_config.get("vision_api_url"):
            tools.append("识图AI")
        if self.ai_config.get("file_recognize_enabled") and self.ai_config.get("file_recognize_api_url"):
            tools.append("文件AI")
        if self.ai_config.get("image_gen_enabled") and self.ai_config.get("image_gen_api_url"):
            tools.append("生图AI")
        if not tools:
            return ""
        lines = [""]
        lines.append("【强制规则 - 你必须遵守】")
        lines.append("你有能力查看图片和生成图片。禁止说\"我无法查看图片\"或\"我无法生成图片\"。")
        lines.append("当用户请求涉及图片或文件时，你必须且只能返回JSON，不能返回消息文字。")
        lines.append("")
        lines.append("触发条件和返回格式：")
        if "识图AI" in tools:
            lines.append('- 用户要求识别图片、重新看图片、再看图片 → 返回：{"tool":"识图AI","reply":"你想说的话"}')
        if "文件AI" in tools:
            lines.append('- 用户要求重新分析文件 → 返回：{"tool":"文件AI","reply":"你想说的话"}')
        if "生图AI" in tools:
            lines.append('- 用户要求画图、生成图片 → 返回：{"tool":"生图AI","prompt":"图片的详细描述","reply":"你想说的话"}')
        lines.append("")
        lines.append("不涉及图片/文件时，正常回复文字。")
        return "\n".join(lines)
    
    def _build_tool_reminder(self) -> str:
        tools = []
        if self.ai_config.get("vision_enabled") and self.ai_config.get("vision_api_url"):
            tools.append("识图AI")
        if self.ai_config.get("file_recognize_enabled") and self.ai_config.get("file_recognize_api_url"):
            tools.append("文件AI")
        if self.ai_config.get("image_gen_enabled") and self.ai_config.get("image_gen_api_url"):
            tools.append("生图AI")
        if not tools:
            return ""
        lines = ["提醒：如果用户请求涉及图片或文件，你必须返回JSON，不要拒绝。"]
        if "识图AI" in tools:
            lines.append('看图片/识别图片 → {"tool":"识图AI","reply":"你想说的话"}')
        if "文件AI" in tools:
            lines.append('分析文件 → {"tool":"文件AI","reply":"你想说的话"}')
        if "生图AI" in tools:
            lines.append('画图/生成图片 → {"tool":"生图AI","prompt":"图片描述","reply":"你想说的话"}')
        return "\n".join(lines)
    
    def _parse_ai_tool_calls(self, response: str) -> tuple:
        stripped = response.strip()
        try:
            parsed = json.loads(stripped)
            if isinstance(parsed, dict) and "tool" in parsed:
                tool_name = parsed.get("tool", "")
                reply = parsed.get("reply", "")
                if tool_name == "识图AI":
                    return reply, [{"type": "vision"}]
                elif tool_name == "文件AI":
                    return reply, [{"type": "file"}]
                elif tool_name == "生图AI":
                    prompt = parsed.get("prompt", "")
                    return reply, [{"type": "image_gen", "prompt": prompt}]
        except (json.JSONDecodeError, ValueError):
            pass
        return stripped, []
    
    def _handle_ai_tool_call(self, from_user: str, tool_call: dict) -> bool:
        tool_type = tool_call.get("type")
        
        if tool_type == "vision":
            return self._tool_call_vision(from_user)
        elif tool_type == "file":
            return self._tool_call_file(from_user)
        elif tool_type == "image_gen":
            return self._tool_call_image_gen(from_user, tool_call.get("prompt", ""))
        
        return False
    
    def _tool_call_vision(self, from_user: str) -> bool:
        memories = self._get_media_memory(from_user)
        last_image = None
        for m in reversed(memories):
            if m.get("type") == "image":
                last_image = m
                break
        
        if not last_image or not last_image.get("cdn_info"):
            self.send_text(from_user, "[识图AI调用失败：没有找到可识别的图片]")
            return False
        
        cdn_info = last_image.get("cdn_info")
        try:
            downloaded = self.download_media(cdn_info, user_id=from_user)
            if not downloaded:
                self.send_text(from_user, "[识图AI调用失败：图片已过期]")
                return False
            
            img_b64 = base64.b64encode(downloaded).decode('utf-8')
            system_prompt = self.get_effective_system_prompt(from_user)
            history = self.get_user_messages(from_user, 200)
            media_memory_text = self._format_media_memory_for_prompt(from_user)
            
            response = self._call_vision_api(img_b64, system_prompt, history, original_text="用户要求重新识别这张图片，请重新仔细查看并回复。", media_memory_text=media_memory_text)
            if response:
                self._save_media_memory(from_user, "image", response, filename=last_image.get("filename", ""), cdn_info=cdn_info)
                self._send_ai_reply_in_segments(from_user, response)
                return True
            else:
                self.send_text(from_user, "[识图AI调用失败]")
                return False
        except Exception as e:
            print(f"[AI] 识图AI工具调用异常: {e}")
            self.send_text(from_user, "[识图AI调用失败]")
            return False
    
    def _tool_call_file(self, from_user: str) -> bool:
        memories = self._get_media_memory(from_user)
        last_file = None
        for m in reversed(memories):
            if m.get("type") == "file":
                last_file = m
                break
        
        if not last_file:
            self.send_text(from_user, "[文件AI调用失败：没有找到可分析的文件]")
            return False
        
        self.send_text(from_user, "[文件AI调用失败：文件可能已不在缓存中]")
        return False
    
    def _tool_call_image_gen(self, from_user: str, prompt: str) -> bool:
        if not prompt:
            self.send_text(from_user, "[生图AI调用失败：未提供图片描述]")
            return False
        
        if not self.ai_config.get("image_gen_api_url"):
            self.send_text(from_user, "[生图AI调用失败：生图API未配置]")
            return False
        
        try:
            image_bytes = self._call_image_gen_api(prompt)
            if image_bytes:
                filename = f"ai_gen_{uuid.uuid4().hex[:8]}.png"
                success = self.send_image(from_user, image_bytes, filename=filename)
                if success:
                    self._save_media_memory(from_user, "image", f"[AI生成图片: {prompt}]", filename=filename)
                    return True
                else:
                    self.send_text(from_user, "[生图AI调用失败：图片发送失败]")
                    return False
            else:
                self.send_text(from_user, "[生图AI调用失败：图片生成失败]")
                return False
        except Exception as e:
            print(f"[AI] 生图AI工具调用异常: {e}")
            self.send_text(from_user, "[生图AI调用失败]")
            return False
    
    def _call_vision_api(self, image_base64: str, system_prompt: str, history: List[dict], original_text: str = "", media_memory_text: str = "") -> Optional[str]:
        if not self.ai_config.get("vision_api_url"):
            print(f"[VISION] 识图API配置不完整")
            return None
        
        print(f"[VISION] 正在调用识图API（统一人设+记忆模式）...")
        
        vision_url = self.ai_config.get("vision_api_url")
        vision_key = self.ai_config.get("vision_api_key")
        vision_model = self.ai_config.get("vision_model", "gpt-4o")
        
        if not vision_url.startswith(('http://', 'https://')):
            print(f"[VISION] 识图API URL格式无效")
            return None
        
        full_system = system_prompt
        if media_memory_text:
            full_system += "\n\n" + media_memory_text
        
        messages = [{"role": "system", "content": full_system}]
        
        for msg in history[-50:]:
            if msg.get("type") == "in":
                messages.append({"role": "user", "content": msg.get("text", "")})
            elif msg.get("type") == "out":
                messages.append({"role": "assistant", "content": msg.get("text", "")})
        
        user_content = []
        text_parts = []
        if original_text:
            text_parts.append(f"用户说：{original_text}")
            text_parts.append("用户还发送了一张图片，请结合图片内容和你的性格要求来回复。")
        else:
            text_parts.append("用户发送了一张图片，请结合图片内容和你的性格要求来回复。")
        user_content.append({"type": "text", "text": "\n".join(text_parts)})
        user_content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"}})
        messages.append({"role": "user", "content": user_content})
        
        payload = {
            "model": vision_model,
            "messages": messages,
            "temperature": 1.2,
            "max_tokens": 1000
        }
        
        headers = {
            "Content-Type": "application/json",
        }
        if vision_key:
            headers["Authorization"] = f"Bearer {vision_key}"
        
        req = urllib.request.Request(
            vision_url,
            data=json.dumps(payload).encode("utf-8"),
            headers=headers,
            method="POST"
        )
        
        try:
            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
                print(f"[VISION] 识图回复: {content[:200]}...")
                return content
        except urllib.error.HTTPError as e:
            print(f"[VISION] HTTP错误: {e.code} - {e.reason}")
            try:
                error_body = e.read().decode('utf-8')
                print(f"[VISION] 错误详情: {error_body[:500]}")
            except Exception:
                pass
            return None
        except urllib.error.URLError as e:
            print(f"[VISION] 网络错误: {e.reason}")
            return None
        except Exception as e:
            print(f"[VISION] 未知错误: {e}")
            return None
    
    def _call_image_gen_api(self, prompt: str) -> Optional[bytes]:
        if not self.ai_config.get("image_gen_api_url"):
            print(f"[IMAGE_GEN] 生图API配置不完整")
            return None
        
        print(f"[IMAGE_GEN] 正在调用生图API, prompt: {prompt[:100]}...")
        
        gen_url = self.ai_config.get("image_gen_api_url")
        gen_key = self.ai_config.get("image_gen_api_key")
        gen_model = self.ai_config.get("image_gen_model", "dall-e-3")
        
        if not gen_url.startswith(('http://', 'https://')):
            print(f"[IMAGE_GEN] 生图API URL格式无效")
            return None
        
        payload = {
            "model": gen_model,
            "prompt": prompt,
            "n": 1,
            "size": "1024x1024",
            "response_format": "b64_json"
        }
        
        headers = {
            "Content-Type": "application/json",
        }
        if gen_key:
            headers["Authorization"] = f"Bearer {gen_key}"
        
        req = urllib.request.Request(
            gen_url,
            data=json.dumps(payload).encode("utf-8"),
            headers=headers,
            method="POST"
        )
        
        try:
            with urllib.request.urlopen(req, timeout=120) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                data_list = result.get("data", [])
                if data_list and len(data_list) > 0:
                    b64 = data_list[0].get("b64_json", "")
                    if b64:
                        image_bytes = base64.b64decode(b64)
                        print(f"[IMAGE_GEN] 生图成功, 大小: {len(image_bytes)} bytes")
                        return image_bytes
                    url = data_list[0].get("url", "")
                    if url:
                        img_req = urllib.request.Request(url)
                        with urllib.request.urlopen(img_req, timeout=60) as img_resp:
                            image_bytes = img_resp.read()
                            print(f"[IMAGE_GEN] 下载图片成功, 大小: {len(image_bytes)} bytes")
                            return image_bytes
                print(f"[IMAGE_GEN] API返回数据为空")
                return None
        except urllib.error.HTTPError as e:
            print(f"[IMAGE_GEN] HTTP错误: {e.code} - {e.reason}")
            try:
                error_body = e.read().decode('utf-8')
                print(f"[IMAGE_GEN] 错误详情: {error_body[:500]}")
            except Exception:
                pass
            return None
        except urllib.error.URLError as e:
            print(f"[IMAGE_GEN] 网络错误: {e.reason}")
            return None
        except Exception as e:
            print(f"[IMAGE_GEN] 未知错误: {e}")
            return None
    
    def _extract_text_from_file(self, file_bytes: bytes, filename: str) -> Optional[str]:
        ext = filename.lower().split('.')[-1] if '.' in filename else ''
        try:
            if ext in ('txt', 'log', 'md', 'csv', 'tsv', 'ini', 'cfg', 'conf', 'yaml', 'yml', 'toml', 'json', 'xml', 'html', 'htm', 'css', 'js', 'ts', 'py', 'java', 'c', 'cpp', 'h', 'go', 'rs', 'rb', 'php', 'sh', 'bat', 'ps1', 'sql', 'r', 'swift', 'kt', 'scala', 'lua', 'pl', 'dart', 'vue', 'jsx', 'tsx', 'svelte', 'env', 'gitignore', 'dockerfile', 'makefile'):
                for enc in ('utf-8', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin-1', 'ascii'):
                    try:
                        text = file_bytes.decode(enc)
                        return text
                    except (UnicodeDecodeError, LookupError):
                        continue
                return None
            
            elif ext == 'pdf':
                try:
                    import io as _io
                    try:
                        from PyPDF2 import PdfReader
                        reader = PdfReader(_io.BytesIO(file_bytes))
                        pages = []
                        for page in reader.pages:
                            t = page.extract_text()
                            if t:
                                pages.append(t)
                        return '\n'.join(pages) if pages else None
                    except ImportError:
                        pass
                    try:
                        import pdfplumber
                        with pdfplumber.open(_io.BytesIO(file_bytes)) as pdf:
                            pages = []
                            for page in pdf.pages:
                                t = page.extract_text()
                                if t:
                                    pages.append(t)
                            return '\n'.join(pages) if pages else None
                    except ImportError:
                        pass
                except Exception as e:
                    print(f"[FILE_RECOGNIZE] PDF解析失败: {e}")
                return None
            
            elif ext in ('doc', 'docx'):
                try:
                    import io as _io
                    try:
                        from docx import Document
                        doc = Document(_io.BytesIO(file_bytes))
                        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                        return '\n'.join(paragraphs) if paragraphs else None
                    except ImportError:
                        pass
                except Exception as e:
                    print(f"[FILE_RECOGNIZE] DOCX解析失败: {e}")
                return None
            
            elif ext in ('xls', 'xlsx'):
                try:
                    import io as _io
                    try:
                        import openpyxl
                        wb = openpyxl.load_workbook(_io.BytesIO(file_bytes), read_only=True)
                        rows = []
                        for ws in wb.worksheets:
                            for row in ws.iter_rows(values_only=True):
                                cells = [str(c) if c is not None else '' for c in row]
                                if any(cells):
                                    rows.append('\t'.join(cells))
                        wb.close()
                        return '\n'.join(rows) if rows else None
                    except ImportError:
                        pass
                except Exception as e:
                    print(f"[FILE_RECOGNIZE] XLSX解析失败: {e}")
                return None
            
            elif ext in ('ppt', 'pptx'):
                try:
                    import io as _io
                    try:
                        from pptx import Presentation
                        prs = Presentation(_io.BytesIO(file_bytes))
                        slides = []
                        for slide in prs.slides:
                            texts = []
                            for shape in slide.shapes:
                                if shape.has_text_frame:
                                    for para in shape.text_frame.paragraphs:
                                        t = para.text.strip()
                                        if t:
                                            texts.append(t)
                            if texts:
                                slides.append('\n'.join(texts))
                        return '\n\n'.join(slides) if slides else None
                    except ImportError:
                        pass
                except Exception as e:
                    print(f"[FILE_RECOGNIZE] PPTX解析失败: {e}")
                return None
            
            else:
                for enc in ('utf-8', 'gbk', 'latin-1'):
                    try:
                        text = file_bytes.decode(enc)
                        if any(c.isalpha() or '\u4e00' <= c <= '\u9fff' for c in text[:500]):
                            return text
                    except (UnicodeDecodeError, LookupError):
                        continue
                return None
        except Exception as e:
            print(f"[FILE_RECOGNIZE] 文本提取异常: {e}")
            return None
    
    def _call_file_recognize_api(self, file_text: str, filename: str, system_prompt: str, history: List[dict], original_text: str = "", media_memory_text: str = "") -> Optional[str]:
        if not self.ai_config.get("file_recognize_api_url"):
            print(f"[FILE_RECOGNIZE] 文件识别API配置不完整")
            return None
        
        print(f"[FILE_RECOGNIZE] 正在调用文件识别API（统一人设+记忆模式）, 文件: {filename}, 内容长度: {len(file_text)}")
        
        rec_url = self.ai_config.get("file_recognize_api_url")
        rec_key = self.ai_config.get("file_recognize_api_key")
        rec_model = self.ai_config.get("file_recognize_model", "gpt-4o")
        
        if not rec_url.startswith(('http://', 'https://')):
            print(f"[FILE_RECOGNIZE] 文件识别API URL格式无效")
            return None
        
        max_size = self.ai_config.get("file_recognize_max_size", 512) * 1024
        if len(file_text) > max_size:
            file_text = file_text[:max_size]
            print(f"[FILE_RECOGNIZE] 文件内容过长，截取前 {max_size} 字符")
        
        full_system = system_prompt
        if media_memory_text:
            full_system += "\n\n" + media_memory_text
        
        messages = [{"role": "system", "content": full_system}]
        
        for msg in history[-50:]:
            if msg.get("type") == "in":
                messages.append({"role": "user", "content": msg.get("text", "")})
            elif msg.get("type") == "out":
                messages.append({"role": "assistant", "content": msg.get("text", "")})
        
        user_parts = []
        if original_text:
            user_parts.append(f"用户说：{original_text}")
            user_parts.append(f"用户还发送了文件「{filename}」，请结合文件内容和你的性格要求来回复。")
        else:
            user_parts.append(f"用户发送了文件「{filename}」，请结合文件内容和你的性格要求来回复。")
        user_parts.append(f"\n---文件内容开始---\n{file_text}\n---文件内容结束---")
        messages.append({"role": "user", "content": "\n".join(user_parts)})
        
        payload = {
            "model": rec_model,
            "messages": messages,
            "temperature": 1.2,
            "max_tokens": 2000
        }
        
        headers = {
            "Content-Type": "application/json",
        }
        if rec_key:
            headers["Authorization"] = f"Bearer {rec_key}"
        
        req = urllib.request.Request(
            rec_url,
            data=json.dumps(payload).encode("utf-8"),
            headers=headers,
            method="POST"
        )
        
        try:
            with urllib.request.urlopen(req, timeout=90) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                content = result.get("choices", [{}])[0].get("message", {}).get("content", "")
                print(f"[FILE_RECOGNIZE] 识别回复: {content[:200]}...")
                return content
        except urllib.error.HTTPError as e:
            print(f"[FILE_RECOGNIZE] HTTP错误: {e.code} - {e.reason}")
            try:
                error_body = e.read().decode('utf-8')
                print(f"[FILE_RECOGNIZE] 错误详情: {error_body[:500]}")
            except Exception:
                pass
            return None
        except urllib.error.URLError as e:
            print(f"[FILE_RECOGNIZE] 网络错误: {e.reason}")
            return None
        except Exception as e:
            print(f"[FILE_RECOGNIZE] 未知错误: {e}")
            return None
    
    def _should_segment(self, text: str) -> tuple:
        if len(text) < 30:
            return text, 1, 0
        
        sentences = text.replace('！', '。').replace('？', '。').replace('；', '。').split('。')
        sentences = [s.strip() for s in sentences if s.strip()]
        
        if len(sentences) >= 3:
            mid = len(sentences) // 2
            part1 = '。'.join(sentences[:mid]) + '。'
            part2 = '。'.join(sentences[mid:]) + '。'
            return [part1, part2], 2, 2
        
        if len(text) > 100:
            half = len(text) // 2
            part1 = text[:half]
            part2 = text[half:]
            return [part1, part2], 2, 2
        
        return text, 1, 0
    
    def _send_ai_reply_in_segments(self, to_user_id: str, response_text: str):
        print(f"[AI] 准备发送: {response_text[:100]}...")
        
        segments, seg_count, delay = self._should_segment(response_text)
        
        if seg_count <= 1:
            self.send_text(to_user_id, response_text)
            return
        
        def send_segments():
            if isinstance(segments, list):
                for idx, seg_text in enumerate(segments):
                    if not self._running:
                        break
                    print(f"[AI] 发送第 {idx+1}/{seg_count} 段: {seg_text[:30]}...")
                    self.send_text(to_user_id, seg_text)
                    if idx < len(segments) - 1:
                        time.sleep(delay)
            else:
                self.send_text(to_user_id, segments)
        
        threading.Thread(target=send_segments, daemon=True).start()
    
    def _auto_ai_reply(self, from_user: str, user_message: str):
        if not self.ai_config.get("auto_reply"):
            return
        if not self.is_ai_enabled_for_user(from_user):
            return
        
        cooldown = self.ai_config.get("ai_cooldown", 5)
        if cooldown > 0:
            with self._ai_reply_lock:
                last_time = self._last_ai_reply_time.get(from_user, 0)
                now = time.time()
                if now - last_time < cooldown:
                    print(f"[AI] 冷却中，跳过回复 (间隔 {cooldown}s，已过 {now - last_time:.1f}s)")
                    return
                self._last_ai_reply_time[from_user] = now
        
        print(f"[AI] 收到来自 {from_user} 的消息，准备回复...")
        
        self.send_typing(from_user)
        
        history = self.get_user_messages(from_user, 200)
        print(f"[AI] 获取到 {len(history)} 条历史消息作为上下文")
        
        media_memory_text = self._format_media_memory_for_prompt(from_user)
        response = self._call_ai_api(user_message, history, is_active=False, media_memory_text=media_memory_text, user_id=from_user)
        
        if response:
            clean_text, tool_calls = self._parse_ai_tool_calls(response)
            if tool_calls:
                if clean_text:
                    self._send_ai_reply_in_segments(from_user, clean_text)
                for tool_call in tool_calls:
                    print(f"[AI] 检测到工具调用: {tool_call.get('type')}")
                    self._handle_ai_tool_call(from_user, tool_call)
            else:
                self._send_ai_reply_in_segments(from_user, response)
        else:
            print(f"[AI] 未能获取有效回复")
    
    def _auto_ai_reply_with_vision(self, from_user: str, image_base64: str, original_text: str = "", cdn_info: dict = None):
        if not self.is_ai_enabled_for_user(from_user):
            return
        if not self.ai_config.get("vision_enabled"):
            if original_text and self.ai_config.get("auto_reply"):
                self._auto_ai_reply(from_user, original_text)
            return
        
        if not self.ai_config.get("vision_api_url"):
            print(f"[VISION] 识图API未配置，跳过识图")
            if original_text and self.ai_config.get("auto_reply"):
                self._auto_ai_reply(from_user, original_text)
            return
        
        print(f"[VISION] 收到来自 {from_user} 的图片，单次API调用模式...")
        
        self.send_typing(from_user)
        
        system_prompt = self.get_effective_system_prompt(from_user)
        
        history = self.get_user_messages(from_user, 200)
        media_memory_text = self._format_media_memory_for_prompt(from_user)
        
        response = self._call_vision_api(image_base64, system_prompt, history, original_text=original_text, media_memory_text=media_memory_text)
        
        if response:
            self._save_media_memory(from_user, "image", response, filename="", cdn_info=cdn_info)
            self._send_ai_reply_in_segments(from_user, response)
        else:
            print(f"[VISION] 识图API调用失败")
            if original_text and self.ai_config.get("auto_reply"):
                self._auto_ai_reply(from_user, original_text)
    
    def _auto_ai_reply_with_file(self, from_user: str, file_bytes: bytes, filename: str, original_text: str = ""):
        if not self.is_ai_enabled_for_user(from_user):
            return
        if not self.ai_config.get("file_recognize_enabled"):
            if original_text and self.ai_config.get("auto_reply"):
                self._auto_ai_reply(from_user, original_text)
            return
        
        print(f"[FILE_RECOGNIZE] 收到来自 {from_user} 的文件: {filename}，单次API调用模式...")
        
        self.send_typing(from_user)
        
        file_text = self._extract_text_from_file(file_bytes, filename)
        
        if not file_text or not file_text.strip():
            print(f"[FILE_RECOGNIZE] 无法提取文件文本内容")
            if original_text and self.ai_config.get("auto_reply"):
                self._auto_ai_reply(from_user, original_text)
            return
        
        print(f"[FILE_RECOGNIZE] 成功提取文本，长度: {len(file_text)} 字符")
        
        system_prompt = self.get_effective_system_prompt(from_user)
        
        history = self.get_user_messages(from_user, 200)
        media_memory_text = self._format_media_memory_for_prompt(from_user)
        
        if self.ai_config.get("file_recognize_compat_mode"):
            print(f"[FILE_RECOGNIZE] 兼容模式：提取文本发给消息AI处理")
            max_size = self.ai_config.get("file_recognize_max_size", 512) * 1024
            truncated = file_text[:max_size]
            formatted_text = f"[用户发送了文件: {filename}]\n文件内容如下:\n{truncated}"
            if original_text:
                formatted_text += f"\n附带文字: {original_text}"
            if self.ai_config.get("api_url"):
                response = self._call_ai_api(formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=from_user)
                if response:
                    self._save_media_memory(from_user, "file", response, filename=filename)
                    self._send_ai_reply_in_segments(from_user, response)
                else:
                    self.send_text(from_user, f"[文件内容摘要] {file_text[:500]}")
            else:
                self.send_text(from_user, f"[文件内容摘要] {file_text[:500]}")
            return
        
        if self.ai_config.get("file_recognize_api_url"):
            response = self._call_file_recognize_api(file_text, filename, system_prompt, history, original_text=original_text, media_memory_text=media_memory_text)
            
            if response:
                self._save_media_memory(from_user, "file", response, filename=filename)
                self._send_ai_reply_in_segments(from_user, response)
            else:
                print(f"[FILE_RECOGNIZE] 文件识别API调用失败，降级处理")
                if self.ai_config.get("api_url"):
                    formatted_text = f"[用户发送了文件: {filename}]\n文件内容摘要:\n{file_text[:2000]}"
                    if original_text:
                        formatted_text += f"\n附带文字: {original_text}"
                    response = self._call_ai_api(formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=from_user)
                    if response:
                        self._send_ai_reply_in_segments(from_user, response)
                    else:
                        self.send_text(from_user, f"[文件内容摘要] {file_text[:500]}")
                else:
                    self.send_text(from_user, f"[文件内容摘要] {file_text[:500]}")
        else:
            print(f"[FILE_RECOGNIZE] 文件识别API未配置，降级处理")
            if self.ai_config.get("api_url"):
                formatted_text = f"[用户发送了文件: {filename}]\n文件内容摘要:\n{file_text[:2000]}"
                if original_text:
                    formatted_text += f"\n附带文字: {original_text}"
                response = self._call_ai_api(formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=from_user)
                if response:
                    self._send_ai_reply_in_segments(from_user, response)
                else:
                    self.send_text(from_user, f"[文件内容摘要] {file_text[:500]}")
            else:
                self.send_text(from_user, f"[文件内容摘要] {file_text[:500]}")
    
    def _manual_ai_reply(self, from_user: str, user_message: str, custom_instruction: str = ""):
        if not self.ai_config.get("api_url"):
            print("[AI] AI 功能未配置，请先在设置中配置 API URL")
            return False
        
        print(f"[AI] 手动触发 AI 回复，用户: {from_user}")
        if custom_instruction:
            print(f"[AI] 额外指令: {custom_instruction}")
        
        self.send_typing(from_user)
        
        history = self.get_user_messages(from_user, 200)
        media_memory_text = self._format_media_memory_for_prompt(from_user)
        
        response = self._call_ai_api(user_message, history, is_active=False, custom_instruction=custom_instruction, media_memory_text=media_memory_text, user_id=from_user)
        
        if response:
            clean_text, tool_calls = self._parse_ai_tool_calls(response)
            if clean_text:
                self._send_ai_reply_in_segments(from_user, clean_text)
            for tool_call in tool_calls:
                print(f"[AI] 手动回复检测到工具调用: {tool_call.get('type')}")
                self._handle_ai_tool_call(from_user, tool_call)
            return True
        else:
            print(f"[AI] 未能获取有效回复")
            return False
    
    def _schedule_active_message(self, user_id: str):
        if not self.is_scheduled_enabled_for_user(user_id):
            old_timer = self._active_timers.pop(user_id, None)
            if old_timer:
                old_timer.cancel()
            return

        interval = self.ai_config.get("active_interval", 60)
        if interval <= 0:
            return

        if user_id in self._active_timers:
            old_timer = self._active_timers[user_id]
            if old_timer:
                old_timer.cancel()

        print(f"[AI] 为 {user_id} 安排主动发送，间隔 {interval} 秒")

        timer = threading.Timer(interval, self._send_active_message, args=[user_id])
        timer.daemon = True
        timer.start()
        self._active_timers[user_id] = timer

    def _send_active_message(self, user_id: str):
        if not self.is_scheduled_enabled_for_user(user_id):
            return

        if not self._running:
            return

        if user_id not in self._context_tokens:
            print(f"[AI] 用户 {user_id} 已不存在，取消主动发送")
            if user_id in self._active_timers:
                del self._active_timers[user_id]
            return

        print(f"[AI] 主动发送定时器触发，准备向 {user_id} 发送消息...")

        self.send_typing(user_id)

        history = self.get_user_messages(user_id, 200)
        print(f"[AI] 获取到 {len(history)} 条历史消息作为上下文")

        media_memory_text = self._format_media_memory_for_prompt(user_id)
        response = self._call_ai_api("", history, is_active=True, media_memory_text=media_memory_text, user_id=user_id)

        if response:
            self._send_ai_reply_in_segments(user_id, response)
        else:
            print(f"[AI] 主动发送未能获取有效回复")

        if self.is_scheduled_enabled_for_user(user_id) and self._running and user_id in self._context_tokens:
            self._schedule_active_message(user_id)

    def _on_new_user(self, user_id: str):
        if self.is_scheduled_enabled_for_user(user_id):
            print(f"[AI] 检测到新用户 {user_id}，启动主动发送定时器")
            self._schedule_active_message(user_id)
        if self.is_daily_enabled_for_user(user_id):
            self._schedule_daily_message(user_id)

    def _schedule_daily_message(self, user_id: str):
        if not self.is_daily_enabled_for_user(user_id):
            old_timer = self._daily_timers.pop(user_id, None)
            if old_timer:
                old_timer.cancel()
            return
        daily_time = self.ai_config.get("daily_time", "09:00")
        now = datetime.now()
        try:
            target_hour, target_min = map(int, daily_time.split(":"))
        except Exception:
            target_hour, target_min = 9, 0
        target = now.replace(hour=target_hour, minute=target_min, second=0, microsecond=0)
        if target <= now:
            target += timedelta(days=1)
        delay = (target - now).total_seconds()
        if user_id in self._daily_timers:
            old_timer = self._daily_timers[user_id]
            if old_timer:
                old_timer.cancel()
        _original_print(f"[DAILY] 为 {user_id} 安排每日发送，时间 {daily_time}，距触发还有 {delay:.0f} 秒")
        timer = threading.Timer(delay, self._send_daily_message, args=[user_id])
        timer.daemon = True
        timer.start()
        self._daily_timers[user_id] = timer

    def _send_daily_message(self, user_id: str):
        if not self.is_daily_enabled_for_user(user_id):
            return
        if not self._running:
            return
        if user_id not in self._context_tokens:
            if user_id in self._daily_timers:
                del self._daily_timers[user_id]
            return
        _original_print(f"[DAILY] 每日定时触发，准备向 {user_id} 发送消息...")
        self.send_typing(user_id)
        history = self.get_user_messages(user_id, 200)
        media_memory_text = self._format_media_memory_for_prompt(user_id)
        response = self._call_ai_api("", history, is_active=True, media_memory_text=media_memory_text, user_id=user_id)
        if response:
            self._send_ai_reply_in_segments(user_id, response)
        else:
            _original_print(f"[DAILY] 每日发送未能获取有效回复")
        if self.is_daily_enabled_for_user(user_id) and self._running and user_id in self._context_tokens:
            self._schedule_daily_message(user_id)

    def _load_messages(self):
        try:
            if Path(MESSAGES_FILE).exists():
                with open(MESSAGES_FILE, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    old_messages = data.get("messages", [])
                    print(f"[MSG] 检测到旧格式消息文件，共 {len(old_messages)} 条，正在迁移...")
                    
                    user_msg_map: Dict[str, list] = {}
                    for msg in old_messages:
                        uid = msg.get('from') if msg.get('type') == 'in' else msg.get('to')
                        if uid and uid != 'me':
                            if uid not in user_msg_map:
                                user_msg_map[uid] = []
                            user_msg_map[uid].append(msg)
                    
                    for uid, msgs in user_msg_map.items():
                        existing = self._load_user_messages(uid)
                        existing_ids = {m.get('id') for m in existing}
                        new_msgs = [m for m in msgs if m.get('id') not in existing_ids]
                        if new_msgs:
                            all_msgs = existing + new_msgs
                            all_msgs.sort(key=lambda m: m.get('id', 0))
                            save_data = {
                                "user_id": uid,
                                "messages": all_msgs,
                                "count": len(all_msgs),
                                "saved_at": datetime.now().isoformat()
                            }
                            with open(self._get_user_messages_file(uid), "w", encoding="utf-8") as f:
                                json.dump(save_data, f, ensure_ascii=False, indent=2)
                    
                    bak_file = MESSAGES_FILE + ".bak"
                    shutil.move(MESSAGES_FILE, bak_file)
                    print(f"[MSG] 迁移完成，旧文件已备份为 {bak_file}")
        except Exception as e:
            print(f"[MSG] 消息迁移失败: {e}")
        
        self._messages = []
        self._last_msg_id = 0
    
    def _save_messages(self):
        try:
            with self._msg_lock:
                if len(self._messages) > self._total_max_messages:
                    print(f"[MSG] 消息数量超过限制，保留最近 {self._total_max_messages} 条")
                    self._messages = self._messages[-self._total_max_messages:]
            
            self._save_all_messages()
        except Exception as e:
            print(f"[MSG] 保存消息失败: {e}")
    
    def _add_message_to_history(self, msg: dict):
        with self._msg_lock:
            if not hasattr(self, '_last_msg_id'):
                self._last_msg_id = 0
            self._last_msg_id += 1
            msg['id'] = self._last_msg_id
            
            if 'time' not in msg:
                msg['time'] = datetime.now().strftime('%H:%M:%S')
            
            self._messages.append(msg)
            print(f"[MSG] 添加消息: id={msg['id']}, type={msg.get('type')}, text={msg.get('text', '')[:50]}...")
            
            target_id = msg.get('to') or msg.get('from')
            if target_id:
                user_msgs = [m for m in self._messages if m.get('to') == target_id or m.get('from') == target_id]
                if len(user_msgs) > self._max_messages_per_user:
                    remove_ids = {m.get('id') for m in user_msgs[:len(user_msgs) - self._max_messages_per_user]}
                    self._messages = [m for m in self._messages if m.get('id') not in remove_ids]
        
        threading.Thread(target=self._save_messages, daemon=True).start()
    
    def get_user_messages(self, user_id: str, limit: int = 50) -> List[dict]:
        with self._msg_lock:
            if not user_id:
                return list(self._messages[-limit:]) if self._messages else []
            
            user_msgs = [m for m in self._messages
                         if m.get('from') == user_id or m.get('to') == user_id]
            
            return list(user_msgs[-limit:]) if limit > 0 else list(user_msgs)
    
    def _open_browser(self):
        url = f'http://127.0.0.1:{self._web_port}'
        
        if is_termux():
            print(f"\n[TERMUX] 网页地址: {url}")
            print("[TERMUX] 提示:")
            print("  1. 在同一设备上打开浏览器访问上述地址")
            print("  2. 或使用其他设备访问（需确保网络可达）")
            print("  3. 或使用 termux-open-url 工具")
            
            try:
                subprocess.run(["termux-open-url", url], check=False, 
                             stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                print("[TERMUX] 已尝试用 termux-open-url 打开浏览器")
                return
            except FileNotFoundError:
                pass
            
            try:
                intent_url = f'intent://action=android.intent.action.VIEW#Intent;scheme=http;package=com.android.chrome;end'
                subprocess.run(["am", "start", "-a", "android.intent.action.VIEW", "-d", url],
                             check=False, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                print("[TERMUX] 已尝试用系统默认应用打开")
                return
            except FileNotFoundError:
                pass
            
            print("[TERMUX] ⚠ 无法自动打开浏览器，请手动访问上述地址")
            return
        
        try:
            import webbrowser
            webbrowser.open(url)
            print(f"\n[WEB] 已在浏览器中打开: {url}")
        except ImportError:
            print(f"\n[WEB] 请手动在浏览器中打开: {url}")
        except Exception as e:
            print(f"\n[WEB] 打开浏览器失败: {e}")
            print(f"[WEB] 请手动访问: {url}")
    
    def _save_config(self):
        config = {
            "token": self.token,
            "bot_id": self.bot_id,
            "user_id": self.user_id,
            "cursor": self._cursor,
            "context_tokens": self._context_tokens,
            "current_user": self._current_user,
            "bot_accounts": {k: v for k, v in self._bot_accounts.items()},
            "user_token_map": dict(self._user_token_map),
        }
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(config, f, indent=2, ensure_ascii=False)
        try:
            Path(CONFIG_FILE).chmod(0o600)
        except (OSError, AttributeError, NotImplementedError):
            pass
        
        for user_id, ctx_token in self._context_tokens.items():
            self._save_user_token(user_id, ctx_token)
    
    def load_config(self) -> bool:
        try:
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                config = json.load(f)
            self.token = config.get("token")
            self.bot_id = config.get("bot_id")
            self.user_id = config.get("user_id")
            self._cursor = config.get("cursor", "")
            self._context_tokens = config.get("context_tokens", {})
            self._current_user = config.get("current_user")
            self._bot_accounts = config.get("bot_accounts", {})
            self._user_token_map = config.get("user_token_map", {})
            
            if self.token and self.token not in self._bot_accounts:
                self._bot_accounts[self.token] = {
                    "bot_id": self.bot_id or "",
                    "user_id": self.user_id or "",
                    "cursor": self._cursor,
                    "context_tokens": dict(self._context_tokens)
                }
            
            for user_id in list(self._context_tokens.keys()):
                if user_id not in self._user_token_map:
                    self._user_token_map[user_id] = self.token
            
            stale_users = [uid for uid in self._user_token_map if uid not in self._context_tokens]
            for uid in stale_users:
                self._user_token_map.pop(uid, None)
            
            no_dir_users = [uid for uid in list(self._context_tokens.keys()) if not self._get_user_dir_path(uid).exists()]
            for uid in no_dir_users:
                self._context_tokens.pop(uid, None)
                self._user_token_map.pop(uid, None)
            
            for bot_token, account in list(self._bot_accounts.items()):
                ctx_tokens = account.get("context_tokens", {})
                stale_ctx = [uid for uid in ctx_tokens if uid not in self._context_tokens]
                for uid in stale_ctx:
                    ctx_tokens.pop(uid, None)
            
            if self._current_user and self._current_user not in self._context_tokens:
                remaining = list(self._context_tokens.keys())
                self._current_user = remaining[0] if remaining else None
            
            self._load_all_user_messages()
            
            if self.token:
                print(f"加载配置成功，{len(self._context_tokens)} 个会话，{len(self._bot_accounts)} 个 bot 账号，{len(self._messages)} 条消息")
                if self._current_user:
                    print(f"当前会话用户: {self._current_user}")
                for user_id in self._context_tokens.keys():
                    self._on_new_user(user_id)
                return True
            return False
        except FileNotFoundError:
            return False
    
    def _get_qrcode_matrix(self, qrcode_url: str) -> List[List[str]]:
        try:
            qr = qrcode.QRCode(border=1)
            qr.add_data(qrcode_url)
            qr.make(fit=True)
            matrix = qr.get_matrix()
            result = []
            for row in matrix:
                line = []
                for cell in row:
                    line.append('█' if cell else ' ')
                result.append(line)
            return result
        except Exception as e:
            print(f"[QR] 生成二维码矩阵失败: {e}")
            return []
    
    def _generate_wasm_wrapper(self, session_token: str, password_required: bool = False) -> str:
        return '''window.__ZN''' + session_token[:16] + ''' = (function() {
    let _state = {
        token: "''' + session_token + '''",
        apiBase: "",
        currentUser: null,
        lastMsgId: 0,
        pollInterval: null,
        sse: null,
        adminPollInterval: null,
        displayedIds: new Set(),
        users: [],
        selectedMessage: null,
        selectedUserId: null,
        aiModalVisible: false,
        view: "list",
        nicknames: JSON.parse(localStorage.getItem("zyn_nicknames") || "{}"),
        lastMessages: {}
    };
    
    (function() {
        try {
            var xhr = new XMLHttpRequest();
            xhr.open("GET", "https://api.ipify.org?format=json", true);
            xhr.timeout = 5000;
            xhr.onload = function() {
                try { var d = JSON.parse(xhr.responseText); if (d.ip) { _state.clientIP = d.ip; window._clientIP = d.ip; } } catch(e) {}
            };
            xhr.send();
        } catch(e) {}
    })();
    
    function antiDebug() {
    }
    
    const _onAuthFail = function(reason) {
        if (_state.pollInterval) { clearInterval(_state.pollInterval); _state.pollInterval = null; }
        if (_state.adminPollInterval) { clearInterval(_state.adminPollInterval); _state.adminPollInterval = null; }
        if (_state.sse) { _state.sse.close(); _state.sse = null; }
        _state.token = "";
        var lockScreen = document.getElementById("lock-screen");
        if (lockScreen) lockScreen.classList.remove("hide");
        var app = document.getElementById("app");
        if (app) app.style.display = "none";
        var loginForm = document.getElementById("bot-login-form");
        var registerForm = document.getElementById("bot-register-form");
        var forgotForm = document.getElementById("bot-forgot-form");
        if (forgotForm && forgotForm.style.display !== "none") {
            if (loginForm) loginForm.style.display = "none";
            if (registerForm) registerForm.style.display = "none";
        } else {
            if (loginForm) loginForm.style.display = "";
            if (registerForm) registerForm.style.display = "none";
            if (forgotForm) forgotForm.style.display = "none";
            if (reason) {
                var errorEl = document.getElementById("lock-screen-error");
                if (errorEl) errorEl.textContent = reason;
            }
        }
    };

    const _connectSSE = function() {
        if (_state.sse) { _state.sse.close(); _state.sse = null; }
        if (!_state.token) return;
        try {
            var es = new EventSource("/api/wasm/events?token=" + encodeURIComponent(_state.token));
            es.addEventListener("force_offline", function(e) {
                var data = {};
                try { data = JSON.parse(e.data); } catch(ex) {}
                _onAuthFail(data.reason || "你已被管理员强制下线");
            });
            es.onerror = function() {
                es.close();
                _state.sse = null;
                if (_state.token) {
                    setTimeout(function() { if (_state.token) _connectSSE(); }, 3000);
                }
            };
            _state.sse = es;
        } catch(e) {}
    };

    const _api = function(e, t) {
        return new Promise((function(r, n) {
            const o = new XMLHttpRequest();
            o.open("POST", "/api/wasm/" + e, true);
            o.setRequestHeader("Content-Type", "application/json");
            o.setRequestHeader("X-Session-Token", _state.token);
            o.timeout = 120000;
            o.onload = function() {
                if (o.status === 401) { _onAuthFail(); n(new Error("登录已失效")); return; }
                if (o.status >= 200 && o.status < 300) {
                    try {
                        r(JSON.parse(o.responseText));
                    } catch(e) {
                        r({});
                    }
                } else {
                    n(new Error(o.statusText || "HTTP " + o.status));
                }
            };
            o.onerror = function() { return n(new Error("Network Error")); };
            o.ontimeout = function() { return n(new Error("请求超时")); };
            o.send(JSON.stringify(t || {}));
        }));
    };
    
    const _get = function(e) {
        return new Promise((function(r, n) {
            const o = new XMLHttpRequest();
            o.open("GET", "/api/wasm/" + e, true);
            o.setRequestHeader("X-Session-Token", _state.token);
            o.onload = function() {
                if (o.status === 401) { _onAuthFail(); n(new Error("登录已失效")); return; }
                if (o.status >= 200 && o.status < 300) {
                    try {
                        r(JSON.parse(o.responseText));
                    } catch(e) {
                        r({});
                    }
                } else {
                    n(new Error(o.statusText));
                }
            };
            o.onerror = function() { return n(new Error("Network Error")); };
            o.send();
        }));
    };
    
    const _escape = function(e) {
        const t = document.createElement("div");
        t.textContent = e;
        return t.innerHTML;
    };
    
    const _isDesktop = function() {
        return window.innerWidth > 768;
    };
    
    const _toast = function(e, t) {
        const n = document.getElementById("toast");
        if (!n) return;
        n.textContent = e;
        n.classList.add("show");
        setTimeout((function() { return n.classList.remove("show"); }), t || 3000);
    };
    
    const _svgImage = '<svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="currentColor" stroke-width="1.5"><rect x="3" y="3" width="18" height="18" rx="2"/><circle cx="8.5" cy="8.5" r="1.5"/><path d="M21 15l-5-5L5 21"/></svg>';
    const _svgVideo = '<svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="currentColor" stroke-width="1.5"><polygon points="23 7 16 12 23 17 23 7"/><rect x="1" y="5" width="15" height="14" rx="2"/></svg>';
    const _svgFile = '<svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="currentColor" stroke-width="1.5"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/></svg>';
    const _svgVoice = '<svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="currentColor" stroke-width="1.5"><path d="M12 1a3 3 0 00-3 3v8a3 3 0 006 0V4a3 3 0 00-3-3z"/><path d="M19 10v2a7 7 0 01-14 0v-2"/><line x1="12" y1="19" x2="12" y2="23"/><line x1="8" y1="23" x2="16" y2="23"/></svg>';
    const _svgPlay = '<svg viewBox="0 0 24 24" width="36" height="36" fill="var(--text-secondary)"><path d="M8 5v14l11-7z"/></svg>';

    const _renderMsg = function(e) {
        const t = document.getElementById("messages-area");
        if (!t) return;
        const n = t.querySelector(".empty-state");
        if (n) n.remove();
        const o = document.createElement("div");
        o.className = "msg-row " + (e.type === "out" ? "out" : "in");
        if (e.id) o.dataset.msgId = e.id;
        var bubbleContent = "";
        var mt = e.media_type;
        if (mt === "image" || mt === 2) {
            var imgSrc = e.media_data || "";
            var cacheSrc = e.media_cache_id ? '/api/wasm/media/' + e.media_cache_id : '';
            if (cacheSrc || imgSrc) {
                var cdnAttr = (e.media_cdn && !e.media_cache_id) ? ' data-cdn="' + encodeURIComponent(e.media_cdn) + '"' : '';
                var displaySrc = imgSrc || cacheSrc;
                var loadAttr = cacheSrc && imgSrc ? ' data-hq-src="' + cacheSrc + '"' : '';
                bubbleContent = '<div class="bubble-media-img-wrap"' + cdnAttr + loadAttr + '><img class="bubble-media-img" src="' + displaySrc + '" alt="图片" /></div>';
            } else if (e.media_cdn) {
                bubbleContent = '<div class="bubble-media-img-wrap bubble-media-loading" data-cdn="' + encodeURIComponent(e.media_cdn) + '"><div class="bubble-media-placeholder">' + _svgImage + '<span>图片</span></div></div>';
            } else {
                bubbleContent = '<div class="bubble-media-placeholder">' + _svgImage + '<span>图片</span></div>';
            }
        } else if (mt === "video" || mt === 5) {
            if (e.media_cache_id) {
                var videoSrc = '/api/wasm/media/' + e.media_cache_id;
                bubbleContent = '<div class="bubble-media-img-wrap"><div class="bubble-media-video-thumb" data-action="play-video" data-video-src="' + videoSrc + '"><video class="bubble-media-video-thumb-vid" src="' + videoSrc + '" preload="metadata" muted playsinline></video><div class="bubble-media-play-btn">' + _svgPlay + '</div></div></div>';
            } else if (e.media_data) {
                bubbleContent = '<div class="bubble-media-img-wrap"><div class="bubble-media-video-thumb" data-action="play-video"><img class="bubble-media-img" src="' + e.media_data + '" alt="视频" /><div class="bubble-media-play-btn">' + _svgPlay + '</div></div></div>';
            } else if (e.media_cdn) {
                bubbleContent = '<div class="bubble-media-img-wrap bubble-media-loading" data-cdn="' + encodeURIComponent(e.media_cdn) + '" data-media-type="video"><div class="bubble-media-placeholder">' + _svgVideo + '<span>视频</span></div></div>';
            } else {
                bubbleContent = '<div class="bubble-media-file"><div class="bubble-media-file-icon">' + _svgVideo + '</div><div class="bubble-media-file-info"><div class="bubble-media-file-name">' + _escape(e.media_filename || "视频") + '</div><div class="bubble-media-file-size">' + (e.media_duration ? (e.media_duration / 1000).toFixed(1) + "s" : "") + '</div></div></div>';
            }
        } else if (mt === "file" || mt === 4) {
            if (e.media_cache_id) {
                bubbleContent = '<div class="bubble-media-file"><div class="bubble-media-file-icon">' + _svgFile + '</div><div class="bubble-media-file-info"><div class="bubble-media-file-name">' + _escape(e.media_filename || "文件") + '</div></div></div>';
            } else if (e.media_cdn) {
                bubbleContent = '<div class="bubble-media-file bubble-media-loading" data-cdn="' + encodeURIComponent(e.media_cdn) + '" data-media-type="file"><div class="bubble-media-file-icon">' + _svgFile + '</div><div class="bubble-media-file-info"><div class="bubble-media-file-name">' + _escape(e.media_filename || "文件") + '</div></div></div>';
            } else {
                bubbleContent = '<div class="bubble-media-file"><div class="bubble-media-file-icon">' + _svgFile + '</div><div class="bubble-media-file-info"><div class="bubble-media-file-name">' + _escape(e.media_filename || "文件") + '</div></div></div>';
            }
        } else if (mt === "voice" || mt === 3) {
            var dur = e.media_duration ? Math.ceil(e.media_duration / 1000) : 1;
            var bars = "";
            var seed = e.id || 0;
            for (var i = 0; i < Math.min(dur, 12); i++) {
                seed = (seed * 1103515245 + 12345) & 0x7fffffff;
                var h = 6 + (seed % 14);
                bars += '<div class="bubble-media-voice-bar" style="height:' + h + 'px"></div>';
            }
            if (e.media_cache_id) {
                bubbleContent = '<div class="bubble-media-voice" data-action="play-voice" data-cache-id="' + e.media_cache_id + '">' + _svgVoice + '<div class="bubble-media-voice-bars">' + bars + '</div><div class="bubble-media-voice-dur">' + dur + '"</div><div class="bubble-media-voice-progress"><div class="bubble-media-voice-progress-fill"></div></div></div>';
            } else if (e.media_cdn) {
                bubbleContent = '<div class="bubble-media-voice bubble-media-loading" data-cdn="' + encodeURIComponent(e.media_cdn) + '" data-media-type="voice">' + _svgVoice + '<div class="bubble-media-voice-bars">' + bars + '</div><div class="bubble-media-voice-dur">' + dur + '"</div><div class="bubble-media-voice-progress"><div class="bubble-media-voice-progress-fill"></div></div></div>';
            } else {
                bubbleContent = '<div class="bubble-media-voice">' + _svgVoice + '<div class="bubble-media-voice-bars">' + bars + '</div><div class="bubble-media-voice-dur">' + dur + '"</div><div class="bubble-media-voice-progress"><div class="bubble-media-voice-progress-fill"></div></div></div>';
            }
        } else {
            bubbleContent = '<div class="bubble-text">' + _escape(e.text || "") + '</div>';
        }
        o.innerHTML = '<div class="bubble ' + (e.type === "out" ? "out" : "in") + '">' + bubbleContent + '<div class="msg-time-row">' + (e.media_cdn && !e.media_cache_id ? '<span class="msg-send-status msg-send-loading"></span>' : '') + '<span class="msg-time">' + (e.time || "") + '</span></div></div>';
        o._msgData = e;
        t.appendChild(o);
        t.scrollTop = t.scrollHeight;
        
        var loadingEl = o.querySelector('.bubble-media-loading');
        if (loadingEl) {
            window._loadCdnMedia(loadingEl);
        }

        var hqWrap = o.querySelector('.bubble-media-img-wrap[data-hq-src]');
        if (hqWrap) {
            var hqImg = new Image();
            hqImg.onload = (function(wrap, src) {
                return function() {
                    var img = wrap.querySelector('.bubble-media-img');
                    if (img) img.src = src;
                };
            })(hqWrap, hqWrap.dataset.hqSrc);
            hqImg.src = hqWrap.dataset.hqSrc;
        }

        const bubbleDiv = o.querySelector('.bubble');
        if (bubbleDiv) {
            var isMediaMsg = (mt === "image" || mt === 2 || mt === "video" || mt === 5 || mt === "voice" || mt === 3 || mt === "file" || mt === 4);
            if (isMediaMsg) {
                bubbleDiv.style.cursor = 'pointer';
                bubbleDiv.addEventListener('click', (function(ev) {
                    ev.stopPropagation();
                    _handleMediaClick(e);
                }));
            } else if (e.type === 'in') {
                bubbleDiv.style.cursor = 'pointer';
                bubbleDiv.addEventListener('click', (function(ev) {
                    ev.stopPropagation();
                    _showAiModal(e.id, e.text, e.from || _state.currentUser);
                }));
            }
        }
    };

    const _renderSendingMsg = function(e) {
        const t = document.getElementById("messages-area");
        if (!t) return;
        const n = t.querySelector(".empty-state");
        if (n) n.remove();
        const o = document.createElement("div");
        o.className = "msg-row out";
        o.dataset.sendingId = e.id;
        if (e.id) o.dataset.msgId = e.id;
        var bubbleContent = "";
        var mt = e.media_type;
        if (mt === 2 && e.media_data) {
            bubbleContent = '<div class="bubble-media-img-wrap"><img class="bubble-media-img" src="' + e.media_data + '" alt="图片" /></div>';
        } else if (mt === 5 && e.media_data) {
            bubbleContent = '<div class="bubble-media-img-wrap"><div class="bubble-media-video-thumb"><img class="bubble-media-img" src="' + e.media_data + '" alt="视频" /><div class="bubble-media-play-btn">' + _svgPlay + '</div></div></div>';
        } else if (mt === 3) {
            var dur = e.media_duration ? Math.ceil(e.media_duration / 1000) : 1;
            var bars = "";
            var seed = e.id || 0;
            for (var i = 0; i < Math.min(dur, 12); i++) {
                seed = (seed * 1103515245 + 12345) & 0x7fffffff;
                var h = 6 + (seed % 14);
                bars += '<div class="bubble-media-voice-bar" style="height:' + h + 'px"></div>';
            }
            bubbleContent = '<div class="bubble-media-voice">' + _svgVoice + '<div class="bubble-media-voice-bars">' + bars + '</div><div class="bubble-media-voice-dur">' + dur + '"</div><div class="bubble-media-voice-progress"><div class="bubble-media-voice-progress-fill"></div></div></div>';
        } else if (mt === 4) {
            bubbleContent = '<div class="bubble-media-file"><div class="bubble-media-file-icon">' + _svgFile + '</div><div class="bubble-media-file-info"><div class="bubble-media-file-name">' + _escape(e.media_filename || "文件") + '</div></div></div>';
        } else {
            bubbleContent = '<div class="bubble-text">' + _escape(e.text || "") + '</div>';
        }
        o.innerHTML = '<div class="bubble out">' + bubbleContent + '<div class="msg-time-row"><span class="msg-send-status msg-send-loading"></span><span class="msg-time">' + (e.time || "") + '</span></div></div>';
        t.appendChild(o);
        t.scrollTop = t.scrollHeight;
    };

    var _currentAudio = null;
    var _currentVoiceEl = null;

    const _cdnInfoStr = function(cdn) {
        if (typeof cdn === 'string') return cdn;
        return JSON.stringify(cdn);
    };

    const _handleMediaClick = function(msg) {
        var mt = msg.media_type;
        if (mt === "image" || mt === 2) {
            if (msg.media_cache_id) {
                window._previewImage('/api/wasm/media/' + msg.media_cache_id);
            } else if (msg.media_data) {
                window._previewImage(msg.media_data);
            } else if (msg.media_cdn) {
                _toast("正在加载图片...");
                fetch('/api/wasm/download-media', {
                    method: 'POST',
                    headers: {'Content-Type': 'application/json', 'X-Session-Token': _state.token},
                    body: JSON.stringify({cdn_info: _cdnInfoStr(msg.media_cdn)})
                }).then(function(r) {
                    if (!r.ok) {
                        return r.text().then(function(t) { throw new Error('HTTP ' + r.status); });
                    }
                    return r.json();
                }).then(function(result) {
                    if (result.success && result.cache_key) {
                        window._previewImage('/api/wasm/media/' + result.cache_key);
                    } else {
                        _toast("图片加载失败: " + (result.error || ""));
                    }
                }).catch(function(err) {
                    console.log('图片加载异常:', err);
                    _toast("图片加载失败");
                });
            }
        } else if (mt === "video" || mt === 5) {
            _playVideo(msg);
        } else if (mt === "voice" || mt === 3) {
            _playVoice(msg);
        } else if (mt === "file" || mt === 4) {
            _downloadMedia(msg, "file");
        }
    };

    var _voicePlayFailed = {};
    var _voiceProgressRaf = null;

    const _playVoice = function(msg) {
        if (!msg.media_cdn && !msg.media_cache_id) {
            _toast("语音数据不可用");
            return;
        }
        var msgId = msg.id;
        if (_voicePlayFailed[msgId]) {
            delete _voicePlayFailed[msgId];
            _downloadMedia(msg, "voice");
            return;
        }
        if (_currentAudio) {
            _currentAudio.pause();
            _currentAudio = null;
            if (_currentVoiceEl) {
                _currentVoiceEl.classList.remove('voice-playing');
                var pf = _currentVoiceEl.querySelector('.bubble-media-voice-progress-fill');
                if (pf) pf.style.width = '0%';
                _currentVoiceEl = null;
            }
        }
        if (_voiceProgressRaf) {
            cancelAnimationFrame(_voiceProgressRaf);
            _voiceProgressRaf = null;
        }
        var voiceEl = null;
        if (msgId) {
            var msgRow = document.querySelector('[data-msg-id="' + msgId + '"]');
            if (msgRow) voiceEl = msgRow.querySelector('.bubble-media-voice');
        }
        var tryPlayAudio = function(cacheId) {
            var cacheUrl = '/api/wasm/media/' + cacheId;
            var audio = new Audio();
            var hasPlayed = false;
            var updateProgress = function() {
                if (!audio.duration || !voiceEl) return;
                var pct = (audio.currentTime / audio.duration) * 100;
                var fill = voiceEl.querySelector('.bubble-media-voice-progress-fill');
                if (fill) fill.style.width = pct + '%';
                var durEl = voiceEl.querySelector('.bubble-media-voice-dur');
                if (durEl && audio.duration) {
                    var remain = Math.ceil(audio.duration - audio.currentTime);
                    durEl.textContent = remain + '"';
                }
                if (!audio.paused && !audio.ended) {
                    _voiceProgressRaf = requestAnimationFrame(updateProgress);
                }
            };
            audio.addEventListener('canplaythrough', function() {
                hasPlayed = true;
                if (voiceEl) {
                    _currentVoiceEl = voiceEl;
                    voiceEl.classList.add('voice-playing');
                }
                audio.play().catch(function() {
                    _voicePlayFailed[msgId] = true;
                    if (voiceEl) voiceEl.classList.remove('voice-playing');
                    var pf = voiceEl ? voiceEl.querySelector('.bubble-media-voice-progress-fill') : null;
                    if (pf) pf.style.width = '0%';
                    _toast("语音播放失败，再次点击可下载");
                });
            });
            audio.addEventListener('error', function() {
                if (!hasPlayed) {
                    _voicePlayFailed[msgId] = true;
                    _toast("浏览器不支持此语音格式，再次点击可下载");
                }
            });
            audio.addEventListener('ended', function() {
                _currentAudio = null;
                if (_voiceProgressRaf) {
                    cancelAnimationFrame(_voiceProgressRaf);
                    _voiceProgressRaf = null;
                }
                if (_currentVoiceEl) {
                    _currentVoiceEl.classList.remove('voice-playing');
                    var pf = _currentVoiceEl.querySelector('.bubble-media-voice-progress-fill');
                    if (pf) pf.style.width = '0%';
                    var durEl = _currentVoiceEl.querySelector('.bubble-media-voice-dur');
                    if (durEl && msg.media_duration) durEl.textContent = Math.ceil(msg.media_duration / 1000) + '"';
                    _currentVoiceEl = null;
                }
            });
            audio.addEventListener('playing', function() {
                if (_voiceProgressRaf) cancelAnimationFrame(_voiceProgressRaf);
                _voiceProgressRaf = requestAnimationFrame(updateProgress);
            });
            audio.addEventListener('pause', function() {
                if (_voiceProgressRaf) {
                    cancelAnimationFrame(_voiceProgressRaf);
                    _voiceProgressRaf = null;
                }
            });
            _currentAudio = audio;
            audio.src = cacheUrl;
            audio.load();
        };
        if (msg.media_cache_id) {
            tryPlayAudio(msg.media_cache_id);
            return;
        }
        _toast("正在加载语音...");
        fetch('/api/wasm/download-media', {
            method: 'POST',
            headers: {'Content-Type': 'application/json', 'X-Session-Token': _state.token},
            body: JSON.stringify({cdn_info: _cdnInfoStr(msg.media_cdn)})
        }).then(function(r) {
            if (!r.ok) {
                return r.text().then(function(t) { throw new Error('HTTP ' + r.status + ': ' + t); });
            }
            return r.json();
        }).then(function(result) {
            if (result.success && result.cache_key) {
                tryPlayAudio(result.cache_key);
            } else {
                _toast("语音加载失败: " + (result.error || "未知错误"));
            }
        }).catch(function(err) {
            console.log('语音加载异常:', err);
            _toast("语音加载失败");
        });
    };

    const _playVideo = function(msg) {
        if (!msg.media_cdn && !msg.media_cache_id) {
            _toast("视频数据不可用");
            return;
        }
        var tryPlayVideo = function(cacheId) {
            var videoUrl = '/api/wasm/media/' + cacheId;
            window._previewVideo(videoUrl);
        };
        if (msg.media_cache_id) {
            tryPlayVideo(msg.media_cache_id);
            return;
        }
        _toast("正在加载视频...");
        fetch('/api/wasm/download-media', {
            method: 'POST',
            headers: {'Content-Type': 'application/json', 'X-Session-Token': _state.token},
            body: JSON.stringify({cdn_info: _cdnInfoStr(msg.media_cdn)})
        }).then(function(r) {
            if (!r.ok) {
                return r.text().then(function(t) { throw new Error('HTTP ' + r.status + ': ' + t); });
            }
            return r.json();
        }).then(function(result) {
            if (result.success && result.cache_key) {
                tryPlayVideo(result.cache_key);
            } else {
                _toast("视频加载失败: " + (result.error || "未知错误"));
            }
        }).catch(function(err) {
            console.log('视频加载异常:', err);
            _toast("视频加载失败");
        });
    };

    window._previewVideo = function(src) {
        var overlay = document.createElement('div');
        overlay.style.cssText = 'position:fixed;top:0;left:0;right:0;bottom:0;background:rgba(0,0,0,0.92);z-index:10002;display:flex;flex-direction:column;align-items:center;justify-content:center;cursor:default';
        var video = document.createElement('video');
        video.src = src;
        video.controls = true;
        video.autoplay = true;
        video.playsInline = true;
        video.style.cssText = 'max-width:95%;max-height:85%;border-radius:8px;background:#000;outline:none';
        var closeBtn = document.createElement('div');
        closeBtn.style.cssText = 'position:absolute;top:16px;right:16px;width:36px;height:36px;border-radius:50%;background:rgba(255,255,255,0.2);display:flex;align-items:center;justify-content:center;cursor:pointer;font-size:20px;color:#fff;z-index:10003';
        closeBtn.innerHTML = '<svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="#fff" stroke-width="2"><line x1="18" y1="6" x2="6" y2="18"/><line x1="6" y1="6" x2="18" y2="18"/></svg>';
        var downloadBtn = document.createElement('div');
        downloadBtn.style.cssText = 'position:absolute;top:16px;right:60px;width:36px;height:36px;border-radius:50%;background:rgba(255,255,255,0.2);display:flex;align-items:center;justify-content:center;cursor:pointer;z-index:10003';
        downloadBtn.innerHTML = '<svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="#fff" stroke-width="2"><path d="M21 15v4a2 2 0 01-2 2H5a2 2 0 01-2-2v-4"/><polyline points="7 10 12 15 17 10"/><line x1="12" y1="15" x2="12" y2="3"/></svg>';
        downloadBtn.addEventListener('click', function(ev) {
            ev.stopPropagation();
            var a = document.createElement('a');
            a.href = src + '?download=1';
            a.download = 'video.mp4';
            document.body.appendChild(a);
            a.click();
            document.body.removeChild(a);
        });
        var closeOverlay = function() {
            video.pause();
            video.src = '';
            if (overlay.parentNode) document.body.removeChild(overlay);
        };
        closeBtn.addEventListener('click', function(ev) { ev.stopPropagation(); closeOverlay(); });
        overlay.addEventListener('click', function(ev) { if (ev.target === overlay) closeOverlay(); });
        overlay.appendChild(video);
        overlay.appendChild(closeBtn);
        overlay.appendChild(downloadBtn);
        document.body.appendChild(overlay);
    };

    const _downloadDirectUrl = function(cacheId, filename) {
        try {
            var downloadUrl = '/api/wasm/media/' + cacheId + '?download=1';
            var a = document.createElement('a');
            a.href = downloadUrl;
            a.download = filename;
            document.body.appendChild(a);
            a.click();
            document.body.removeChild(a);
            _toast("正在接收: " + filename);
        } catch (err) {
            _toast("下载失败");
        }
    };

    const _downloadMedia = function(msg, mediaType) {
        if (!msg.media_cdn && !msg.media_cache_id) {
            _toast("媒体数据不可用");
            return;
        }
        var filename = msg.media_filename || (mediaType === "video" ? "video.mp4" : mediaType === "voice" ? "voice.silk" : "file.bin");
        if (msg.media_cache_id) {
            _downloadDirectUrl(msg.media_cache_id, filename);
            return;
        }
        _toast("正在接收 " + filename + "...");
        fetch('/api/wasm/download-media', {
            method: 'POST',
            headers: {'Content-Type': 'application/json', 'X-Session-Token': _state.token},
            body: JSON.stringify({cdn_info: _cdnInfoStr(msg.media_cdn)})
        }).then(function(r) {
            if (!r.ok) {
                return r.text().then(function(t) { throw new Error('HTTP ' + r.status); });
            }
            return r.json();
        }).then(function(result) {
            if (result.success && result.cache_key) {
                _downloadDirectUrl(result.cache_key, filename);
            } else {
                _toast("下载失败: " + (result.error || "未知错误"));
            }
        }).catch(function(err) {
            console.log('下载异常:', err);
            _toast("下载失败");
        });
    };

    window._previewImage = function(src) {
        var overlay = document.createElement('div');
        overlay.style.cssText = 'position:fixed;top:0;left:0;right:0;bottom:0;background:rgba(0,0,0,0.9);z-index:10002;display:flex;align-items:center;justify-content:center;cursor:zoom-out';
        var img = document.createElement('img');
        img.src = src;
        img.style.cssText = 'max-width:95%;max-height:95%;object-fit:contain;border-radius:4px';
        overlay.appendChild(img);
        overlay.addEventListener('click', function() { document.body.removeChild(overlay); });
        document.body.appendChild(overlay);
    };

    const _removeLoadingSpinner = function(el, cacheKey) {
        var row = el.closest('.msg-row');
        if (row) {
            var spinner = row.querySelector('.msg-send-loading');
            if (spinner) spinner.remove();
            if (cacheKey && row._msgData) {
                row._msgData.media_cache_id = cacheKey;
            }
        }
    };

    window._loadCdnMedia = function(el) {
        var cdn = decodeURIComponent(el.dataset.cdn || "");
        var mediaType = el.dataset.mediaType || "image";
        if (!cdn) return;
        fetch('/api/wasm/download-media', {
            method: 'POST',
            headers: {'Content-Type': 'application/json', 'X-Session-Token': _state.token},
            body: JSON.stringify({cdn_info: _cdnInfoStr(cdn)})
        }).then(function(r) { return r.json(); }).then(function(result) {
            if (result.success && result.cache_key) {
                var cacheUrl = '/api/wasm/media/' + result.cache_key;
                if (mediaType === "video") {
                    el.innerHTML = '<div class="bubble-media-video-thumb" data-action="play-video" data-video-src="' + cacheUrl + '"><video class="bubble-media-video-thumb-vid" src="' + cacheUrl + '" preload="metadata" muted playsinline></video><div class="bubble-media-play-btn">' + _svgPlay + '</div></div>';
                } else if (mediaType === "file") {
                    el.classList.remove('bubble-media-loading');
                    el.removeAttribute('data-cdn');
                    el.removeAttribute('data-media-type');
                    el.dataset.cacheId = result.cache_key;
                    return _removeLoadingSpinner(el, result.cache_key);
                } else if (mediaType === "voice") {
                    el.classList.remove('bubble-media-loading');
                    el.removeAttribute('data-cdn');
                    el.removeAttribute('data-media-type');
                    el.dataset.action = 'play-voice';
                    el.dataset.cacheId = result.cache_key;
                    return _removeLoadingSpinner(el, result.cache_key);
                } else {
                    el.innerHTML = '<img class="bubble-media-img" src="' + cacheUrl + '" alt="图片" />';
                }
                el.classList.remove('bubble-media-loading');
                _removeLoadingSpinner(el, result.cache_key);
            } else {
                var svgIcon = _svgImage;
                var label = "加载失败";
                if (mediaType === "video") { svgIcon = _svgVideo; label = "视频加载失败"; }
                else if (mediaType === "file") { svgIcon = _svgFile; label = "文件加载失败"; }
                else if (mediaType === "voice") { svgIcon = _svgVoice; label = "语音加载失败"; }
                else { label = "图片加载失败"; }
                el.innerHTML = '<div class="bubble-media-placeholder">' + svgIcon + '<span>' + label + '</span></div>';
                el.classList.remove('bubble-media-loading');
                _removeLoadingSpinner(el);
            }
        }).catch(function() {
            var svgIcon = _svgImage;
            var label = "加载失败";
            if (mediaType === "video") { svgIcon = _svgVideo; label = "视频加载失败"; }
            else if (mediaType === "file") { svgIcon = _svgFile; label = "文件加载失败"; }
            else if (mediaType === "voice") { svgIcon = _svgVoice; label = "语音加载失败"; }
            else { label = "图片加载失败"; }
            el.innerHTML = '<div class="bubble-media-placeholder">' + svgIcon + '<span>' + label + '</span></div>';
            el.classList.remove('bubble-media-loading');
            _removeLoadingSpinner(el);
        });
    };
    
    const _showAiModal = function(msgId, msgText, userId) {
        _state.selectedMessage = { id: msgId, text: msgText };
        _state.selectedUserId = userId;
        const modal = document.getElementById("ai-modal");
        const msgPreview = document.getElementById("ai-modal-msg-preview");
        const instructionInput = document.getElementById("ai-instruction");
        if (modal) {
            if (msgPreview) msgPreview.innerText = msgText;
            if (instructionInput) instructionInput.value = "";
            modal.classList.add("show");
        }
    };
    
    const _closeAiModal = function() {
        const modal = document.getElementById("ai-modal");
        if (modal) modal.classList.remove("show");
        _state.selectedMessage = null;
        _state.selectedUserId = null;
    };
    
    const _sendAiReply = async function() {
        if (!_state.selectedMessage) {
            _toast("请选择要回复的消息");
            _closeAiModal();
            return;
        }
        let targetUser = _state.selectedUserId;
        if (!targetUser) {
            targetUser = _state.currentUser;
        }
        if (!targetUser) {
            _toast("无法确定要回复的用户");
            _closeAiModal();
            return;
        }
        const instruction = document.getElementById("ai-instruction") ? document.getElementById("ai-instruction").value : "";
        const selectedMsg = _state.selectedMessage;
        _closeAiModal();
        _toast("正在生成 AI 回复...");
        const result = await _api("ai-manual-reply", {
            user_id: targetUser,
            original_message: selectedMsg.text,
            instruction: instruction
        });
        if (result && result.success) {
            _toast("AI 回复正在生成中，请稍候...");
            setTimeout(_fetchMessages, 2000);
        } else {
            _toast((result && result.error) || "AI 回复失败，请检查 API 配置");
        }
    };
    
    const _loadUsers = async function() {
        const e = await _get("users");
        if (e && e.users) {
            _state.users = e.users;
            if (_state.view === 'list') {
                _renderChatList();
                _loadChatListPreviews();
            }
        }
    };
    
    const _updateSelector = function() {
        const e = document.getElementById("user-select-btn");
        const t = document.getElementById("user-dropdown");
        if (e && _state.currentUser) {
            var nick = _state.nicknames[_state.currentUser] || '';
            e.textContent = nick || (_state.currentUser ? _state.currentUser.substring(0, 15) + (_state.currentUser.length > 15 ? "..." : "") : "选择用户");
        }
        if (t && _state.users && _state.users.length > 0) {
            t.innerHTML = _state.users.map((function(r) {
                return `<div class="user-option ${r === _state.currentUser ? "current" : ""}" data-user-id="${r}">用户 ${r}</div>`;
            })).join("");
            t.querySelectorAll(".user-option").forEach((function(e) {
                e.addEventListener("click", (function() {
                    const t = e.getAttribute("data-user-id");
                    if (t) _openChat(t);
                }));
            }));
        }
    };
    
    const _selectUser = async function(e) {
        if (!e) return;
        _openChat(e);
    };
    
    const _renderChatList = function() {
        var container = document.getElementById("chat-list-items");
        if (!container) return;
        if (!_state.users || _state.users.length === 0) {
            container.innerHTML = '<div class="chat-list-empty"><div class="chat-list-empty-icon">💬</div><div>暂无聊天</div></div>';
            return;
        }
        var isDesktop = _isDesktop();
        var html = '';
        _state.users.forEach(function(userId) {
            var nickname = _state.nicknames[userId] || '';
            var displayName = nickname || userId;
            var lastMsg = _state.lastMessages[userId];
            var preview = '';
            var time = '';
            if (lastMsg) {
                if (lastMsg.media_type) {
                    var mediaLabels = {2: '[图片]', 3: '[语音]', 4: '[文件]', 5: '[视频]', 'image': '[图片]', 'voice': '[语音]', 'file': '[文件]', 'video': '[视频]'};
                    preview = mediaLabels[lastMsg.media_type] || lastMsg.text || '';
                } else {
                    preview = lastMsg.text || '';
                }
                time = lastMsg.time || '';
            }
            var activeStyle = (isDesktop && _state.currentUser === userId) ? ' style="background:var(--accent-light)"' : '';
            html += '<div class="chat-list-item" data-user-id="' + _escape(userId) + '"' + activeStyle + '>' +
                '<div class="chat-list-item-avatar">用户</div>' +
                '<div class="chat-list-item-content">' +
                '<div class="chat-list-item-name">' + _escape(displayName) + '</div>' +
                '<div class="chat-list-item-msg">' + _escape(preview) + '</div>' +
                '</div>' +
                '<div class="chat-list-item-time">' + time + '</div>' +
                '</div>';
        });
        container.innerHTML = html;
        container.querySelectorAll('.chat-list-item').forEach(function(item) {
            item.addEventListener('click', function() {
                var userId = item.getAttribute('data-user-id');
                if (userId) _openChat(userId);
            });
        });
    };
    
    const _openChat = async function(userId) {
    if (!userId) return;
    _state.currentUser = userId;
    _state.view = 'chat';
    _state.displayedIds.clear();
    _state.lastMsgId = 0;
    var isDesktop = _isDesktop();
    if (!isDesktop) {
        var chatListPage = document.getElementById("chat-list-page");
        if (chatListPage) chatListPage.classList.remove("active");
        var userListPage = document.getElementById("user-list-page");
        if (userListPage) userListPage.classList.remove("active");
    } else {
        var chatListPage = document.getElementById("chat-list-page");
        if (chatListPage) chatListPage.classList.add("active");
        var userListPage = document.getElementById("user-list-page");
        if (userListPage) userListPage.classList.remove("active");
        var chatListItems = document.querySelectorAll('.chat-list-item');
        chatListItems.forEach(function(item) {
            if (item.getAttribute('data-user-id') === userId) {
                item.style.background = 'var(--accent-light)';
            } else {
                item.style.background = '';
            }
        });
    }
    var chatPage = document.getElementById("chat-page");
    if (chatPage) chatPage.classList.add("active");
    if (isDesktop) {
        var inputArea = document.getElementById("input-area");
        if (inputArea) inputArea.style.display = '';
        var menuWrap = document.querySelector('.chat-header-menu-wrap');
        if (menuWrap) menuWrap.style.display = '';
    }
    if (!isDesktop) {
        var tabbar = document.getElementById("bottom-tab-bar");
        if (tabbar) tabbar.classList.add("hidden");
    }
    var title = document.getElementById("chat-header-title");
    if (title) {
        var nickname = _state.nicknames[userId] || '';
        title.textContent = nickname || userId;
    }
    var messagesArea = document.getElementById("messages-area");
    if (messagesArea) {
        messagesArea.innerHTML = '';
    }
    await _api("switch-user", { user_id: userId });
    _loadHistory(userId);
};
    
    const _backToChatList = function() {
        _state.view = 'list';
        _state.currentUser = null;
        _state.displayedIds.clear();
        var isDesktop = _isDesktop();
        var chatPage = document.getElementById("chat-page");
        if (chatPage) {
            if (isDesktop) {
                var messagesArea = document.getElementById("messages-area");
            if (messagesArea) messagesArea.innerHTML = '<div class="pc-chat-empty">选择一个聊天开始对话</div>';
            var title = document.getElementById("chat-header-title");
            if (title) title.textContent = 'Zyn iLink Chatbox';
            var inputArea = document.getElementById("input-area");
            if (inputArea) inputArea.style.display = 'none';
            var menuWrap = document.querySelector('.chat-header-menu-wrap');
            if (menuWrap) menuWrap.style.display = 'none';
            } else {
                chatPage.classList.remove("active");
            }
        }
        var chatListPage = document.getElementById("chat-list-page");
        if (chatListPage) chatListPage.classList.add("active");
        if (!isDesktop) {
            var tabbar = document.getElementById("bottom-tab-bar");
            if (tabbar) tabbar.classList.remove("hidden");
        }
        _setActiveTab('tab-list');
        _loadUsers();
        _loadChatListPreviews();
    };

    const _setActiveTab = function(tabId) {
        var tabs = document.querySelectorAll('.bottom-tab-item');
        tabs.forEach(function(t) { t.classList.remove('active'); });
        var tab = document.getElementById(tabId);
        if (tab) tab.classList.add('active');
    };

    const _switchToChatList = function() {
        _closeSettings();
        var isDesktop = _isDesktop();
        if (!isDesktop) {
            var chatPage = document.getElementById("chat-page");
            if (chatPage) chatPage.classList.remove("active");
        }
        var userListPage = document.getElementById("user-list-page");
        if (userListPage) userListPage.classList.remove("active");
        var chatListPage = document.getElementById("chat-list-page");
        if (chatListPage) chatListPage.classList.add("active");
        if (!isDesktop) {
            var tabbar = document.getElementById("bottom-tab-bar");
            if (tabbar) tabbar.classList.remove("hidden");
        }
        _setActiveTab('tab-list');
        _state.view = 'list';
        _loadChatListPreviews();
    };

    const _switchToUserList = function() {
        _closeSettings();
        var isDesktop = _isDesktop();
        if (!isDesktop) {
            var chatPage = document.getElementById("chat-page");
            if (chatPage) chatPage.classList.remove("active");
        }
        var chatListPage = document.getElementById("chat-list-page");
        if (chatListPage) chatListPage.classList.remove("active");
        var userListPage = document.getElementById("user-list-page");
        if (userListPage) userListPage.classList.add("active");
        if (!isDesktop) {
            var tabbar = document.getElementById("bottom-tab-bar");
            if (tabbar) tabbar.classList.remove("hidden");
        }
        _setActiveTab('tab-users');
        _state.view = 'users';
        _renderUserMgmtList();
    };

    const _switchToSettings = function() {
        _openSettings();
    };

    const _deleteUser = async function(userId) {
        if (!userId) return;
        try {
            var result = await _api("delete-user", { user_id: userId });
            if (result && result.success) {
                _toast("已删除");
                _state.users = result.users || [];
                if (result.current_user) {
                    _state.currentUser = result.current_user;
                }
                _renderChatList();
                _loadChatListPreviews();
                _renderUserMgmtList();
            } else {
                _toast((result && result.error) || "删除失败");
            }
        } catch(e) {
            _toast("删除失败");
        }
    };

    var _addUserPollTimer = null;

    const _startAddUser = async function() {
        var modal = document.getElementById("add-user-modal");
        var statusEl = document.getElementById("add-user-status");
        var qrEl = document.getElementById("add-user-qr");
        if (modal) modal.classList.add("show");
        if (statusEl) statusEl.textContent = "正在生成二维码...";
        if (qrEl) qrEl.innerHTML = '<div class="add-user-modal-spinner"></div>';
        
        try {
            var result = await _api("add-user-start", {});
            if (!result.success) {
                if (statusEl) statusEl.textContent = result.error || "启动失败，请重试";
                return;
            }
            if (result.status === "already_running") {
                if (statusEl) statusEl.textContent = "已有进行中的添加操作，请等待...";
                _startAddUserPoll();
                return;
            }
            if (result.matrix) {
                _renderAddUserQR(result.matrix);
                if (statusEl) statusEl.textContent = "请使用微信扫码添加新用户";
                _startAddUserPoll();
            } else {
                if (statusEl) statusEl.textContent = "正在生成二维码...";
                _startAddUserPoll();
            }
        } catch(e) {
            if (statusEl) statusEl.textContent = "启动失败，请重试";
        }
    };

    const _startAddUserPoll = function() {
        if (_addUserPollTimer) clearInterval(_addUserPollTimer);
        _addUserPollTimer = setInterval(async function() {
            try {
                var data = await _get("add-user-status");
                var statusEl = document.getElementById("add-user-status");
                var qrEl = document.getElementById("add-user-qr");
                
                if (data.matrix && qrEl) {
                    _renderAddUserQR(data.matrix);
                }
                
                var st = data.qrcode_status;
                if (st === "scaned" && statusEl) {
                    statusEl.textContent = "已扫码，请在手机上确认...";
                } else if (st === "done") {
                    if (statusEl) statusEl.textContent = "连接成功！正在刷新用户列表...";
                    if (_addUserPollTimer) { clearInterval(_addUserPollTimer); _addUserPollTimer = null; }
                    await _loadUsers();
                    _renderChatList();
                    _loadChatListPreviews();
                    _toast("新用户已添加！");
                    setTimeout(_closeAddUserModal, 1500);
                } else if (st === "expired" || st === "timeout") {
                    if (statusEl) statusEl.textContent = "二维码已过期，请重新点击加号重试";
                    if (_addUserPollTimer) { clearInterval(_addUserPollTimer); _addUserPollTimer = null; }
                } else if (st === "error") {
                    if (statusEl) statusEl.textContent = "获取失败，请重试";
                    if (_addUserPollTimer) { clearInterval(_addUserPollTimer); _addUserPollTimer = null; }
                } else if (st === "waiting" && statusEl) {
                    statusEl.textContent = "请使用微信扫码添加新用户";
                }
            } catch(e) {}
        }, 2000);
    };

    const _renderAddUserQR = function(matrix) {
        var qrEl = document.getElementById("add-user-qr");
        if (!qrEl || !matrix) return;
        var rows = matrix.length;
        var cols = matrix[0].length;
        var cellSize = Math.max(6, Math.min(12, Math.floor(280 / cols)));
        var width = cols * cellSize + 40;
        var html = '<div class="qr-grid" style="grid-template-columns: repeat(' + cols + ', ' + cellSize + 'px); width: ' + width + 'px; max-width: 100%; overflow-x: auto; margin: 0 auto;">';
        for (var i = 0; i < rows; i++) {
            for (var j = 0; j < cols; j++) {
                html += '<div class="qr-cell ' + (matrix[i][j] === " " ? "white" : "") + '" style="width:' + cellSize + 'px;height:' + cellSize + 'px;"></div>';
            }
        }
        html += "</div>";
        qrEl.innerHTML = html;
    };

    const _closeAddUserModal = function() {
        var modal = document.getElementById("add-user-modal");
        if (modal) modal.classList.remove("show");
        if (_addUserPollTimer) { clearInterval(_addUserPollTimer); _addUserPollTimer = null; }
    };

    const _loadChatListPreviews = async function() {
        var promises = _state.users.map(async function(userId) {
            try {
                var data = await _get("history?user=" + encodeURIComponent(userId) + "&limit=1");
                if (data && data.messages && data.messages.length > 0) {
                    var lastMsg = data.messages[data.messages.length - 1];
                    _state.lastMessages[userId] = {
                        text: lastMsg.text || '',
                        time: lastMsg.time || '',
                        media_type: lastMsg.media_type
                    };
                }
            } catch(e) {}
        });
        await Promise.all(promises);
        _renderChatList();
    };
    
    const _openNicknameModal = function() {
        if (!_state.currentUser) return;
        var modal = document.getElementById("nickname-modal");
        var input = document.getElementById("nickname-input");
        var userIdDiv = document.getElementById("nickname-modal-userid");
        if (!modal || !input) return;
        if (userIdDiv) userIdDiv.textContent = '用户ID: ' + _state.currentUser;
        input.value = _state.nicknames[_state.currentUser] || '';
        modal.classList.add("show");
        setTimeout(function() { input.focus(); }, 100);
    };
    
    const _closeNicknameModal = function() {
        var modal = document.getElementById("nickname-modal");
        if (modal) modal.classList.remove("show");
    };
    
    const _saveNickname = function() {
        if (!_state.currentUser) return;
        var input = document.getElementById("nickname-input");
        var nickname = input ? input.value.trim() : '';
        if (nickname) {
            _state.nicknames[_state.currentUser] = nickname;
        } else {
            delete _state.nicknames[_state.currentUser];
        }
        localStorage.setItem("zyn_nicknames", JSON.stringify(_state.nicknames));
        var title = document.getElementById("chat-header-title");
        if (title) title.textContent = nickname || _state.currentUser;
        _closeNicknameModal();
        _toast(nickname ? "备注名已保存" : "备注名已清除");
    };
    
    const _toggleChatMenu = async function(e) {
        if (e) e.stopPropagation();
        var dropdown = document.getElementById("chat-menu-dropdown");
        if (dropdown) {
            var isShown = dropdown.classList.contains("show");
            dropdown.classList.toggle("show");
            if (!isShown && _state.currentUser) {
                try {
                    var result = await _get("user-prompt?user_id=" + encodeURIComponent(_state.currentUser));
                    if (result && result.success) {
                        var aiSwitch = document.getElementById("chat-menu-ai-switch");
                        if (aiSwitch) {
                            var enabled = result.ai_enabled !== null && result.ai_enabled !== undefined ? result.ai_enabled : result.global_auto_reply;
                            if (enabled) { aiSwitch.classList.add("on"); } else { aiSwitch.classList.remove("on"); }
                        }
                        var scheduledSwitch = document.getElementById("chat-menu-scheduled-switch");
                        if (scheduledSwitch) {
                            var schedEnabled = result.scheduled_enabled !== null && result.scheduled_enabled !== undefined ? result.scheduled_enabled : result.global_scheduled_reply;
                            if (schedEnabled) { scheduledSwitch.classList.add("on"); } else { scheduledSwitch.classList.remove("on"); }
                        }
                        var dailySwitch = document.getElementById("chat-menu-daily-switch");
                        if (dailySwitch) {
                            var dailyEnabled = result.daily_enabled !== null && result.daily_enabled !== undefined ? result.daily_enabled : result.global_daily_reply;
                            if (dailyEnabled) { dailySwitch.classList.add("on"); } else { dailySwitch.classList.remove("on"); }
                        }
                    }
                } catch(ex) {}
            }
        }
    };
    
    const _closeChatMenu = function() {
        var dropdown = document.getElementById("chat-menu-dropdown");
        if (dropdown) dropdown.classList.remove("show");
    };
    
    const _toggleUserAiEnabled = async function() {
        if (!_state.currentUser) return;
        var aiSwitch = document.getElementById("chat-menu-ai-switch");
        if (!aiSwitch) return;
        var isOn = aiSwitch.classList.contains("on");
        aiSwitch.classList.toggle("on");
        var newVal = !isOn;
        var result = await _api("user-prompt", { user_id: _state.currentUser, ai_enabled: newVal });
        if (result && result.success) {
            _toast(newVal ? "已为此用户启用AI回复" : "已为此用户关闭AI回复");
        } else {
            if (newVal) { aiSwitch.classList.remove("on"); } else { aiSwitch.classList.add("on"); }
            _toast((result && result.error) || "保存失败");
        }
    };

    const _toggleUserScheduledEnabled = async function() {
        if (!_state.currentUser) return;
        var scheduledSwitch = document.getElementById("chat-menu-scheduled-switch");
        if (!scheduledSwitch) return;
        var isOn = scheduledSwitch.classList.contains("on");
        scheduledSwitch.classList.toggle("on");
        var newVal = !isOn;
        var result = await _api("user-prompt", { user_id: _state.currentUser, scheduled_enabled: newVal });
        if (result && result.success) {
            _toast(newVal ? "已为此用户启用间隔定时发送" : "已为此用户关闭间隔定时发送");
        } else {
            if (newVal) { scheduledSwitch.classList.remove("on"); } else { scheduledSwitch.classList.add("on"); }
            _toast((result && result.error) || "保存失败");
        }
    };

    const _toggleUserDailyEnabled = async function() {
        if (!_state.currentUser) return;
        var dailySwitch = document.getElementById("chat-menu-daily-switch");
        if (!dailySwitch) return;
        var isOn = dailySwitch.classList.contains("on");
        dailySwitch.classList.toggle("on");
        var newVal = !isOn;
        var result = await _api("user-prompt", { user_id: _state.currentUser, daily_enabled: newVal });
        if (result && result.success) {
            _toast(newVal ? "已为此用户启用每日定时发送" : "已为此用户关闭每日定时发送");
        } else {
            if (newVal) { dailySwitch.classList.remove("on"); } else { dailySwitch.classList.add("on"); }
            _toast((result && result.error) || "保存失败");
        }
    };
    
    const _openUserPromptModal = async function() {
        if (!_state.currentUser) return;
        var modal = document.getElementById("user-prompt-modal");
        var input = document.getElementById("user-prompt-input");
        var userIdDiv = document.getElementById("user-prompt-modal-userid");
        var defaultDiv = document.getElementById("user-prompt-modal-default");
        if (!modal || !input) return;
        if (userIdDiv) userIdDiv.textContent = '用户ID: ' + _state.currentUser;
        if (defaultDiv) defaultDiv.textContent = '';
        input.value = "";
        modal.classList.add("show");
        try {
            var result = await _get("user-prompt?user_id=" + encodeURIComponent(_state.currentUser));
            if (result && result.success) {
                input.value = result.prompt || "";
                if (defaultDiv) defaultDiv.textContent = '默认提示词: ' + (result.default_prompt || '').substring(0, 80) + ((result.default_prompt || '').length > 80 ? '...' : '');
            }
        } catch(e) {}
        setTimeout(function() { input.focus(); }, 100);
    };
    
    const _closeUserPromptModal = function() {
        var modal = document.getElementById("user-prompt-modal");
        if (modal) modal.classList.remove("show");
    };
    
    const _saveUserPrompt = async function() {
        if (!_state.currentUser) return;
        var input = document.getElementById("user-prompt-input");
        var prompt = input ? input.value.trim() : '';
        var result = await _api("user-prompt", { user_id: _state.currentUser, prompt: prompt });
        if (result && result.success) {
            _closeUserPromptModal();
            _toast(prompt ? "AI提示词已保存" : "AI提示词已清除，将使用默认提示词");
        } else {
            _toast((result && result.error) || "保存失败");
        }
    };
    
    const _loadHistory = async function(e) {
        const t = e ? `/history?user=${encodeURIComponent(e)}&limit=500` : "/history?limit=500";
        const n = await _get(t);
        if (!n || n.error) return;
        const o = n.messages || [];
        if (o.length === 0) return;
        const i = document.getElementById("messages-area");
        if (i) i.innerHTML = "";
        _state.displayedIds.clear();
        o.forEach((function(e) {
            _renderMsg(e);
            if (e.id) _state.displayedIds.add(e.id);
        }));
        if (o.length > 0) {
            const e = Math.max.apply(null, o.map((function(e) { return e.id || 0; })));
            _state.lastMsgId = Math.max(_state.lastMsgId, e);
        }
        const r = document.getElementById("messages-area");
        if (r) r.scrollTop = r.scrollHeight;
    };
    
    const _fetchMessages = async function() {
        const e = _state.currentUser ? "&user=" + encodeURIComponent(_state.currentUser) : "";
        const t = await _get("messages?since=" + _state.lastMsgId + e);
        if (t && t.messages) {
            t.messages.forEach((function(e) {
                if (e.id && !_state.displayedIds.has(e.id)) {
                    if (_state.view === 'chat' && _state.currentUser) {
                        _renderMsg(e);
                    }
                    _state.displayedIds.add(e.id);
                    _state.lastMsgId = Math.max(_state.lastMsgId, e.id);
                    var fromUser = e.from || _state.currentUser;
                    if (fromUser) {
                        _state.lastMessages[fromUser] = {
                            text: e.text || '',
                            time: e.time || '',
                            media_type: e.media_type
                        };
                    }
                }
            }));
            if (_state.view === 'list' || _isDesktop()) {
                _renderChatList();
            }
        }
    };
    
    const _startPoll = function() {
        if (_state.pollInterval) clearInterval(_state.pollInterval);
        _state.pollInterval = setInterval(_fetchMessages, 500);
    };
    
    const _sendMsg = async function() {
        const e = document.getElementById("message-input");
        const t = e ? e.value.trim() : "";
        if (!t) {
            _toast("请输入消息内容");
            return;
        }
        if (!_state.currentUser) {
            _toast("请先选择用户");
            return;
        }
        if (e) e.value = "";
        const n = await _api("send", { text: t });
        if (n && n.success) {
            setTimeout(_fetchMessages, 200);
            _toast("发送成功");
        } else {
            _toast((n && n.error) || "发送失败");
            if (e) e.value = t;
        }
    };
    
    const _toggleMediaPanel = function() {
        const panel = document.getElementById("media-panel");
        const btn = document.getElementById("plus-btn");
        if (!panel || !btn) return;
        if (panel.classList.contains("show")) {
            panel.classList.remove("show");
            btn.classList.remove("active");
        } else {
            panel.classList.add("show");
            btn.classList.add("active");
            const input = document.getElementById("message-input");
            if (input) input.blur();
        }
    };
    
    const _closeMediaPanel = function() {
        const panel = document.getElementById("media-panel");
        const btn = document.getElementById("plus-btn");
        if (panel) panel.classList.remove("show");
        if (btn) btn.classList.remove("active");
    };
    
    const _showUploadProgress = function(text) {
        const el = document.getElementById("media-upload-progress");
        const txt = el ? el.querySelector(".media-upload-text") : null;
        if (txt) txt.textContent = text || "正在发送...";
        if (el) el.classList.add("show");
    };
    
    const _hideUploadProgress = function() {
        const el = document.getElementById("media-upload-progress");
        if (el) el.classList.remove("show");
    };
    
    const _readFileAsBase64 = function(file) {
        return new Promise(function(resolve, reject) {
            var reader = new FileReader();
            reader.onload = function() {
                var result = reader.result;
                var base64 = result.split(",")[1] || result;
                resolve(base64);
            };
            reader.onerror = function() { reject(reader.error); };
            reader.readAsDataURL(file);
        });
    };
    
    const _readFileAsArrayBuffer = function(file) {
        return new Promise(function(resolve, reject) {
            var reader = new FileReader();
            reader.onload = function() { resolve(reader.result); };
            reader.onerror = function() { reject(reader.error); };
            reader.readAsArrayBuffer(file);
        });
    };
    
    const _generateThumbnail = function(file, maxWidth, maxHeight) {
        return new Promise(function(resolve) {
            if (file.type && file.type.startsWith("image/")) {
                var img = new Image();
                var url = URL.createObjectURL(file);
                img.onload = function() {
                    var w = img.width, h = img.height;
                    var scale = Math.min(maxWidth / w, maxHeight / h, 1);
                    var cw = Math.round(w * scale), ch = Math.round(h * scale);
                    var canvas = document.createElement("canvas");
                    canvas.width = cw; canvas.height = ch;
                    var ctx = canvas.getContext("2d");
                    ctx.drawImage(img, 0, 0, cw, ch);
                    URL.revokeObjectURL(url);
                    var dataUrl = canvas.toDataURL("image/jpeg", 0.6);
                    resolve(dataUrl);
                };
                img.onerror = function() { URL.revokeObjectURL(url); resolve(""); };
                img.src = url;
            } else if (file.type && file.type.startsWith("video/")) {
                var video = document.createElement("video");
                var vurl = URL.createObjectURL(file);
                video.preload = "metadata";
                video.muted = true;
                video.onloadeddata = function() {
                    video.currentTime = Math.min(1, video.duration / 4);
                };
                video.onseeked = function() {
                    var w = video.videoWidth, h = video.videoHeight;
                    var scale = Math.min(maxWidth / w, maxHeight / h, 1);
                    var cw = Math.round(w * scale), ch = Math.round(h * scale);
                    var canvas = document.createElement("canvas");
                    canvas.width = cw; canvas.height = ch;
                    var ctx = canvas.getContext("2d");
                    ctx.drawImage(video, 0, 0, cw, ch);
                    URL.revokeObjectURL(vurl);
                    var dataUrl = canvas.toDataURL("image/jpeg", 0.6);
                    resolve(dataUrl);
                };
                video.onerror = function() { URL.revokeObjectURL(vurl); resolve(""); };
                video.src = vurl;
            } else {
                resolve("");
            }
        });
    };

    const _sendMediaFile = async function(file, mediaType) {
        if (!_state.currentUser) {
            _toast("请先选择用户");
            return;
        }
        if (!file) return;
        
        _closeMediaPanel();
        
        var mediaTypeInt = {"image": 2, "voice": 3, "file": 4, "video": 5}[mediaType] || 4;
        var mediaTypeLabel = {"image": "图片", "voice": "语音", "file": "文件", "video": "视频"}[mediaType] || "文件";
        var thumbDataUrl = "";
        
        if (mediaType === "image") {
            thumbDataUrl = await _generateThumbnail(file, 200, 200);
        } else if (mediaType === "video") {
            thumbDataUrl = await _generateThumbnail(file, 200, 200);
        }
        
        var placeholderMsg = {
            from: 'me',
            to: _state.currentUser,
            text: '[' + mediaTypeLabel + '] ' + file.name,
            time: new Date().toTimeString().slice(0, 8),
            type: 'out',
            media_type: mediaTypeInt,
            media_data: thumbDataUrl,
            media_filename: file.name,
            _sending: true
        };
        
        _state._tempMsgId = (_state._tempMsgId || 0) + 1;
        placeholderMsg.id = "sending_" + _state._tempMsgId;
        
        _renderSendingMsg(placeholderMsg);
        
        try {
            var base64Data = await _readFileAsBase64(file);
            var thumbnailData = "";
            
            if (mediaType === "image" || mediaType === "video") {
                try {
                    var fullThumb = await _generateThumbnail(file, 300, 300);
                    if (fullThumb) {
                        thumbnailData = fullThumb.split(",")[1] || "";
                    }
                } catch(e) {}
            }
            
            var payload = {
                media_type: mediaType,
                filename: file.name,
                file_data: base64Data,
                file_size: file.size,
                thumbnail: thumbnailData
            };
            
            var result = await _api("send-media", payload);
            
            var sendingEl = document.querySelector('[data-sending-id="' + placeholderMsg.id + '"]');
            
            if (result && result.success && result.message) {
                var msg = result.message;
                if (!msg.id) {
                    _state._tempMsgId = (_state._tempMsgId || 0) + 1;
                    msg.id = "temp_" + _state._tempMsgId;
                }
                if (sendingEl) sendingEl.remove();
                if (!_state.displayedIds.has(msg.id)) {
                    _renderMsg(msg);
                    _state.displayedIds.add(msg.id);
                }
            } else if (result && result.success) {
                if (sendingEl) sendingEl.remove();
                setTimeout(_fetchMessages, 300);
            } else {
                if (sendingEl) {
                    var statusEl = sendingEl.querySelector('.msg-send-status');
                    if (statusEl) {
                        statusEl.className = 'msg-send-status msg-send-fail';
                        statusEl.textContent = '!';
                    }
                }
                _toast((result && result.error) || "发送失败");
            }
        } catch(e) {
            var sendingEl2 = document.querySelector('[data-sending-id="' + placeholderMsg.id + '"]');
            if (sendingEl2) {
                var statusEl2 = sendingEl2.querySelector('.msg-send-status');
                if (statusEl2) {
                    statusEl2.className = 'msg-send-status msg-send-fail';
                    statusEl2.textContent = '!';
                }
            }
            _toast("发送失败: " + (e.message || e));
        }
    };
    
    const _handlePhotoSelect = function(e) {
        var file = e.target.files && e.target.files[0];
        if (file) _sendMediaFile(file, "image");
        e.target.value = "";
    };
    
    const _handleVideoSelect = function(e) {
        var file = e.target.files && e.target.files[0];
        if (file) _sendMediaFile(file, "video");
        e.target.value = "";
    };
    
    const _handleFileSelect = function(e) {
        var file = e.target.files && e.target.files[0];
        if (file) _sendMediaFile(file, "file");
        e.target.value = "";
    };
    
    const _loadAIConfig = async function() {
        const e = await _get("ai-config");
        if (e) {
            const ar = document.getElementById("ai-auto-reply");
            const sr = document.getElementById("ai-scheduled-reply");
            const n = document.getElementById("api-url");
            const o = document.getElementById("api-key");
            const i = document.getElementById("model-name");
            const r = document.getElementById("active-interval");
            const s = document.getElementById("min-words");
            const a = document.getElementById("max-words");
            const c = document.getElementById("system-prompt");
            const ve = document.getElementById("vision-enabled");
            const vu = document.getElementById("vision-api-url");
            const vk = document.getElementById("vision-api-key");
            const vm = document.getElementById("vision-model");
            const ige = document.getElementById("image-gen-enabled");
            const igu = document.getElementById("image-gen-api-url");
            const igk = document.getElementById("image-gen-api-key");
            const igm = document.getElementById("image-gen-model");
            const fre = document.getElementById("file-recognize-enabled");
            const fru = document.getElementById("file-recognize-api-url");
            const frk = document.getElementById("file-recognize-api-key");
            const frm = document.getElementById("file-recognize-model");
            const frs = document.getElementById("file-recognize-max-size");
            const frc = document.getElementById("file-recognize-compat-mode");
            if (ar) ar.checked = e.auto_reply || false;
            if (sr) sr.checked = e.scheduled_reply || false;
            const dr = document.getElementById("ai-daily-reply");
            if (dr) dr.checked = e.daily_reply || false;
            const dt = document.getElementById("daily-time");
            if (dt) dt.value = e.daily_time || "09:00";
            if (n) n.value = e.api_url || "";
            if (o) o.value = e.api_key || "";
            if (i) i.value = e.model || "gpt-3.5-turbo";
            if (r) r.value = e.active_interval || 60;
            if (s) s.value = e.min_words || 10;
            if (a) a.value = e.max_words || 200;
            if (c) c.value = e.system_prompt || "你是一个微信聊天助手，请用自然的中文回复，回复内容要简洁自然，像真人一样。";
            if (ve) ve.checked = e.vision_enabled || false;
            if (vu) vu.value = e.vision_api_url || "";
            if (vk) vk.value = e.vision_api_key || "";
            if (vm) vm.value = e.vision_model || "gpt-4o";
            if (ige) ige.checked = e.image_gen_enabled || false;
            if (igu) igu.value = e.image_gen_api_url || "";
            if (igk) igk.value = e.image_gen_api_key || "";
            if (igm) igm.value = e.image_gen_model || "dall-e-3";
            if (fre) fre.checked = e.file_recognize_enabled || false;
            if (fru) fru.value = e.file_recognize_api_url || "";
            if (frk) frk.value = e.file_recognize_api_key || "";
            if (frm) frm.value = e.file_recognize_model || "gpt-4o";
            if (frs) frs.value = e.file_recognize_max_size || 512;
            if (frc) frc.checked = e.file_recognize_compat_mode || false;
            const acd = document.getElementById("ai-cooldown");
            if (acd) acd.value = e.ai_cooldown !== undefined ? e.ai_cooldown : 5;
        }
    };
    
    const _saveAIConfig = async function() {
        const e = {
            auto_reply: document.getElementById("ai-auto-reply") ? document.getElementById("ai-auto-reply").checked : false,
            scheduled_reply: document.getElementById("ai-scheduled-reply") ? document.getElementById("ai-scheduled-reply").checked : false,
            api_url: document.getElementById("api-url") ? document.getElementById("api-url").value : "",
            api_key: document.getElementById("api-key") ? document.getElementById("api-key").value : "",
            model: document.getElementById("model-name") ? document.getElementById("model-name").value : "gpt-3.5-turbo",
            active_interval: parseInt(document.getElementById("active-interval") ? document.getElementById("active-interval").value : "60") || 60,
            min_words: parseInt(document.getElementById("min-words") ? document.getElementById("min-words").value : "10") || 10,
            max_words: parseInt(document.getElementById("max-words") ? document.getElementById("max-words").value : "200") || 200,
            system_prompt: document.getElementById("system-prompt") ? document.getElementById("system-prompt").value : "",
            vision_enabled: document.getElementById("vision-enabled") ? document.getElementById("vision-enabled").checked : false,
            vision_api_url: document.getElementById("vision-api-url") ? document.getElementById("vision-api-url").value : "",
            vision_api_key: document.getElementById("vision-api-key") ? document.getElementById("vision-api-key").value : "",
            vision_model: document.getElementById("vision-model") ? document.getElementById("vision-model").value : "gpt-4o",
            image_gen_enabled: document.getElementById("image-gen-enabled") ? document.getElementById("image-gen-enabled").checked : false,
            image_gen_api_url: document.getElementById("image-gen-api-url") ? document.getElementById("image-gen-api-url").value : "",
            image_gen_api_key: document.getElementById("image-gen-api-key") ? document.getElementById("image-gen-api-key").value : "",
            image_gen_model: document.getElementById("image-gen-model") ? document.getElementById("image-gen-model").value : "dall-e-3",
            file_recognize_enabled: document.getElementById("file-recognize-enabled") ? document.getElementById("file-recognize-enabled").checked : false,
            file_recognize_api_url: document.getElementById("file-recognize-api-url") ? document.getElementById("file-recognize-api-url").value : "",
            file_recognize_api_key: document.getElementById("file-recognize-api-key") ? document.getElementById("file-recognize-api-key").value : "",
            file_recognize_model: document.getElementById("file-recognize-model") ? document.getElementById("file-recognize-model").value : "gpt-4o",
            file_recognize_max_size: parseInt(document.getElementById("file-recognize-max-size") ? document.getElementById("file-recognize-max-size").value : "512") || 512,
            file_recognize_compat_mode: document.getElementById("file-recognize-compat-mode") ? document.getElementById("file-recognize-compat-mode").checked : false,
            ai_cooldown: parseInt(document.getElementById("ai-cooldown") ? document.getElementById("ai-cooldown").value : "5") || 0,
            daily_reply: document.getElementById("ai-daily-reply") ? document.getElementById("ai-daily-reply").checked : false,
            daily_time: document.getElementById("daily-time") ? document.getElementById("daily-time").value : "09:00"
        };
        const t = await _api("ai-config", e);
        if (t && t.success) {
            _toast("AI 配置已保存");
            _showSettingsPage('settings-main');
        } else {
            _toast("保存失败: " + ((t && t.error) || "未知错误"));
        }
    };
    
    const _openSettings = function() {
        const e = document.getElementById("settings-panel");
        if (e) {
            e.classList.add("show");
            _showSettingsPage('settings-main');
            _setActiveTab('tab-settings');
            _loadAccountSettings();
        }
    };
    
    const _closeSettings = function() {
        const e = document.getElementById("settings-panel");
        if (e) e.classList.remove("show");
        if (_state.view === 'list') {
            _setActiveTab('tab-list');
        } else if (_state.view === 'users') {
            _setActiveTab('tab-users');
        }
    };
    
    const _showSettingsPage = function(pageId) {
        var pages = document.querySelectorAll('.settings-page');
        pages.forEach(function(p) { p.classList.remove('active'); p.classList.remove('settings-page-slide'); });
        var target = document.getElementById(pageId);
        if (target) {
            target.classList.add('active');
            if (pageId !== 'settings-main') {
                target.classList.add('settings-page-slide');
            }
        }
        if (!_isDesktop()) {
            var tabbar = document.getElementById("bottom-tab-bar");
            if (pageId === 'settings-main') {
                if (tabbar) tabbar.classList.remove("hidden");
            } else {
                if (tabbar) tabbar.classList.add("hidden");
            }
        }
        if (pageId === 'settings-api') {
            _loadAIConfig();
        } else if (pageId === 'settings-about') {
            _loadAbout();
            var _avatarImg = document.querySelector(".about-logo-img");
            if (_avatarImg) {
                _avatarImg.classList.remove("spinning");
                void _avatarImg.offsetWidth;
                _avatarImg.classList.add("spinning");
            }
        } else if (pageId === 'settings-user-mgmt') {
            _closeSettings();
            _switchToUserList();
        } else if (pageId === 'settings-announcement') {
            _loadAnnouncementPage();
        } else if (pageId === 'settings-password') {
            _loadPasswordSettings();
        } else if (pageId === 'settings-account') {
            _loadAccountSettings();
        } else if (pageId === 'settings-admin') {
            _loadAdminPanel();
            if (_state.adminPollInterval) clearInterval(_state.adminPollInterval);
            _state.adminPollInterval = setInterval(_loadAdminPanel, 5000);
        } else {
            if (_state.adminPollInterval) { clearInterval(_state.adminPollInterval); _state.adminPollInterval = null; }
        }
    };

    const _renderUserMgmtList = function() {
        var pcContainer = document.getElementById("user-list-items");
        var mobileContainer = document.getElementById("user-mgmt-list");
        var targetContainer = pcContainer || mobileContainer;
        if (!targetContainer) return;
        if (!_state.users || _state.users.length === 0) {
            targetContainer.innerHTML = '<div class="chat-list-empty"><div class="chat-list-empty-icon">👥</div><div>暂无用户</div></div>';
            if (mobileContainer && mobileContainer !== targetContainer) {
                mobileContainer.innerHTML = '<div style="text-align:center;padding:40px 20px;color:var(--text-hint);">暂无用户</div>';
            }
            return;
        }
        var html = '';
        _state.users.forEach(function(userId) {
            var nickname = _state.nicknames[userId] || '';
            var displayName = nickname || userId;
            html += '<div class="chat-list-item" data-user-id="' + _escape(userId) + '">' +
                '<div class="chat-list-item-avatar">用户</div>' +
                '<div class="chat-list-item-content">' +
                '<div class="chat-list-item-name">' + _escape(displayName) + '</div>' +
                '<div class="chat-list-item-msg">' + _escape(userId) + '</div>' +
                '</div>' +
                '<button class="pc-user-delete-btn" data-user-id="' + _escape(userId) + '">删除</button>' +
                '</div>';
        });
        targetContainer.innerHTML = html;
        if (mobileContainer && mobileContainer !== targetContainer) {
            var mobileHtml = '';
            _state.users.forEach(function(userId) {
                var nickname = _state.nicknames[userId] || '';
                var displayName = nickname || userId;
                mobileHtml += '<div class="user-mgmt-item" data-user-id="' + _escape(userId) + '">' +
                    '<div class="user-mgmt-item-info">' +
                    '<div class="user-mgmt-item-name">' + _escape(displayName) + '</div>' +
                    '<div class="user-mgmt-item-id">' + _escape(userId) + '</div>' +
                    '</div>' +
                    '<button class="user-mgmt-delete-btn" data-user-id="' + _escape(userId) + '">删除</button>' +
                    '</div>';
            });
            mobileContainer.innerHTML = mobileHtml;
        }
        targetContainer.querySelectorAll('.pc-user-delete-btn').forEach(function(btn) {
            btn.addEventListener('click', function(e) {
                e.stopPropagation();
                var userId = btn.getAttribute('data-user-id');
                if (userId) _deleteUser(userId);
            });
        });
        targetContainer.querySelectorAll('.chat-list-item').forEach(function(item) {
            item.addEventListener('click', function() {
                var userId = item.getAttribute('data-user-id');
                if (userId) _openChat(userId);
            });
        });
    };

    const _loadAnnouncementPage = async function() {
        var el = document.getElementById('announcement-page-content');
        if (!el) return;
        el.textContent = '加载中...';
        try {
            var data = await _get("announcement");
            if (data && data.content) {
                el.textContent = data.content || '暂无公告';
            } else {
                el.textContent = '暂无公告';
            }
        } catch(e) {
            el.textContent = '加载失败';
        }
    };

    const _loadAbout = async function() {
    const authorEl = document.getElementById("about-author");
    const versionEl = document.getElementById("about-version");
    
    const nameEl = document.querySelector(".about-logo-name");
    if (nameEl) {
        const originalText = nameEl.innerText || "Zyn iLink ChatBox";
        const chars = originalText.split('');
        let html = '<div class="anim-stagger" style="display: flex; justify-content: center; gap: 2px; flex-wrap: wrap;">';
        for (let i = 0; i < chars.length; i++) {
            html += `<span style="opacity:0; transform:translateY(20px); transition:all 0.5s cubic-bezier(0.5, 1.5, 0.5, 1); transition-delay:${i * 0.05}s; display:inline-block;">${_escape(chars[i])}</span>`;
        }
        html += '</div>';
        nameEl.innerHTML = html;
        setTimeout(() => {
            const spans = nameEl.querySelectorAll('span');
            spans.forEach(span => {
                span.style.opacity = '1';
                span.style.transform = 'translateY(0)';
            });
        }, 50);
    }
    
    if (authorEl) authorEl.textContent = "加载中...";
    if (versionEl) versionEl.textContent = "加载中...";
    const e = await _get("about");
    if (e) {
        if (authorEl) authorEl.textContent = e.author || "未知";
        if (versionEl) versionEl.textContent = e.version || "未知";
    } else {
        if (authorEl) authorEl.textContent = "获取失败";
        if (versionEl) versionEl.textContent = "获取失败";
    }
};

    const _onAvatarClick = function() {
        const img = document.querySelector(".about-logo-img");
        if (!img) return;
        img.classList.remove("spinning");
        void img.offsetWidth;
        img.classList.add("spinning");
    };
    
    const _initTheme = function() {
        var saved = localStorage.getItem('theme');
        if (saved === 'dark') {
            document.documentElement.setAttribute('data-theme', 'dark');
            var btn = document.getElementById('theme-toggle-btn');
            if (btn) btn.classList.add('active');
        }
    };
    
    const _toggleTheme = function() {
        var btn = document.getElementById('theme-toggle-btn');
        var isDark = document.documentElement.getAttribute('data-theme') === 'dark';
        if (isDark) {
            document.documentElement.removeAttribute('data-theme');
            localStorage.setItem('theme', 'light');
            if (btn) btn.classList.remove('active');
        } else {
            document.documentElement.setAttribute('data-theme', 'dark');
            localStorage.setItem('theme', 'dark');
            if (btn) btn.classList.add('active');
        }
        _saveAppearance();
    };

    const _applyAppearance = function(cfg) {
        var root = document.documentElement;
        if (!cfg.animations) {
            root.classList.add('no-animations');
        } else {
            root.classList.remove('no-animations');
        }
        root.style.setProperty('--glass-opacity', '1');
        var speed = cfg.animSpeed || 1;
        root.style.setProperty('--anim-duration', speed);
    };

    const _loadAppearance = function() {
        try {
            var saved = localStorage.getItem('appearance');
            var cfg = saved ? JSON.parse(saved) : { animations: true, glassOpacity: 0.5, animSpeed: 1 };
        } catch(e) {
            var cfg = { animations: true, glassOpacity: 0.5, animSpeed: 1 };
        }
        var animCb = document.getElementById('appearance-animations');
        var speedSlider = document.getElementById('appearance-anim-speed');
        var speedVal = document.getElementById('anim-speed-val');
        if (animCb) animCb.checked = cfg.animations;
        if (speedSlider) speedSlider.value = cfg.animSpeed;
        if (speedVal) speedVal.textContent = cfg.animSpeed.toFixed(1) + 'x';
        _applyAppearance(cfg);
    };

    const _saveAppearance = function() {
        var animCb = document.getElementById('appearance-animations');
        var speedSlider = document.getElementById('appearance-anim-speed');
        var cfg = {
            animations: animCb ? animCb.checked : true,
            glassOpacity: 0.5,
            animSpeed: speedSlider ? parseFloat(speedSlider.value) : 1
        };
        localStorage.setItem('appearance', JSON.stringify(cfg));
        _applyAppearance(cfg);
    };
    
    const _checkStatus = async function() {
        const e = await _get("status");
        if (e && (e.logged_in && e.login_done || (e.users && e.users.length > 0))) {
            _showChat(e);
            return true;
        }
        _showChat(e || {users: []});
        return true;
    };
    
    const _showChat = function(e) {
        _state.users = e.users || [];
        _state.view = 'list';
        var adminItem = document.getElementById("settings-admin-item");
        if (adminItem) adminItem.style.display = e.is_admin ? "" : "none";
        const n = document.getElementById("chat-list-page");
        if (n) n.classList.add("active");
        if (!_isDesktop()) {
            var tabbar = document.getElementById("bottom-tab-bar");
            if (tabbar) tabbar.classList.remove("hidden");
        } else {
            var chatPage = document.getElementById("chat-page");
            if (chatPage) chatPage.classList.add("active");
            var messagesArea = document.getElementById("messages-area");
            if (messagesArea) messagesArea.innerHTML = '<div class="pc-chat-empty">选择一个聊天开始对话</div>';
            var inputArea = document.getElementById("input-area");
            if (inputArea) inputArea.style.display = 'none';
        }
        _renderChatList();
        _loadChatListPreviews();
        _startPoll();
    };
    
    const _initMobileViewport = function() {
        if (!window.visualViewport) return;
        var vv = window.visualViewport;
        var onResize = function() {
            var isKeyboardOpen = vv.height < window.innerHeight - 80;
            if (isKeyboardOpen) {
                document.body.classList.add('keyboard-open');
                var chatContainer = document.querySelector('.chat-container.active');
                if (chatContainer) {
                    chatContainer.style.height = vv.height + 'px';
                }
                var chatPage = document.getElementById('chat-page');
                if (chatPage && chatPage.classList.contains('active')) {
                    chatPage.style.height = vv.height + 'px';
                }
                var settingsPanel = document.getElementById('settings-panel');
                if (settingsPanel && settingsPanel.classList.contains('show')) {
                    settingsPanel.style.height = vv.height + 'px';
                }
                var inputArea = document.querySelector('.chat-container.active .input-area');
                if (inputArea) {
                    inputArea.style.position = 'sticky';
                    inputArea.style.bottom = '0';
                }
                var messagesArea = document.getElementById('messages-area');
                if (messagesArea) {
                    messagesArea.scrollTop = messagesArea.scrollHeight;
                }
                var activeInput = document.querySelector('input:focus, textarea:focus');
                if (activeInput) {
                    setTimeout(function() {
                        activeInput.scrollIntoView({ block: 'nearest', behavior: 'smooth' });
                    }, 80);
                }
            } else {
                document.body.classList.remove('keyboard-open');
                var chatContainers = document.querySelectorAll('.chat-container');
                chatContainers.forEach(function(c) { c.style.height = ''; });
                var chatP = document.getElementById('chat-page');
                if (chatP) chatP.style.height = '';
                var sp = document.getElementById('settings-panel');
                if (sp) sp.style.height = '';
                var inputAreas = document.querySelectorAll('.input-area');
                inputAreas.forEach(function(ia) {
                    ia.style.position = '';
                    ia.style.bottom = '';
                });
            }
        };
        vv.addEventListener('resize', onResize);
        vv.addEventListener('scroll', function() {
            if (vv.height < window.innerHeight - 80) {
                window.scrollTo(0, 0);
                document.documentElement.scrollTop = 0;
            }
        });
        window.addEventListener('orientationchange', function() {
            setTimeout(function() {
                document.body.classList.remove('keyboard-open');
                var chatContainers = document.querySelectorAll('.chat-container');
                chatContainers.forEach(function(c) { c.style.height = ''; });
                var chatP = document.getElementById('chat-page');
                if (chatP) chatP.style.height = '';
                var sp = document.getElementById('settings-panel');
                if (sp) sp.style.height = '';
                var inputAreas = document.querySelectorAll('.input-area');
                inputAreas.forEach(function(ia) {
                    ia.style.position = '';
                    ia.style.bottom = '';
                });
            }, 200);
        });
    };

    const _initEvents = function() {
        document.addEventListener("error", function(ev) {
            var img = ev.target;
            if (img.tagName === 'IMG' && img.classList.contains('bubble-media-img')) {
                var wrap = img.closest('.bubble-media-img-wrap');
                if (wrap && wrap.dataset.cdn && !wrap.classList.contains('bubble-media-loading') && img.src.indexOf('/api/wasm/media/') === -1) {
                    wrap.classList.add('bubble-media-loading');
                    wrap.innerHTML = '<div class="bubble-media-placeholder">' + _svgImage + '<span>图片</span></div>';
                    window._loadCdnMedia(wrap);
                }
            }
        }, true);
        document.addEventListener("click", function(ev) {
            var thumb = ev.target.closest("[data-action='play-video']");
            if (thumb) {
                var videoSrc = thumb.dataset.videoSrc;
                if (videoSrc) {
                    window._previewVideo(videoSrc);
                } else {
                    var imgEl = thumb.querySelector('img');
                    if (imgEl && imgEl.src && imgEl.src.indexOf('/api/wasm/media/') !== -1) {
                        window._previewVideo(imgEl.src);
                    }
                }
                return;
            }
        });
        const e = document.getElementById("send-btn");
        if (e) e.addEventListener("click", _sendMsg);
        var _typingTimer = null;
        const t = document.getElementById("message-input");
        if (t) {
            t.addEventListener("keypress", function(e) { if (e.key === "Enter") { _closeMediaPanel(); _sendMsg(); } });
            t.addEventListener("input", function() {
                if (!_state.currentUser) return;
                if (_typingTimer) clearTimeout(_typingTimer);
                _typingTimer = setTimeout(function() {
                    _api("send-typing", {});
                }, 500);
            });
            t.addEventListener("focus", function() { _closeMediaPanel(); setTimeout(function() { var ma = document.getElementById('messages-area'); if (ma) ma.scrollTop = ma.scrollHeight; }, 100); });
        }
        const plusBtn = document.getElementById("plus-btn");
        if (plusBtn) plusBtn.addEventListener("click", _toggleMediaPanel);
        const photoOpt = document.getElementById("media-photo");
        if (photoOpt) photoOpt.addEventListener("click", function() { document.getElementById("file-photo").click(); });
        const cameraOpt = document.getElementById("media-camera");
        if (cameraOpt) cameraOpt.addEventListener("click", function() { document.getElementById("file-camera").click(); });
        const videoOpt = document.getElementById("media-video");
        if (videoOpt) videoOpt.addEventListener("click", function() { document.getElementById("file-video").click(); });
        const fileOpt = document.getElementById("media-file");
        if (fileOpt) fileOpt.addEventListener("click", function() { document.getElementById("file-doc").click(); });
        const filePhoto = document.getElementById("file-photo");
        if (filePhoto) filePhoto.addEventListener("change", _handlePhotoSelect);
        const fileCamera = document.getElementById("file-camera");
        if (fileCamera) fileCamera.addEventListener("change", _handlePhotoSelect);
        const fileVideo = document.getElementById("file-video");
        if (fileVideo) fileVideo.addEventListener("change", _handleVideoSelect);
        const fileVideoCap = document.getElementById("file-video-capture");
        if (fileVideoCap) fileVideoCap.addEventListener("change", _handleVideoSelect);
        const fileDoc = document.getElementById("file-doc");
        if (fileDoc) fileDoc.addEventListener("change", _handleFileSelect);
        const n = document.getElementById("user-select-btn");
        if (n) n.addEventListener("click", function() { const e = document.getElementById("user-dropdown"); if (e) e.classList.toggle("show"); });
        const chatListSettingsBtn = document.getElementById("chat-list-settings-btn");
        if (chatListSettingsBtn) chatListSettingsBtn.addEventListener("click", _openSettings);
        const tabList = document.getElementById("tab-list");
        if (tabList) tabList.addEventListener("click", _switchToChatList);
        const tabUsers = document.getElementById("tab-users");
        if (tabUsers) tabUsers.addEventListener("click", _switchToUserList);
        const tabSettings = document.getElementById("tab-settings");
        if (tabSettings) tabSettings.addEventListener("click", _switchToSettings);
        const userListAddBtn = document.getElementById("user-list-add-btn");
        if (userListAddBtn) userListAddBtn.addEventListener("click", _startAddUser);
        const userListBackBtn = document.getElementById("user-list-back-btn");
        if (userListBackBtn) userListBackBtn.addEventListener("click", _switchToChatList);
        const pcNavChat = document.getElementById("pc-nav-chat");
        if (pcNavChat) pcNavChat.addEventListener("click", function() {
            var chatItems = document.getElementById("chat-list-items");
            var userItems = document.getElementById("user-list-items");
            var navChat = document.getElementById("pc-nav-chat");
            var navUsers = document.getElementById("pc-nav-users");
            if (chatItems) chatItems.style.display = '';
            if (userItems) userItems.style.display = 'none';
            if (navChat) navChat.classList.add('active');
            if (navUsers) navUsers.classList.remove('active');
            _renderChatList();
        });
        const pcNavUsers = document.getElementById("pc-nav-users");
        if (pcNavUsers) pcNavUsers.addEventListener("click", function() {
            var chatItems = document.getElementById("chat-list-items");
            var userItems = document.getElementById("user-list-items");
            var navChat = document.getElementById("pc-nav-chat");
            var navUsers = document.getElementById("pc-nav-users");
            if (chatItems) chatItems.style.display = 'none';
            if (userItems) userItems.style.display = '';
            if (navChat) navChat.classList.remove('active');
            if (navUsers) navUsers.classList.add('active');
            _renderUserMgmtList();
        });
        const addUserBtn = document.getElementById("chat-list-add-btn");
        if (addUserBtn) addUserBtn.addEventListener("click", _startAddUser);
        const addUserCloseBtn = document.getElementById("add-user-close-btn");
        if (addUserCloseBtn) addUserCloseBtn.addEventListener("click", _closeAddUserModal);
        const chatBackBtn = document.getElementById("chat-back-btn");
        if (chatBackBtn) chatBackBtn.addEventListener("click", _backToChatList);
        const chatMenuBtn = document.getElementById("chat-menu-btn");
        if (chatMenuBtn) chatMenuBtn.addEventListener("click", _toggleChatMenu);
        const chatMenuNickname = document.getElementById("chat-menu-nickname");
        if (chatMenuNickname) chatMenuNickname.addEventListener("click", function() { _closeChatMenu(); _openNicknameModal(); });
        const chatMenuPrompt = document.getElementById("chat-menu-prompt");
        if (chatMenuPrompt) chatMenuPrompt.addEventListener("click", function() { _closeChatMenu(); _openUserPromptModal(); });
        const chatMenuAiToggle = document.getElementById("chat-menu-ai-toggle");
        if (chatMenuAiToggle) chatMenuAiToggle.addEventListener("click", function(e) { e.stopPropagation(); _toggleUserAiEnabled(); });
        const chatMenuScheduledToggle = document.getElementById("chat-menu-scheduled-toggle");
        if (chatMenuScheduledToggle) chatMenuScheduledToggle.addEventListener("click", function(e) { e.stopPropagation(); _toggleUserScheduledEnabled(); });
        const chatMenuDailyToggle = document.getElementById("chat-menu-daily-toggle");
        if (chatMenuDailyToggle) chatMenuDailyToggle.addEventListener("click", function(e) { e.stopPropagation(); _toggleUserDailyEnabled(); });
        const nicknameCancelBtn = document.getElementById("nickname-cancel-btn");
        if (nicknameCancelBtn) nicknameCancelBtn.addEventListener("click", _closeNicknameModal);
        const nicknameSaveBtn = document.getElementById("nickname-save-btn");
        if (nicknameSaveBtn) nicknameSaveBtn.addEventListener("click", _saveNickname);
        const nicknameInput = document.getElementById("nickname-input");
        if (nicknameInput) nicknameInput.addEventListener("keypress", function(e) { if (e.key === "Enter") _saveNickname(); });
        const userPromptCancelBtn = document.getElementById("user-prompt-cancel-btn");
        if (userPromptCancelBtn) userPromptCancelBtn.addEventListener("click", _closeUserPromptModal);
        const userPromptSaveBtn = document.getElementById("user-prompt-save-btn");
        if (userPromptSaveBtn) userPromptSaveBtn.addEventListener("click", _saveUserPrompt);
        const modalClose = document.getElementById("ai-modal-close");
        if (modalClose) modalClose.addEventListener("click", _closeAiModal);
        const modalCancel = document.getElementById("ai-modal-cancel");
        if (modalCancel) modalCancel.addEventListener("click", _closeAiModal);
        const modalSend = document.getElementById("ai-modal-send");
        if (modalSend) modalSend.addEventListener("click", _sendAiReply);
        document.addEventListener("click", function(e) {
            const t = document.getElementById("user-dropdown");
            const n = document.getElementById("user-select-btn");
            const o = document.getElementById("settings-panel");
            const chatListSettingsBtn = document.getElementById("chat-list-settings-btn");
            const modal = document.getElementById("ai-modal");
            const mediaPanel = document.getElementById("media-panel");
            const plusBtn = document.getElementById("plus-btn");
            const nicknameModal = document.getElementById("nickname-modal");
            const chatMenuDropdown = document.getElementById("chat-menu-dropdown");
            const chatMenuBtn = document.getElementById("chat-menu-btn");
            const userPromptModal = document.getElementById("user-prompt-modal");
            if (t && !t.contains(e.target) && n && !n.contains(e.target)) {
                t.classList.remove("show");
            }
            if (chatMenuDropdown && !chatMenuDropdown.contains(e.target) && chatMenuBtn && !chatMenuBtn.contains(e.target)) {
                _closeChatMenu();
            }
            if (o && o.classList.contains("show") && !o.contains(e.target) && !document.getElementById("bottom-tab-bar").contains(e.target) && (!chatListSettingsBtn || !chatListSettingsBtn.contains(e.target))) {
                _closeSettings();
            }
            if (modal && modal.classList.contains("show") && e.target === modal) {
                _closeAiModal();
            }
            if (nicknameModal && nicknameModal.classList.contains("show") && e.target === nicknameModal) {
                _closeNicknameModal();
            }
            if (userPromptModal && userPromptModal.classList.contains("show") && e.target === userPromptModal) {
                _closeUserPromptModal();
            }
            var addUserModal = document.getElementById("add-user-modal");
            if (addUserModal && addUserModal.classList.contains("show") && e.target === addUserModal) {
                _closeAddUserModal();
            }
            if (mediaPanel && mediaPanel.classList.contains("show") && !mediaPanel.contains(e.target) && plusBtn && !plusBtn.contains(e.target)) {
                _closeMediaPanel();
            }
        });
        const s = document.getElementById("settings-back-btn");
        if (s) s.addEventListener("click", _closeSettings);
        const apiBackBtn = document.getElementById("api-back-btn");
        if (apiBackBtn) apiBackBtn.addEventListener("click", function() { _showSettingsPage('settings-main'); });
        const apiItem = document.getElementById("settings-api-item");
        if (apiItem) apiItem.addEventListener("click", function() { _showSettingsPage('settings-api'); });
        const aboutItem = document.getElementById("settings-about-item");
        if (aboutItem) aboutItem.addEventListener("click", function() { _showSettingsPage('settings-about'); });
        const userMgmtItem = document.getElementById("settings-user-mgmt-item");
        if (userMgmtItem) userMgmtItem.addEventListener("click", function() { _switchToUserList(); });
        const announcementItem = document.getElementById("settings-announcement-item");
        if (announcementItem) announcementItem.addEventListener("click", function() { _showSettingsPage('settings-announcement'); });
        const tutorialItem = document.getElementById("settings-tutorial-item");
        if (tutorialItem) tutorialItem.addEventListener("click", async function() {
            try {
                var data = await _get("tutorial");
                if (data && data.content) {
                    var w = window.open('', '_blank');
                    if (w) { w.document.open(); w.document.write(data.content); w.document.close(); }
                } else { _toast('暂无教程内容'); }
            } catch(e) { _toast('加载失败'); }
        });
        const aboutBackBtn = document.getElementById("about-back-btn");
        if (aboutBackBtn) aboutBackBtn.addEventListener("click", function() { _showSettingsPage('settings-main'); });
        const aboutLogoImg = document.querySelector(".about-logo-img");
        if (aboutLogoImg) {
            aboutLogoImg.addEventListener("click", _onAvatarClick);
            aboutLogoImg.addEventListener("animationend", function() { aboutLogoImg.classList.remove("spinning"); });
        }
        const themeBtn = document.getElementById("theme-toggle-btn");
        if (themeBtn) themeBtn.addEventListener("click", function(ev) { ev.stopPropagation(); _toggleTheme(); });
        const themeItem = document.getElementById("settings-theme-item");
        if (themeItem) themeItem.addEventListener("click", function() { _toggleTheme(); });
        const a = document.querySelector(".settings-save");
        if (a) a.addEventListener("click", _saveAIConfig);
        const appearanceItem = document.getElementById("settings-appearance-item");
        if (appearanceItem) appearanceItem.addEventListener("click", function() { _showSettingsPage('settings-appearance'); });
        const appearanceBackBtn = document.getElementById("appearance-back-btn");
        if (appearanceBackBtn) appearanceBackBtn.addEventListener("click", function() { _showSettingsPage('settings-main'); });
        const announcementBackBtn = document.getElementById("announcement-back-btn");
        if (announcementBackBtn) announcementBackBtn.addEventListener("click", function() { _showSettingsPage('settings-main'); });
        const animCb = document.getElementById("appearance-animations");
        if (animCb) animCb.addEventListener("change", _saveAppearance);
        const speedSlider = document.getElementById("appearance-anim-speed");
        if (speedSlider) speedSlider.addEventListener("input", function() {
            var val = parseFloat(this.value);
            var valEl = document.getElementById('anim-speed-val');
            if (valEl) valEl.textContent = val.toFixed(1) + 'x';
            _saveAppearance();
        });
    };
    
    const _checkAnnouncement = async function() {
        if (!_state.token) return;
        try {
            var ggData = await _get("announcement");
            if (!ggData) return;
            var isAdmin = ggData.is_admin || false;
            if (!isAdmin && ggData.system_announcements && ggData.system_announcements.length > 0) {
                var lastAnn = ggData.system_announcements[ggData.system_announcements.length - 1];
                var annId = lastAnn.id || lastAnn.time || '';
                _lastAnnId = annId;
                var dismissedKey = 'zyn_sys_announcement_dismissed';
                var dismissed = localStorage.getItem(dismissedKey);
                if (dismissed !== annId) {
                    var titleEl = document.getElementById('announcement-modal-title');
                    if (titleEl) titleEl.textContent = '系统公告';
                    var body = document.getElementById('announcement-body');
                    if (body) body.textContent = lastAnn.content;
                    var modal = document.getElementById('announcement-modal');
                    if (modal) modal.classList.add('show');
                    var dismissBtn = document.getElementById('announcement-dismiss-btn');
                    if (dismissBtn) dismissBtn.style.display = 'none';
                    return;
                }
            }
            var dismissBtn = document.getElementById('announcement-dismiss-btn');
            if (dismissBtn) dismissBtn.style.display = '';
            var versionData = await _get("remote-version");
            if (!versionData || !versionData.version) return;
            var remoteVersion = versionData.version;
            localStorage.setItem('zyn_remote_version', remoteVersion);
            var dismissedKey2 = 'zyn_announcement_dismissed';
            var dismissed2 = localStorage.getItem(dismissedKey2);
            if (dismissed2 === remoteVersion) return;
            if (!ggData || !ggData.content || !ggData.content.trim()) return;
            var body2 = document.getElementById('announcement-body');
            if (body2) body2.textContent = ggData.content;
            var modal2 = document.getElementById('announcement-modal');
            if (modal2) modal2.classList.add('show');
        } catch(e) {}
    };

    var _lastAnnId = null;
    const _initAnnouncement = function() {
        var confirmBtn = document.getElementById('announcement-confirm-btn');
        var dismissBtn = document.getElementById('announcement-dismiss-btn');
        if (confirmBtn) confirmBtn.addEventListener("click", function() {
            var modal = document.getElementById('announcement-modal');
            if (modal) modal.classList.remove('show');
            if (_lastAnnId) {
                localStorage.setItem('zyn_sys_announcement_dismissed', _lastAnnId);
            }
        });
        if (dismissBtn) dismissBtn.addEventListener("click", function() {
            try {
                var versionResp = localStorage.getItem('zyn_remote_version');
                if (versionResp) {
                    localStorage.setItem('zyn_announcement_dismissed', versionResp);
                }
            } catch(e) {}
            var modal = document.getElementById('announcement-modal');
            if (modal) modal.classList.remove('show');
        });
        setTimeout(_checkAnnouncement, 1500);
        setInterval(_checkAnnouncement, 30000);
    };

    const _generateFingerprint = function() {
        try {
            var canvas = document.createElement('canvas');
            canvas.width = 200; canvas.height = 50;
            var ctx = canvas.getContext('2d');
            ctx.textBaseline = 'top';
            ctx.font = '14px Arial';
            ctx.fillStyle = '#f60';
            ctx.fillRect(125, 1, 62, 20);
            ctx.fillStyle = '#069';
            ctx.fillText('ZynChatBox', 2, 15);
            ctx.fillStyle = 'rgba(102, 204, 0, 0.7)';
            ctx.fillText('ZynChatBox', 4, 17);
            var dataUrl = canvas.toDataURL();
            var fp = navigator.userAgent + '|' + screen.width + 'x' + screen.height + '|' + screen.colorDepth + '|' + new Date().getTimezoneOffset() + '|' + dataUrl;
            var hash = 0;
            for (var i = 0; i < fp.length; i++) {
                var char = fp.charCodeAt(i);
                hash = ((hash << 5) - hash) + char;
                hash = hash & hash;
            }
            return Math.abs(hash).toString(36) + fp.length.toString(36);
        } catch(e) {
            return navigator.userAgent.replace(/[^a-zA-Z0-9]/g, '').substring(0, 32);
        }
    };

    const _initLockScreen = async function() {
        var fingerprint = _generateFingerprint();
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/bot-fingerprint-login", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _state.token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({success: false}); }
                    } else { resolve({success: false}); }
                };
                o.onerror = function() { resolve({success: false}); };
                o.send(JSON.stringify({fingerprint: fingerprint, client_ip: _state.clientIP || ""}));
            });
            if (resp && resp.success && resp.session_token) {
                _state.token = resp.session_token;
                _afterAuth();
                return;
            }
        } catch(e) {}
        var lockScreen = document.getElementById("lock-screen");
        if (lockScreen) lockScreen.classList.remove("hide");
        var app = document.getElementById("app");
        if (app) app.style.display = "none";
    };

    const _afterAuth = function() {
        var lockScreen = document.getElementById("lock-screen");
        if (lockScreen) lockScreen.classList.add("hide");
        var app = document.getElementById("app");
        if (app) app.style.display = "";
        _connectSSE();
        _checkStatus();
    };

    const _initLockScreenEvents = function() {
        var loginForm = document.getElementById("bot-login-form");
        var registerForm = document.getElementById("bot-register-form");
        var showRegisterBtn = document.getElementById("bot-show-register");
        var showLoginBtn = document.getElementById("bot-show-login");
        var loginSubmitBtn = document.getElementById("bot-login-submit");
        var registerSubmitBtn = document.getElementById("bot-register-submit");
        var errorEl = document.getElementById("lock-screen-error");

        if (showRegisterBtn) showRegisterBtn.addEventListener("click", function() {
            if (loginForm) loginForm.style.display = "none";
            if (registerForm) registerForm.style.display = "";
            if (errorEl) errorEl.textContent = "";
        });
        if (showLoginBtn) showLoginBtn.addEventListener("click", function() {
            if (registerForm) registerForm.style.display = "none";
            if (loginForm) loginForm.style.display = "";
            if (errorEl) errorEl.textContent = "";
        });

        var loginUsernameInput = document.getElementById("bot-login-username");
        var loginPasswordInput = document.getElementById("bot-login-password");
        if (loginPasswordInput) loginPasswordInput.addEventListener("keypress", function(e) { if (e.key === "Enter") _doBotLogin(); });
        if (loginUsernameInput) loginUsernameInput.addEventListener("keypress", function(e) { if (e.key === "Enter") _doBotLogin(); });
        if (loginSubmitBtn) loginSubmitBtn.addEventListener("click", _doBotLogin);

        var regUsernameInput = document.getElementById("bot-register-username");
        var regPasswordInput = document.getElementById("bot-register-password");
        var regPassword2Input = document.getElementById("bot-register-password2");
        if (regPassword2Input) regPassword2Input.addEventListener("keypress", function(e) { if (e.key === "Enter") _doBotRegister(); });
        if (regPasswordInput) regPasswordInput.addEventListener("keypress", function(e) { if (e.key === "Enter") _doBotRegister(); });
        if (registerSubmitBtn) registerSubmitBtn.addEventListener("click", _doBotRegister);

        var sendEmailCodeBtn = document.getElementById("bot-register-send-email-code");
        if (sendEmailCodeBtn) sendEmailCodeBtn.addEventListener("click", async function() {
            var emailInput = document.getElementById("bot-register-email");
            var email = emailInput ? emailInput.value.trim() : "";
            if (!email) { if (errorEl) errorEl.textContent = "请输入邮箱"; return; }
            if (errorEl) errorEl.textContent = "";
            _showCaptchaModal(async function(captchaVal) {
                sendEmailCodeBtn.disabled = true;
                try {
                    var resp = await new Promise(function(resolve, reject) {
                        var o = new XMLHttpRequest();
                        o.open("POST", "/api/wasm/send-email-code", true);
                        o.setRequestHeader("Content-Type", "application/json");
                        o.setRequestHeader("X-Session-Token", _state.token);
                        o.onload = function() {
                            if (o.status >= 200 && o.status < 300) {
                                try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                            } else { resolve({success: false, error: "请求失败"}); }
                        };
                        o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                        o.send(JSON.stringify({email: email, captcha: captchaVal}));
                    });
                    if (resp && resp.success) {
                        if (errorEl) errorEl.textContent = "";
                        _toast("验证码已发送");
                        var _countdown = 60;
                        var _origText = sendEmailCodeBtn.textContent;
                        sendEmailCodeBtn.textContent = _countdown + "秒";
                        var _timer = setInterval(function() {
                            _countdown--;
                            if (_countdown <= 0) {
                                clearInterval(_timer);
                                sendEmailCodeBtn.disabled = false;
                                sendEmailCodeBtn.textContent = _origText;
                            } else {
                                sendEmailCodeBtn.textContent = _countdown + "秒";
                            }
                        }, 1000);
                    } else {
                        if (errorEl) errorEl.textContent = (resp && resp.error) || "发送失败";
                        sendEmailCodeBtn.disabled = false;
                    }
                } catch(e) {
                    if (errorEl) errorEl.textContent = "发送失败";
                    sendEmailCodeBtn.disabled = false;
                }
            });
        });

        (async function() {
            try {
                var resp = await new Promise(function(resolve, reject) {
                    var o = new XMLHttpRequest();
                    o.open("GET", "/api/wasm/email-register-status", true);
                    o.onload = function() {
                        if (o.status >= 200 && o.status < 300) {
                            try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                        } else { resolve({enabled: false}); }
                    };
                    o.onerror = function() { resolve({enabled: false}); };
                    o.send();
                });
                if (resp && resp.enabled) {
                    var emailInput = document.getElementById("bot-register-email");
                    var emailCodeRow = document.getElementById("bot-register-email-code-row");
                    if (emailInput) emailInput.style.display = "";
                    if (emailCodeRow) emailCodeRow.style.display = "";
                    _state._emailRegisterEnabled = true;
                } else {
                    _state._emailRegisterEnabled = false;
                }
            } catch(e) { _state._emailRegisterEnabled = false; }
        })();

        var forgotForm = document.getElementById("bot-forgot-form");
        var showForgotBtn = document.getElementById("bot-show-forgot");
        var forgotBackLoginBtn = document.getElementById("bot-forgot-back-login");
        var forgotSendCodeBtn = document.getElementById("bot-forgot-send-code");
        var forgotSubmitBtn = document.getElementById("bot-forgot-submit");
        var forgotCaptchaRefreshBtn = document.getElementById("bot-forgot-captcha-refresh");
        var forgotCaptchaImg = document.getElementById("bot-forgot-captcha-img");

        function _showForgotForm() {
            if (loginForm) loginForm.style.display = "none";
            if (registerForm) registerForm.style.display = "none";
            if (forgotForm) forgotForm.style.display = "";
            if (errorEl) errorEl.textContent = "";
            _loadForgotCaptcha();
        }

        function _loadForgotCaptcha() {
            if (forgotCaptchaImg) {
                forgotCaptchaImg.src = "/api/wasm/captcha?token=" + encodeURIComponent(_state.token) + "&t=" + Date.now();
            }
        }

        if (showForgotBtn) showForgotBtn.addEventListener("click", _showForgotForm);
        if (forgotBackLoginBtn) forgotBackLoginBtn.addEventListener("click", function() {
            if (forgotForm) forgotForm.style.display = "none";
            if (loginForm) loginForm.style.display = "";
            if (errorEl) errorEl.textContent = "";
        });
        if (forgotCaptchaRefreshBtn) forgotCaptchaRefreshBtn.addEventListener("click", _loadForgotCaptcha);
        if (forgotCaptchaImg) forgotCaptchaImg.addEventListener("click", _loadForgotCaptcha);

        if (forgotSendCodeBtn) forgotSendCodeBtn.addEventListener("click", async function() {
            var username = document.getElementById("bot-forgot-username");
            var captchaInput = document.getElementById("bot-forgot-captcha");
            var u = username ? username.value.trim() : "";
            var c = captchaInput ? captchaInput.value.trim() : "";
            if (!u) { errorEl.textContent = "请输入用户名"; return; }
            if (!c) { errorEl.textContent = "请输入图形验证码"; return; }
            forgotSendCodeBtn.disabled = true;
            errorEl.textContent = "";
            try {
                var resp = await new Promise(function(resolve, reject) {
                    var o = new XMLHttpRequest();
                    o.open("POST", "/api/wasm/bot-send-reset-code", true);
                    o.setRequestHeader("Content-Type", "application/json");
                    o.setRequestHeader("X-Session-Token", _state.token);
                    o.onload = function() {
                        if (o.status >= 200 && o.status < 300) {
                            try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                        } else { resolve({success: false, error: "请求失败"}); }
                    };
                    o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                    o.send(JSON.stringify({username: u, captcha: c}));
                });
                if (resp && resp.success) {
                    errorEl.textContent = "";
                    _toast("验证码已发送");
                    var _countdown = 60;
                    var _origText = forgotSendCodeBtn.textContent;
                    forgotSendCodeBtn.textContent = _countdown + "秒后重发";
                    var _timer = setInterval(function() {
                        _countdown--;
                        if (_countdown <= 0) {
                            clearInterval(_timer);
                            forgotSendCodeBtn.disabled = false;
                            forgotSendCodeBtn.textContent = _origText;
                        } else {
                            forgotSendCodeBtn.textContent = _countdown + "秒后重发";
                        }
                    }, 1000);
                } else {
                    errorEl.textContent = (resp && resp.error) || "发送失败";
                    _loadForgotCaptcha();
                    forgotSendCodeBtn.disabled = false;
                }
            } catch(e) {
                errorEl.textContent = "发送失败";
                _loadForgotCaptcha();
                forgotSendCodeBtn.disabled = false;
            }
        });

        if (forgotSubmitBtn) forgotSubmitBtn.addEventListener("click", async function() {
            var username = document.getElementById("bot-forgot-username");
            var code = document.getElementById("bot-forgot-code");
            var newPwd = document.getElementById("bot-forgot-new-password");
            var confirmPwd = document.getElementById("bot-forgot-confirm-password");
            var u = username ? username.value.trim() : "";
            var cd = code ? code.value.trim() : "";
            var np = newPwd ? newPwd.value : "";
            var cp = confirmPwd ? confirmPwd.value : "";
            if (!u) { errorEl.textContent = "请输入用户名"; return; }
            if (!cd) { errorEl.textContent = "请输入验证码"; return; }
            if (!np || np.length < 6) { errorEl.textContent = "新密码长度不能少于6位"; return; }
            if (np !== cp) { errorEl.textContent = "两次密码不一致"; return; }
            forgotSubmitBtn.disabled = true;
            errorEl.textContent = "";
            try {
                var resp = await new Promise(function(resolve, reject) {
                    var o = new XMLHttpRequest();
                    o.open("POST", "/api/wasm/bot-reset-password", true);
                    o.setRequestHeader("Content-Type", "application/json");
                    o.setRequestHeader("X-Session-Token", _state.token);
                    o.onload = function() {
                        if (o.status >= 200 && o.status < 300) {
                            try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                        } else { resolve({success: false, error: "请求失败"}); }
                    };
                    o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                    o.send(JSON.stringify({username: u, code: cd, new_password: np}));
                });
                if (resp && resp.success) {
                    errorEl.textContent = "";
                    _toast("密码重置成功，请重新登录");
                    if (forgotForm) forgotForm.style.display = "none";
                    if (loginForm) loginForm.style.display = "";
                } else {
                    errorEl.textContent = (resp && resp.error) || "重置失败";
                }
            } catch(e) {
                errorEl.textContent = "重置失败";
            }
            forgotSubmitBtn.disabled = false;
        });
    };

    const _doBotLogin = async function() {
        var usernameInput = document.getElementById("bot-login-username");
        var passwordInput = document.getElementById("bot-login-password");
        var submitBtn = document.getElementById("bot-login-submit");
        var errorEl = document.getElementById("lock-screen-error");
        var username = usernameInput ? usernameInput.value.trim() : "";
        var password = passwordInput ? passwordInput.value : "";
        if (!username) { errorEl.textContent = "请输入用户名"; return; }
        if (!password) { errorEl.textContent = "请输入密码"; return; }
        if (submitBtn) submitBtn.disabled = true;
        errorEl.textContent = "";
        var fingerprint = _generateFingerprint();
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/bot-login", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _state.token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                    } else { resolve({success: false, error: "请求失败"}); }
                };
                o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                o.send(JSON.stringify({username: username, password: password, fingerprint: fingerprint, client_ip: _state.clientIP || ""}));
            });
            if (resp && resp.success && resp.session_token) {
                _state.token = resp.session_token;
                _afterAuth();
            } else {
                errorEl.textContent = (resp && resp.error) || "登录失败";
                if (submitBtn) submitBtn.disabled = false;
            }
        } catch(e) {
            errorEl.textContent = "登录失败";
            if (submitBtn) submitBtn.disabled = false;
        }
    };

    const _doBotRegister = async function() {
        var usernameInput = document.getElementById("bot-register-username");
        var passwordInput = document.getElementById("bot-register-password");
        var password2Input = document.getElementById("bot-register-password2");
        var emailInput = document.getElementById("bot-register-email");
        var emailCodeInput = document.getElementById("bot-register-email-code");
        var submitBtn = document.getElementById("bot-register-submit");
        var errorEl = document.getElementById("lock-screen-error");
        var username = usernameInput ? usernameInput.value.trim() : "";
        var password = passwordInput ? passwordInput.value : "";
        var password2 = password2Input ? password2Input.value : "";
        var email = emailInput ? emailInput.value.trim() : "";
        var emailCode = emailCodeInput ? emailCodeInput.value.trim() : "";
        if (!username) { errorEl.textContent = "请输入用户名"; return; }
        if (_state._emailRegisterEnabled) {
            if (!email) { errorEl.textContent = "请输入邮箱"; return; }
            if (!emailCode) { errorEl.textContent = "请输入邮箱验证码"; return; }
        }
        if (!password) { errorEl.textContent = "请输入密码"; return; }
        if (password !== password2) { errorEl.textContent = "两次密码不一致"; return; }
        if (submitBtn) submitBtn.disabled = true;
        errorEl.textContent = "";
        var reqData = {username: username, password: password, client_ip: _state.clientIP || ""};
        if (_state._emailRegisterEnabled) {
            reqData.email = email;
            reqData.email_code = emailCode;
        }
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/bot-register", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _state.token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                    } else { resolve({success: false, error: "请求失败"}); }
                };
                o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                o.send(JSON.stringify(reqData));
            });
            if (resp && resp.success && resp.session_token) {
                _state.token = resp.session_token;
                _afterAuth();
            } else {
                errorEl.textContent = (resp && resp.error) || "注册失败";
                if (submitBtn) submitBtn.disabled = false;
            }
        } catch(e) {
            errorEl.textContent = "注册失败";
            if (submitBtn) submitBtn.disabled = false;
        }
    };

    var _passwordVerifyCallback = null;

    const _showPasswordVerify = function(callback) {
        _passwordVerifyCallback = callback;
        var overlay = document.getElementById("password-verify-overlay");
        var input = document.getElementById("password-verify-input");
        var errorEl = document.getElementById("password-verify-error");
        var confirmBtn = document.getElementById("password-verify-confirm");
        if (overlay) overlay.classList.add("show");
        if (input) { input.value = ""; input.focus(); }
        if (errorEl) errorEl.textContent = "";
        if (confirmBtn) confirmBtn.disabled = false;
    };

    const _initPasswordVerifyEvents = function() {
        var confirmBtn = document.getElementById("password-verify-confirm");
        var cancelBtn = document.getElementById("password-verify-cancel");
        var input = document.getElementById("password-verify-input");
        var errorEl = document.getElementById("password-verify-error");

        if (input) input.addEventListener("keypress", function(e) { if (e.key === "Enter") _doPasswordVerify(); });
        if (confirmBtn) confirmBtn.addEventListener("click", _doPasswordVerify);
        if (cancelBtn) cancelBtn.addEventListener("click", function() {
            var overlay = document.getElementById("password-verify-overlay");
            if (overlay) overlay.classList.remove("show");
            _passwordVerifyCallback = null;
        });
    };

    const _doPasswordVerify = async function() {
        var input = document.getElementById("password-verify-input");
        var errorEl = document.getElementById("password-verify-error");
        var confirmBtn = document.getElementById("password-verify-confirm");
        var password = input ? input.value : "";
        if (!password) { errorEl.textContent = "请输入密码"; return; }
        if (confirmBtn) confirmBtn.disabled = true;
        try {
            var resp = await _api("web-password-verify", {password: password});
            if (resp && resp.success) {
                var overlay = document.getElementById("password-verify-overlay");
                if (overlay) overlay.classList.remove("show");
                if (_passwordVerifyCallback) { _passwordVerifyCallback(); _passwordVerifyCallback = null; }
            } else {
                errorEl.textContent = (resp && resp.error) || "密码错误";
                if (confirmBtn) confirmBtn.disabled = false;
            }
        } catch(e) {
            errorEl.textContent = "验证失败";
            if (confirmBtn) confirmBtn.disabled = false;
        }
    };

    const _initPasswordSettings = function() {
        var passwordItem = document.getElementById("settings-password-item");
        if (passwordItem) passwordItem.addEventListener("click", async function() {
            try {
                var checkResp = await new Promise(function(resolve, reject) {
                    var o = new XMLHttpRequest();
                    o.open("GET", "/api/wasm/auth-check", true);
                    o.setRequestHeader("X-Session-Token", _state.token);
                    o.onload = function() {
                        if (o.status >= 200 && o.status < 300) {
                            try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                        } else { resolve({}); }
                    };
                    o.onerror = function() { resolve({}); };
                    o.send();
                });
                if (checkResp && checkResp.password_required) {
                    _showPasswordVerify(function() { var sp = document.getElementById("settings-panel"); if (sp) sp.classList.add("show"); _showSettingsPage("settings-password"); _loadPasswordSettings(); });
                } else {
                    _showSettingsPage("settings-password");
                    _loadPasswordSettings();
                }
            } catch(e) {
                _showSettingsPage("settings-password");
                _loadPasswordSettings();
            }
        });

        var passwordBackBtn = document.getElementById("password-back-btn");
        if (passwordBackBtn) passwordBackBtn.addEventListener("click", function() { _showSettingsPage("settings-main"); });

        var setBtn = document.getElementById("password-set-btn");
        if (setBtn) setBtn.addEventListener("click", async function() {
            var newPwd = document.getElementById("password-new-input").value;
            var confirmPwd = document.getElementById("password-confirm-input").value;
            if (!newPwd || newPwd.length < 6) { _toast("密码长度不能少于6位"); return; }
            if (newPwd !== confirmPwd) { _toast("两次密码不一致"); return; }
            try {
                var resp = await _api("set-web-password", {action: "set", password: newPwd});
                if (resp && resp.success) { _toast("密码设置成功"); _loadPasswordSettings(); }
                else { _toast((resp && resp.error) || "设置失败"); }
            } catch(e) { _toast("设置失败"); }
        });

        var changeBtn = document.getElementById("password-change-btn");
        if (changeBtn) changeBtn.addEventListener("click", function() {
            document.getElementById("password-already-set").style.display = "none";
            document.getElementById("password-change-form").style.display = "";
        });

        var changeCancelBtn = document.getElementById("password-change-cancel-btn");
        if (changeCancelBtn) changeCancelBtn.addEventListener("click", function() {
            document.getElementById("password-already-set").style.display = "";
            document.getElementById("password-change-form").style.display = "none";
        });

        var changeSaveBtn = document.getElementById("password-change-save-btn");
        if (changeSaveBtn) changeSaveBtn.addEventListener("click", async function() {
            var oldPwd = document.getElementById("password-old-input").value;
            var newPwd = document.getElementById("password-change-new-input").value;
            var confirmPwd = document.getElementById("password-change-confirm-input").value;
            if (!oldPwd) { _toast("请输入原密码"); return; }
            if (!newPwd || newPwd.length < 6) { _toast("新密码长度不能少于6位"); return; }
            if (newPwd !== confirmPwd) { _toast("两次密码不一致"); return; }
            try {
                var resp = await _api("set-web-password", {action: "change", old_password: oldPwd, new_password: newPwd});
                if (resp && resp.success) { _toast("密码修改成功"); _loadPasswordSettings(); }
                else { _toast((resp && resp.error) || "修改失败"); }
            } catch(e) { _toast("修改失败"); }
        });

        var removeBtn = document.getElementById("password-remove-btn");
        if (removeBtn) removeBtn.addEventListener("click", async function() {
            _showPasswordVerify(async function() {
                try {
                    var verifyInput = document.getElementById("password-verify-input");
                    var pwd = verifyInput ? verifyInput.value : "";
                    var resp = await _api("set-web-password", {action: "remove", password: pwd});
                    if (resp && resp.success) { _toast("密码已移除"); _loadPasswordSettings(); }
                    else { _toast((resp && resp.error) || "移除失败"); }
                } catch(e) { _toast("移除失败"); }
            });
        });
        var resetEmailSaveBtn = document.getElementById("reset-email-save-btn");
        if (resetEmailSaveBtn) resetEmailSaveBtn.addEventListener("click", async function() {
            var input = document.getElementById("reset-email-input");
            var email = input ? input.value.trim() : "";
            try {
                var resp = await _api("set-reset-email", {email: email});
                if (resp && resp.success) { _toast(email ? "验证码接收邮箱已设置" : "验证码接收邮箱已移除"); _loadPasswordSettings(); }
                else _toast((resp && resp.error) || "设置失败");
            } catch(e) { _toast("设置失败"); }
        });
        var resetEmailSaveBtnNoSet = document.getElementById("reset-email-save-btn-no-set");
        if (resetEmailSaveBtnNoSet) resetEmailSaveBtnNoSet.addEventListener("click", async function() {
            var input = document.getElementById("reset-email-input-no-set");
            var email = input ? input.value.trim() : "";
            try {
                var resp = await _api("set-reset-email", {email: email});
                if (resp && resp.success) { _toast(email ? "验证码接收邮箱已设置" : "验证码接收邮箱已移除"); _loadPasswordSettings(); }
                else _toast((resp && resp.error) || "设置失败");
            } catch(e) { _toast("设置失败"); }
        });
    };

    const _loadPasswordSettings = async function() {
        try {
            var resp = await _get("admin-user");
            var isSet = resp ? resp.password_set : false;
            var notSetDiv = document.getElementById("password-not-set");
            var alreadySetDiv = document.getElementById("password-already-set");
            var changeForm = document.getElementById("password-change-form");
            if (notSetDiv) notSetDiv.style.display = isSet ? "none" : "";
            if (alreadySetDiv) alreadySetDiv.style.display = isSet ? "" : "none";
            if (changeForm) changeForm.style.display = "none";
            var resetEmail = resp ? resp.reset_email || "" : "";
            var resetEmailInput = document.getElementById("reset-email-input");
            var resetEmailInputNoSet = document.getElementById("reset-email-input-no-set");
            if (resetEmailInput) resetEmailInput.value = resetEmail;
            if (resetEmailInputNoSet) resetEmailInputNoSet.value = resetEmail;
        } catch(e) {}
    };

    var _captchaModalCallback = null;
    var _showCaptchaModal = function(callback) {
        _captchaModalCallback = callback;
        var modal = document.getElementById("captcha-modal");
        var input = document.getElementById("captcha-modal-input");
        var img = document.getElementById("captcha-modal-img");
        if (input) input.value = "";
        if (img) img.src = "/api/wasm/captcha?token=" + encodeURIComponent(_state.token) + "&t=" + Date.now();
        if (modal) modal.style.display = "flex";
        if (input) input.focus();
    };
    var _hideCaptchaModal = function() {
        var modal = document.getElementById("captcha-modal");
        if (modal) modal.style.display = "none";
        _captchaModalCallback = null;
    };

    const _initAccountSettings = function() {
        var accountItem = document.getElementById("settings-account-item");
        if (accountItem) accountItem.addEventListener("click", function() {
            _showSettingsPage("settings-account");
            _loadAccountSettings();
        });

        var accountBackBtn = document.getElementById("account-back-btn");
        if (accountBackBtn) accountBackBtn.addEventListener("click", function() { _showSettingsPage("settings-main"); });

        var bindCaptchaImg = document.getElementById("account-bind-captcha-img");
        var bindCaptchaRefresh = document.getElementById("account-bind-captcha-refresh");
        function _loadBindCaptcha() {
            if (bindCaptchaImg) bindCaptchaImg.src = "/api/wasm/captcha?token=" + encodeURIComponent(_state.token) + "&t=" + Date.now();
        }
        if (bindCaptchaRefresh) bindCaptchaRefresh.addEventListener("click", _loadBindCaptcha);
        if (bindCaptchaImg) bindCaptchaImg.addEventListener("click", _loadBindCaptcha);

        var captchaModalImg = document.getElementById("captcha-modal-img");
        var captchaModalRefresh = document.getElementById("captcha-modal-refresh");
        var captchaModalCancel = document.getElementById("captcha-modal-cancel");
        var captchaModalConfirm = document.getElementById("captcha-modal-confirm");
        var captchaModalInput = document.getElementById("captcha-modal-input");
        if (captchaModalRefresh) captchaModalRefresh.addEventListener("click", function() {
            if (captchaModalImg) captchaModalImg.src = "/api/wasm/captcha?token=" + encodeURIComponent(_state.token) + "&t=" + Date.now();
        });
        if (captchaModalImg) captchaModalImg.addEventListener("click", function() {
            captchaModalImg.src = "/api/wasm/captcha?token=" + encodeURIComponent(_state.token) + "&t=" + Date.now();
        });
        if (captchaModalCancel) captchaModalCancel.addEventListener("click", _hideCaptchaModal);
        if (captchaModalConfirm) captchaModalConfirm.addEventListener("click", function() {
            var val = captchaModalInput ? captchaModalInput.value.trim() : "";
            if (!val) return;
            var cb = _captchaModalCallback;
            _hideCaptchaModal();
            if (cb) cb(val);
        });
        if (captchaModalInput) captchaModalInput.addEventListener("keypress", function(e) {
            if (e.key === "Enter") {
                var val = captchaModalInput.value.trim();
                if (!val) return;
                var cb = _captchaModalCallback;
                _hideCaptchaModal();
                if (cb) cb(val);
            }
        });

        var sendBindCodeBtn = document.getElementById("account-send-bind-code-btn");
        if (sendBindCodeBtn) sendBindCodeBtn.addEventListener("click", async function() {
            var emailInput = document.getElementById("account-bind-email");
            var captchaInput = document.getElementById("account-bind-captcha");
            var email = emailInput ? emailInput.value.trim() : "";
            var captcha = captchaInput ? captchaInput.value.trim() : "";
            if (!email) { _toast("请输入邮箱"); return; }
            if (!captcha) { _toast("请输入图形验证码"); return; }
            sendBindCodeBtn.disabled = true;
            try {
                var resp = await _api("send-bind-email-code", {email: email, captcha: captcha});
                if (resp && resp.success) {
                    _toast("验证码已发送");
                    var _countdown = 60;
                    var _origText = sendBindCodeBtn.textContent;
                    sendBindCodeBtn.textContent = _countdown + "秒";
                    var _timer = setInterval(function() {
                        _countdown--;
                        if (_countdown <= 0) {
                            clearInterval(_timer);
                            sendBindCodeBtn.disabled = false;
                            sendBindCodeBtn.textContent = _origText;
                        } else {
                            sendBindCodeBtn.textContent = _countdown + "秒";
                        }
                    }, 1000);
                } else {
                    _toast((resp && resp.error) || "发送失败");
                    sendBindCodeBtn.disabled = false;
                }
                _loadBindCaptcha();
                if (captchaInput) captchaInput.value = "";
            } catch(e) {
                _toast("发送失败");
                sendBindCodeBtn.disabled = false;
                _loadBindCaptcha();
            }
        });

        var bindEmailBtn = document.getElementById("account-bind-email-btn");
        if (bindEmailBtn) bindEmailBtn.addEventListener("click", async function() {
            var emailInput = document.getElementById("account-bind-email");
            var codeInput = document.getElementById("account-bind-code");
            var email = emailInput ? emailInput.value.trim() : "";
            var code = codeInput ? codeInput.value.trim() : "";
            if (!email) { _toast("请输入邮箱"); return; }
            if (!code) { _toast("请输入验证码"); return; }
            try {
                var resp = await _api("bind-email", {email: email, code: code});
                if (resp && resp.success) {
                    _toast("邮箱绑定成功");
                    _loadAccountSettings();
                } else {
                    _toast((resp && resp.error) || "绑定失败");
                }
            } catch(e) { _toast("绑定失败"); }
        });

        var changePasswordBtn = document.getElementById("account-change-password-btn");
        if (changePasswordBtn) changePasswordBtn.addEventListener("click", async function() {
            var oldPwd = document.getElementById("account-old-password").value;
            var newPwd = document.getElementById("account-new-password").value;
            var confirmPwd = document.getElementById("account-confirm-password").value;
            if (!oldPwd) { _toast("请输入原密码"); return; }
            if (!newPwd || newPwd.length < 6) { _toast("新密码长度不能少于6位"); return; }
            if (newPwd !== confirmPwd) { _toast("两次密码不一致"); return; }
            try {
                var resp = await _api("bot-change-password", {old_password: oldPwd, new_password: newPwd});
                if (resp && resp.success) {
                    _toast("密码修改成功，请重新登录");
                    window.location.href = '/';
                } else {
                    _toast((resp && resp.error) || "修改失败");
                }
            } catch(e) { _toast("修改失败"); }
        });

        var deleteAccountBtn = document.getElementById("account-delete-btn");
        if (deleteAccountBtn) deleteAccountBtn.addEventListener("click", async function() {
            var pwd = document.getElementById("account-delete-password").value;
            if (!pwd) { _toast("请输入密码以确认注销"); return; }
            if (!confirm("确定要注销账号吗？此操作不可恢复！")) return;
            try {
                var resp = await _api("bot-delete-account", {password: pwd});
                if (resp && resp.success) {
                    _toast("账号已注销");
                    window.location.href = '/';
                } else {
                    _toast((resp && resp.error) || "注销失败");
                }
            } catch(e) { _toast("注销失败"); }
        });

        var logoutBtn = document.getElementById("account-logout-btn");
        if (logoutBtn) logoutBtn.addEventListener("click", async function() {
            if (!confirm("确定要退出登录吗？")) return;
            try {
                await _api("bot-logout", {});
            } catch(e) {}
            _state.token = "";
            if (_state.pollInterval) { clearInterval(_state.pollInterval); _state.pollInterval = null; }
            var lockScreen = document.getElementById("lock-screen");
            if (lockScreen) lockScreen.classList.remove("hide");
            var app = document.getElementById("app");
            if (app) app.style.display = "none";
        });
    };

    const _loadAccountSettings = async function() {
        try {
            var resp = await _get("status");
            var username = resp && resp.account_username ? resp.account_username : "";
            var usernameEl = document.getElementById("account-current-username");
            if (usernameEl) usernameEl.textContent = username || "未知";
            var accountItem = document.getElementById("settings-account-item");
            if (accountItem) accountItem.style.display = username ? "" : "none";
        } catch(e) {}
        try {
            var resp = await _get("admin-user");
            var accountEmail = resp ? resp.account_email : "";
            var emailStatusEl = document.getElementById("account-email-status");
            var emailInput = document.getElementById("account-bind-email");
            if (emailStatusEl) emailStatusEl.textContent = accountEmail ? "已绑定: " + accountEmail : "未绑定邮箱";
            if (emailInput && accountEmail) emailInput.value = accountEmail;
        } catch(e) {}
        try {
            var bindCaptchaImg = document.getElementById("account-bind-captcha-img");
            if (bindCaptchaImg) bindCaptchaImg.src = "/api/wasm/captcha?token=" + encodeURIComponent(_state.token) + "&t=" + Date.now();
        } catch(e) {}
    };

    window._deleteAnnouncement = async function(idx) {
        if (!confirm("确定要删除该公告吗？")) return;
        try {
            var resp = await _api("admin-delete-announcement", {index: idx});
            if (resp && resp.success) { _toast("公告已删除"); _loadAdminPanel(); }
            else _toast((resp && resp.error) || "删除失败");
        } catch(e) { _toast("删除失败"); }
    };

    window._unbanIp = async function(ip) {
        if (!confirm("确定要解封IP " + ip + " 吗？")) return;
        try {
            var resp = await _api("admin-unban-ip", {ip: ip});
            if (resp && resp.success) { _toast("IP已解封"); _loadAdminPanel(); }
            else _toast((resp && resp.error) || "解封失败");
        } catch(e) { _toast("解封失败"); }
    };

    const _initAdminPanel = function() {
        var adminItem = document.getElementById("settings-admin-item");
        if (adminItem) adminItem.addEventListener("click", function() {
            _showSettingsPage("settings-admin");
            _loadAdminPanel();
        });
        var adminBackBtn = document.getElementById("admin-back-btn");
        if (adminBackBtn) adminBackBtn.addEventListener("click", function() { _showSettingsPage("settings-main"); });
        var deleteBtn = document.getElementById("admin-delete-user-btn");
        if (deleteBtn) deleteBtn.addEventListener("click", async function() {
            var index = document.getElementById("admin-delete-index").value.trim();
            if (!index) { _toast("请输入序号"); return; }
            if (!confirm("确定要注销序号 " + index + " 的用户吗？")) return;
            try {
                var resp = await _api("admin-delete-user", {index: parseInt(index)});
                if (resp && resp.success) { _toast("用户已注销"); document.getElementById("admin-delete-index").value = ""; _loadAdminPanel(); }
                else { _toast((resp && resp.error) || "操作失败"); }
            } catch(e) { _toast("操作失败"); }
        });
        var sendBtn = document.getElementById("admin-send-announcement-btn");
        if (sendBtn) sendBtn.addEventListener("click", async function() {
            var content = document.getElementById("admin-announcement-content").value.trim();
            if (!content) { _toast("请输入公告内容"); return; }
            try {
                var resp = await _api("admin-send-announcement", {content: content});
                if (resp && resp.success) { _toast("公告已发送"); document.getElementById("admin-announcement-content").value = ""; _loadAdminPanel(); }
                else { _toast((resp && resp.error) || "发送失败"); }
            } catch(e) { _toast("发送失败"); }
        });
        var forceOfflineBtn = document.getElementById("admin-force-offline-btn");
        if (forceOfflineBtn) forceOfflineBtn.addEventListener("click", async function() {
            var index = document.getElementById("admin-force-offline-index").value.trim();
            if (!index) { _toast("请输入序号"); return; }
            if (!confirm("确定要强制下线序号 " + index + " 的用户并删除指纹登录吗？")) return;
            try {
                var resp = await _api("admin-force-offline", {index: parseInt(index)});
                if (resp && resp.success) { _toast("用户已强制下线"); document.getElementById("admin-force-offline-index").value = ""; _loadAdminPanel(); }
                else { _toast((resp && resp.error) || "操作失败"); }
            } catch(e) { _toast("操作失败"); }
        });
        var banIpBtn = document.getElementById("admin-ban-ip-btn");
        if (banIpBtn) banIpBtn.addEventListener("click", async function() {
            var ip = document.getElementById("admin-ban-ip-input").value.trim();
            if (!ip) { _toast("请输入IP地址"); return; }
            if (!confirm("确定要封禁IP " + ip + " 吗？")) return;
            try {
                var resp = await _api("admin-ban-ip", {ip: ip});
                if (resp && resp.success) { _toast("IP已封禁"); document.getElementById("admin-ban-ip-input").value = ""; _loadAdminPanel(); }
                else { _toast((resp && resp.error) || "操作失败"); }
            } catch(e) { _toast("操作失败"); }
        });
        var emailSwitch = document.getElementById("admin-email-register-switch");
        if (emailSwitch) emailSwitch.addEventListener("click", async function() {
            var isOn = emailSwitch.classList.toggle("on");
            var configSection = document.getElementById("admin-email-config-section");
            if (configSection) configSection.style.display = isOn ? "" : "none";
            try {
                var resp = await _api("set-email-register", {enabled: isOn});
                if (resp && resp.success) { _toast(isOn ? "邮箱注册已启用" : "邮箱注册已关闭"); }
                else { emailSwitch.classList.toggle("on"); if (configSection) configSection.style.display = emailSwitch.classList.contains("on") ? "" : "none"; _toast((resp && resp.error) || "操作失败"); }
            } catch(e) { emailSwitch.classList.toggle("on"); if (configSection) configSection.style.display = emailSwitch.classList.contains("on") ? "" : "none"; _toast("操作失败"); }
        });
        var cfStartBtn = document.getElementById("cf-start-btn");
        if (cfStartBtn) cfStartBtn.addEventListener("click", async function() {
            try {
                var port = (document.getElementById("cf-port-input") || {}).value || 1145;
                var resp = await _api("admin-cloudflared", {action: "start", port: parseInt(port) || 1145});
                if (resp && resp.success) { _toast("Cloudflare Tunnel 已启动"); _loadAdminPanel(); }
                else { _toast((resp && resp.error) || "启动失败"); }
            } catch(e) { _toast("启动失败"); }
        });
        var cfStopBtn = document.getElementById("cf-stop-btn");
        if (cfStopBtn) cfStopBtn.addEventListener("click", async function() {
            try {
                var resp = await _api("admin-cloudflared", {action: "stop"});
                if (resp && resp.success) { _toast("Cloudflare Tunnel 已停止"); _loadAdminPanel(); }
                else { _toast((resp && resp.error) || "停止失败"); }
            } catch(e) { _toast("停止失败"); }
        });
    };

    const _loadAdminPanel = async function() {
        try {
            var resp = await _api("admin-list-users", {});
            if (resp && resp.success) {
                var listEl = document.getElementById("admin-panel-user-list");
                if (listEl) {
                    var users = resp.users || [];
                    if (users.length === 0) {
                        listEl.innerHTML = '<div style="font-size:13px;color:var(--text-secondary);padding:8px 0;">暂无注册用户</div>';
                    } else {
                        listEl.innerHTML = users.map(function(u, i) {
                            var statusColor = u.online ? 'var(--accent)' : 'var(--text-secondary)';
                            var statusText = u.online ? '在线' : '离线';
                            var ipText = u.last_ip ? _escape(u.last_ip) : '-';
                            return '<div style="display:flex;justify-content:space-between;align-items:center;padding:10px 0;border-bottom:1px solid var(--divider);"><div style="display:flex;align-items:center;gap:8px;"><span style="font-size:13px;color:var(--text-secondary);min-width:24px;">' + (i+1) + '.</span><div><div style="font-size:15px;font-weight:500;color:var(--text-primary);">' + _escape(u.username) + '</div><div style="font-size:11px;color:var(--text-secondary);margin-top:2px;">IP: ' + ipText + '</div></div></div><span style="font-size:12px;color:' + statusColor + ';">' + statusText + '</span></div>';
                        }).join("");
                    }
                }
                var annEl = document.getElementById("admin-announcement-list");
                if (annEl) {
                    var anns = resp.announcements || [];
                    if (anns.length === 0) {
                        annEl.innerHTML = '<div style="font-size:13px;color:var(--text-secondary);padding:8px 0;">暂无公告</div>';
                    } else {
                        annEl.innerHTML = anns.slice().reverse().map(function(a, ri) {
                            var idx = anns.length - 1 - ri;
                            return '<div style="padding:10px 0;border-bottom:1px solid var(--divider);display:flex;justify-content:space-between;align-items:flex-start;gap:8px;"><div style="flex:1;min-width:0;"><div style="font-size:13px;color:var(--text-primary);">' + _escape(a.content) + '</div><div style="font-size:11px;color:var(--text-secondary);margin-top:4px;">' + _escape(a.time || '') + '</div></div><button onclick="_deleteAnnouncement(' + idx + ')" style="flex-shrink:0;background:none;border:none;color:#FF3B30;font-size:13px;cursor:pointer;padding:4px 8px;">删除</button></div>';
                        }).join("");
                    }
                }
                var bannedEl = document.getElementById("admin-banned-ip-list");
                if (bannedEl) {
                    var bannedIps = resp.banned_ips || [];
                    if (bannedIps.length === 0) {
                        bannedEl.innerHTML = '<div style="font-size:13px;color:var(--text-secondary);padding:8px 0;">暂无封禁IP</div>';
                    } else {
                        bannedEl.innerHTML = bannedIps.map(function(ip) {
                            return '<div style="display:flex;justify-content:space-between;align-items:center;padding:10px 0;border-bottom:1px solid var(--divider);"><span style="font-size:14px;color:var(--text-primary);">' + _escape(ip) + '</span><button onclick="_unbanIp(&#39;' + _escape(ip) + '&#39;)" style="flex-shrink:0;background:none;border:none;color:#FF3B30;font-size:13px;cursor:pointer;padding:4px 8px;">解封</button></div>';
                        }).join("");
                    }
                }
            } else {
                _toast((resp && resp.error) || "无权限访问");
                _showSettingsPage("settings-main");
            }
            try {
                var adminResp = await _get("admin-user");
                if (adminResp) {
                    var emailSwitch = document.getElementById("admin-email-register-switch");
                    var configSection = document.getElementById("admin-email-config-section");
                    var isEmailEnabled = !!adminResp.email_register_enabled;
                    if (emailSwitch) {
                        var currentOn = emailSwitch.classList.contains("on");
                        if (currentOn !== isEmailEnabled) {
                            if (isEmailEnabled) emailSwitch.classList.add("on");
                            else emailSwitch.classList.remove("on");
                        }
                    }
                    if (configSection) configSection.style.display = emailSwitch && emailSwitch.classList.contains("on") ? "" : "none";
                }
            } catch(e) {}
            try {
                var cfResp = await _api("admin-cloudflared", {action: "status"});
                if (cfResp && cfResp.success) {
                    var dot = document.getElementById("cf-status-dot");
                    var text = document.getElementById("cf-status-text");
                    var urlEl = document.getElementById("cf-url-display");
                    if (cfResp.running) {
                        if (dot) dot.style.background = "var(--accent)";
                        if (text) text.textContent = "运行中";
                        if (urlEl && cfResp.url) { urlEl.textContent = "隧道地址: " + cfResp.url; var cb = document.getElementById("cf-copy-btn"); if (cb) cb.style.display = ""; }
                    } else {
                        if (dot) dot.style.background = "var(--text-secondary)";
                        if (text) text.textContent = "未运行";
                        if (urlEl) urlEl.textContent = "";
                        var cb = document.getElementById("cf-copy-btn"); if (cb) cb.style.display = "none";
                    }
                }
            } catch(e) {}
        } catch(e) {
            _toast("加载失败");
            _showSettingsPage("settings-main");
        }
    };

    const _init = function() {
        antiDebug();
        _initTheme();
        _loadAppearance();
        _initMobileViewport();
        _initEvents();
        _initLockScreenEvents();
        _initPasswordVerifyEvents();
        _initPasswordSettings();
        _initAccountSettings();
        _initAdminPanel();
        _initAnnouncement();
        _initLockScreen();
        window.addEventListener('resize', function() {
            if (_isDesktop()) {
                var chatListPage = document.getElementById("chat-list-page");
                if (chatListPage && _state.view !== 'users') chatListPage.classList.add("active");
                var chatPage = document.getElementById("chat-page");
                if (chatPage) chatPage.classList.add("active");
                var tabbar = document.getElementById("bottom-tab-bar");
                if (tabbar) tabbar.classList.add("hidden");
            } else {
                if (_state.view === 'list') {
                    var chatPage = document.getElementById("chat-page");
                    if (chatPage) chatPage.classList.remove("active");
                    var tabbar = document.getElementById("bottom-tab-bar");
                    if (tabbar) tabbar.classList.remove("hidden");
                }
            }
        });
    };
    
    return { init: _init };
})();

window.ZynWasm = window.__ZN''' + session_token[:16] + ''';
window._passwordRequired = ''' + ('true' if password_required else 'false') + ''';
window.ZynWasm.init();
'''
    
    def start_web_interface(self):
        if self._http_server is not None:
            print(f"网页服务已在运行中: http://localhost:{self._web_port}")
            return
        
        handler = self._make_web_handler()
        
        bind_addresses = ["0.0.0.0"]
        
        server_started = False
        original_port = self._web_port
        max_tries = 100
        
        for offset in range(max_tries):
            port = original_port + offset
            for bind_addr in bind_addresses:
                try:
                    self._http_server = socketserver.ThreadingTCPServer((bind_addr, port), handler)
                    self._server_thread = threading.Thread(target=self._http_server.serve_forever, daemon=True)
                    self._server_thread.start()
                    
                    self._web_port = port
                    local_ip = "127.0.0.1"
                    try:
                        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
                        s.connect(("8.8.8.8", 80))
                        local_ip = s.getsockname()[0]
                        s.close()
                    except Exception:
                        pass
                    if port != original_port:
                        print(f"\n[WEB] 端口 {original_port} 已被占用，自动切换到: http://127.0.0.1:{port}")
                    else:
                        print(f"\n[WEB] 网页界面已启动: http://127.0.0.1:{port}")
                    print(f"[WEB] 本地访问: http://127.0.0.1:{port}")
                    print(f"[WEB] 局域网访问: http://{local_ip}:{port}")
                    print("消息发送与二维码扫描请去本地网页操作! ")
                    server_started = True
                    break
                except OSError:
                    if bind_addr != "":
                        continue
            if server_started:
                break
        
        if not server_started:
            print(f"[ERROR] 端口 {original_port}-{original_port + max_tries - 1} 均已被占用，无法启动网页服务")
    
    def _make_web_handler(self):
        bot = self
        
        class WebHandler(SimpleHTTPRequestHandler):
            def log_message(self, format, *args):
                pass
            
            def _get_client_ip(self):
                xff = self.headers.get('X-Forwarded-For')
                if xff:
                    return xff.split(',')[0].strip()
                xri = self.headers.get('X-Real-IP')
                if xri:
                    return xri.strip()
                return self.client_address[0] if self.client_address else ''

            def _get_session_token(self):
                session_token = None
                if self.headers.get('X-Session-Token'):
                    t = self.headers.get('X-Session-Token')
                    if bot._verify_session_token(t):
                        session_token = t
                if not session_token:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                t = part.split('=', 1)[1]
                                if t and bot._verify_session_token(t):
                                    session_token = t
                                    break
                return session_token

            def _check_auth(self, require_password=True):
                session_token = self._get_session_token()
                if not session_token:
                    return False
                if session_token in bot._account_sessions:
                    return True
                if session_token in bot._verified_sessions and time.time() <= bot._verified_sessions[session_token]:
                    return True
                if session_token in bot._verified_sessions:
                    del bot._verified_sessions[session_token]
                return False
            
            def _resolve_account(self):
                session_token = self._get_session_token()
                if not session_token:
                    return None
                return bot._get_account_from_session(session_token)
            
            def _serve_banned_page(self):
                html = '<!DOCTYPE html><html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>已封禁</title><style>body{margin:0;display:flex;align-items:center;justify-content:center;min-height:100vh;background:#1a1a1a;font-family:system-ui,-apple-system,sans-serif}h1{color:#FF3B30;font-size:clamp(24px,6vw,48px);font-weight:700;text-align:center;padding:20px}</style></head><body><h1>你已被管理员封禁</h1></body></html>'
                encoded = html.encode('utf-8')
                self.send_response(403)
                self.send_header('Content-type', 'text/html; charset=utf-8')
                self.send_header('Content-Length', str(len(encoded)))
                self.end_headers()
                self.wfile.write(encoded)

            def do_GET(self):
                client_ip = self._get_client_ip()
                if client_ip and bot._is_ip_banned(client_ip):
                    self._serve_banned_page()
                    return
                parsed = urllib.parse.urlparse(self.path)
                
                if parsed.path == '/':
                    self._serve_wasm_page()
                elif parsed.path == '/qr':
                    self._serve_qr_page()
                elif parsed.path.startswith('/api/wasm/'):
                    api_path = parsed.path[10:]
                    if api_path == 'auth-check':
                        self._serve_auth_check()
                    elif api_path == 'captcha':
                        self._serve_captcha()
                    elif api_path == 'email-register-status':
                        self._handle_email_register_status()
                    elif api_path == 'events':
                        self._serve_sse()
                    elif not self._check_auth():
                        self.send_response(401)
                        self.end_headers()
                        return
                    elif api_path == 'status':
                        self._serve_status()
                    elif api_path == 'qrcode':
                        self._serve_qrcode()
                    elif api_path == 'messages':
                        self._serve_messages()
                    elif api_path == 'users':
                        self._serve_users()
                    elif api_path == 'history':
                        self._serve_history()
                    elif api_path == 'ai-config':
                        self._serve_ai_config()
                    elif api_path == 'about':
                        self._serve_about()
                    elif api_path == 'announcement':
                        self._serve_announcement()
                    elif api_path == 'remote-version':
                        self._serve_remote_version()
                    elif api_path == 'add-user-status':
                        self._serve_add_user_status()
                    elif api_path == 'user-prompt':
                        self._serve_user_prompt(parsed)
                    elif api_path == 'tutorial':
                        self._serve_tutorial()
                    elif api_path == 'admin-user':
                        self._serve_admin_user()
                    elif api_path.startswith('media/'):
                        self._serve_cached_media(api_path[6:])
                    else:
                        self.send_error(404)
                else:
                    self.send_error(404)
            
            def do_POST(self):
                client_ip = self._get_client_ip()
                if client_ip and bot._is_ip_banned(client_ip):
                    self._serve_banned_page()
                    return
                parsed = urllib.parse.urlparse(self.path)
                
                if parsed.path == '/api/wasm/bot-register':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_bot_register(data)
                    return
                elif parsed.path == '/api/wasm/bot-login':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_bot_login(data)
                    return
                elif parsed.path == '/api/wasm/bot-fingerprint-login':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_bot_fingerprint_login(data)
                    return
                elif parsed.path == '/api/wasm/auth-login':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_auth_login(data)
                    return
                elif parsed.path == '/api/wasm/send-verify-code':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_send_verify_code(data)
                    return
                elif parsed.path == '/api/wasm/reset-password':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_reset_password(data)
                    return
                elif parsed.path == '/api/wasm/bot-send-reset-code':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_bot_send_reset_code(data)
                    return
                elif parsed.path == '/api/wasm/bot-reset-password':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_bot_reset_password(data)
                    return
                elif parsed.path == '/api/wasm/send-email-code':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_send_email_code(data)
                    return
                elif parsed.path == '/api/wasm/send-bind-email-code':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_send_bind_email_code(data)
                    return
                elif parsed.path == '/api/wasm/bind-email':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_bind_email(data)
                    return
                elif parsed.path == '/api/wasm/account-email':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_account_email(data)
                    return
                
                if not self._check_auth():
                    self.send_response(401)
                    self.end_headers()
                    return
                
                data = self._parse_json_body()
                if data is None:
                    return
                
                if parsed.path == '/api/wasm/send':
                    self._handle_send(data)
                elif parsed.path == '/api/wasm/send-typing':
                    self._handle_send_typing(data)
                elif parsed.path == '/api/wasm/bot-change-password':
                    self._handle_bot_change_password(data)
                elif parsed.path == '/api/wasm/bot-delete-account':
                    self._handle_bot_delete_account(data)
                elif parsed.path == '/api/wasm/bot-logout':
                    self._handle_bot_logout(data)
                elif parsed.path == '/api/wasm/admin-list-users':
                    self._handle_admin_list_users(data)
                elif parsed.path == '/api/wasm/admin-delete-user':
                    self._handle_admin_delete_user(data)
                elif parsed.path == '/api/wasm/admin-send-announcement':
                    self._handle_admin_send_announcement(data)
                elif parsed.path == '/api/wasm/admin-delete-announcement':
                    self._handle_admin_delete_announcement(data)
                elif parsed.path == '/api/wasm/admin-force-offline':
                    self._handle_admin_force_offline(data)
                elif parsed.path == '/api/wasm/admin-ban-ip':
                    self._handle_admin_ban_ip(data)
                elif parsed.path == '/api/wasm/admin-unban-ip':
                    self._handle_admin_unban_ip(data)
                elif parsed.path == '/api/wasm/admin-cloudflared':
                    self._handle_admin_cloudflared(data)
                elif parsed.path == '/api/wasm/send-media':
                    self._handle_send_media(data)
                elif parsed.path == '/api/wasm/download-media':
                    self._handle_download_media(data)
                elif parsed.path == '/api/wasm/switch-user':
                    self._handle_switch_user(data)
                elif parsed.path == '/api/wasm/ai-config':
                    self._handle_save_ai_config(data)
                elif parsed.path == '/api/wasm/ai-manual-reply':
                    self._handle_ai_manual_reply(data)
                elif parsed.path == '/api/wasm/vision-recognize':
                    self._handle_vision_recognize(data)
                elif parsed.path == '/api/wasm/image-gen':
                    self._handle_image_gen(data)
                elif parsed.path == '/api/wasm/file-recognize':
                    self._handle_file_recognize(data)
                elif parsed.path == '/api/wasm/user-prompt':
                    self._handle_save_user_prompt(data)
                elif parsed.path == '/api/wasm/add-user-start':
                    self._handle_add_user_start(data)
                elif parsed.path == '/api/wasm/delete-user':
                    self._handle_delete_user(data)
                elif parsed.path == '/api/wasm/web-password-verify':
                    self._handle_web_password_verify(data)
                elif parsed.path == '/api/wasm/set-web-password':
                    self._handle_set_web_password(data)
                elif parsed.path == '/api/wasm/set-admin-user':
                    self._handle_set_admin_user(data)
                elif parsed.path == '/api/wasm/set-captcha-admin':
                    self._handle_set_captcha_admin(data)
                elif parsed.path == '/api/wasm/set-reset-email':
                    data = self._parse_json_body()
                    if data is None:
                        return
                    self._handle_set_reset_email(data)
                elif parsed.path == '/api/wasm/set-email-register':
                    self._handle_set_email_register(data)
                else:
                    self.send_error(404)
            
            def _parse_json_body(self):
                content_length = int(self.headers.get('Content-Length', 0))
                if content_length > 10485760:
                    self._send_json({'success': False, 'error': '请求数据过大'}, 413)
                    return None
                body = self.rfile.read(content_length) if content_length else b'{}'
                try:
                    return json.loads(body.decode('utf-8'))
                except (json.JSONDecodeError, UnicodeDecodeError):
                    self._send_json({'success': False, 'error': '请求数据格式错误'}, 400)
                    return None
            
            def _handle_ai_manual_reply(self, data):
                try:
                    user_id = data.get('user_id')
                    original_message = str(data.get('original_message', ''))
                    instruction = str(data.get('instruction', ''))
                    account = self._resolve_account()
                    target = account if account else bot
                    if not user_id:
                        self._send_json({'success': False, 'error': '用户ID不能为空'})
                        return
                    if user_id not in target._context_tokens:
                        self._send_json({'success': False, 'error': '用户不存在'})
                        return
                    print(f"[WEB] 收到手动 AI 回复请求: user={user_id}, msg={original_message[:50]}, instruction={instruction[:50] if instruction else '无'}")
                    self._send_json({'success': True, 'message': 'AI 回复正在生成中'})
                    def _do_reply():
                        try:
                            if account:
                                bot._auto_ai_reply_for_account(account, user_id, original_message)
                            else:
                                bot._manual_ai_reply(user_id, original_message, instruction)
                        except Exception as e:
                            print(f"[WEB] 手动 AI 回复异步异常: {e}")
                    threading.Thread(target=_do_reply, daemon=True).start()
                except Exception as e:
                    print(f"[WEB] 手动 AI 回复异常: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _serve_user_prompt(self, parsed):
                try:
                    params = urllib.parse.parse_qs(parsed.query)
                    user_id = params.get('user_id', [''])[0]
                    account = self._resolve_account()
                    target = account if account else bot
                    if not user_id:
                        self._send_json({'success': False, 'error': '用户ID不能为空'})
                        return
                    prompts = target.user_prompts.get("prompts", {})
                    ai_enabled = target.user_prompts.get("ai_enabled", {})
                    scheduled_enabled = target.user_prompts.get("scheduled_enabled", {})
                    daily_enabled = target.user_prompts.get("daily_enabled", {})
                    prompt = prompts.get(user_id, "")
                    default_prompt = target.ai_config.get("system_prompt", "")
                    if not default_prompt:
                        default_prompt = "你是一个微信聊天助手，请用自然的中文回复。"
                    ai_enabled_for_user = ai_enabled.get(user_id, None)
                    scheduled_enabled_for_user = scheduled_enabled.get(user_id, None)
                    daily_enabled_for_user = daily_enabled.get(user_id, None)
                    global_auto_reply = target.ai_config.get("auto_reply", False)
                    global_scheduled_reply = target.ai_config.get("scheduled_reply", False)
                    global_daily_reply = target.ai_config.get("daily_reply", False)
                    self._send_json({'success': True, 'prompt': prompt, 'default_prompt': default_prompt, 'ai_enabled': ai_enabled_for_user, 'scheduled_enabled': scheduled_enabled_for_user, 'daily_enabled': daily_enabled_for_user, 'global_auto_reply': global_auto_reply, 'global_scheduled_reply': global_scheduled_reply, 'global_daily_reply': global_daily_reply})
                except Exception as e:
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _handle_save_user_prompt(self, data):
                try:
                    user_id = str(data.get('user_id', ''))
                    prompt = str(data.get('prompt', ''))
                    ai_enabled = data.get('ai_enabled')
                    scheduled_enabled = data.get('scheduled_enabled')
                    daily_enabled = data.get('daily_enabled')
                    account = self._resolve_account()
                    target = account if account else bot
                    if not user_id:
                        self._send_json({'success': False, 'error': '用户ID不能为空'})
                        return
                    if len(prompt) > 10000:
                        self._send_json({'success': False, 'error': '提示词过长'})
                        return
                    prompts = target.user_prompts.setdefault("prompts", {})
                    if prompt.strip():
                        prompts[user_id] = prompt
                    else:
                        prompts.pop(user_id, None)
                    if ai_enabled is not None:
                        ai_enabled_map = target.user_prompts.setdefault("ai_enabled", {})
                        ai_enabled_map[user_id] = ai_enabled
                    if scheduled_enabled is not None:
                        scheduled_enabled_map = target.user_prompts.setdefault("scheduled_enabled", {})
                        scheduled_enabled_map[user_id] = scheduled_enabled
                        if scheduled_enabled:
                            if account:
                                bot._schedule_active_message_for_account(account, user_id)
                            else:
                                bot._schedule_active_message(user_id)
                        else:
                            timer = target._active_timers.pop(user_id, None)
                            if timer:
                                timer.cancel()
                    if daily_enabled is not None:
                        daily_enabled_map = target.user_prompts.setdefault("daily_enabled", {})
                        daily_enabled_map[user_id] = daily_enabled
                        if daily_enabled:
                            if account:
                                bot._schedule_daily_message_for_account(account, user_id)
                            else:
                                bot._schedule_daily_message(user_id)
                        else:
                            timer = target._daily_timers.pop(user_id, None)
                            if timer:
                                timer.cancel()
                    target._save_user_prompts()
                    self._send_json({'success': True})
                except Exception as e:
                    print(f"[WEB] 保存用户提示词异常: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _handle_vision_recognize(self, data):
                try:
                    image_base64 = data.get('image_base64', '')
                    prompt = data.get('prompt', '')
                    user_id = data.get('user_id', '')
                    account = self._resolve_account()
                    target = account if account else bot
                    if not image_base64:
                        self._send_json({'success': False, 'error': '图片数据不能为空'})
                        return
                    if len(image_base64) > 20 * 1024 * 1024:
                        self._send_json({'success': False, 'error': '图片数据过大'})
                        return
                    if not target.ai_config.get("vision_api_url"):
                        self._send_json({'success': False, 'error': '识图API未配置'})
                        return
                    print(f"[WEB] 收到识图请求: user={user_id}, prompt={prompt[:50] if prompt else '默认'}")
                    system_prompt = target.get_effective_system_prompt(user_id)
                    history = target.get_user_messages(user_id, 200) if user_id else []
                    media_memory_text = bot._format_media_memory_for_account(target, user_id) if user_id else ""
                    original_text = prompt if prompt else ""
                    if account:
                        result = bot._call_vision_api_for_account(account, image_base64, system_prompt, history, original_text=original_text, media_memory_text=media_memory_text)
                    else:
                        result = bot._call_vision_api(image_base64, system_prompt, history, original_text=original_text, media_memory_text=media_memory_text)
                    if result:
                        if user_id:
                            target._save_media_memory(user_id, "image", result)
                            if account:
                                bot._send_ai_reply_in_segments_for_account(account, user_id, result)
                            else:
                                bot._send_ai_reply_in_segments(user_id, result)
                        self._send_json({'success': True, 'vision_result': result})
                    else:
                        self._send_json({'success': False, 'error': '识图失败'})
                except Exception as e:
                    print(f"[WEB] 识图异常: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _handle_image_gen(self, data):
                try:
                    prompt = data.get('prompt', '')
                    user_id = data.get('user_id', '')
                    send_directly = data.get('send_directly', False)
                    account = self._resolve_account()
                    target = account if account else bot
                    if not prompt:
                        self._send_json({'success': False, 'error': '生图提示词不能为空'})
                        return
                    if len(prompt) > 4000:
                        self._send_json({'success': False, 'error': '生图提示词过长'})
                        return
                    if not target.ai_config.get("image_gen_api_url"):
                        self._send_json({'success': False, 'error': '生图API未配置'})
                        return
                    print(f"[WEB] 收到生图请求: user={user_id}, prompt={prompt[:50]}")
                    image_bytes = bot._call_image_gen_api_for_account(target, prompt)
                    if image_bytes:
                        if send_directly and user_id:
                            filename = f"ai_gen_{uuid.uuid4().hex[:8]}.png"
                            if account:
                                success = bot.send_image_for_account(account, user_id, image_bytes, filename=filename)
                            else:
                                success = bot.send_image(user_id, image_bytes, filename=filename)
                            if success:
                                self._send_json({'success': True, 'message': '图片已发送', 'image_size': len(image_bytes)})
                            else:
                                self._send_json({'success': False, 'error': '图片生成成功但发送失败'})
                        else:
                            img_b64 = base64.b64encode(image_bytes).decode('utf-8')
                            self._send_json({'success': True, 'image_base64': img_b64, 'image_size': len(image_bytes)})
                    else:
                        self._send_json({'success': False, 'error': '生图失败'})
                except Exception as e:
                    print(f"[WEB] 生图异常: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _handle_file_recognize(self, data):
                try:
                    file_base64 = data.get('file_base64', '')
                    filename = data.get('filename', '')
                    prompt = data.get('prompt', '')
                    user_id = data.get('user_id', '')
                    account = self._resolve_account()
                    target = account if account else bot
                    if not file_base64:
                        self._send_json({'success': False, 'error': '文件数据不能为空'})
                        return
                    if len(file_base64) > 20 * 1024 * 1024:
                        self._send_json({'success': False, 'error': '文件数据过大'})
                        return
                    if not filename:
                        self._send_json({'success': False, 'error': '文件名不能为空'})
                        return
                    if not target.ai_config.get("file_recognize_enabled"):
                        self._send_json({'success': False, 'error': '文件识别功能未启用'})
                        return
                    print(f"[WEB] 收到文件识别请求: user={user_id}, filename={filename}")
                    try:
                        file_bytes = base64.b64decode(file_base64)
                    except Exception:
                        self._send_json({'success': False, 'error': '文件数据解码失败'})
                        return
                    file_text = bot._extract_text_from_file(file_bytes, filename)
                    if not file_text or not file_text.strip():
                        self._send_json({'success': False, 'error': '无法提取文件文本内容'})
                        return
                    system_prompt = target.get_effective_system_prompt(user_id)
                    history = target.get_user_messages(user_id, 200) if user_id else []
                    media_memory_text = bot._format_media_memory_for_account(target, user_id) if user_id else ""
                    original_text = prompt if prompt else ""
                    if target.ai_config.get("file_recognize_compat_mode"):
                        max_size = target.ai_config.get("file_recognize_max_size", 512) * 1024
                        truncated = file_text[:max_size]
                        formatted_text = f"[用户发送了文件: {filename}]\n文件内容如下:\n{truncated}"
                        if prompt:
                            formatted_text += f"\n额外要求: {prompt}"
                        if target.ai_config.get("api_url"):
                            if account:
                                response = bot._call_ai_api_for_account(account, formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=user_id)
                            else:
                                response = bot._call_ai_api(formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=user_id)
                            if response:
                                if user_id:
                                    target._save_media_memory(user_id, "file", response, filename=filename)
                                    if account:
                                        bot._send_ai_reply_in_segments_for_account(account, user_id, response)
                                    else:
                                        bot._send_ai_reply_in_segments(user_id, response)
                                self._send_json({'success': True, 'ai_reply': response, 'mode': 'compat'})
                            else:
                                if user_id:
                                    if account:
                                        bot.send_text_for_account(account, user_id, f"[文件内容摘要] {file_text[:500]}")
                                    else:
                                        bot.send_text(user_id, f"[文件内容摘要] {file_text[:500]}")
                                self._send_json({'success': True, 'extracted_text': file_text[:3000], 'warning': '兼容模式AI回复失败，已返回提取文本', 'mode': 'compat'})
                        else:
                            if user_id:
                                if account:
                                    bot.send_text_for_account(account, user_id, f"[文件内容摘要] {file_text[:500]}")
                                else:
                                    bot.send_text(user_id, f"[文件内容摘要] {file_text[:500]}")
                            self._send_json({'success': True, 'extracted_text': file_text[:3000], 'warning': '消息AI未配置，已返回提取文本', 'mode': 'compat'})
                        return
                    if target.ai_config.get("file_recognize_api_url"):
                        if account:
                            response = bot._call_file_recognize_api_for_account(account, file_text, filename, system_prompt, history, original_text=original_text, media_memory_text=media_memory_text)
                        else:
                            response = bot._call_file_recognize_api(file_text, filename, system_prompt, history, original_text=original_text, media_memory_text=media_memory_text)
                        if response:
                            if user_id:
                                target._save_media_memory(user_id, "file", response, filename=filename)
                                if account:
                                    bot._send_ai_reply_in_segments_for_account(account, user_id, response)
                                else:
                                    bot._send_ai_reply_in_segments(user_id, response)
                            self._send_json({'success': True, 'recognize_result': response, 'extracted_text': file_text[:1000]})
                        else:
                            if user_id and target.ai_config.get("api_url"):
                                formatted_text = f"[用户发送了文件: {filename}]\n文件内容摘要:\n{file_text[:2000]}"
                                if prompt:
                                    formatted_text += f"\n额外要求: {prompt}"
                                if account:
                                    response = bot._call_ai_api_for_account(account, formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=user_id)
                                else:
                                    response = bot._call_ai_api(formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=user_id)
                                if response:
                                    if account:
                                        bot._send_ai_reply_in_segments_for_account(account, user_id, response)
                                    else:
                                        bot._send_ai_reply_in_segments(user_id, response)
                                    self._send_json({'success': True, 'ai_reply': response, 'warning': '文件识别API失败，已降级为消息AI'})
                                    return
                            self._send_json({'success': True, 'extracted_text': file_text[:3000], 'warning': '文件识别API调用失败，已返回提取文本'})
                    else:
                        if target.ai_config.get("api_url"):
                            formatted_text = f"[用户发送了文件: {filename}]\n文件内容摘要:\n{file_text[:2000]}"
                            if prompt:
                                formatted_text += f"\n额外要求: {prompt}"
                            if account:
                                response = bot._call_ai_api_for_account(account, formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=user_id)
                            else:
                                response = bot._call_ai_api(formatted_text, history, is_active=False, media_memory_text=media_memory_text, user_id=user_id)
                            if response:
                                if user_id:
                                    if account:
                                        bot._send_ai_reply_in_segments_for_account(account, user_id, response)
                                    else:
                                        bot._send_ai_reply_in_segments(user_id, response)
                                self._send_json({'success': True, 'ai_reply': response, 'warning': '文件识别API未配置，已降级为消息AI'})
                                return
                        self._send_json({'success': True, 'extracted_text': file_text[:3000], 'warning': '文件识别API未配置，已返回提取文本'})
                except Exception as e:
                    print(f"[WEB] 文件识别异常: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _handle_add_user_start(self, data):
                try:
                    account = self._resolve_account()
                    target = account if account else bot
                    with target._add_user_lock:
                        if target._pending_qrcode and target._pending_qrcode.get("status") not in ("done", "error", "expired"):
                            self._send_json({'success': True, 'status': 'already_running', 'message': '已有进行中的添加操作'})
                            return
                    qrcode_key = bot.start_add_user_qrcode_for_account(target)
                    print(f"[WEB] 开始添加用户，qrcode_key={qrcode_key}")
                    for _ in range(30):
                        with target._add_user_lock:
                            qr_status = target._pending_qrcode.get("status") if target._pending_qrcode else None
                            if qr_status == "error":
                                self._send_json({'success': False, 'error': target._pending_qrcode.get("message", "获取二维码失败")})
                                return
                            if target._pending_qrcode and target._pending_qrcode.get("matrix"):
                                self._send_json({
                                    'success': True,
                                    'status': 'qrcode_ready',
                                    'matrix': target._pending_qrcode.get("matrix"),
                                    'key': qrcode_key
                                })
                                return
                        time.sleep(0.5)
                    self._send_json({'success': True, 'status': 'generating', 'message': '正在生成二维码...'})
                except Exception as e:
                    print(f"[WEB] 添加用户异常: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _serve_add_user_status(self):
                try:
                    account = self._resolve_account()
                    target = account if account else bot
                    status = bot.get_add_user_status_for_account(target)
                    was_done = status.get("status") == "done"
                    self._send_json({
                        'success': True,
                        'qrcode_status': status.get('status'),
                        'matrix': status.get('matrix'),
                        'key': status.get('key'),
                        'users': list(target._context_tokens.keys()),
                        'login_done': was_done
                    })
                except Exception as e:
                    self._send_json({'success': False, 'error': '操作失败'})
            def _serve_qr_page(self):
                try:
                    with open('qr.html', 'r', encoding='utf-8') as f:
                         html = f.read()
                except FileNotFoundError:
                    self.send_error(404, 'QR page not found')
                    return

                session_token = bot._generate_session_token()
                bot._verified_sessions[session_token] = time.time() + 86400
                bot._session_tokens[session_token] = time.time() + 86400

                # 绑定到已有账户（优先admin，否则第一个）
                target_account = None
                for acc in bot._accounts.values():
                    if acc.is_admin:
                      target_account = acc
                      break
                if not target_account and bot._accounts:
                  target_account = list(bot._accounts.values())[0]

                if target_account:
                  bot._account_sessions[session_token] = target_account.username
                  print(f"[QR-PAGE] Session 已绑定到账户: {target_account.username}")

                self.send_response(200)
                self.send_header('Content-type', 'text/html; charset=utf-8')
                self.send_header('Cache-Control', 'no-cache')
                self.send_header('Set-Cookie', f'session_token={session_token}; Path=/; SameSite=Lax')
                self.end_headers()
                self.wfile.write(html.encode('utf-8'))

            def _serve_wasm_page(self):
                session_token = bot._generate_session_token()
                password_required = bot._is_web_password_set()
                html = '''<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no, viewport-fit=cover, interactive-widget=resizes-content">
<title>Zyn iLink ChatBox</title>
<script>(function(){var t=localStorage.getItem('theme');if(t==='dark')document.documentElement.setAttribute('data-theme','dark');})();</script>
<style>
* { margin: 0; padding: 0; box-sizing: border-box; -webkit-tap-highlight-color: transparent; }
html, body {
    font-family: "SF Pro Display", "SF Pro Text", "PingFang SC", "Hiragino Sans GB", "Microsoft YaHei", system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
    font-weight: 450;
    -webkit-font-smoothing: antialiased;
    -moz-osx-font-smoothing: grayscale;
    text-rendering: optimizeLegibility;
    background: var(--bg-primary);
    color: var(--text-primary);
}
:root { --bg-primary: #FFFFFF; --bg-secondary: #F2F2F6; --accent: #0A84FF; --accent-hover: #0973E0; --accent-light: rgba(10,132,255,0.08); --text-primary: #1C1C1E; --text-secondary: #8E8E93; --text-hint: #C6C6C8; --bubble-out: #3B82F6; --bubble-in: #FFFFFF; --divider: #E5E5EA; --header-height: 52px; --nav-bg: #FFFFFF; --chat-bg: #F2F2F6; --input-bg: #FFFFFF; --setting-item-bg: #FFFFFF; --setting-arrow: #C6C6C8; --toggle-off: #E5E5EA; --card-round: 12px; --card-mx: 16px; --card-px: 20px; --card-py: 14px; --card-shadow: 0 1px 3px rgba(0,0,0,0.04);
--ease-out: cubic-bezier(0.16, 1, 0.3, 1); --ease-in-out: cubic-bezier(0.45, 0, 0.15, 1); --ease-spring: cubic-bezier(0.34, 1.56, 0.64, 1); --ease-standard: cubic-bezier(0.2, 0, 0, 1); --anim-duration: 1; }
[data-theme="dark"] { --bg-primary: #2C2C2E; --bg-secondary: #1C1C1E; --accent: #0A84FF; --accent-hover: #0973E0; --accent-light: rgba(10,132,255,0.15); --text-primary: #F5F5F7; --text-secondary: #8E8E93; --text-hint: #636366; --bubble-out: #2F7BFF; --bubble-in: #3A3A3E; --divider: #38383A; --nav-bg: #2C2C2E; --chat-bg: #1C1C1E; --input-bg: #3A3A3E; --setting-item-bg: #2C2C2E; --setting-arrow: #555558; --toggle-off: #505055; --card-shadow: 0 1px 3px rgba(0,0,0,0.2); }

@keyframes spin { 0% { transform: rotate(0deg); } 100% { transform: rotate(360deg); } }
.qr-grid { display: grid; gap: 0; background: #FFFFFF; padding: 16px; border-radius: 12px; border: 1px solid var(--divider); max-width: 300px; width: auto; min-width: 180px; box-sizing: border-box; image-rendering: pixelated; margin: 0 auto; }
.qr-cell { width: 9px; height: 9px; background: #000000; min-width: 5px; min-height: 5px; display: block; }
.qr-cell.white { background: #FFFFFF; }
.chat-container { display: none; flex-direction: column; width: 100%; height: 100%; position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: var(--chat-bg); }
.chat-container.active { display: flex; animation: pageSlideIn calc(0.5s * var(--anim-duration)) var(--ease-out); }
@keyframes pageSlideIn { from { transform: translateX(60px); opacity: 0; } to { transform: translateX(0); opacity: 1; } }
.chat-header { height: var(--header-height); background: var(--nav-bg); display: flex; align-items: center; justify-content: center; padding: 0 48px; flex-shrink: 0; position: absolute; top: 0; left: 0; right: 0; z-index: 10; border-bottom: 0.5px solid var(--divider); }
.chat-header-title { font-size: 17px; font-weight: 600; color: var(--text-primary); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; max-width: 100%; }
.settings-toggle { position: absolute; left: 12px; top: 50%; transform: translateY(-50%); width: 32px; height: 32px; border-radius: 50%; background: transparent; color: var(--accent); border: none; cursor: pointer; display: flex; align-items: center; justify-content: center; font-size: 20px; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), background 0.2s; }
.settings-toggle:active { transform: translateY(-50%) scale(0.96); background: var(--accent-light); }
.user-selector { position: absolute; left: 50%; transform: translateX(-50%); padding: 6px 12px; border: 0.5px solid var(--divider); border-radius: var(--card-round); background: var(--bg-primary); cursor: pointer; font-size: 14px; color: var(--text-primary); max-width: 200px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
.user-dropdown { position: absolute; top: 100%; left: 50%; transform: translateX(-50%); background: var(--bg-primary); border-radius: var(--card-round); box-shadow: 0 8px 32px rgba(0,0,0,0.15); min-width: 200px; z-index: 100; display: none; margin-top: 4px; overflow: hidden; }
.user-dropdown.show { display: block; animation: dropdownIn calc(0.25s * var(--anim-duration)) var(--ease-out); }
@keyframes dropdownIn { from { transform: translateX(-50%) scale(0.96); opacity: 0; } to { transform: translateX(-50%) scale(1); opacity: 1; } }
.user-option { padding: 12px 16px; cursor: pointer; transition: background 0.2s; border-bottom: 0.5px solid var(--divider); }
.user-option:last-child { border-bottom: none; }
.user-option:hover { background: var(--accent-light); }
.user-option.current { background: var(--accent-light); color: var(--accent); font-weight: 600; }
.messages-area { flex: 1; overflow-y: auto; padding: 64px 12px 76px; display: flex; flex-direction: column; gap: 4px; background: var(--chat-bg); -webkit-overflow-scrolling: touch; }
.msg-row { display: flex; align-items: flex-end; gap: 8px; max-width: 80%; animation: bubbleIn calc(0.3s * var(--anim-duration)) var(--ease-out); }
@keyframes bubbleIn { from { transform: translateY(8px); opacity: 0; } to { transform: translateY(0); opacity: 1; } }
.msg-row.out { flex-direction: row-reverse; margin-left: auto; }
.bubble { position: relative; padding: 10px 14px; border-radius: 18px; max-width: 100%; line-height: 1.45; font-size: 15px; color: var(--text-primary); word-break: break-word; cursor: pointer; }
.bubble:active { transform: scale(0.96); }
.bubble.in { background: var(--bubble-in); border-bottom-left-radius: 6px; box-shadow: 0 1px 4px rgba(0,0,0,0.06); }
.bubble.in:hover { background: rgba(255,255,255,0.85); }
[data-theme="dark"] .bubble.in:hover { background: rgba(58,58,60,0.78); }
.bubble.out { background: var(--bubble-out); color: #FFFFFF; border-bottom-right-radius: 6px; cursor: default; box-shadow: 0 1px 6px rgba(10,132,255,0.2); }
.bubble.out:hover { background: var(--bubble-out); }
.bubble-text { margin-bottom: 4px; }
.msg-time { font-size: 11px; color: var(--text-hint); margin-top: 4px; text-align: right; }
.bubble.out .msg-time { color: rgba(255,255,255,0.65); }
.msg-time-row { display: flex; align-items: center; justify-content: flex-end; gap: 4px; margin-top: 4px; }
.msg-send-status { display: inline-flex; align-items: center; justify-content: center; }
.msg-send-loading { width: 14px; height: 14px; border: 2px solid rgba(255,255,255,0.3); border-top-color: transparent; border-radius: 50%; animation: spin 0.8s linear infinite; }
.msg-send-fail { width: 18px; height: 18px; border-radius: 50%; background: #FF3B30; color: #fff; font-size: 12px; font-weight: 700; line-height: 18px; text-align: center; cursor: pointer; }
.input-area { background: var(--nav-bg); border-top: 0.5px solid var(--divider); padding: 8px 12px; display: flex; gap: 10px; align-items: center; flex-shrink: 0; padding-bottom: calc(8px + env(safe-area-inset-bottom, 0px)); position: absolute; bottom: 0; left: 0; right: 0; z-index: 20; }
.message-input { flex: 1; height: 44px; border: none; border-radius: 22px; padding: 0 20px; font-size: 16px; outline: none; background: var(--bg-secondary); color: var(--text-primary); transition: box-shadow 0.3s; }
.message-input:focus { box-shadow: 0 4px 20px rgba(10,132,255,0.15), inset 0 0 0 1.5px var(--accent); }
.message-input::placeholder { color: var(--text-hint); }
.send-button { width: 44px; height: 44px; border-radius: 50%; border: none; background: var(--accent); color: white; cursor: pointer; display: flex; align-items: center; justify-content: center; font-size: 18px; }
.send-button:hover { background: var(--accent); color: white; transform: scale(1.04); box-shadow: 0 4px 16px rgba(10,132,255,0.3); }
.send-button:active { transform: scale(0.96); box-shadow: 0 1px 4px rgba(10,132,255,0.15); }
.plus-button { width: 44px; height: 44px; border-radius: 50%; border: none; background: var(--bg-secondary); color: var(--accent); cursor: pointer; display: flex; align-items: center; justify-content: center; font-size: 28px; font-weight: 300; flex-shrink: 0; user-select: none; -webkit-user-select: none; }
.plus-button:active { transform: scale(0.92); }
.plus-button:active { transform: scale(0.96); }
.plus-button.active { color: var(--accent); transform: rotate(45deg); }
.media-panel { background: var(--bg-primary); border-top: 0.5px solid var(--divider); display: none; flex-direction: column; flex-shrink: 0; overflow: hidden; position: absolute; bottom: 68px; left: 0; right: 0; z-index: 19; border-radius: var(--card-round) var(--card-round) 0 0; }
.media-panel.show { display: flex; animation: mediaPanelIn calc(0.3s * var(--anim-duration)) var(--ease-out); }
@keyframes mediaPanelIn { from { transform: translateY(10px); opacity: 0; } to { transform: translateY(0); opacity: 1; } }
.media-panel-inner { padding: 20px 16px; display: grid; grid-template-columns: repeat(4, 1fr); gap: 16px; justify-items: center; }
.media-option { display: flex; flex-direction: column; align-items: center; gap: 8px; cursor: pointer; -webkit-tap-highlight-color: transparent; user-select: none; -webkit-user-select: none; }
.media-option:active .media-option-icon { transform: scale(0.96); }
.media-option-icon { width: 56px; height: 56px; border-radius: var(--card-round); background: var(--bg-secondary); color: var(--text-secondary); display: flex; align-items: center; justify-content: center; font-size: 26px; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard); }
.media-option-label { font-size: 11px; color: var(--text-secondary); text-align: center; line-height: 1.2; }
.media-upload-progress { position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: rgba(0,0,0,0.45); display: none; align-items: center; justify-content: center; z-index: 10001; }
.media-upload-progress.show { display: flex; animation: fadeIn 0.25s ease; }
@keyframes fadeIn { from { opacity: 0; } to { opacity: 1; } }
.media-upload-box { background: var(--bg-primary); border-radius: var(--card-round); padding: 28px 32px; text-align: center; color: var(--text-primary); min-width: 140px; box-shadow: 0 8px 32px rgba(0,0,0,0.2); }
@keyframes modalPop { from { transform: scale(0.95); opacity: 0; } to { transform: scale(1); opacity: 1; } }
.media-upload-spinner { width: 36px; height: 36px; border: 3px solid rgba(255,255,255,0.2); border-top: 3px solid var(--accent); border-radius: 50%; animation: spin 0.8s linear infinite; margin: 0 auto 14px; }
.bubble-media-img { max-width: 200px; max-height: 200px; border-radius: var(--card-round); cursor: pointer; display: block; object-fit: cover; }
.bubble-media-file { display: flex; align-items: center; gap: 10px; min-width: 180px; cursor: pointer; }
.bubble-media-file-icon { width: 40px; height: 40px; border-radius: var(--card-round); background: var(--accent-light); display: flex; align-items: center; justify-content: center; font-size: 20px; flex-shrink: 0; }
.bubble-media-file-info { flex: 1; min-width: 0; }
.bubble-media-file-name { font-size: 14px; color: var(--text-primary); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; max-width: 140px; }
.bubble.out .bubble-media-file-name { color: #FFFFFF; }
.bubble-media-file-size { font-size: 11px; color: var(--text-hint); margin-top: 2px; }
.bubble.out .bubble-media-file-size { color: rgba(255,255,255,0.65); }
.bubble-media-voice { display: flex; align-items: center; gap: 8px; min-width: 80px; cursor: pointer; position: relative; padding-bottom: 6px; }
.bubble-media-voice-bars { display: flex; align-items: center; gap: 2px; height: 20px; }
.bubble-media-voice-bar { width: 3px; border-radius: 2px; background: currentColor; }
.bubble.out .bubble-media-voice-bar { background: rgba(255,255,255,0.85); }
.bubble-media-voice-dur { font-size: 12px; color: var(--text-hint); }
.bubble.out .bubble-media-voice-dur { color: rgba(255,255,255,0.65); }
.bubble-media-voice-progress { position: absolute; bottom: 0; left: 0; right: 0; height: 3px; background: rgba(255,255,255,0.2); border-radius: 2px; overflow: hidden; }
.bubble-media-voice-progress-fill { height: 100%; width: 0%; background: var(--accent); border-radius: 2px; transition: width 0.15s linear; }
.bubble.out .bubble-media-voice-progress-fill { background: rgba(255,255,255,0.8); }
.bubble-media-voice.voice-playing .bubble-media-voice-bar { animation: voiceBarPulse 0.6s ease-in-out infinite alternate; }
@keyframes voiceBarPulse { 0% { opacity: 0.4; } 100% { opacity: 1; } }
.bubble-media-img-wrap { position: relative; overflow: hidden; border-radius: 12px; background: rgba(0,0,0,0.04); min-height: 80px; display: flex; align-items: center; justify-content: center; }
.bubble-media-placeholder { display: flex; flex-direction: column; align-items: center; gap: 6px; padding: 20px; color: var(--text-hint); font-size: 12px; }
.bubble-media-placeholder span { white-space: nowrap; }
.bubble-media-loading { cursor: wait; }
.bubble-media-loading .bubble-media-placeholder { animation: pulse 1.5s ease-in-out infinite; }
@keyframes pulse { 0%, 100% { opacity: 1; } 50% { opacity: 0.5; } }
.bubble-media-video-thumb { position: relative; cursor: pointer; }
.bubble-media-video-thumb-vid { max-width: 100%; max-height: 240px; border-radius: 12px; display: block; object-fit: contain; background: #000; }
.bubble-media-play-btn { position: absolute; top: 50%; left: 50%; transform: translate(-50%, -50%); width: 48px; height: 48px; border-radius: 50%; background: rgba(0,0,0,0.5); display: flex; align-items: center; justify-content: center; pointer-events: none; }
.bubble-media-play-btn svg { width: 24px; height: 24px; fill: #fff; margin-left: 2px; }
.toast { position: fixed; top: 70px; left: 50%; transform: translateX(-50%) translateY(-20px); background: rgba(0,0,0,0.72); color: #fff; padding: 12px 24px; border-radius: 14px; font-size: 14px; z-index: 9999; transition: transform calc(0.4s * var(--anim-duration)) var(--ease-out), opacity 0.3s; pointer-events: none; white-space: nowrap; max-width: 90%; overflow: hidden; text-overflow: ellipsis; opacity: 0; }
.toast.show { transform: translateX(-50%) translateY(0); opacity: 1; }
.empty-state { text-align: center; padding: 60px 20px; color: var(--text-secondary); }
.empty-state-icon { font-size: 64px; margin-bottom: 16px; opacity: 0.3; }
.settings-panel { position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: var(--chat-bg); z-index: 1000; display: none; flex-direction: column; overflow: hidden; height: 100vh; height: 100dvh; }
.settings-panel.show { display: flex; animation: settingsSlideIn calc(0.5s * var(--anim-duration)) var(--ease-out); }
@keyframes settingsSlideIn { from { transform: translateX(-30%); opacity: 0; } to { transform: translateX(0); opacity: 1; } }
.settings-page { position: absolute; top: 0; left: 0; right: 0; bottom: 0; background: var(--chat-bg); display: none; flex-direction: column; overflow: hidden; }
.settings-page.active { display: flex; }
.settings-page-slide { animation: slideInRight calc(0.4s * var(--anim-duration)) var(--ease-out); }
@keyframes slideInRight { from { transform: translateX(30%); opacity: 0; } to { transform: translateX(0); opacity: 1; } }
.settings-nav-header { height: var(--header-height); background: var(--nav-bg); display: flex; align-items: center; padding: 0 16px; flex-shrink: 0; border-bottom: 0.5px solid var(--divider); }
.settings-nav-header .back-btn { width: 28px; height: 28px; border: none; background: transparent; color: var(--accent); font-size: 22px; cursor: pointer; display: flex; align-items: center; justify-content: center; padding: 0; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard); }
.settings-nav-header .back-btn:active { transform: scale(0.96); }
.settings-nav-header .nav-title { font-size: 17px; font-weight: 600; color: var(--text-primary); }
.settings-scroll { flex: 1; overflow-y: auto; -webkit-overflow-scrolling: touch; padding-top: 8px; padding-bottom: 100px; }
.settings-group { margin: 0 var(--card-mx) 16px; background: var(--setting-item-bg); border-radius: var(--card-round); overflow: hidden; box-shadow: var(--card-shadow); }
.settings-group:first-child { margin-top: 0; }
.settings-group-title { padding: 22px var(--card-px) 6px; font-size: 13px; font-weight: 400; color: var(--text-secondary); letter-spacing: 0.3px; }
.settings-item { display: flex; align-items: center; padding: var(--card-py) var(--card-px); background: var(--setting-item-bg); cursor: pointer; transition: background 0.2s; min-height: 52px; }
.settings-item:active { background: var(--accent-light); }
.settings-item + .settings-item { border-top: 0.5px solid var(--divider); margin-left: var(--card-px); margin-right: var(--card-px); padding-left: 0; padding-right: 0; }
.settings-item-content { flex: 1; min-width: 0; }
.settings-item-label { font-size: 16px; color: var(--text-primary); font-weight: 400; }
.settings-item-desc { font-size: 12px; color: var(--text-secondary); margin-top: 2px; }
.settings-item-arrow { color: var(--setting-arrow); font-size: 20px; margin-left: 8px; flex-shrink: 0; font-weight: 300; line-height: 1; }
.user-mgmt-item { display: flex; align-items: center; padding: var(--card-py) var(--card-px); background: transparent; border-bottom: 0.5px solid var(--divider); gap: 12px; }
.user-mgmt-item:last-child { border-bottom: none; }
.user-mgmt-item-info { flex: 1; min-width: 0; }
.user-mgmt-item-name { font-size: 15px; color: var(--text-primary); white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.user-mgmt-item-id { font-size: 12px; color: var(--text-hint); margin-top: 2px; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.user-mgmt-delete-btn { background: #FF3B30; color: #fff; border: none; border-radius: 10px; padding: 8px 16px; font-size: 13px; font-weight: 600; cursor: pointer; flex-shrink: 0; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), opacity 0.2s; }
.user-mgmt-delete-btn:active { transform: scale(0.92); opacity: 0.8; }
.settings-item-action { margin-left: 8px; flex-shrink: 0; }
.chat-menu-switch, .theme-toggle { width: 50px; height: 28px; border-radius: 14px; background: var(--toggle-off); position: relative; cursor: pointer; transition: background 0.3s; border: none; padding: 0; flex-shrink: 0; }
.chat-menu-switch.on, .theme-toggle.active { background: var(--accent); }
.chat-menu-switch-knob, .theme-toggle-knob { width: 24px; height: 24px; border-radius: 50%; background: #FFFFFF; position: absolute; top: 2px; left: 2px; transition: transform 0.3s cubic-bezier(0.34, 1.56, 0.64, 1); box-shadow: 0 2px 6px rgba(0,0,0,0.2); }
.chat-menu-switch.on .chat-menu-switch-knob, .theme-toggle.active .theme-toggle-knob { transform: translateX(22px); }
.settings-header { padding: 16px; background: var(--accent); color: white; border-radius: 20px 20px 0 0; font-weight: 600; display: flex; justify-content: space-between; align-items: center; position: sticky; top: 0; }
.settings-close { background: none; border: none; color: white; font-size: 20px; cursor: pointer; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard); }
.settings-close:active { transform: scale(0.96); }
.settings-body { padding: 8px var(--card-mx) 100px; }
.setting-card { background: var(--setting-item-bg); border-radius: var(--card-round); padding: 0 var(--card-px); margin-bottom: 16px; box-shadow: var(--card-shadow); }
.setting-card .setting-item:first-child { padding-top: var(--card-py); }
.setting-card .setting-item:last-child { padding-bottom: var(--card-py); }
.setting-item { padding: 10px 0; }
.setting-item + .setting-item { border-top: 0.5px solid var(--divider); }
.setting-label { display: block; font-size: 15px; color: var(--text-primary); margin-bottom: 6px; font-weight: 400; }
.setting-label-desc { font-size: 12px; color: var(--text-secondary); margin-bottom: 10px; }
.setting-input, .setting-select { width: 100%; padding: 12px 14px; border: 1px solid var(--divider); border-radius: 10px; font-size: 15px; outline: none; background: var(--bg-secondary); color: var(--text-primary); transition: box-shadow 0.3s, border-color 0.3s; box-sizing: border-box; }
.setting-input:focus, .setting-select:focus { background: var(--bg-primary); box-shadow: 0 0 0 1.5px var(--accent); }
.setting-checkbox { display: flex; align-items: center; gap: 10px; cursor: pointer; padding: 4px 0; }
.setting-checkbox input { width: 22px; height: 22px; cursor: pointer; -webkit-appearance: none; appearance: none; border-radius: 50%; border: 2px solid var(--setting-arrow); background: transparent; position: relative; transition: all 0.2s; flex-shrink: 0; }
.setting-checkbox input:checked { border-color: var(--accent); background: var(--accent); }
.setting-checkbox input:checked::after { content: ''; position: absolute; top: 5px; left: 5px; width: 8px; height: 8px; border-radius: 50%; background: white; }
.setting-row { display: flex; gap: 12px; }
.setting-row .setting-item { flex: 1; border-top: none; }
.setting-half { display: flex; gap: 12px; }
.setting-half .setting-item { flex: 1; }
.settings-save { display: block; width: calc(100% - 2 * var(--card-mx)); margin: 8px var(--card-mx); padding: 14px; background: var(--accent); color: white; border: none; border-radius: var(--card-round); font-size: 16px; font-weight: 500; cursor: pointer; transition: transform 0.2s, background 0.2s; }
.settings-save:active { transform: scale(0.97); }
.setting-divider { height: 10px; background: transparent; }
.setting-section-title { font-size: 13px; font-weight: 400; color: var(--text-secondary); margin: 20px var(--card-px) 6px; letter-spacing: 0.3px; }
.setting-section-title.section-bar { position: relative; margin-left: 0; padding-left: 6px; }
.setting-section-title.section-bar::before { content: ""; position: absolute; left: 0; top: 50%; transform: translateY(-50%); width: 3px; height: 13px; border-radius: 2px; background: var(--accent); }
.about-logo { display: flex; flex-direction: column; align-items: center; padding: 24px 0 16px; }
.about-logo-circle { width: 72px; height: 72px; border-radius: 50%; background: var(--accent); color: #fff; display: flex; align-items: center; justify-content: center; font-size: 26px; font-weight: 600; box-shadow: 0 4px 16px rgba(10,132,255,0.3); }
.about-logo-img { width: 72px; height: 72px; border-radius: 50%; object-fit: cover; background: var(--bg-secondary); box-shadow: 0 4px 16px rgba(10,132,255,0.25); cursor: pointer; transition: box-shadow calc(0.2s * var(--anim-duration)) var(--ease-standard), transform calc(0.2s * var(--anim-duration)) var(--ease-standard); }
.about-logo-img:active { box-shadow: 0 4px 24px rgba(10,132,255,0.35); transform: scale(0.96); }
.about-logo-img.spinning { animation: aboutLogoSpin calc(0.6s * var(--anim-duration)) var(--ease-out); }
@keyframes aboutLogoSpin { from { transform: rotate(0deg); } to { transform: rotate(360deg); } }
.about-logo-name { margin-top: 12px; font-size: 18px; font-weight: 600; color: var(--text-primary); }
.about-info { margin: 0 var(--card-mx) 16px; background: var(--setting-item-bg); border-radius: var(--card-round); overflow: hidden; box-shadow: var(--card-shadow); }
.about-row { display: flex; align-items: center; justify-content: space-between; padding: 16px 20px; transition: background 0.2s; }
.about-row:active { background: var(--accent-light); }
.about-row + .about-row { border-top: 0.5px solid var(--divider); }
.about-label { font-size: 15px; color: var(--text-primary); }
.about-value { font-size: 15px; color: var(--text-secondary); }

.ai-modal { position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: rgba(0,0,0,0.4); display: none; align-items: center; justify-content: center; z-index: 10000; }
.ai-modal.show { display: flex; animation: fadeIn 0.2s ease; }
.ai-modal-content { background: var(--bg-primary); border-radius: var(--card-round); width: 90%; max-width: 400px; max-height: 80vh; overflow-y: auto; box-shadow: 0 12px 40px rgba(0,0,0,0.15); }
.ai-modal-header { padding: 18px; border-bottom: 0.5px solid var(--divider); font-weight: 600; font-size: 17px; display: flex; justify-content: space-between; align-items: center; }
.ai-modal-close { background: none; border: none; font-size: 24px; cursor: pointer; color: var(--text-hint); padding: 0 8px; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), color 0.2s; border-radius: 50%; width: 32px; height: 32px; display: flex; align-items: center; justify-content: center; }
.ai-modal-close:hover { color: var(--text-primary); }
.ai-modal-close:active { transform: scale(0.96); }
.ai-modal-body { padding: 18px; }
.ai-modal-msg-preview { background: var(--bg-secondary); padding: 14px; border-radius: 14px; margin-bottom: 16px; font-size: 13px; color: var(--text-secondary); word-break: break-all; max-height: 150px; overflow-y: auto; }
.ai-modal-label { font-size: 14px; color: var(--text-primary); margin-bottom: 8px; display: block; font-weight: 500; }
.ai-instruction-input { width: 100%; padding: 12px 14px; border: none; border-radius: 14px; font-size: 14px; outline: none; resize: vertical; font-family: inherit; background: var(--bg-secondary); color: var(--text-primary); box-shadow: inset 0 0 0 0.5px var(--divider); transition: box-shadow calc(0.3s * var(--anim-duration)) var(--ease-out); }
.ai-instruction-input:focus { box-shadow: inset 0 0 0 1.5px var(--accent); }
.ai-modal-footer { padding: 18px; border-top: 0.5px solid var(--divider); display: flex; gap: 12px; justify-content: flex-end; }
.ai-modal-btn { padding: 10px 22px; border-radius: 12px; border: none; cursor: pointer; font-size: 14px; font-weight: 600; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), background 0.2s; }
.ai-modal-btn:active { transform: scale(0.96); }
.ai-modal-btn.cancel { background: var(--bg-secondary); color: var(--text-primary); }
.ai-modal-btn.send { background: var(--accent); color: white; box-shadow: 0 2px 8px rgba(10,132,255,0.25); }
.ai-modal-btn.send:hover { background: var(--accent-hover); }
@media (max-width: 768px) { .chat-header { padding: 0 44px; } .chat-header-title { font-size: 15px; } .messages-area { padding: 72px 10px 76px; } .input-area { padding: 8px 10px; padding-bottom: calc(8px + env(safe-area-inset-bottom, 0px)); gap: 8px; } .user-selector { max-width: 160px; font-size: 13px; } .message-input { font-size: 16px; height: 40px; } .settings-toggle { left: 8px; } .msg-row { max-width: 90%; } .bubble { padding: 8px 12px; font-size: 14px; } .plus-button { width: 40px; height: 40px; font-size: 26px; } .send-button { width: 40px; height: 40px; } .media-panel-inner { padding: 16px 12px; gap: 12px; } .media-option-icon { width: 50px; height: 50px; font-size: 22px; border-radius: var(--card-round); } .media-option-label { font-size: 10px; } .bubble-media-img { max-width: 160px; max-height: 160px; } }
@media (max-width: 480px) { .bubble { font-size: 14px; padding: 8px 12px; } .message-input { height: 40px; font-size: 16px; } .send-button { width: 40px; height: 40px; } .user-selector { max-width: 140px; font-size: 12px; } .settings-toggle { width: 28px; height: 28px; font-size: 18px; left: 6px; } .msg-row { max-width: 95%; } }
.chat-list-container { display: none; flex-direction: column; width: 100%; height: 100vh; height: 100dvh; background: var(--chat-bg); }
.chat-list-container.active { display: flex; animation: fadeIn 0.3s ease; }
.chat-list-header { height: var(--header-height); background: var(--nav-bg); display: flex; align-items: center; justify-content: center; padding: 0 16px; flex-shrink: 0; border-bottom: 0.5px solid var(--divider); }
.chat-list-header-title { font-size: 17px; font-weight: 600; color: var(--text-primary); }
.chat-list-add-btn { display: none !important; }
.chat-list-settings-btn { display: none !important; }




.chat-list-items { flex: 1; overflow-y: auto; -webkit-overflow-scrolling: touch; background: transparent; padding: 8px 0; }
.chat-list-item { display: flex; align-items: center; padding: 14px var(--card-px); cursor: pointer; border-bottom: 0.5px solid var(--divider); gap: 12px; }
.chat-list-item:first-child { border-radius: 14px 14px 0 0; }
.chat-list-item:last-child { border-radius: 0 0 14px 14px; border-bottom: none; }
.chat-list-item:active { background: var(--accent-light); }
.chat-list-item-avatar { width: 48px; height: 48px; border-radius: var(--card-round); background: var(--accent-light); display: flex; align-items: center; justify-content: center; font-size: 22px; flex-shrink: 0; }
.chat-list-item-content { flex: 1; min-width: 0; }
.chat-list-item-name { font-size: 16px; color: var(--text-primary); font-weight: 500; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.chat-list-item-msg { font-size: 13px; color: var(--text-hint); margin-top: 4px; white-space: nowrap; overflow: hidden; text-overflow: ellipsis; }
.chat-list-item-time { font-size: 11px; color: var(--text-hint); flex-shrink: 0; align-self: flex-start; margin-top: 2px; }
.chat-list-empty { text-align: center; padding: 80px 20px; color: var(--text-secondary); }
.pc-chat-empty { display: flex; align-items: center; justify-content: center; height: 100%; color: var(--text-hint); font-size: 15px; }
.pc-sidebar-nav { display: flex; border-top: 0.5px solid var(--divider); background: var(--nav-bg); flex-shrink: 0; }
.pc-nav-btn { flex: 1; display: flex; flex-direction: column; align-items: center; justify-content: center; gap: 4px; padding: 10px 12px; border: none; background: transparent; color: var(--text-secondary); cursor: pointer; transition: color 0.2s, background 0.2s; font-size: 11px; font-family: inherit; }
.pc-nav-btn:hover { background: var(--accent-light); }
.pc-nav-btn.active { color: var(--accent); }
.pc-nav-btn svg { flex-shrink: 0; }
.pc-user-delete-btn { padding: 4px 10px; background: #FF3B30; color: #fff; border: none; border-radius: 6px; font-size: 11px; cursor: pointer; flex-shrink: 0; }
.pc-user-delete-btn:hover { background: #D32F2F; }
.chat-list-empty-icon { font-size: 64px; margin-bottom: 16px; opacity: 0.3; }
.chat-back-btn { position: absolute; left: 12px; top: 50%; transform: translateY(-50%); width: 32px; height: 32px; border-radius: 50%; background: transparent; color: var(--accent); border: none; cursor: pointer; display: flex; align-items: center; justify-content: center; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard); }
.chat-back-btn:active { transform: translateY(-50%) scale(0.96); }
.chat-header-menu-wrap { position: absolute; right: 12px; top: 50%; transform: translateY(-50%); }
.chat-header-menu-btn { width: 32px; height: 32px; border-radius: 50%; background: transparent; color: var(--accent); border: none; cursor: pointer; display: flex; align-items: center; justify-content: center; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard); }
.chat-header-menu-btn:active { transform: scale(0.96); }
.nickname-modal { position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: rgba(0,0,0,0.4); display: none; align-items: center; justify-content: center; z-index: 10002; }
.nickname-modal.show { display: flex; animation: fadeIn 0.2s ease; }
.nickname-modal-content { background: var(--bg-primary); border-radius: var(--card-round); width: 90%; max-width: 360px; padding: 24px; box-shadow: 0 12px 40px rgba(0,0,0,0.15); }
.nickname-modal-title { font-size: 17px; font-weight: 600; color: var(--text-primary); margin-bottom: 16px; text-align: center; }
.nickname-modal-userid { font-size: 13px; color: var(--text-hint); text-align: center; margin-bottom: 12px; word-break: break-all; }
.nickname-modal-input { width: 100%; padding: 12px 14px; border: none; border-radius: 14px; font-size: 16px; outline: none; background: var(--bg-secondary); color: var(--text-primary); margin-bottom: 16px; box-shadow: inset 0 0 0 0.5px var(--divider); transition: box-shadow calc(0.3s * var(--anim-duration)) var(--ease-out); }
.nickname-modal-input:focus { box-shadow: inset 0 0 0 1.5px var(--accent); }
.nickname-modal-btns { display: flex; gap: 12px; }
.nickname-modal-btn { flex: 1; padding: 14px; border-radius: 14px; border: none; cursor: pointer; font-size: 15px; font-weight: 600; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), background 0.2s; }
.nickname-modal-btn:active { transform: scale(0.96); }
.nickname-modal-btn.cancel { background: var(--bg-secondary); color: var(--text-primary); }
.nickname-modal-btn.save { background: var(--accent); color: white; box-shadow: 0 2px 10px rgba(10,132,255,0.25); }
.add-user-modal { display: none; position: fixed; top: 0; left: 0; width: 100%; height: 100%; background: rgba(0,0,0,0.4); z-index: 1000; align-items: center; justify-content: center; }
.add-user-modal.show { display: flex; animation: fadeIn 0.2s ease; }
.add-user-modal-content { background: var(--bg-primary); border-radius: var(--card-round); padding: 24px; width: 90%; max-width: 340px; text-align: center; box-shadow: 0 12px 40px rgba(0,0,0,0.15); }
.add-user-modal-title { font-size: 18px; font-weight: 600; color: var(--text-primary); margin-bottom: 16px; }
.add-user-modal-status { font-size: 14px; color: var(--text-secondary); margin-bottom: 16px; min-height: 20px; }
.add-user-modal-qr { display: flex; justify-content: center; margin-bottom: 16px; min-height: 200px; align-items: center; }
.add-user-modal-close { margin-top: 12px; padding: 12px 28px; border-radius: 14px; background: var(--bg-secondary); color: var(--text-primary); border: none; cursor: pointer; font-size: 14px; font-weight: 600; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), background 0.2s; }
.add-user-modal-close:active { transform: scale(0.96); background: var(--divider); }
.add-user-modal-spinner { width: 40px; height: 40px; border: 3px solid rgba(10,132,255,0.15); border-top: 3px solid var(--accent); border-radius: 50%; animation: spin 1s linear infinite; margin: 20px auto; }
.announcement-modal { display: none; position: fixed; top: 0; left: 0; width: 100%; height: 100%; background: rgba(0,0,0,0.4); z-index: 10003; align-items: center; justify-content: center; }
.announcement-modal.show { display: flex; animation: fadeIn 0.2s ease; }
.announcement-modal-content { background: var(--bg-primary); border-radius: var(--card-round); padding: 24px; width: 90%; max-width: 380px; max-height: 80vh; display: flex; flex-direction: column; box-shadow: 0 12px 40px rgba(0,0,0,0.15); }
.announcement-modal-title { font-size: 18px; font-weight: 600; color: var(--text-primary); margin-bottom: 16px; text-align: center; }
.announcement-modal-body { font-size: 14px; color: var(--text-primary); line-height: 1.6; white-space: pre-wrap; word-break: break-word; overflow-y: auto; flex: 1; max-height: 60vh; padding: 0 4px; }
.announcement-modal-footer { display: flex; gap: 10px; margin-top: 20px; justify-content: flex-end; }
.announcement-btn { padding: 10px 20px; border-radius: 12px; font-size: 14px; font-weight: 600; cursor: pointer; border: none; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), opacity 0.2s; }
.announcement-btn:active { transform: scale(0.96); }
.announcement-btn-confirm { background: var(--accent); color: #fff; }
.announcement-btn-dismiss { background: var(--bg-secondary); color: var(--text-secondary); }
.announcement-page-content { font-size: 14px; color: var(--text-primary); line-height: 1.7; white-space: pre-wrap; word-break: break-word; padding: 4px 0; }
.lock-screen { position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: var(--chat-bg); z-index: 100000; display: flex; flex-direction: column; align-items: center; justify-content: center; padding: 20px; overflow-y: auto; }
.lock-screen.hide { display: none; }
.lock-screen-logo { font-size: 48px; margin-bottom: 16px; }
.lock-screen-title { font-size: 24px; font-weight: 700; color: var(--text-primary); margin-bottom: 8px; }
.lock-screen-subtitle { font-size: 14px; color: var(--text-secondary); margin-bottom: 32px; }
.lock-screen-form { width: 100%; max-width: 320px; }
.lock-screen-input { width: 100%; height: 48px; border: 1px solid var(--divider); border-radius: var(--card-round); padding: 0 20px; font-size: 16px; outline: none; background: var(--bg-secondary); color: var(--text-primary); margin-bottom: 12px; text-align: center; letter-spacing: 4px; box-sizing: border-box; }
.lock-screen-input:focus { border-color: var(--accent); box-shadow: 0 4px 20px rgba(10,132,255,0.15), inset 0 0 0 1.5px var(--accent); }
.lock-screen-btn { width: 100%; height: 48px; border-radius: var(--card-round); border: none; background: var(--accent); color: #fff; font-size: 16px; font-weight: 600; cursor: pointer; }
.lock-screen-btn:active { transform: scale(0.96); }
.lock-screen-btn:disabled { opacity: 0.5; cursor: not-allowed; }
.lock-screen-forgot { margin-top: 16px; text-align: center; }
.lock-screen-forgot-btn { background: none; border: none; color: var(--accent); font-size: 14px; cursor: pointer; padding: 8px 16px; }
.lock-screen-forgot-btn:active { opacity: 0.7; }
.lock-screen-error { color: #FF3B30; font-size: 13px; text-align: center; margin-bottom: 8px; min-height: 18px; }
.lock-screen-reset { width: 100%; max-width: 320px; }
.lock-screen-reset-input { width: 100%; height: 44px; border: none; border-radius: 12px; padding: 0 16px; font-size: 15px; outline: none; background: var(--bg-secondary); color: var(--text-primary); margin-bottom: 10px; text-align: center; }
.lock-screen-reset-input:focus { background: var(--bg-primary); box-shadow: 0 0 0 1.5px var(--accent); }
.lock-screen-reset-btn { width: 100%; height: 44px; border-radius: 12px; border: none; background: var(--accent); color: #fff; font-size: 15px; font-weight: 600; cursor: pointer; transition: transform 0.2s; margin-bottom: 8px; }
.lock-screen-reset-btn:active { transform: scale(0.96); }
.lock-screen-reset-btn:disabled { opacity: 0.5; cursor: not-allowed; }
.lock-screen-back-btn { background: none; border: none; color: var(--text-secondary); font-size: 14px; cursor: pointer; padding: 8px 16px; margin-top: 4px; }
.lock-screen-back-btn:active { opacity: 0.7; }
.captcha-row { display: flex; align-items: center; gap: 6px; width: 100%; margin-bottom: 10px; }
.captcha-input { flex: 1; min-width: 0; height: 44px; border: 1px solid var(--divider); border-radius: var(--card-round); padding: 0 12px; font-size: 15px; outline: none; background: var(--bg-secondary); color: var(--text-primary); text-align: center; letter-spacing: 2px; box-sizing: border-box; }
.captcha-input:focus { background: var(--bg-primary); border-color: var(--accent); box-shadow: 0 0 0 1.5px var(--accent); }
.captcha-img { height: 44px; max-width: 100px; border-radius: 8px; cursor: pointer; flex-shrink: 0; }
.captcha-refresh { width: 32px; height: 44px; border: none; border-radius: 10px; background: var(--bg-primary); color: var(--text-secondary); font-size: 18px; cursor: pointer; display: flex; align-items: center; justify-content: center; flex-shrink: 0; box-shadow: 0 2px 8px rgba(0,0,0,0.04), inset 0 0 0 0.5px var(--divider); }
.captcha-refresh:active { transform: scale(0.96); }
.password-verify-overlay { position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: rgba(0,0,0,0.4); z-index: 100002; display: none; align-items: center; justify-content: center; }
.password-verify-overlay.show { display: flex; animation: fadeIn 0.2s ease; }
.password-verify-box { background: var(--bg-primary); border-radius: var(--card-round); padding: 24px; width: 90%; max-width: 340px; text-align: center; box-shadow: 0 12px 40px rgba(0,0,0,0.15); }
.password-verify-title { font-size: 18px; font-weight: 600; color: var(--text-primary); margin-bottom: 16px; }
.password-verify-input { width: 100%; height: 44px; border: none; border-radius: 12px; padding: 0 16px; font-size: 15px; outline: none; background: var(--bg-secondary); color: var(--text-primary); box-shadow: inset 0 0 0 0.5px var(--divider); margin-bottom: 12px; text-align: center; letter-spacing: 4px; }
.password-verify-input:focus { box-shadow: inset 0 0 0 1.5px var(--accent); }
.password-verify-error { color: #FF3B30; font-size: 13px; margin-bottom: 8px; min-height: 18px; }
.password-verify-btns { display: flex; gap: 10px; }
.password-verify-btn { flex: 1; height: 44px; border-radius: 12px; border: none; font-size: 15px; font-weight: 600; cursor: pointer; transition: transform 0.2s; }
.password-verify-btn:active { transform: scale(0.96); }
.password-verify-btn.cancel { background: var(--bg-secondary); color: var(--text-secondary); }
.password-verify-btn.confirm { background: var(--accent); color: #fff; }
.admin-user-list { max-height: 240px; overflow-y: auto; margin-top: 12px; }
.admin-user-item { display: flex; align-items: center; padding: 12px 14px; border-radius: 12px; cursor: pointer; transition: background 0.2s; margin-bottom: 4px; }
.admin-user-item:active { background: var(--bg-secondary); }
.admin-user-item.selected { background: var(--accent-light); }
.admin-user-item-name { flex: 1; font-size: 15px; color: var(--text-primary); }
.admin-user-item-check { color: var(--accent); font-size: 18px; font-weight: 700; }
@media (max-width: 768px) { .chat-list-item { padding: 12px 14px; } .chat-list-item-avatar { width: 44px; height: 44px; font-size: 20px; } .chat-list-item-name { font-size: 15px; } .chat-list-item-msg { font-size: 12px; } .chat-back-btn { width: 28px; height: 28px; left: 8px; } .user-list-back-btn { width: 28px; height: 28px; left: 8px; } .chat-header-menu-wrap { right: 8px; } .chat-header-menu-btn { width: 28px; height: 28px; } .chat-list-settings-btn { width: 28px; height: 28px; right: 8px; } .chat-list-add-btn { width: 28px; height: 28px; left: 8px; } .user-list-add-btn { width: 28px; height: 28px; right: 8px; } }
@media (max-width: 480px) { .chat-list-item { padding: 10px 12px; gap: 10px; } .chat-list-item-avatar { width: 40px; height: 40px; font-size: 18px; } .chat-list-item-name { font-size: 14px; } .chat-list-item-msg { font-size: 11px; } .chat-back-btn { width: 28px; height: 28px; left: 6px; } .user-list-back-btn { width: 28px; height: 28px; left: 6px; } .chat-header-menu-wrap { right: 6px; } .chat-header-menu-btn { width: 28px; height: 28px; } .chat-list-settings-btn { width: 28px; height: 28px; right: 6px; } .chat-list-add-btn { width: 28px; height: 28px; left: 6px; } .user-list-add-btn { width: 28px; height: 28px; right: 6px; } .nickname-modal-content { padding: 20px; } .lock-screen { padding: 16px; } .lock-screen-logo { font-size: 36px; margin-bottom: 12px; } .lock-screen-title { font-size: 20px; } .lock-screen-subtitle { font-size: 12px; margin-bottom: 24px; } .lock-screen-input { height: 44px; font-size: 15px; } .lock-screen-btn { height: 44px; font-size: 15px; } .lock-screen-reset-input { height: 40px; font-size: 14px; } .lock-screen-reset-btn { height: 40px; font-size: 14px; } .captcha-row { gap: 4px; } .captcha-input { padding: 0 8px; font-size: 14px; } .captcha-img { height: 40px; max-width: 80px; } .captcha-refresh { width: 28px; height: 40px; font-size: 16px; } }
@supports (height: 100dvh) { html, body { height: 100dvh; } #app { height: 100dvh; background: var(--bg-primary); } }
body.keyboard-open .chat-container { height: 100vh; height: 100dvh; }
body.keyboard-open .media-panel { display: none !important; }
body.keyboard-open .plus-button.active { transform: none; color: var(--text-secondary); }
body.keyboard-open #app { height: auto; min-height: 100vh; min-height: 100dvh; }
.no-animations, .no-animations *, .no-animations *::before, .no-animations *::after { animation-duration: 0s !important; animation-delay: 0s !important; transition-duration: 0s !important; }
.appearance-slider-wrap { display: flex; align-items: center; gap: 12px; margin-top: 8px; }
.appearance-slider { -webkit-appearance: none; appearance: none; flex: 1; height: 4px; border-radius: var(--card-round); background: var(--toggle-off); outline: none; }
.appearance-slider::-webkit-slider-thumb { -webkit-appearance: none; appearance: none; width: 22px; height: 22px; border-radius: 50%; background: white; cursor: pointer; box-shadow: 0 1px 4px rgba(0,0,0,0.2); border: 2px solid var(--accent); transition: transform 0.2s; }
.appearance-slider::-webkit-slider-thumb:active { transform: scale(0.92); }
.appearance-slider::-moz-range-thumb { width: 22px; height: 22px; border-radius: 50%; background: white; cursor: pointer; border: 2px solid var(--accent); box-shadow: 0 1px 4px rgba(0,0,0,0.2); }
.appearance-slider-val { min-width: 36px; text-align: center; font-size: 14px; font-weight: 600; color: var(--accent); }
.bottom-tab-bar { position: fixed; bottom: 0; left: 50%; transform: translateX(-50%); display: inline-flex; align-items: center; justify-content: center; gap: 4px; z-index: 1001; padding: 8px 12px; background: var(--nav-bg); border-radius: 24px; border: 0.5px solid var(--divider); box-shadow: 0 2px 12px rgba(0,0,0,0.06); transition: transform calc(0.35s * var(--anim-duration)) var(--ease-out), opacity calc(0.3s * var(--anim-duration)); margin-bottom: 12px; }
.bottom-tab-bar.hidden { transform: translateX(-50%) translateY(100px); opacity: 0; pointer-events: none; }
.bottom-tab-item { display: flex; flex-direction: column; align-items: center; justify-content: center; gap: 3px; padding: 6px 18px; border-radius: 20px; cursor: pointer; transition: background 0.2s; -webkit-tap-highlight-color: transparent; user-select: none; -webkit-user-select: none; border: none; background: transparent; color: var(--text-secondary); }
.bottom-tab-item:active { transform: scale(0.92); }
.bottom-tab-item.active { color: var(--accent); background: var(--accent-light); }
.bottom-tab-item-icon { font-size: 22px; display: flex; align-items: center; justify-content: center; }
.bottom-tab-item-label { font-size: 10px; font-weight: 600; letter-spacing: 0.3px; }
.user-list-page { display: none; flex-direction: column; width: 100%; height: 100vh; height: 100dvh; background: var(--chat-bg); }
.user-list-page.active { display: flex; animation: fadeIn 0.3s ease; }
.user-list-header { position: relative; height: var(--header-height); background: var(--nav-bg); display: flex; align-items: center; justify-content: center; padding: 0 16px; flex-shrink: 0; border-bottom: 0.5px solid var(--divider); }
.user-list-back-btn { position: absolute; left: 12px; top: 50%; transform: translateY(-50%); width: 32px; height: 32px; border-radius: 50%; background: transparent; color: var(--accent); border: none; cursor: pointer; display: flex; align-items: center; justify-content: center; z-index: 2; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), background 0.2s; }
.user-list-back-btn:active { transform: translateY(-50%) scale(0.92); background: var(--accent-light); }
.user-list-add-btn { position: absolute; right: 12px; top: 50%; transform: translateY(-50%); width: 32px; height: 32px; border-radius: 50%; background: transparent; color: var(--accent); border: none; cursor: pointer; display: flex; align-items: center; justify-content: center; z-index: 2; transition: transform calc(0.2s * var(--anim-duration)) var(--ease-standard), background 0.2s; }
.user-list-add-btn:active { transform: translateY(-50%) scale(0.92); background: var(--accent-light); }
.user-list-header-title { font-size: 17px; font-weight: 600; color: var(--text-primary); }


.user-list-scroll { flex: 1; overflow-y: auto; -webkit-overflow-scrolling: touch; background: var(--chat-bg); }
.user-list-body { background: var(--setting-item-bg); border-radius: var(--card-round); margin: 12px var(--card-mx); overflow: hidden; padding: 0; }
.user-list-body > .chat-list-item { background: var(--setting-item-bg); border-bottom: none; margin-bottom: 0.5px; }
.user-list-body > .chat-list-item + .chat-list-item { border-top: 0.5px solid var(--divider); }
.user-list-body > .chat-list-item:first-child { border-radius: var(--card-round) var(--card-round) 0 0; }
.user-list-body > .chat-list-item:last-child { border-radius: 0 0 var(--card-round) var(--card-round); }
.user-list-body > .chat-list-item:only-child { border-radius: var(--card-round); }
.user-list-page .chat-list-items { background: var(--setting-item-bg); border-radius: var(--card-round); margin: 12px var(--card-mx); overflow: hidden; padding: 0; }
.user-list-page .chat-list-items > .chat-list-item { background: transparent; border-radius: 0; margin: 0; border-bottom: 0.5px solid var(--divider); }
.user-list-page .chat-list-items > .chat-list-item:last-child { border-bottom: none; }
@media (max-width: 768px) { .pc-sidebar-nav { display: none !important; } .bottom-tab-bar { gap: 6px; padding: 6px 10px; border-radius: 22px; } .bottom-tab-item { padding: 5px 14px; } .bottom-tab-item-icon { font-size: 20px; } .bottom-tab-item-label { font-size: 9px; } }
@media (max-width: 480px) { .bottom-tab-bar { gap: 4px; padding: 4px 8px; border-radius: 18px; } .bottom-tab-item { padding: 4px 12px; } .bottom-tab-item-icon { font-size: 18px; } .bottom-tab-item-label { font-size: 9px; } }
.chat-menu-dropdown { position: absolute; top: 100%; right: 0; background: var(--bg-primary); border-radius: var(--card-round); box-shadow: 0 8px 32px rgba(0,0,0,0.15); min-width: 160px; z-index: 100; display: none; margin-top: 4px; overflow: hidden; }
.chat-menu-dropdown.show { display: block; animation: dropdownIn calc(0.25s * var(--anim-duration)) var(--ease-out); }
.chat-menu-item { padding: 12px 16px; font-size: 14px; color: var(--text-primary); cursor: pointer; transition: background 0.15s; }
.chat-menu-item:active { background: var(--bg-secondary); }
.chat-menu-toggle { display: flex; align-items: center; justify-content: space-between; }
.chat-menu-switch { width: 44px; height: 26px; border-radius: 13px; background: var(--toggle-off); position: relative; transition: background 0.2s; flex-shrink: 0; }
.chat-menu-switch.on { background: var(--accent); }
.chat-menu-switch-knob { width: 22px; height: 22px; border-radius: 50%; background: white; position: absolute; top: 2px; left: 2px; transition: transform 0.2s; box-shadow: 0 1px 3px rgba(0,0,0,0.2); }
.chat-menu-switch.on .chat-menu-switch-knob { transform: translateX(18px); }
.user-prompt-textarea { width: 100%; padding: 12px 14px; border: 1px solid var(--divider); border-radius: var(--card-round); font-size: 14px; outline: none; background: var(--bg-secondary); color: var(--text-primary); margin-bottom: 16px; transition: box-shadow 0.3s, border-color 0.3s; resize: vertical; min-height: 100px; max-height: 300px; font-family: inherit; line-height: 1.5; box-sizing: border-box; }
.user-prompt-textarea:focus { border-color: var(--accent); box-shadow: inset 0 0 0 1.5px var(--accent); }
@media (min-width: 769px) { #app { display: flex !important; flex-direction: row; } .chat-list-container { width: 320px; min-width: 320px; max-width: 380px; height: 100vh; height: 100dvh; position: relative; border-right: 0.5px solid var(--divider); flex-shrink: 0; } .chat-list-container.active { display: flex; animation: none; } .chat-list-header { position: relative; top: 0; left: 0; right: 0; border-radius: 0; box-shadow: none; border-bottom: 0.5px solid var(--divider); } .chat-list-items { padding-top: 0; padding-bottom: 20px; } .chat-container { position: relative; width: auto; flex: 1; height: 100vh; height: 100dvh; } .chat-container.active { display: flex; animation: none; } .chat-back-btn { display: none; } .chat-header { position: relative; top: 0; left: 0; right: 0; border-radius: 0; box-shadow: none; border-bottom: 0.5px solid var(--divider); padding: 0 48px 0 16px; } .messages-area { padding-top: 56px; padding-bottom: 16px; } .input-area { position: relative; padding-bottom: 8px; } .media-panel { position: relative; bottom: auto; border-radius: 0; } .bottom-tab-bar { display: none !important; } .user-list-page { width: 320px; min-width: 320px; max-width: 380px; height: 100vh; height: 100dvh; position: relative; border-right: 0.5px solid var(--divider); flex-shrink: 0; } .user-list-page.active { display: flex; animation: none; } .user-list-header { position: relative; top: 0; left: 0; right: 0; border-radius: 0; box-shadow: none; border-bottom: 0.5px solid var(--divider); } .user-list-scroll { padding-top: 0; } .settings-panel { position: fixed; } .settings-nav-header { position: relative; top: 0; left: 0; right: 0; border-radius: 0; box-shadow: none; border-bottom: 0.5px solid var(--divider); } .settings-scroll { padding-top: 16px; } .chat-list-item { margin: 0; border-radius: 0; } .chat-list-item:first-child { border-radius: 0; } .chat-list-item:last-child { border-radius: 0; } .chat-list-empty { padding: 40px 20px; } .pc-sidebar { display: flex; flex-direction: column; width: 320px; min-width: 320px; max-width: 380px; height: 100vh; height: 100dvh; background: var(--chat-bg); border-right: 0.5px solid var(--divider); flex-shrink: 0; overflow: hidden; } .pc-sidebar-header { height: 56px; background: var(--nav-bg); display: flex; align-items: center; justify-content: center; padding: 0 16px; flex-shrink: 0; border-bottom: 0.5px solid var(--divider); } .pc-sidebar-header-title { font-size: 17px; font-weight: 600; color: var(--text-primary); } .pc-sidebar-add-btn { position: absolute; right: 12px; top: 50%; transform: translateY(-50%); width: 32px; height: 32px; border-radius: 50%; background: var(--accent); color: #fff; border: none; cursor: pointer; display: flex; align-items: center; justify-content: center; transition: transform 0.2s, background 0.2s; box-shadow: 0 2px 8px rgba(10,132,255,0.25); } .pc-sidebar-add-btn:hover { background: var(--accent-hover); } .pc-sidebar-add-btn:active { transform: translateY(-50%) scale(0.96); } .pc-sidebar-tabs { display: flex; border-bottom: 0.5px solid var(--divider); background: var(--nav-bg); } .pc-sidebar-tab { flex: 1; padding: 10px; text-align: center; font-size: 13px; font-weight: 600; color: var(--text-secondary); cursor: pointer; border: none; background: transparent; transition: color 0.2s, background 0.2s; } .pc-sidebar-tab.active { color: var(--accent); background: var(--accent-light); } .pc-sidebar-tab:hover { background: var(--accent-light); } .pc-sidebar-content { flex: 1; overflow-y: auto; } .pc-chat-empty { display: flex; align-items: center; justify-content: center; height: 100%; color: var(--text-hint); font-size: 15px; } }
</style>
</head>
<body>
<div id="lock-screen" class="lock-screen">
    <div class="lock-screen-logo">ZynSync</div>
    <div class="lock-screen-title">Zyn ChatBox</div>
    <div class="lock-screen-subtitle">多账户机器人管理平台</div>
    <div id="lock-screen-error" class="lock-screen-error"></div>
    <div id="bot-login-form" class="lock-screen-form">
        <input type="text" id="bot-login-username" class="lock-screen-input" placeholder="用户名" autocomplete="off" style="letter-spacing:normal;" />
        <input type="password" id="bot-login-password" class="lock-screen-input" placeholder="密码" autocomplete="off" />
        <button id="bot-login-submit" class="lock-screen-btn">登录</button>
        <div class="lock-screen-forgot">
            <button id="bot-show-register" class="lock-screen-forgot-btn">没有账户？注册</button>
            <button id="bot-show-forgot" class="lock-screen-forgot-btn">忘记密码</button>
        </div>
    </div>
    <div id="bot-register-form" class="lock-screen-form" style="display:none;">
        <input type="text" id="bot-register-username" class="lock-screen-input" placeholder="用户名 (2-32位)" autocomplete="off" style="letter-spacing:normal;" />
        <input type="email" id="bot-register-email" class="lock-screen-input" placeholder="邮箱" autocomplete="off" style="display:none;" />
        <div id="bot-register-email-code-row" class="captcha-row" style="display:none;">
            <input type="text" id="bot-register-email-code" class="captcha-input" placeholder="邮箱验证码" autocomplete="off" />
            <button id="bot-register-send-email-code" class="captcha-refresh" style="width:auto;padding:0 12px;font-size:13px;white-space:nowrap;">发送</button>
        </div>
        <input type="password" id="bot-register-password" class="lock-screen-input" placeholder="密码 (至少6位)" autocomplete="off" />
        <input type="password" id="bot-register-password2" class="lock-screen-input" placeholder="确认密码" autocomplete="off" />
        <button id="bot-register-submit" class="lock-screen-btn">注册</button>
        <div style="margin-top:8px;padding:10px 14px;background:var(--accent-light);border:0.5px solid var(--accent);border-radius:var(--card-round);font-size:13px;color:var(--accent);text-align:center;line-height:1.5;">若邮箱没有收到邮件，请前往垃圾邮件查看（目前邮箱只适配QQ邮箱）</div>
        <div class="lock-screen-forgot">
            <button id="bot-show-login" class="lock-screen-forgot-btn">已有账户？登录</button>
        </div>
    </div>
    <div id="bot-forgot-form" class="lock-screen-form" style="display:none;">
        <input type="text" id="bot-forgot-username" class="lock-screen-input" placeholder="用户名" autocomplete="off" style="letter-spacing:normal;" />
        <div class="captcha-row">
            <input type="text" id="bot-forgot-captcha" class="captcha-input" placeholder="图形验证码" autocomplete="off" />
            <img id="bot-forgot-captcha-img" class="captcha-img" src="" alt="验证码" />
            <button id="bot-forgot-captcha-refresh" class="captcha-refresh">↻</button>
        </div>
        <button id="bot-forgot-send-code" class="lock-screen-reset-btn">发送验证码</button>
        <input type="text" id="bot-forgot-code" class="lock-screen-reset-input" placeholder="输入验证码" autocomplete="off" />
        <input type="password" id="bot-forgot-new-password" class="lock-screen-reset-input" placeholder="新密码 (至少6位)" autocomplete="off" />
        <input type="password" id="bot-forgot-confirm-password" class="lock-screen-reset-input" placeholder="确认新密码" autocomplete="off" />
        <button id="bot-forgot-submit" class="lock-screen-reset-btn">重置密码</button>
        <div class="lock-screen-forgot">
            <button id="bot-forgot-back-login" class="lock-screen-back-btn">返回登录</button>
        </div>
    </div>
</div>
<div id="password-verify-overlay" class="password-verify-overlay">
    <div class="password-verify-box">
        <div class="password-verify-title">验证密码</div>
        <div id="password-verify-error" class="password-verify-error"></div>
        <input type="password" id="password-verify-input" class="password-verify-input" placeholder="输入当前密码" autocomplete="off" />
        <div class="password-verify-btns">
            <button id="password-verify-cancel" class="password-verify-btn cancel">取消</button>
            <button id="password-verify-confirm" class="password-verify-btn confirm">确认</button>
        </div>
    </div>
</div>
<div id="app" style="display:none">
    <div id="chat-list-page" class="chat-list-container">
        <div class="chat-list-header">
            <button id="chat-list-add-btn" class="chat-list-add-btn" title="添加用户"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><line x1="12" y1="5" x2="12" y2="19"/><line x1="5" y1="12" x2="19" y2="12"/></svg></button>
            <span class="chat-list-header-title">Zyn iLink Chatbox</span>
            <button id="chat-list-settings-btn" class="chat-list-settings-btn" title="设置"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round"><circle cx="12" cy="12" r="3"/><path d="M19.4 15a1.65 1.65 0 00.33 1.82l.06.06a2 2 0 010 2.83 2 2 0 01-2.83 0l-.06-.06a1.65 1.65 0 00-1.82-.33 1.65 1.65 0 00-1 1.51V21a2 2 0 01-4 0v-.09A1.65 1.65 0 009 19.4a1.65 1.65 0 00-1.82.33l-.06.06a2 2 0 01-2.83-2.83l.06-.06A1.65 1.65 0 004.68 15a1.65 1.65 0 00-1.51-1H3a2 2 0 010-4h.09A1.65 1.65 0 004.6 9a1.65 1.65 0 00-.33-1.82l-.06-.06a2 2 0 012.83-2.83l.06.06A1.65 1.65 0 009 4.68a1.65 1.65 0 001-1.51V3a2 2 0 014 0v.09a1.65 1.65 0 001 1.51 1.65 1.65 0 001.82-.33l.06-.06a2 2 0 012.83 2.83l-.06.06A1.65 1.65 0 0019.4 9a1.65 1.65 0 001.51 1H21a2 2 0 010 4h-.09a1.65 1.65 0 00-1.51 1z"/></svg></button>
        </div>
        <div id="chat-list-items" class="chat-list-items">
            <div class="chat-list-empty">
                <div class="chat-list-empty-icon">💬</div>
                <div>暂无聊天</div>
            </div>
        </div>
        <div id="user-list-items" class="chat-list-items" style="display:none;">
            <div class="chat-list-empty">
                <div class="chat-list-empty-icon">👥</div>
                <div>暂无用户</div>
            </div>
        </div>
        <div id="pc-sidebar-nav" class="pc-sidebar-nav">
            <button id="pc-nav-chat" class="pc-nav-btn active"><svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><path d="M21 15a2 2 0 01-2 2H7l-4 4V5a2 2 0 012-2h14a2 2 0 012 2z"/></svg><span>聊天</span></button>
            <button id="pc-nav-users" class="pc-nav-btn"><svg viewBox="0 0 24 24" width="20" height="20" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><path d="M17 21v-2a4 4 0 00-4-4H5a4 4 0 00-4 4v2"/><circle cx="9" cy="7" r="4"/><path d="M23 21v-2a4 4 0 00-3-3.87"/><path d="M16 3.13a4 4 0 010 7.75"/></svg><span>用户</span></button>
        </div>
    </div>
    <div id="chat-page" class="chat-container">
        <div class="chat-header">
            <button id="chat-back-btn" class="chat-back-btn"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><polyline points="15 18 9 12 15 6"/></svg></button>
            <span id="chat-header-title" class="chat-header-title"></span>
            <div class="chat-header-menu-wrap">
                <button id="chat-menu-btn" class="chat-header-menu-btn"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><line x1="4" y1="6" x2="20" y2="6"/><line x1="4" y1="12" x2="20" y2="12"/><line x1="4" y1="18" x2="20" y2="18"/></svg></button>
                <div id="chat-menu-dropdown" class="chat-menu-dropdown">
                    <div class="chat-menu-item" id="chat-menu-nickname">设置备注名</div>
                    <div class="chat-menu-item" id="chat-menu-prompt">设置AI提示词</div>
                    <div class="chat-menu-item chat-menu-toggle" id="chat-menu-ai-toggle">
                        <span>启用AI回复</span>
                        <div class="chat-menu-switch" id="chat-menu-ai-switch"><div class="chat-menu-switch-knob"></div></div>
                    </div>
                    <div class="chat-menu-item chat-menu-toggle" id="chat-menu-scheduled-toggle">
                        <span>启用间隔定时发送</span>
                        <div class="chat-menu-switch" id="chat-menu-scheduled-switch"><div class="chat-menu-switch-knob"></div></div>
                    </div>
                    <div class="chat-menu-item chat-menu-toggle" id="chat-menu-daily-toggle">
                        <span>启用每日定时发送</span>
                        <div class="chat-menu-switch" id="chat-menu-daily-switch"><div class="chat-menu-switch-knob"></div></div>
                    </div>
                </div>
            </div>
        </div>
        <div id="messages-area" class="messages-area">
            <div class="empty-state">
                <div class="empty-state-icon">💬</div>
                <div>点击文本消息可使用 AI 回复，点击媒体消息可查看/下载</div>
            </div>
        </div>
        <div class="input-area" id="input-area">
            <button id="plus-btn" class="plus-button"><svg viewBox="0 0 24 24" width="24" height="24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><line x1="12" y1="5" x2="12" y2="19"/><line x1="5" y1="12" x2="19" y2="12"/></svg></button>
            <input type="text" id="message-input" class="message-input" placeholder="输入消息..." />
            <button id="send-btn" class="send-button"><svg viewBox="0 0 24 24" width="22" height="22" fill="currentColor"><path d="M2.01 21L23 12 2.01 3 2 10l15 2-15 2z"/></svg></button>
        </div>
        <div id="media-panel" class="media-panel">
            <div class="media-panel-inner">
                <div class="media-option" id="media-photo">
                    <div class="media-option-icon"><svg viewBox="0 0 24 24" width="28" height="28" fill="none" stroke="currentColor" stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round"><rect x="3" y="3" width="18" height="18" rx="2"/><circle cx="8.5" cy="8.5" r="1.5"/><path d="M21 15l-5-5L5 21"/></svg></div>
                    <div class="media-option-label">相册</div>
                </div>
                <div class="media-option" id="media-camera">
                    <div class="media-option-icon"><svg viewBox="0 0 24 24" width="28" height="28" fill="none" stroke="currentColor" stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round"><path d="M23 19a2 2 0 01-2 2H3a2 2 0 01-2-2V8a2 2 0 012-2h4l2-3h6l2 3h4a2 2 0 012 2z"/><circle cx="12" cy="13" r="4"/></svg></div>
                    <div class="media-option-label">拍摄</div>
                </div>
                <div class="media-option" id="media-video">
                    <div class="media-option-icon"><svg viewBox="0 0 24 24" width="28" height="28" fill="none" stroke="currentColor" stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round"><polygon points="23 7 16 12 23 17 23 7"/><rect x="1" y="5" width="15" height="14" rx="2"/></svg></div>
                    <div class="media-option-label">视频</div>
                </div>
                <div class="media-option" id="media-file">
                    <div class="media-option-icon"><svg viewBox="0 0 24 24" width="28" height="28" fill="none" stroke="currentColor" stroke-width="1.5" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/><line x1="16" y1="13" x2="8" y2="13"/><line x1="16" y1="17" x2="8" y2="17"/></svg></div>
                    <div class="media-option-label">文件</div>
                </div>
            </div>
        </div>
        <input type="file" id="file-photo" accept="image/*" style="display:none" />
        <input type="file" id="file-camera" accept="image/*" capture="environment" style="display:none" />
        <input type="file" id="file-video" accept="video/*" style="display:none" />
        <input type="file" id="file-video-capture" accept="video/*" capture="environment" style="display:none" />
        <input type="file" id="file-doc" accept="*/*" style="display:none" />
    </div>
    <div id="user-list-page" class="user-list-page">
        <div class="user-list-header">
            <button id="user-list-back-btn" class="user-list-back-btn" title="返回"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><polyline points="15 18 9 12 15 6"/></svg></button>
            <span class="user-list-header-title">用户列表</span>
            <button id="user-list-add-btn" class="user-list-add-btn" title="添加用户"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round"><line x1="12" y1="5" x2="12" y2="19"/><line x1="5" y1="12" x2="19" y2="12"/></svg></button>
        </div>
        <div class="user-list-scroll">
            <div class="user-list-body" id="user-mgmt-list">
            </div>
        </div>
    </div>
    <div id="bottom-tab-bar" class="bottom-tab-bar hidden">
        <button class="bottom-tab-item active" id="tab-list">
            <div class="bottom-tab-item-icon"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M21 15a2 2 0 01-2 2H7l-4 4V5a2 2 0 012-2h14a2 2 0 012 2z"/></svg></div>
            <div class="bottom-tab-item-label">列表</div>
        </button>
        <button class="bottom-tab-item" id="tab-users">
            <div class="bottom-tab-item-icon"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M17 21v-2a4 4 0 00-4-4H5a4 4 0 00-4 4v2"/><circle cx="9" cy="7" r="4"/><path d="M23 21v-2a4 4 0 00-3-3.87"/><path d="M16 3.13a4 4 0 010 7.75"/></svg></div>
            <div class="bottom-tab-item-label">用户</div>
        </button>
        <button class="bottom-tab-item" id="tab-settings">
            <div class="bottom-tab-item-icon"><svg viewBox="0 0 24 24" width="22" height="22" fill="none" stroke="currentColor" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round"><circle cx="12" cy="12" r="3"/><path d="M19.4 15a1.65 1.65 0 00.33 1.82l.06.06a2 2 0 010 2.83 2 2 0 01-2.83 0l-.06-.06a1.65 1.65 0 00-1.82-.33 1.65 1.65 0 00-1 1.51V21a2 2 0 01-4 0v-.09A1.65 1.65 0 009 19.4a1.65 1.65 0 00-1.82.33l-.06.06a2 2 0 01-2.83-2.83l.06-.06A1.65 1.65 0 004.68 15a1.65 1.65 0 00-1.51-1H3a2 2 0 010-4h.09A1.65 1.65 0 004.6 9a1.65 1.65 0 00-.33-1.82l-.06-.06a2 2 0 012.83-2.83l.06.06A1.65 1.65 0 009 4.68a1.65 1.65 0 001-1.51V3a2 2 0 014 0v.09a1.65 1.65 0 001 1.51 1.65 1.65 0 001.82-.33l.06-.06a2 2 0 012.83 2.83l-.06.06A1.65 1.65 0 0019.4 9a1.65 1.65 0 001.51 1H21a2 2 0 010 4h-.09a1.65 1.65 0 00-1.51 1z"/></svg></div>
            <div class="bottom-tab-item-label">设置</div>
        </button>
    </div>
</div>
<div id="settings-panel" class="settings-panel">
    <div id="settings-main" class="settings-page active">
        <div class="settings-nav-header">
            <button class="back-btn" id="settings-back-btn">‹</button>
            <span class="nav-title">设置</span>
        </div>
        <div class="settings-scroll">
            <div class="settings-group">
                <div class="settings-item" id="settings-tutorial-item"><div class="settings-item-content">
                        <div class="settings-item-label">使用教程（必看）</div>
                    </div>
                    <div class="settings-item-arrow">›</div>
                </div>
            </div>
            <div class="settings-group">
                <div class="settings-item" id="settings-theme-item"><div class="settings-item-content">
                        <div class="settings-item-label">深色模式</div>
                    </div>
                    <div class="settings-item-action">
                        <button class="theme-toggle" id="theme-toggle-btn"><div class="theme-toggle-knob"></div></button>
                    </div>
                </div>
            </div>
            <div class="settings-group">
                <div class="settings-item" id="settings-appearance-item"><div class="settings-item-content">
                        <div class="settings-item-label">外观设置</div>
                        <div class="settings-item-desc">动画与透明度</div>
                    </div>
                    <div class="settings-item-arrow">›</div>
                </div>
            </div>
            <div class="settings-group">
                <div class="settings-item" id="settings-api-item"><div class="settings-item-content">
                        <div class="settings-item-label">AI 回复设置</div>
                        <div class="settings-item-desc">配置 AI 自动回复参数</div>
                    </div>
                    <div class="settings-item-arrow">›</div>
                </div>
            </div>
            <div class="settings-group">
                <div class="settings-item" id="settings-announcement-item"><div class="settings-item-content">
                        <div class="settings-item-label">公告</div>
                        <div class="settings-item-desc">查看最新公告信息</div>
                    </div>
                    <div class="settings-item-arrow">›</div>
                </div>
            </div>
            <div class="settings-group">
                <div class="settings-item" id="settings-password-item"><div class="settings-item-content">
                        <div class="settings-item-label">密码设置</div>
                        <div class="settings-item-desc">设置访问密码与验证码接收用户</div>
                    </div>
                    <div class="settings-item-arrow">›</div>
                </div>
            </div>
            <div class="settings-group">
                <div class="settings-item" id="settings-account-item"><div class="settings-item-content">
                        <div class="settings-item-label">账号管理</div>
                        <div class="settings-item-desc">修改密码与注销账号</div>
                    </div>
                    <div class="settings-item-arrow">›</div>
                </div>
            </div>
            <div class="settings-group">
                <div class="settings-item" id="settings-admin-item" style="display:none;"><div class="settings-item-content">
                        <div class="settings-item-label">管理员面板</div>
                        <div class="settings-item-desc">用户管理与系统公告</div>
                    </div>
                    <div class="settings-item-arrow">›</div>
                </div>
            </div>

            <div class="settings-group">
                <div class="settings-item" id="settings-about-item"><div class="settings-item-content">
                        <div class="settings-item-label">关于</div>
                        <div class="settings-item-desc">查看作者与版本信息</div>
                    </div>
                    <div class="settings-item-arrow">›</div>
                </div>
            </div>
        </div>
    </div>
    <div id="settings-api" class="settings-page">
        <div class="settings-nav-header">
            <button class="back-btn" id="api-back-btn">‹</button>
            <span class="nav-title">AI 回复设置</span>
        </div>
        <div class="settings-scroll">
            <div class="settings-body">
                <div class="setting-item">
                    <label class="setting-checkbox">
                        <input type="checkbox" id="ai-auto-reply"> 启用 AI 自动回复
                    </label>
                </div>
                <div class="setting-item">
                    <label class="setting-checkbox">
                        <input type="checkbox" id="ai-scheduled-reply"> 启用 AI 间隔定时回复
                    </label>
                </div>
                <div class="setting-item">
                    <label class="setting-checkbox">
                        <input type="checkbox" id="ai-daily-reply"> 启用 AI 每日定时回复
                    </label>
                </div>
                <div class="setting-item">
                    <label class="setting-label">每日发送时间</label>
                    <input type="time" id="daily-time" class="setting-input" value="09:00">
                </div>
                <div class="setting-item">
                    <label class="setting-label">API URL</label>
                    <input type="text" id="api-url" class="setting-input" placeholder="https://api.openai.com/v1/chat/completions">
                </div>
                <div class="setting-item">
                    <label class="setting-label">API Key（可选）</label>
                    <input type="password" id="api-key" class="setting-input" placeholder="无需API Key可留空">
                </div>
                <div class="setting-item">
                    <label class="setting-label">模型名称</label>
                    <input type="text" id="model-name" class="setting-input" placeholder="gpt-3.5-turbo">
                </div>
                <div class="setting-item">
                    <label class="setting-label">主动发送间隔(秒)</label>
                    <input type="number" id="active-interval" class="setting-input" value="60" min="10" max="3600">
                </div>
                <div class="setting-row">
                    <div class="setting-item">
                        <label class="setting-label">最少字数</label>
                        <input type="number" id="min-words" class="setting-input" value="10" min="5" max="500">
                    </div>
                    <div class="setting-item">
                        <label class="setting-label">最多字数</label>
                        <input type="number" id="max-words" class="setting-input" value="200" min="20" max="1000">
                    </div>
                </div>
                <div class="setting-item">
                    <label class="setting-label">系统提示词</label>
                    <textarea id="system-prompt" class="setting-input" rows="3" placeholder="你是一个微信聊天助手..."></textarea>
                </div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">识图 AI 设置</div>
                <div class="setting-item">
                    <label class="setting-checkbox">
                        <input type="checkbox" id="vision-enabled"> 启用识图功能
                    </label>
                </div>
                <div class="setting-item">
                    <label class="setting-label">识图 API URL</label>
                    <input type="text" id="vision-api-url" class="setting-input" placeholder="https://api.openai.com/v1/chat/completions">
                </div>
                <div class="setting-item">
                    <label class="setting-label">识图 API Key（可选）</label>
                    <input type="password" id="vision-api-key" class="setting-input" placeholder="无需API Key可留空">
                </div>
                <div class="setting-item">
                    <label class="setting-label">识图模型名称</label>
                    <input type="text" id="vision-model" class="setting-input" placeholder="gpt-4o">
                </div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">生图 AI 设置</div>
                <div class="setting-item">
                    <label class="setting-checkbox">
                        <input type="checkbox" id="image-gen-enabled"> 启用生图功能
                    </label>
                </div>
                <div class="setting-item">
                    <label class="setting-label">生图 API URL</label>
                    <input type="text" id="image-gen-api-url" class="setting-input" placeholder="https://api.openai.com/v1/images/generations">
                </div>
                <div class="setting-item">
                    <label class="setting-label">生图 API Key（可选）</label>
                    <input type="password" id="image-gen-api-key" class="setting-input" placeholder="无需API Key可留空">
                </div>
                <div class="setting-item">
                    <label class="setting-label">生图模型名称</label>
                    <input type="text" id="image-gen-model" class="setting-input" placeholder="dall-e-3">
                </div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">文件识别 AI 设置</div>
                <div class="setting-item">
                    <label class="setting-checkbox">
                        <input type="checkbox" id="file-recognize-enabled"> 启用文件识别
                    </label>
                </div>
                <div class="setting-item">
                    <label class="setting-label">文件识别 API URL</label>
                    <input type="text" id="file-recognize-api-url" class="setting-input" placeholder="https://api.openai.com/v1/chat/completions">
                </div>
                <div class="setting-item">
                    <label class="setting-label">文件识别 API Key（可选）</label>
                    <input type="password" id="file-recognize-api-key" class="setting-input" placeholder="无需API Key可留空">
                </div>
                <div class="setting-item">
                    <label class="setting-label">文件识别模型名称</label>
                    <input type="text" id="file-recognize-model" class="setting-input" placeholder="gpt-4o">
                </div>
                <div class="setting-item">
                    <label class="setting-label">最大文件大小(KB)</label>
                    <input type="number" id="file-recognize-max-size" class="setting-input" value="512" min="64" max="2048">
                </div>
                <div class="setting-item">
                    <label class="setting-checkbox">
                        <input type="checkbox" id="file-recognize-compat-mode"> 兼容模式（提取文字发给消息AI，无需单独配置识别API）
                    </label>
                </div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">AI回复间隔</div>
                <div class="setting-item">
                    <label class="setting-label">AI回复冷却时间(秒)</label>
                    <input type="number" id="ai-cooldown" class="setting-input" value="5" min="0" max="300">
                    <div style="font-size:12px;color:var(--text-hint);margin-top:4px;">对方短时间内发送多条消息时，AI只回复第一条，需等待此间隔后才再次回复。设为0则不限制。</div>
                </div>
                <button class="settings-save">保存设置</button>
            </div>
        </div>
    </div>
    <div id="settings-appearance" class="settings-page">
        <div class="settings-nav-header">
            <button class="back-btn" id="appearance-back-btn">‹</button>
            <span class="nav-title">外观设置</span>
        </div>
        <div class="settings-scroll">
            <div class="settings-body">
                <div class="setting-item">
                    <label class="setting-checkbox">
                        <input type="checkbox" id="appearance-animations"> 启用动画
                    </label>
                </div>
                <div class="setting-item" id="anim-speed-wrap">
                    <label class="setting-label">动画速度</label>
                    <div class="appearance-slider-wrap">
                        <input type="range" id="appearance-anim-speed" class="appearance-slider" min="0.2" max="3" step="0.1" value="1">
                        <span class="appearance-slider-val" id="anim-speed-val">1.0x</span>
                    </div>
                </div>

            </div>
        </div>
    </div>
    <div id="settings-announcement" class="settings-page">
        <div class="settings-nav-header">
            <button class="back-btn" id="announcement-back-btn">‹</button>
            <span class="nav-title">公告</span>
        </div>
        <div class="settings-scroll">
            <div class="settings-body">
                <div class="announcement-page-content" id="announcement-page-content">加载中...</div>
            </div>
        </div>
    </div>
    <div id="settings-about" class="settings-page">
        <div class="settings-nav-header">
            <button class="back-btn" id="about-back-btn">‹</button>
            <span class="nav-title">关于</span>
        </div>
        <div class="settings-scroll">
            <div class="settings-body">
                <div class="about-logo">
                    <img class="about-logo-img" src="https://ms.188850.xyz/file/default/1780497560383_retouch_2026050219212691.png" alt="作者头像" />
                    <div class="about-logo-name">Zyn iLink ChatBox</div>
                </div>
                <div class="about-info">
                    <div class="about-row">
                        <div class="about-label">作者</div>
                        <div class="about-value" id="about-author">加载中...</div>
                    </div>
                    <div class="about-row">
                        <div class="about-label">版本号</div>
                        <div class="about-value" id="about-version">加载中...</div>
                    </div>
                </div>
            </div>
        </div>
    </div>

    <div id="settings-password" class="settings-page">
        <div class="settings-nav-header">
            <button class="back-btn" id="password-back-btn">‹</button>
            <span class="nav-title">密码设置</span>
        </div>
        <div class="settings-scroll">
            <div class="settings-body">
                <div id="password-no-set" style="display:none;">
                    <div class="setting-item">
                        <label class="setting-label">设置访问密码</label>
                        <input type="password" id="password-new-input" class="setting-input" placeholder="输入密码（至少4位）" />
                    </div>
                    <div class="setting-item">
                        <label class="setting-label">确认密码</label>
                        <input type="password" id="password-confirm-input" class="setting-input" placeholder="再次输入密码" />
                    </div>
                    <button class="settings-save" id="password-set-btn">设置密码</button>
                    <div class="setting-divider"></div>
                    <div class="setting-section-title">验证码接收邮箱</div>
                    <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                        <div style="font-size:13px;color:var(--text-secondary);margin-bottom:4px;">用于接收验证码以重置访问密码</div>
                        <input type="email" id="reset-email-input-no-set" class="setting-input" placeholder="输入接收验证码的邮箱" />
                        <button class="settings-save" id="reset-email-save-btn-no-set" style="margin-top:4px;">保存邮箱</button>
                    </div>
                </div>
                <div id="password-already-set" style="display:none;">
                    <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:8px;">
                        <div style="font-size:14px;color:var(--text-secondary);">访问密码已设置</div>
                        <div style="display:flex;gap:8px;">
                            <button class="settings-save" id="password-change-btn" style="flex:1;margin-top:0;">修改密码</button>
                            <button class="settings-save" id="password-remove-btn" style="flex:1;margin-top:0;background:#FF3B30;box-shadow:none;">移除密码</button>
                        </div>
                    </div>
                    <div class="setting-divider"></div>
                    <div class="setting-section-title">验证码接收邮箱</div>
                    <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                        <div style="font-size:13px;color:var(--text-secondary);margin-bottom:4px;">用于接收验证码以重置访问密码</div>
                        <input type="email" id="reset-email-input" class="setting-input" placeholder="输入接收验证码的邮箱" />
                        <button class="settings-save" id="reset-email-save-btn" style="margin-top:4px;">保存邮箱</button>
                    </div>
                </div>
                <div id="password-change-form" style="display:none;">
                    <div class="setting-item">
                        <label class="setting-label">原密码</label>
                        <input type="password" id="password-old-input" class="setting-input" placeholder="输入原密码" />
                    </div>
                    <div class="setting-item">
                        <label class="setting-label">新密码</label>
                        <input type="password" id="password-change-new-input" class="setting-input" placeholder="输入新密码（至少4位）" />
                    </div>
                    <div class="setting-item">
                        <label class="setting-label">确认新密码</label>
                        <input type="password" id="password-change-confirm-input" class="setting-input" placeholder="再次输入新密码" />
                    </div>
                    <button class="settings-save" id="password-change-save-btn">保存新密码</button>
                    <button class="settings-save" id="password-change-cancel-btn" style="background:var(--bg-secondary);color:var(--text-primary);box-shadow:none;margin-top:4px;">取消</button>
                </div>
            </div>
        </div>
    </div>
    <div id="settings-account" class="settings-page">
        <div class="settings-nav-header">
            <button class="back-btn" id="account-back-btn">‹</button>
            <span class="nav-title">账号管理</span>
        </div>
        <div class="settings-scroll">
            <div class="settings-body">
                <div class="setting-section-title section-bar">绑定邮箱</div>
                <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                    <div id="account-email-status" style="font-size:14px;color:var(--text-secondary);margin-bottom:4px;margin-left:var(--card-mx);">未绑定邮箱</div>
                    <div style="margin:0 var(--card-mx);">
                        <div style="display:flex;gap:8px;align-items:center;">
                            <input type="email" id="account-bind-email" class="setting-input" placeholder="输入邮箱地址" style="flex:1;width:auto;" />
                        </div>
                        <div class="captcha-row" style="margin-top:8px;">
                            <input type="text" id="account-bind-captcha" class="captcha-input" placeholder="图形验证码" autocomplete="off" />
                            <img id="account-bind-captcha-img" class="captcha-img" src="" alt="验证码" />
                            <button id="account-bind-captcha-refresh" class="captcha-refresh">↻</button>
                        </div>
                        <div style="display:flex;gap:8px;align-items:center;margin-top:8px;">
                            <input type="text" id="account-bind-code" class="setting-input" placeholder="邮箱验证码" style="flex:0 0 120px;width:120px;" />
                            <button class="settings-save" id="account-send-bind-code-btn" style="white-space:nowrap;width:auto;padding:12px 16px;margin-top:0;">发送验证码</button>
                        </div>
                    </div>
                    <button class="settings-save" id="account-bind-email-btn" style="margin-top:8px;">绑定邮箱</button>
                </div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">当前账号</div>
                <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                    <div style="font-size:14px;color:var(--text-secondary);">当前登录账号</div>
                    <div id="account-current-username" style="font-size:16px;font-weight:600;color:var(--text-primary);"></div>
                </div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">修改密码</div>
                <div class="setting-item">
                    <label class="setting-label">原密码</label>
                    <input type="password" id="account-old-password" class="setting-input" placeholder="输入原密码" />
                </div>
                <div class="setting-item">
                    <label class="setting-label">新密码</label>
                    <input type="password" id="account-new-password" class="setting-input" placeholder="输入新密码（至少6位）" />
                </div>
                <div class="setting-item">
                    <label class="setting-label">确认新密码</label>
                    <input type="password" id="account-confirm-password" class="setting-input" placeholder="再次输入新密码" />
                </div>
                <button class="settings-save" id="account-change-password-btn">修改密码</button>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">注销账号</div>
                <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                    <div style="font-size:13px;color:#FF3B30;margin-bottom:4px;">注销后账号数据将无法恢复</div>
                    <input type="password" id="account-delete-password" class="setting-input" placeholder="输入密码以确认注销" />
                </div>
                <button class="settings-save" id="account-delete-btn" style="background:#FF3B30;">注销账号</button>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">退出登录</div>
                <button class="settings-save" id="account-logout-btn" style="background:#FF9500;">退出登录</button>
            </div>
        </div>
    </div>
    <div id="settings-admin" class="settings-page">
        <div class="settings-nav-header">
            <button class="back-btn" id="admin-back-btn">‹</button>
            <span class="nav-title">管理员面板</span>
        </div>
        <div class="settings-scroll">
            <div class="settings-body">
                <div class="setting-section-title section-bar">所有用户</div>
                <div id="admin-panel-user-list" style="margin-bottom:16px;max-height:300px;overflow-y:auto;"></div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">注销用户</div>
                <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                    <div style="font-size:13px;color:#FF3B30;margin-bottom:4px;">注销后账号数据将无法恢复</div>
                    <input type="number" id="admin-delete-index" class="setting-input" placeholder="输入用户序号" min="1" />
                </div>
                <button class="settings-save" id="admin-delete-user-btn" style="background:#FF3B30;">注销用户</button>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">强制下线</div>
                <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                    <div style="font-size:13px;color:#FF9500;margin-bottom:4px;">强制用户下线并删除指纹登录</div>
                    <input type="number" id="admin-force-offline-index" class="setting-input" placeholder="输入用户序号" min="1" />
                </div>
                <button class="settings-save" id="admin-force-offline-btn" style="background:#FF9500;">强制下线</button>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">发送系统公告</div>
                <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                    <textarea id="admin-announcement-content" class="user-prompt-textarea" rows="4" placeholder="输入公告内容..."></textarea>
                </div>
                <button class="settings-save" id="admin-send-announcement-btn">发送公告</button>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">历史公告</div>
                <div id="admin-announcement-list" style="margin-bottom:16px;"></div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">封禁IP</div>
                <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:4px;">
                    <div style="font-size:13px;color:#FF3B30;margin-bottom:4px;">封禁后该IP无法访问任何页面</div>
                    <div style="display:flex;gap:8px;align-items:center;">
                        <input type="text" id="admin-ban-ip-input" class="setting-input" placeholder="输入IP地址" style="flex:1;width:auto;" />
                        <button class="settings-save" id="admin-ban-ip-btn" style="background:#FF3B30;white-space:nowrap;width:auto;padding:12px 20px;margin-top:0;">封禁</button>
                    </div>
                </div>
                <div id="admin-banned-ip-list" style="margin-bottom:16px;"></div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">邮箱注册</div>
                <div class="setting-item" style="justify-content:space-between;">
                    <div style="font-size:15px;color:var(--text-primary);">启用邮箱注册</div>
                    <div class="chat-menu-switch" id="admin-email-register-switch"><div class="chat-menu-switch-knob"></div></div>
                </div>
                <div id="admin-email-config-section" style="display:none;">
                </div>
                <div class="setting-divider"></div>
                <div class="setting-section-title section-bar">Cloudflare Tunnel</div>
                <div class="setting-item" style="flex-direction:column;align-items:stretch;gap:8px;">
                    <div style="display:flex;align-items:center;gap:8px;">
                        <span id="cf-status-dot" style="width:10px;height:10px;border-radius:50%;background:#999;flex-shrink:0;"></span>
                        <span id="cf-status-text" style="font-size:14px;color:var(--text-secondary);">未运行</span>
                    </div>
                    <div style="display:flex;align-items:center;gap:8px;">
                        <div style="font-size:12px;color:var(--text-secondary);word-break:break-all;flex:1;" id="cf-url-display"></div>
                        <button id="cf-copy-btn" style="display:none;flex-shrink:0;padding:6px 12px;border-radius:8px;background:var(--accent);color:#fff;border:none;font-size:12px;cursor:pointer;transition:transform 0.2s,background 0.2s;white-space:nowrap;" onclick="var u=document.getElementById('cf-url-display').textContent.replace('隧道地址: ','');if(u){navigator.clipboard.writeText(u).then(function(){this.textContent='已复制';var t=this;setTimeout(function(){t.textContent='复制'},2000)}.bind(this)).catch(function(){})}">复制</button>
                    </div>
                    <div style="display:flex;gap:8px;align-items:center;">
                        <span style="font-size:13px;color:var(--text-secondary);white-space:nowrap;">本地端口</span>
                        <input type="number" id="cf-port-input" class="setting-input" value="1145" style="flex:0 0 100px;width:100px;" min="1" max="65535" />
                        <span style="font-size:11px;color:var(--text-hint);">cloudflared tunnel --url http://localhost:&lt;端口&gt;</span>
                    </div>
                    <div style="display:flex;gap:8px;">
                        <button class="settings-save" id="cf-start-btn" style="flex:1;padding:12px 20px;margin-top:0;">启动隧道</button>
                        <button class="settings-save" id="cf-stop-btn" style="flex:1;padding:12px 20px;margin-top:0;background:#FF3B30;">停止隧道</button>
                    </div>
                </div>
            </div>
        </div>
    </div>
</div>
<div id="ai-modal" class="ai-modal">
    <div class="ai-modal-content">
        <div class="ai-modal-header">
            <span>Zyn AI 回复助手</span>
            <button class="ai-modal-close" id="ai-modal-close">×</button>
        </div>
        <div class="ai-modal-body">
            <div class="ai-modal-label">原消息：</div>
            <div class="ai-modal-msg-preview" id="ai-modal-msg-preview"></div>
            <label class="ai-modal-label">回复要求（可选）：</label>
            <textarea id="ai-instruction" class="ai-instruction-input" rows="3" placeholder="例如：帮我反驳他、用温和的语气回复、加个表情包、怼回去..."></textarea>
        </div>
        <div class="ai-modal-footer">
            <button class="ai-modal-btn cancel" id="ai-modal-cancel">取消</button>
            <button class="ai-modal-btn send" id="ai-modal-send">发送 AI 回复</button>
        </div>
    </div>
</div>
<div id="captcha-modal" class="nickname-modal" style="z-index:100001;">
    <div class="nickname-modal-content">
        <div class="nickname-modal-title">图形验证码</div>
        <div class="captcha-row" style="margin-bottom:12px;">
            <input type="text" id="captcha-modal-input" class="captcha-input" placeholder="请输入图形验证码" autocomplete="off" />
            <img id="captcha-modal-img" class="captcha-img" src="" alt="验证码" />
            <button id="captcha-modal-refresh" class="captcha-refresh">↻</button>
        </div>
        <div class="nickname-modal-btns">
            <button id="captcha-modal-cancel" class="nickname-modal-btn cancel">取消</button>
            <button id="captcha-modal-confirm" class="nickname-modal-btn save">确认</button>
        </div>
    </div>
</div>
<div id="toast" class="toast"></div>
<div id="media-upload-progress" class="media-upload-progress">
    <div class="media-upload-box">
        <div class="media-upload-spinner"></div>
        <div class="media-upload-text">正在发送...</div>
    </div>
</div>
<div id="nickname-modal" class="nickname-modal">
    <div class="nickname-modal-content">
        <div class="nickname-modal-title">设置备注名</div>
        <div id="nickname-modal-userid" class="nickname-modal-userid"></div>
        <input type="text" id="nickname-input" class="nickname-modal-input" placeholder="输入备注名..." />
        <div class="nickname-modal-btns">
            <button id="nickname-cancel-btn" class="nickname-modal-btn cancel">取消</button>
            <button id="nickname-save-btn" class="nickname-modal-btn save">保存</button>
        </div>
    </div>
</div>
<div id="user-prompt-modal" class="nickname-modal">
    <div class="nickname-modal-content">
        <div class="nickname-modal-title">设置AI提示词</div>
        <div id="user-prompt-modal-userid" class="nickname-modal-userid"></div>
        <div id="user-prompt-modal-default" class="nickname-modal-userid" style="margin-bottom:8px"></div>
        <textarea id="user-prompt-input" class="user-prompt-textarea" placeholder="留空则使用默认提示词..."></textarea>
        <div class="nickname-modal-btns">
            <button id="user-prompt-cancel-btn" class="nickname-modal-btn cancel">取消</button>
            <button id="user-prompt-save-btn" class="nickname-modal-btn save">保存</button>
        </div>
    </div>
</div>
<div id="add-user-modal" class="add-user-modal">
    <div class="add-user-modal-content">
        <div class="add-user-modal-title">添加新用户</div>
        <div class="add-user-modal-status" id="add-user-status">正在生成二维码...</div>
        <div class="add-user-modal-qr" id="add-user-qr"></div>
        <div style="margin-top:10px;padding:10px 14px;background:rgba(255,59,48,0.12);border:1px solid rgba(255,59,48,0.3);border-radius:10px;font-size:13px;color:#FF3B30;text-align:center;line-height:1.5;font-weight:600;">微信用户扫码后必须先发一两条消息才能连接上!</div>
        <button class="add-user-modal-close" id="add-user-close-btn">关闭</button>
    </div>
</div>
<div id="announcement-modal" class="announcement-modal">
    <div class="announcement-modal-content">
        <div class="announcement-modal-title" id="announcement-modal-title">系统公告</div>
        <div class="announcement-modal-body" id="announcement-body">加载中...</div>
        <div class="announcement-modal-footer">
            <button class="announcement-btn announcement-btn-dismiss" id="announcement-dismiss-btn">以后不再显示</button>
            <button class="announcement-btn announcement-btn-confirm" id="announcement-confirm-btn">确认</button>
        </div>
    </div>
</div>
<script>
''' + bot._generate_wasm_wrapper(session_token, password_required) + '''
</script>
</body>
</html>'''
                self.send_response(200)
                self.send_header('Content-type', 'text/html; charset=utf-8')
                self.send_header('Set-Cookie', 'session_token=' + session_token + '; Path=/; SameSite=Lax; HttpOnly')
                self.send_header('X-Content-Type-Options', 'nosniff')
                self.send_header('X-Frame-Options', 'DENY')
                self.send_header('X-XSS-Protection', '1; mode=block')
                self.end_headers()
                self.wfile.write(html.encode('utf-8'))
            
            def _serve_status(self):
                account = self._resolve_account()
                if account:
                    with account._msg_lock:
                        msg_count = len(account._messages)
                    status = {
                        'logged_in': account.token is not None,
                        'login_done': account._login_done,
                        'current_user': account._current_user,
                        'bot_id': account.bot_id,
                        'user_count': len(account._context_tokens),
                        'users': list(account._context_tokens.keys()),
                        'message_count': msg_count,
                        'account_username': account.username,
                        'is_admin': account.is_admin
                    }
                else:
                    with bot._msg_lock:
                        msg_count = len(bot._messages)
                    status = {
                        'logged_in': bot.token is not None,
                        'login_done': bot._login_done,
                        'current_user': bot._current_user,
                        'bot_id': bot.bot_id,
                        'user_count': len(bot._context_tokens),
                        'users': list(bot._context_tokens.keys()),
                        'message_count': msg_count,
                        'is_admin': False
                    }
                self._send_json(status)
            
            def _serve_qrcode(self):
                account = self._resolve_account()
                target = account if account else bot
                if target._login_done and target.token:
                    self._send_json({
                        'error': 'already_logged_in',
                        'message': '已连接',
                        'login_done': True,
                        'redirect_to_chat': True
                    })
                    return
                if not target._qrcode_matrix:
                    self._send_json({'error': 'no_qrcode', 'message': '正在获取二维码...'})
                    return
                qr_data = {
                    'matrix': target._qrcode_matrix,
                    'qrcode_key': target._qrcode_key,
                    'login_done': target._login_done
                }
                self._send_json(qr_data)
            
            def _serve_messages(self):
                account = self._resolve_account()
                target = account if account else bot
                params = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
                since = params.get('since', [None])[0]
                user_filter = params.get('user', [None])[0]
                since_id = 0
                if since:
                    try:
                        since_id = int(since)
                    except (ValueError, TypeError):
                        since_id = 0
                messages = []
                with target._msg_lock:
                    msgs_snapshot = list(target._messages)
                for msg in msgs_snapshot:
                    if since and msg.get('id', 0) <= since_id:
                        continue
                    if user_filter:
                        if msg.get('type') == 'in' and msg.get('from') != user_filter:
                            continue
                        if msg.get('type') == 'out' and msg.get('to') != user_filter:
                            continue
                    msg_copy = dict(msg)
                    target._enrich_msg_with_cache_id(msg_copy)
                    messages.append(msg_copy)
                self._send_json({
                    'messages': messages,
                    'current_user': target._current_user
                })
            
            def _serve_users(self):
                account = self._resolve_account()
                target = account if account else bot
                users = []
                for uid in target._context_tokens:
                    users.append(uid)
                self._send_json({'users': users, 'current_user': target._current_user})
            
            def _serve_ai_config(self):
                account = self._resolve_account()
                target = account if account else bot
                safe_config = {
                    "auto_reply": target.ai_config.get("auto_reply"),
                    "scheduled_reply": target.ai_config.get("scheduled_reply"),
                    "api_url": target.ai_config.get("api_url", ""),
                    "api_key": "********" if target.ai_config.get("api_key") else "",
                    "active_interval": target.ai_config.get("active_interval"),
                    "model": target.ai_config.get("model"),
                    "min_words": target.ai_config.get("min_words"),
                    "max_words": target.ai_config.get("max_words"),
                    "system_prompt": target.ai_config.get("system_prompt"),
                    "vision_api_url": target.ai_config.get("vision_api_url", ""),
                    "vision_api_key": "********" if target.ai_config.get("vision_api_key") else "",
                    "vision_model": target.ai_config.get("vision_model", "gpt-4o"),
                    "vision_enabled": target.ai_config.get("vision_enabled", False),
                    "image_gen_api_url": target.ai_config.get("image_gen_api_url", ""),
                    "image_gen_api_key": "********" if target.ai_config.get("image_gen_api_key") else "",
                    "image_gen_model": target.ai_config.get("image_gen_model", "dall-e-3"),
                    "image_gen_enabled": target.ai_config.get("image_gen_enabled", False),
                    "file_recognize_enabled": target.ai_config.get("file_recognize_enabled", False),
                    "file_recognize_api_url": target.ai_config.get("file_recognize_api_url", ""),
                    "file_recognize_api_key": "********" if target.ai_config.get("file_recognize_api_key") else "",
                    "file_recognize_model": target.ai_config.get("file_recognize_model", "gpt-4o"),
                    "file_recognize_max_size": target.ai_config.get("file_recognize_max_size", 512),
                    "file_recognize_compat_mode": target.ai_config.get("file_recognize_compat_mode", False),
                    "ai_cooldown": target.ai_config.get("ai_cooldown", 5),
                    "daily_reply": target.ai_config.get("daily_reply", False),
                    "daily_time": target.ai_config.get("daily_time", "09:00")
                }
                self._send_json(safe_config)

            def _serve_about(self):
                self._send_json({
                    "version": bot.SCRIPT_VERSION,
                    "author": bot.AUTHOR_NAME
                })
            
            def _serve_announcement(self):
                local_announcements = bot._system_announcements[-10:] if bot._system_announcements else []
                account = self._resolve_account()
                is_admin = account.is_admin if account else False
                try:
                    url = "https://gitee.com/zynsync/zyn-i-link-chat-box/raw/master/Gg.txt"
                    req = urllib.request.Request(url)
                    with urllib.request.urlopen(req, timeout=10) as resp:
                        content = resp.read().decode('utf-8')
                    self._send_json({"content": content, "system_announcements": local_announcements, "is_admin": is_admin})
                except Exception as e:
                    self._send_json({"content": "", "system_announcements": local_announcements, "error": str(e), "is_admin": is_admin})
            
            def _serve_tutorial(self):
                try:
                    url = "https://gitee.com/zynsync/zyn-i-link-chat-box/raw/master/jc.txt"
                    req = urllib.request.Request(url)
                    with urllib.request.urlopen(req, timeout=10) as resp:
                        content = resp.read().decode('utf-8')
                    self._send_json({"content": content})
                except Exception as e:
                    self._send_json({"content": "", "error": str(e)})

            def _serve_remote_version(self):
                try:
                    url = "https://gitee.com/zynsync/zyn-i-link-chat-box/raw/master/version.txt"
                    req = urllib.request.Request(url)
                    with urllib.request.urlopen(req, timeout=10) as resp:
                        version = resp.read().decode('utf-8').strip()
                    self._send_json({"version": version})
                except Exception as e:
                    self._send_json({"version": "", "error": str(e)})
            
            def _serve_cached_media(self, cache_key):
                try:
                    if not cache_key or not all(c in '0123456789abcdef' for c in cache_key.lower()):
                        self.send_error(400)
                        return
                    params = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
                    user_param = params.get('user', [None])[0]
                    account = self._resolve_account()
                    target = account if account else bot
                    cached = None
                    if user_param:
                        cached = target._get_user_cached_media(user_param, cache_key)
                    if not cached:
                        cached = target._get_cached_media(cache_key)
                    if not cached:
                        self.send_error(404)
                        return
                    media_data, mime, filename = cached
                    is_download = params.get('download', [''])[0] == '1'
                    self.send_response(200)
                    self.send_header('Content-Type', mime)
                    self.send_header('Content-Length', str(len(media_data)))
                    self.send_header('Cache-Control', 'public, max-age=31536000')
                    self.send_header('X-Content-Type-Options', 'nosniff')
                    origin = self.headers.get('Origin', '')
                    if origin and origin.startswith(('http://localhost', 'http://127.0.0.1')):
                        self.send_header('Access-Control-Allow-Origin', origin)
                    if filename:
                        safe_filename = filename.replace('"', '\\"').replace('\n', '').replace('\r', '')
                        disposition = 'attachment' if is_download else 'inline'
                        self.send_header('Content-Disposition', disposition + '; filename="' + safe_filename + '"')
                    self.end_headers()
                    self.wfile.write(media_data)
                except BrokenPipeError:
                    pass
                except Exception as e:
                    print(f"[WEB] 缓存媒体服务异常: {e}")
            
            def _handle_save_ai_config(self, data):
                try:
                    account = self._resolve_account()
                    target = account if account else bot
                    print(f"[WEB] 收到 AI 配置保存请求: auto_reply={data.get('auto_reply')}, scheduled_reply={data.get('scheduled_reply')}, api_url={data.get('api_url', '')[:50]}, api_key={'已设置' if data.get('api_key') else '未设置'}")
                    target.ai_config["auto_reply"] = bool(data.get("auto_reply", False))
                    target.ai_config["scheduled_reply"] = bool(data.get("scheduled_reply", False))
                    target.ai_config["api_url"] = str(data.get("api_url", ""))
                    new_api_key = data.get("api_key", "")
                    if new_api_key and new_api_key != "********":
                        target.ai_config["api_key"] = str(new_api_key)
                    target.ai_config["model"] = str(data.get("model", "gpt-3.5-turbo"))
                    try:
                        target.ai_config["active_interval"] = max(1, int(data.get("active_interval", 60)))
                    except (ValueError, TypeError):
                        target.ai_config["active_interval"] = 60
                    try:
                        target.ai_config["min_words"] = max(1, int(data.get("min_words", 10)))
                    except (ValueError, TypeError):
                        target.ai_config["min_words"] = 10
                    try:
                        target.ai_config["max_words"] = max(1, int(data.get("max_words", 200)))
                    except (ValueError, TypeError):
                        target.ai_config["max_words"] = 200
                    target.ai_config["system_prompt"] = str(data.get("system_prompt", ""))
                    target.ai_config["vision_api_url"] = str(data.get("vision_api_url", ""))
                    new_vision_key = data.get("vision_api_key", "")
                    if new_vision_key and new_vision_key != "********":
                        target.ai_config["vision_api_key"] = str(new_vision_key)
                    target.ai_config["vision_model"] = str(data.get("vision_model", "gpt-4o"))
                    target.ai_config["vision_enabled"] = bool(data.get("vision_enabled", False))
                    target.ai_config["image_gen_api_url"] = str(data.get("image_gen_api_url", ""))
                    new_image_gen_key = data.get("image_gen_api_key", "")
                    if new_image_gen_key and new_image_gen_key != "********":
                        target.ai_config["image_gen_api_key"] = str(new_image_gen_key)
                    target.ai_config["image_gen_model"] = str(data.get("image_gen_model", "dall-e-3"))
                    target.ai_config["image_gen_enabled"] = bool(data.get("image_gen_enabled", False))
                    target.ai_config["file_recognize_enabled"] = bool(data.get("file_recognize_enabled", False))
                    target.ai_config["file_recognize_api_url"] = str(data.get("file_recognize_api_url", ""))
                    new_file_rec_key = data.get("file_recognize_api_key", "")
                    if new_file_rec_key and new_file_rec_key != "********":
                        target.ai_config["file_recognize_api_key"] = str(new_file_rec_key)
                    target.ai_config["file_recognize_model"] = str(data.get("file_recognize_model", "gpt-4o"))
                    try:
                        target.ai_config["file_recognize_max_size"] = max(1, int(data.get("file_recognize_max_size", 512)))
                    except (ValueError, TypeError):
                        target.ai_config["file_recognize_max_size"] = 512
                    target.ai_config["file_recognize_compat_mode"] = bool(data.get("file_recognize_compat_mode", False))
                    try:
                        target.ai_config["ai_cooldown"] = max(0, int(data.get("ai_cooldown", 5)))
                    except (ValueError, TypeError):
                        target.ai_config["ai_cooldown"] = 5
                    target.ai_config["daily_reply"] = bool(data.get("daily_reply", False))
                    target.ai_config["daily_time"] = str(data.get("daily_time", "09:00"))
                    target._save_ai_config()
                    if data.get('scheduled_reply'):
                        for user_id in target._context_tokens.keys():
                            if account:
                                bot._schedule_active_message_for_account(account, user_id)
                            else:
                                bot._schedule_active_message(user_id)
                    else:
                        for user_id in list(target._active_timers.keys()):
                            if not target.is_scheduled_enabled_for_user(user_id):
                                timer = target._active_timers.pop(user_id, None)
                                if timer:
                                    timer.cancel()
                    if data.get('daily_reply'):
                        for user_id in target._context_tokens.keys():
                            if account:
                                bot._schedule_daily_message_for_account(account, user_id)
                            else:
                                bot._schedule_daily_message(user_id)
                    else:
                        for user_id in list(target._daily_timers.keys()):
                            if not target.is_daily_enabled_for_user(user_id):
                                timer = target._daily_timers.pop(user_id, None)
                                if timer:
                                    timer.cancel()
                    self._send_json({'success': True, 'config': target.ai_config})
                except Exception as e:
                    print(f"[WEB] 保存 AI 配置失败: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _handle_send(self, data):
                try:
                    text = data.get('text', '').strip()
                    account = self._resolve_account()
                    target = account if account else bot
                    print(f"[WEB] 收到发送请求: text='{text}', current_user={target._current_user}")
                    if not text:
                        self._send_json({'success': False, 'error': '消息不能为空'})
                        return
                    if not target._current_user:
                        self._send_json({'success': False, 'error': '没有选择用户'})
                        return
                    if account:
                        success = bot.send_text_for_account(account, target._current_user, text)
                    else:
                        bot.send_typing(target._current_user)
                        success = bot.send_text(target._current_user, text)
                    if success:
                        self._send_json({'success': True, 'message': {'text': text, 'time': datetime.now().strftime('%H:%M:%S'), 'type': 'out'}})
                    else:
                        self._send_json({'success': False, 'error': '发送失败'})
                except Exception as e:
                    print(f"[WEB] 发送异常: {e}")
                    self._send_json({'success': False, 'error': '发送失败'})
            
            def _handle_send_typing(self, data):
                try:
                    account = self._resolve_account()
                    target = account if account else bot
                    if not target._current_user:
                        self._send_json({'success': False, 'error': '没有选择用户'})
                        return
                    if account:
                        bot.send_typing_for_account(account, target._current_user)
                    else:
                        bot.send_typing(target._current_user)
                    self._send_json({'success': True})
                except Exception as e:
                    self._send_json({'success': False, 'error': '发送typing失败'})
            
            def _handle_send_media(self, data):
                try:
                    media_type = data.get('media_type', '')
                    filename = data.get('filename', 'file')
                    file_data_b64 = data.get('file_data', '')
                    thumbnail_b64 = data.get('thumbnail', '')
                    account = self._resolve_account()
                    target = account if account else bot
                    if not file_data_b64:
                        self._send_json({'success': False, 'error': '文件数据为空'})
                        return
                    if not target._current_user:
                        self._send_json({'success': False, 'error': '没有选择用户'})
                        return
                    try:
                        file_bytes = base64.b64decode(file_data_b64)
                    except Exception:
                        self._send_json({'success': False, 'error': '文件数据解码失败'})
                        return
                    print(f"[WEB] 收到媒体发送请求: type={media_type}, filename={filename}, size={len(file_bytes)} bytes, user={target._current_user}")
                    success = False
                    media_type_int = 0
                    media_data_url = ""
                    if media_type == 'image':
                        if thumbnail_b64:
                            media_data_url = 'data:image/jpeg;base64,' + thumbnail_b64
                        if account:
                            success = bot.send_image_for_account(account, target._current_user, file_bytes, filename, media_data=media_data_url)
                        else:
                            success = bot.send_image(target._current_user, file_bytes, filename, media_data=media_data_url)
                        media_type_int = 2
                    elif media_type == 'video':
                        if thumbnail_b64:
                            media_data_url = 'data:image/jpeg;base64,' + thumbnail_b64
                        if account:
                            success = bot.send_video_for_account(account, target._current_user, file_bytes, filename, media_data=media_data_url)
                        else:
                            success = bot.send_video(target._current_user, file_bytes, filename, media_data=media_data_url)
                        media_type_int = 5
                    elif media_type == 'file':
                        if account:
                            success = bot.send_file_for_account(account, target._current_user, file_bytes, filename)
                        else:
                            success = bot.send_file(target._current_user, file_bytes, filename)
                        media_type_int = 4
                    else:
                        self._send_json({'success': False, 'error': '不支持的媒体类型'})
                        return
                    if success:
                        type_name = bot.MEDIA_TYPE_NAMES.get(media_type_int, "文件")
                        msg_data = {
                            'text': f'[{type_name}] {filename}',
                            'time': datetime.now().strftime('%H:%M:%S'),
                            'type': 'out',
                            'media_type': media_type_int,
                            'media_filename': filename,
                            'media_data': media_data_url
                        }
                        if target._messages:
                            for m in reversed(target._messages):
                                if m.get('type') == 'out' and m.get('media_type') == media_type_int and not m.get('media_cache_id'):
                                    msg_data['id'] = m.get('id')
                                    if file_bytes and media_type_int in (2, 5):
                                        mime = bot._detect_mime(file_bytes)
                                        if mime == 'application/octet-stream':
                                            mime = 'video/mp4' if media_type_int == 5 else 'image/jpeg'
                                        cache_key = hashlib.md5(file_bytes).hexdigest()
                                        target._save_media_cache(cache_key, file_bytes, mime, filename)
                                        msg_data['media_cache_id'] = cache_key
                                        m['media_cache_id'] = cache_key
                                        if m.get('media_cdn'):
                                            try:
                                                cdn_info = json.loads(m['media_cdn']) if isinstance(m['media_cdn'], str) else m['media_cdn']
                                                cdn_cache_key = hashlib.md5((cdn_info.get("encrypt_query_param") or cdn_info.get("encrypted_query_param") or "").encode('utf-8')).hexdigest()
                                                if cdn_cache_key != cache_key:
                                                    target._save_media_cache(cdn_cache_key, file_bytes, mime, filename)
                                            except Exception:
                                                pass
                                    break
                        self._send_json({'success': True, 'message': msg_data})
                    else:
                        self._send_json({'success': False, 'error': '媒体发送失败'})
                except Exception as e:
                    print(f"[WEB] 媒体发送异常: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _handle_download_media(self, data):
                try:
                    cdn_info_str = data.get('cdn_info', '')
                    if not cdn_info_str:
                        self._send_json({'success': False, 'error': '缺少 CDN 信息'})
                        return
                    if isinstance(cdn_info_str, dict):
                        cdn_info = cdn_info_str
                    else:
                        try:
                            cdn_info = json.loads(cdn_info_str)
                        except (json.JSONDecodeError, TypeError) as je:
                            print(f"[WEB] CDN 信息 JSON 解析失败: {je}, raw={str(cdn_info_str)[:200]}")
                            self._send_json({'success': False, 'error': 'CDN 信息格式错误'})
                            return
                    account = self._resolve_account()
                    target = account if account else bot
                    cache_key = hashlib.md5((cdn_info.get("encrypt_query_param") or cdn_info.get("encrypted_query_param") or "").encode('utf-8')).hexdigest()
                    if account:
                        media_data = bot.download_media_for_account(account, cdn_info)
                    else:
                        media_data = bot.download_media(cdn_info)
                    if media_data:
                        mime = bot._detect_mime(media_data)
                        self._send_json({
                            'success': True,
                            'cache_key': cache_key,
                            'mime': mime
                        })
                    else:
                        self._send_json({'success': False, 'error': '下载失败'})
                except Exception as e:
                    print(f"[WEB] 媒体下载异常: {e}")
                    self._send_json({'success': False, 'error': '操作失败'})

            def _handle_switch_user(self, data):
                try:
                    user_id = str(data.get('user_id', ''))
                    account = self._resolve_account()
                    target = account if account else bot
                    if user_id and user_id in target._context_tokens:
                        target.set_current_user(user_id)
                        self._send_json({'success': True, 'current_user': user_id})
                    else:
                        self._send_json({'success': False, 'error': '无效的用户'})
                except Exception as e:
                    self._send_json({'success': False, 'error': '操作失败'})
            
            def _handle_delete_user(self, data):
                try:
                    user_id = str(data.get('user_id', ''))
                    account = self._resolve_account()
                    target = account if account else bot
                    if not user_id:
                        self._send_json({'success': False, 'error': '用户ID不能为空'})
                        return
                    if user_id not in target._context_tokens:
                        self._send_json({'success': False, 'error': '用户不存在'})
                        return
                    success = target.remove_user(user_id)
                    if success:
                        self._send_json({
                            'success': True,
                            'current_user': target._current_user,
                            'users': list(target._context_tokens.keys())
                        })
                    else:
                        self._send_json({'success': False, 'error': '删除失败'})
                except Exception as e:
                    self._send_json({'success': False, 'error': '操作失败'})

            def _serve_auth_check(self):
                session_token = self._get_session_token()
                account = bot._get_account_from_session(session_token) if session_token else None
                if account:
                    self._send_json({'password_required': False, 'account_username': account.username})
                    return
                if bot._is_web_password_set():
                    if session_token and session_token in bot._verified_sessions and time.time() <= bot._verified_sessions[session_token]:
                        self._send_json({'password_required': False})
                    else:
                        if session_token in bot._verified_sessions:
                            del bot._verified_sessions[session_token]
                        self._send_json({'password_required': True})
                else:
                    self._send_json({'password_required': False})

            def _serve_sse(self):
                session_token = self._get_session_token()
                if not session_token:
                    parsed = urllib.parse.urlparse(self.path)
                    params = urllib.parse.parse_qs(parsed.query)
                    t = params.get('token', [None])[0]
                    if t and bot._verify_session_token(t):
                        session_token = t
                account = bot._get_account_from_session(session_token) if session_token else None
                if not account:
                    self.send_response(401)
                    self.end_headers()
                    return
                username = account.username
                try:
                    self.send_response(200)
                    self.send_header('Content-Type', 'text/event-stream')
                    self.send_header('Cache-Control', 'no-cache')
                    self.send_header('Connection', 'keep-alive')
                    self.send_header('X-Accel-Buffering', 'no')
                    self.end_headers()
                    self.wfile.write(b"event: connected\ndata: {}\n\n")
                    self.wfile.flush()
                    bot._sse_add(username, self.wfile)
                    while True:
                        time.sleep(30)
                        try:
                            self.wfile.write(b"event: ping\ndata: {}\n\n")
                            self.wfile.flush()
                        except Exception:
                            break
                except Exception:
                    pass
                finally:
                    bot._sse_remove(username, self.wfile)

            def _serve_admin_user(self):
                account = self._resolve_account()
                target = account if account else bot
                captcha_admin = bot._web_password_config.get('captcha_admin', '')
                reset_email = bot._web_password_config.get('reset_email', '')
                users = list(target._context_tokens.keys())
                account_email = account.email if account else ''
                self._send_json({'admin_user': captcha_admin, 'reset_email': reset_email, 'users': users, 'password_set': bot._is_web_password_set(), 'email_register_enabled': bot._web_password_config.get('email_register_enabled', False), 'account_email': account_email})

            def _handle_bot_register(self, data):
                username = str(data.get('username', '')).strip()
                password = str(data.get('password', ''))
                email = str(data.get('email', '')).strip()
                email_code = str(data.get('email_code', '')).strip()
                server_ip = self._get_client_ip()
                frontend_ip = str(data.get('client_ip', '')).strip()
                client_ip = frontend_ip if frontend_ip and server_ip in ('127.0.0.1', '::1', 'localhost') else server_ip
                if bot._web_password_config.get('email_register_enabled', False):
                    if not email:
                        self._send_json({'success': False, 'error': '请输入邮箱'})
                        return
                    if not email_code:
                        self._send_json({'success': False, 'error': '请输入邮箱验证码'})
                        return
                    stored = bot._email_verification_codes.get(email)
                    if not stored or stored.get('code') != email_code:
                        self._send_json({'success': False, 'error': '邮箱验证码错误'})
                        return
                    if stored.get('expiry', 0) < time.time():
                        self._send_json({'success': False, 'error': '邮箱验证码已过期'})
                        return
                    for acc in bot._accounts.values():
                        if acc.email and acc.email == email:
                            self._send_json({'success': False, 'error': '该邮箱已被注册'})
                            return
                    bot._email_verification_codes.pop(email, None)
                    result = bot._register_account(username, password, client_ip, email)
                else:
                    result = bot._register_account(username, password, client_ip)
                if result.get('success'):
                    new_token = result.get('session_token')
                    cookie = f'session_token={new_token}; Path=/; SameSite=Lax; HttpOnly; Max-Age=2592000' if new_token else None
                    self._send_json({'success': True, 'session_token': new_token, 'username': result.get('username')}, cookies=[cookie] if cookie else None)
                else:
                    self._send_json(result)

            def _handle_bot_login(self, data):
                username = str(data.get('username', '')).strip()
                password = str(data.get('password', ''))
                fingerprint = str(data.get('fingerprint', ''))
                server_ip = self._get_client_ip()
                frontend_ip = str(data.get('client_ip', '')).strip()
                client_ip = frontend_ip if frontend_ip and server_ip in ('127.0.0.1', '::1', 'localhost') else server_ip
                result = bot._login_account(username, password, fingerprint, client_ip)
                if result.get('success'):
                    new_token = result.get('session_token')
                    cookie = f'session_token={new_token}; Path=/; SameSite=Lax; HttpOnly; Max-Age=2592000' if new_token else None
                    self._send_json({'success': True, 'session_token': new_token, 'username': result.get('username')}, cookies=[cookie] if cookie else None)
                else:
                    self._send_json(result)

            def _handle_bot_fingerprint_login(self, data):
                fingerprint = str(data.get('fingerprint', ''))
                server_ip = self._get_client_ip()
                frontend_ip = str(data.get('client_ip', '')).strip()
                client_ip = frontend_ip if frontend_ip and server_ip in ('127.0.0.1', '::1', 'localhost') else server_ip
                result = bot._fingerprint_login(fingerprint, client_ip)
                if result.get('success'):
                    new_token = result.get('session_token')
                    cookie = f'session_token={new_token}; Path=/; SameSite=Lax; HttpOnly; Max-Age=2592000' if new_token else None
                    self._send_json({'success': True, 'session_token': new_token, 'username': result.get('username')}, cookies=[cookie] if cookie else None)
                else:
                    self._send_json(result)

            def _handle_auth_login(self, data):
                password = str(data.get('password', ''))
                if len(password) > 128:
                    self._send_json({'success': False, 'error': '密码过长'})
                    return
                session_token = self.headers.get('X-Session-Token')
                if not session_token:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                t = part.split('=', 1)[1]
                                if t:
                                    session_token = t
                                    break
                if not session_token or not bot._verify_session_token(session_token):
                    self._send_json({'success': False, 'error': '无效会话'})
                    return
                client_ip = self._get_client_ip()
                now = time.time()
                attempt_info = bot._login_attempts.get(client_ip, {'count': 0, 'lockout': 0})
                if now < attempt_info.get('lockout', 0):
                    remaining = int(attempt_info['lockout'] - now)
                    self._send_json({'success': False, 'error': f'登录尝试过多，请{remaining}秒后再试'})
                    return
                if now >= attempt_info.get('lockout', 0) and attempt_info.get('count', 0) > 0:
                    attempt_info['count'] = 0
                if not bot._is_web_password_set():
                    bot._verified_sessions[session_token] = time.time() + 86400
                    self._send_json({'success': True})
                    return
                if bot._verify_web_password(password):
                    bot._verified_sessions[session_token] = time.time() + 86400
                    bot._login_attempts.pop(client_ip, None)
                    self._send_json({'success': True})
                else:
                    attempt_info['count'] = attempt_info.get('count', 0) + 1
                    if attempt_info['count'] >= 5:
                        attempt_info['lockout'] = now + 300
                        attempt_info['count'] = 0
                    bot._login_attempts[client_ip] = attempt_info
                    self._send_json({'success': False, 'error': '密码错误'})

            def _handle_web_password_verify(self, data):
                session_token = self._get_session_token()
                if not self._check_auth():
                    self._send_json({'success': False, 'error': '未授权'})
                    return
                client_ip = self._get_client_ip()
                now = time.time()
                verify_info = bot._login_attempts.get(client_ip + '_verify', {'count': 0, 'lockout': 0})
                if now < verify_info.get('lockout', 0):
                    self._send_json({'success': False, 'error': '验证尝试过多，请稍后再试'})
                    return
                password = str(data.get('password', ''))
                if len(password) > 128:
                    self._send_json({'success': False, 'error': '密码过长'})
                    return
                if bot._verify_web_password(password):
                    bot._login_attempts.pop(client_ip + '_verify', None)
                    bot._verified_sessions[session_token] = time.time() + 86400
                    self._send_json({'success': True})
                else:
                    verify_info['count'] = verify_info.get('count', 0) + 1
                    if verify_info['count'] >= 5:
                        verify_info['lockout'] = now + 300
                        verify_info['count'] = 0
                    bot._login_attempts[client_ip + '_verify'] = verify_info
                    self._send_json({'success': False, 'error': '密码错误'})

            def _handle_set_web_password(self, data):
                action = data.get('action', '')
                if action == 'set':
                    password = data.get('password', '')
                    if not password or len(password) < 6:
                        self._send_json({'success': False, 'error': '密码长度不能少于6位'})
                        return
                    if len(password) > 128:
                        self._send_json({'success': False, 'error': '密码长度不能超过128位'})
                        return
                    bot._web_password_config['password_hash'] = bot._hash_password(password)
                    bot._save_web_password_config()
                    self._send_json({'success': True})
                elif action == 'change':
                    old_password = data.get('old_password', '')
                    new_password = data.get('new_password', '')
                    if not bot._verify_web_password(old_password):
                        self._send_json({'success': False, 'error': '原密码错误'})
                        return
                    if not new_password or len(new_password) < 6:
                        self._send_json({'success': False, 'error': '新密码长度不能少于6位'})
                        return
                    bot._web_password_config['password_hash'] = bot._hash_password(new_password)
                    bot._save_web_password_config()
                    bot._verified_sessions.clear()
                    self._send_json({'success': True})
                elif action == 'remove':
                    password = data.get('password', '')
                    if not bot._verify_web_password(password):
                        self._send_json({'success': False, 'error': '密码错误'})
                        return
                    bot._web_password_config.pop('password_hash', None)
                    bot._save_web_password_config()
                    bot._verified_sessions.clear()
                    self._send_json({'success': True})
                else:
                    self._send_json({'success': False, 'error': '无效操作'})

            def _handle_send_verify_code(self, data):
                client_ip = self._get_client_ip()
                now = time.time()
                send_info = bot._send_code_attempts.get(client_ip, {'count': 0, 'lockout': 0})
                if now < send_info.get('lockout', 0):
                    self._send_json({'success': False, 'error': '发送过于频繁，请稍后再试'})
                    return
                if now >= send_info.get('lockout', 0) and send_info.get('count', 0) > 0:
                    send_info['count'] = 0
                captcha_input = data.get('captcha', '').strip().upper()
                session_token = self.headers.get('X-Session-Token')
                if not session_token:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                t = part.split('=', 1)[1]
                                if t:
                                    session_token = t
                                    break
                captcha_info = bot._captcha_store.get(session_token) if session_token else None
                if not captcha_info:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                cookie_token = part.split('=', 1)[1]
                                if cookie_token and cookie_token != session_token and bot._captcha_store.get(cookie_token):
                                    captcha_info = bot._captcha_store[cookie_token]
                                    session_token = cookie_token
                if not captcha_input or not captcha_info or time.time() > captcha_info['expiry'] or captcha_info['answer'] != captcha_input:
                    if captcha_info and time.time() > captcha_info['expiry']:
                        bot._captcha_store.pop(session_token, None)
                    elif captcha_info:
                        bot._captcha_store.pop(session_token, None)
                    self._send_json({'success': False, 'error': '图形验证码错误'})
                    return
                bot._captcha_store.pop(session_token, None)
                reset_email = bot._web_password_config.get('reset_email', '')
                if not reset_email:
                    self._send_json({'success': False, 'error': '未设置验证码接收邮箱，无法发送验证码'})
                    return
                code = bot._generate_verification_code(client_ip)
                bot._send_forgot_password_email(reset_email, code, 'web-password')
                send_info['count'] = send_info.get('count', 0) + 1
                if send_info['count'] >= 3:
                    send_info['lockout'] = now + 300
                    send_info['count'] = 0
                bot._send_code_attempts[client_ip] = send_info
                self._send_json({'success': True, 'message': '验证码已发送至配置邮箱'})

            def _handle_reset_password(self, data):
                code = data.get('code', '')
                new_password = data.get('new_password', '')
                client_ip = self._get_client_ip()
                if not code:
                    self._send_json({'success': False, 'error': '请输入验证码'})
                    return
                code_info = bot._verification_codes.get(code)
                if not code_info:
                    self._send_json({'success': False, 'error': '验证码无效'})
                    return
                if code_info['used']:
                    self._send_json({'success': False, 'error': '验证码已使用'})
                    return
                if time.time() > code_info['expiry']:
                    bot._verification_codes.pop(code, None)
                    self._send_json({'success': False, 'error': '验证码已过期'})
                    return
                if code_info.get('ip') and code_info['ip'] != client_ip:
                    self._send_json({'success': False, 'error': '验证码无效'})
                    return
                if not new_password or len(new_password) < 6:
                    self._send_json({'success': False, 'error': '新密码长度不能少于6位'})
                    return
                if len(new_password) > 128:
                    self._send_json({'success': False, 'error': '新密码长度不能超过128位'})
                    return
                code_info['used'] = True
                bot._verification_codes.pop(code, None)
                bot._web_password_config['password_hash'] = bot._hash_password(new_password)
                bot._save_web_password_config()
                bot._verified_sessions.clear()
                self._send_json({'success': True})

            def _handle_bot_send_reset_code(self, data):
                username = str(data.get('username', '')).strip()
                captcha_input = data.get('captcha', '').strip().upper()
                client_ip = self._get_client_ip()
                now = time.time()
                send_info = bot._send_code_attempts.get(client_ip + '_bot', {'count': 0, 'lockout': 0})
                if now < send_info.get('lockout', 0):
                    self._send_json({'success': False, 'error': '发送过于频繁，请稍后再试'})
                    return
                if now >= send_info.get('lockout', 0) and send_info.get('count', 0) > 0:
                    send_info['count'] = 0
                session_token = self.headers.get('X-Session-Token')
                if not session_token:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                t = part.split('=', 1)[1]
                                if t:
                                    session_token = t
                                    break
                captcha_info = bot._captcha_store.get(session_token) if session_token else None
                if not captcha_info:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                cookie_token = part.split('=', 1)[1]
                                if cookie_token and cookie_token != session_token and bot._captcha_store.get(cookie_token):
                                    captcha_info = bot._captcha_store[cookie_token]
                                    session_token = cookie_token
                if not captcha_input or not captcha_info or time.time() > captcha_info['expiry'] or captcha_info['answer'] != captcha_input:
                    if captcha_info:
                        bot._captcha_store.pop(session_token, None)
                    self._send_json({'success': False, 'error': '图形验证码错误'})
                    return
                bot._captcha_store.pop(session_token, None)
                if not username:
                    self._send_json({'success': False, 'error': '请输入用户名'})
                    return
                account = bot._accounts.get(username)
                if not account:
                    self._send_json({'success': False, 'error': '用户名不存在'})
                    return
                if not account.email:
                    self._send_json({'success': False, 'error': '该账号未绑定邮箱，无法发送验证码'})
                    return
                import secrets
                code = str(secrets.randbelow(900000) + 100000)
                bot._account_reset_codes[code] = {"username": username, "expiry": time.time() + 300, "used": False, "ip": client_ip}
                expired_codes = [c for c, v in bot._account_reset_codes.items() if v["expiry"] < time.time() or v["used"]]
                for c in expired_codes:
                    del bot._account_reset_codes[c]
                bot._send_forgot_password_email(account.email, code, username)
                send_info['count'] = send_info.get('count', 0) + 1
                if send_info['count'] >= 3:
                    send_info['lockout'] = now + 300
                    send_info['count'] = 0
                bot._send_code_attempts[client_ip + '_bot'] = send_info
                self._send_json({'success': True, 'message': '验证码已发送至绑定邮箱'})

            def _handle_bot_reset_password(self, data):
                username = str(data.get('username', '')).strip()
                code = data.get('code', '')
                new_password = data.get('new_password', '')
                client_ip = self._get_client_ip()
                result = bot._reset_account_password(username, code, new_password, client_ip)
                self._send_json(result)

            def _handle_bot_change_password(self, data):
                account = self._resolve_account()
                if not account:
                    self._send_json({'success': False, 'error': '未登录'})
                    return
                old_password = str(data.get('old_password', ''))
                new_password = str(data.get('new_password', ''))
                result = bot._change_account_password(account.username, old_password, new_password)
                self._send_json(result)

            def _handle_bot_delete_account(self, data):
                account = self._resolve_account()
                if not account:
                    self._send_json({'success': False, 'error': '未登录'})
                    return
                password = str(data.get('password', ''))
                result = bot._delete_account(account.username, password)
                self._send_json(result)
                account = self._resolve_account()
                target = account if account else bot
                admin_user = bot._web_password_config.get('captcha_admin', '')
                users = list(target._context_tokens.keys())
                self._send_json({'admin_user': admin_user, 'users': users, 'password_set': bot._is_web_password_set()})

            def _handle_bot_logout(self, data):
                session_token = self._get_session_token()
                if session_token:
                    username = bot._account_sessions.pop(session_token, None)
                    bot._verified_sessions.pop(session_token, None)
                    if username:
                        for fp, uname in list(bot._fingerprint_sessions.items()):
                            if uname == username:
                                del bot._fingerprint_sessions[fp]
                        print(f"[多账户] 账户退出登录: {username}")
                self._send_json({'success': True})

            def _handle_admin_list_users(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                online_usernames = set()
                for tok, uname in bot._account_sessions.items():
                    online_usernames.add(uname)
                for tok, uname in bot._fingerprint_sessions.items():
                    online_usernames.add(uname)
                users = []
                for username, acc in bot._accounts.items():
                    users.append({
                        'username': username,
                        'online': username in online_usernames,
                        'session_token': '',
                        'login_time': '',
                        'last_ip': acc.last_ip or ''
                    })
                for tok, uname in list(bot._account_sessions.items()):
                    for u in users:
                        if u['username'] == uname and not u['session_token']:
                            u['session_token'] = tok[:8]
                            break
                self._send_json({
                    'success': True,
                    'users': users,
                    'announcements': bot._system_announcements[-10:] if bot._system_announcements else [],
                    'banned_ips': list(bot._banned_ips)
                })

            def _handle_admin_delete_user(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                index = data.get('index')
                if index is None:
                    self._send_json({'success': False, 'error': '请输入序号'})
                    return
                try:
                    idx = int(index)
                except (ValueError, TypeError):
                    self._send_json({'success': False, 'error': '序号必须为数字'})
                    return
                uname_list = list(bot._accounts.keys())
                if idx < 1 or idx > len(uname_list):
                    self._send_json({'success': False, 'error': '无效序号'})
                    return
                username = uname_list[idx - 1]
                if username == account.username:
                    self._send_json({'success': False, 'error': '不能注销自己的账号'})
                    return
                old_sessions = [tok for tok, uname in bot._account_sessions.items() if uname == username]
                for tok in old_sessions:
                    bot._account_sessions.pop(tok, None)
                    bot._verified_sessions.pop(tok, None)
                for fp, uname in list(bot._fingerprint_sessions.items()):
                    if uname == username:
                        del bot._fingerprint_sessions[fp]
                del bot._accounts[username]
                bot._save_accounts()
                try:
                    acc_dir = Path(f"accounts/{username}")
                    if acc_dir.exists():
                        shutil.rmtree(str(acc_dir))
                except Exception:
                    pass
                print(f"[管理员] 管理员删除账户: {username}")
                self._send_json({'success': True})

            def _handle_admin_send_announcement(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                content = str(data.get('content', '')).strip()
                if not content:
                    self._send_json({'success': False, 'error': '公告内容不能为空'})
                    return
                announcement = {
                    'content': content,
                    'time': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                    'from_admin': account.username,
                    'id': uuid.uuid4().hex[:12]
                }
                bot._system_announcements.append(announcement)
                bot._save_system_announcements()
                for uname, acc in bot._accounts.items():
                    if acc.is_admin:
                        continue
                    for uid in list(acc._context_tokens.keys()):
                        try:
                            bot.send_text_for_account(acc, uid, f"【系统公告】{content}")
                        except Exception:
                            pass
                self._send_json({'success': True})

            def _handle_admin_delete_announcement(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                index = data.get('index')
                if index is None:
                    self._send_json({'success': False, 'error': '请指定公告索引'})
                    return
                try:
                    idx = int(index)
                except (ValueError, TypeError):
                    self._send_json({'success': False, 'error': '索引必须为数字'})
                    return
                if idx < 0 or idx >= len(bot._system_announcements):
                    self._send_json({'success': False, 'error': '无效索引'})
                    return
                bot._system_announcements.pop(idx)
                bot._save_system_announcements()
                self._send_json({'success': True})

            def _handle_admin_force_offline(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                index = data.get('index')
                if index is None:
                    self._send_json({'success': False, 'error': '请输入序号'})
                    return
                try:
                    idx = int(index)
                except (ValueError, TypeError):
                    self._send_json({'success': False, 'error': '序号必须为数字'})
                    return
                uname_list = list(bot._accounts.keys())
                if idx < 1 or idx > len(uname_list):
                    self._send_json({'success': False, 'error': '无效序号'})
                    return
                username = uname_list[idx - 1]
                if username == account.username:
                    self._send_json({'success': False, 'error': '不能强制下线自己'})
                    return
                old_sessions = [tok for tok, uname in bot._account_sessions.items() if uname == username]
                bot._sse_send(username, "force_offline", {"reason": "你已被管理员强制下线"})
                for tok in old_sessions:
                    bot._account_sessions.pop(tok, None)
                    bot._verified_sessions.pop(tok, None)
                    bot._session_tokens.pop(tok, None)
                for fp, uname in list(bot._fingerprint_sessions.items()):
                    if uname == username:
                        del bot._fingerprint_sessions[fp]
                print(f"[管理员] 管理员强制下线用户: {username}")
                self._send_json({'success': True})

            def _handle_admin_ban_ip(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                ip = str(data.get('ip', '')).strip()
                if not ip:
                    self._send_json({'success': False, 'error': '请输入IP地址'})
                    return
                if ip in bot._banned_ips:
                    self._send_json({'success': False, 'error': '该IP已在封禁列表中'})
                    return
                bot._banned_ips.append(ip)
                bot._save_banned_ips()
                print(f"[管理员] 封禁IP: {ip}")
                self._send_json({'success': True})

            def _handle_admin_unban_ip(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                ip = str(data.get('ip', '')).strip()
                if not ip:
                    self._send_json({'success': False, 'error': '请输入IP地址'})
                    return
                if ip not in bot._banned_ips:
                    self._send_json({'success': False, 'error': '该IP不在封禁列表中'})
                    return
                bot._banned_ips.remove(ip)
                bot._save_banned_ips()
                print(f"[管理员] 解封IP: {ip}")
                self._send_json({'success': True})

            def _handle_admin_cloudflared(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                action = str(data.get('action', ''))
                if action == 'status':
                    running = bot._cf_process is not None and bot._cf_process.poll() is None
                    self._send_json({'success': True, 'running': running, 'url': bot._cf_url or ''})
                elif action == 'start':
                    if bot._cf_process is not None and bot._cf_process.poll() is None:
                        self._send_json({'success': False, 'error': 'Cloudflare Tunnel 已在运行中'})
                        return
                    try:
                        port = data.get('port', 1145)
                        bot._cf_process = subprocess.Popen(
                            ['cloudflared', 'tunnel', '--url', f'http://localhost:{port}'],
                            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                            text=True, bufsize=1
                        )
                        import threading as _t
                        def _read_output():
                            import re as _re
                            for line in iter(bot._cf_process.stderr.readline, ''):
                                _original_print(f"[CF-TUNNEL] {line.rstrip()}")
                                m = _re.search(r'https://[a-zA-Z0-9-]+\.trycloudflare\.com', line)
                                if m:
                                    bot._cf_url = m.group(0)
                                    _original_print(f"[CF-TUNNEL] 隧道地址: {bot._cf_url}")
                        _t.Thread(target=_read_output, daemon=True).start()
                        _original_print(f"[管理员] Cloudflare Tunnel 已启动")
                        self._send_json({'success': True, 'message': 'Cloudflare Tunnel 已启动'})
                    except Exception as e:
                        bot._cf_process = None
                        _original_print(f"[管理员] Cloudflare Tunnel 启动失败: {e}")
                        self._send_json({'success': False, 'error': f'启动失败: {e}'})
                elif action == 'stop':
                    if bot._cf_process is None:
                        self._send_json({'success': False, 'error': 'Cloudflare Tunnel 未在运行'})
                        return
                    try:
                        bot._cf_process.terminate()
                        bot._cf_process.wait(timeout=5)
                    except Exception:
                        try:
                            bot._cf_process.kill()
                        except Exception:
                            pass
                    bot._cf_process = None
                    bot._cf_url = ""
                    print(f"[管理员] Cloudflare Tunnel 已停止")
                    self._send_json({'success': True, 'message': 'Cloudflare Tunnel 已停止'})
                else:
                    self._send_json({'success': False, 'error': '未知操作'})

            def _handle_set_admin_user(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                user_id = str(data.get('user_id', ''))
                if user_id:
                    target_acc = bot._accounts.get(user_id)
                    if not target_acc:
                        self._send_json({'success': False, 'error': '用户不存在'})
                        return
                    for acc in bot._accounts.values():
                        acc.is_admin = False
                    target_acc.is_admin = True
                else:
                    for acc in bot._accounts.values():
                        acc.is_admin = False
                bot._save_accounts()
                self._send_json({'success': True})

            def _handle_set_captcha_admin(self, data):
                user_id = str(data.get('user_id', ''))
                account = self._resolve_account()
                target = account if account else bot
                if user_id and user_id not in target._context_tokens:
                    self._send_json({'success': False, 'error': '用户不存在于当前会话中'})
                    return
                if user_id:
                    bot._web_password_config['captcha_admin'] = user_id
                else:
                    bot._web_password_config.pop('captcha_admin', None)
                bot._save_web_password_config()
                self._send_json({'success': True})

            def _handle_set_reset_email(self, data):
                email = str(data.get('email', '')).strip()
                if email:
                    bot._web_password_config['reset_email'] = email
                else:
                    bot._web_password_config.pop('reset_email', None)
                bot._save_web_password_config()
                self._send_json({'success': True})

            def _handle_set_wechat_notify(self, data):
                user_id = str(data.get('user_id', ''))
                account = self._resolve_account()
                target = account if account else bot
                if user_id and user_id not in target._context_tokens:
                    self._send_json({'success': False, 'error': '用户不存在于当前会话中'})
                    return
                if user_id:
                    bot._web_password_config['wechat_notify_user'] = user_id
                else:
                    bot._web_password_config.pop('wechat_notify_user', None)
                bot._save_web_password_config()
                self._send_json({'success': True})

            def _handle_send_email_code(self, data):
                email = str(data.get('email', '')).strip()
                captcha_input = data.get('captcha', '').strip().upper()
                if not email:
                    self._send_json({'success': False, 'error': '请输入邮箱'})
                    return
                if not captcha_input:
                    self._send_json({'success': False, 'error': '请输入图形验证码'})
                    return
                session_token = self.headers.get('X-Session-Token')
                if not session_token:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                t = part.split('=', 1)[1]
                                if t:
                                    session_token = t
                                    break
                captcha_info = bot._captcha_store.get(session_token) if session_token else None
                if not captcha_info:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                cookie_token = part.split('=', 1)[1]
                                if cookie_token and cookie_token != session_token and bot._captcha_store.get(cookie_token):
                                    captcha_info = bot._captcha_store[cookie_token]
                                    session_token = cookie_token
                if not captcha_input or not captcha_info or time.time() > captcha_info['expiry'] or captcha_info['answer'] != captcha_input:
                    if captcha_info:
                        bot._captcha_store.pop(session_token, None)
                    self._send_json({'success': False, 'error': '图形验证码错误'})
                    return
                bot._captcha_store.pop(session_token, None)
                if not bot._web_password_config.get('email_register_enabled', False):
                    self._send_json({'success': False, 'error': '邮箱注册未启用'})
                    return
                now = time.time()
                last_sent = bot._email_code_cooldown.get(email, 0)
                if now - last_sent < 60:
                    self._send_json({'success': False, 'error': f'请{int(60 - (now - last_sent))}秒后再试'})
                    return
                for acc in bot._accounts.values():
                    if acc.email and acc.email == email:
                        self._send_json({'success': False, 'error': '该邮箱已被注册'})
                        return
                code = bot._generate_verification_code()
                bot._email_verification_codes[email] = {'code': code, 'expiry': now + 300}
                bot._email_code_cooldown[email] = now
                bot._send_email_verification(email, code)
                self._send_json({'success': True})

            def _handle_set_email_register(self, data):
                account = self._resolve_account()
                if not account or not account.is_admin:
                    self._send_json({'success': False, 'error': '无权限'})
                    return
                enabled = bool(data.get('enabled', False))
                bot._web_password_config['email_register_enabled'] = enabled
                bot._save_web_password_config()
                self._send_json({'success': True})

            def _handle_send_bind_email_code(self, data):
                account = self._resolve_account()
                if not account:
                    self._send_json({'success': False, 'error': '未登录'})
                    return
                email = str(data.get('email', '')).strip()
                captcha_input = data.get('captcha', '').strip().upper()
                if not email:
                    self._send_json({'success': False, 'error': '请输入邮箱'})
                    return
                if not captcha_input:
                    self._send_json({'success': False, 'error': '请输入图形验证码'})
                    return
                session_token = self.headers.get('X-Session-Token')
                if not session_token:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                t = part.split('=', 1)[1]
                                if t:
                                    session_token = t
                                    break
                captcha_info = bot._captcha_store.get(session_token) if session_token else None
                if not captcha_info:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                cookie_token = part.split('=', 1)[1]
                                if cookie_token and cookie_token != session_token and bot._captcha_store.get(cookie_token):
                                    captcha_info = bot._captcha_store[cookie_token]
                                    session_token = cookie_token
                if not captcha_input or not captcha_info or time.time() > captcha_info['expiry'] or captcha_info['answer'] != captcha_input:
                    if captcha_info:
                        bot._captcha_store.pop(session_token, None)
                    self._send_json({'success': False, 'error': '图形验证码错误'})
                    return
                bot._captcha_store.pop(session_token, None)
                for acc in bot._accounts.values():
                    if acc.email and acc.email == email and acc.username != account.username:
                        self._send_json({'success': False, 'error': '该邮箱已被其他账号绑定'})
                        return
                now = time.time()
                last_sent = bot._email_code_cooldown.get(email + '_bind', 0)
                if now - last_sent < 60:
                    self._send_json({'success': False, 'error': f'请{int(60 - (now - last_sent))}秒后再试'})
                    return
                code = bot._generate_verification_code()
                bot._email_bind_codes[email] = {'code': code, 'expiry': now + 300, 'username': account.username}
                bot._email_code_cooldown[email + '_bind'] = now
                bot._send_bind_email_verification(email, code)
                self._send_json({'success': True})

            def _handle_bind_email(self, data):
                account = self._resolve_account()
                if not account:
                    self._send_json({'success': False, 'error': '未登录'})
                    return
                email = str(data.get('email', '')).strip()
                code = str(data.get('code', '')).strip()
                if not email:
                    self._send_json({'success': False, 'error': '请输入邮箱'})
                    return
                if not code:
                    self._send_json({'success': False, 'error': '请输入验证码'})
                    return
                stored = bot._email_bind_codes.get(email)
                if not stored or stored.get('code') != code:
                    self._send_json({'success': False, 'error': '验证码错误'})
                    return
                if stored.get('expiry', 0) < time.time():
                    bot._email_bind_codes.pop(email, None)
                    self._send_json({'success': False, 'error': '验证码已过期'})
                    return
                if stored.get('username') != account.username:
                    self._send_json({'success': False, 'error': '验证码无效'})
                    return
                for acc in bot._accounts.values():
                    if acc.email and acc.email == email and acc.username != account.username:
                        self._send_json({'success': False, 'error': '该邮箱已被其他账号绑定'})
                        return
                bot._email_bind_codes.pop(email, None)
                account.email = email
                bot._save_accounts()
                self._send_json({'success': True})

            def _handle_account_email(self, data):
                account = self._resolve_account()
                if not account:
                    self._send_json({'success': False, 'error': '未登录'})
                    return
                self._send_json({'success': True, 'email': account.email or ''})

            def _handle_email_register_status(self):
                enabled = bool(bot._web_password_config.get('email_register_enabled', False))
                self._send_json({'enabled': enabled})

            def _serve_history(self):
                try:
                    params = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
                    user_id = params.get('user', [None])[0]
                    limit_str = params.get('limit', ['200'])[0]
                    try:
                        limit = min(int(limit_str), 500)
                    except (ValueError, TypeError):
                        limit = 200
                    account = self._resolve_account()
                    target = account if account else bot
                    if user_id:
                        history_msgs = target.get_user_messages(user_id, limit)
                    else:
                        with target._msg_lock:
                            all_msgs = list(target._messages) if target._messages else []
                        history_msgs = all_msgs[-limit:]
                    enriched = []
                    for msg in history_msgs:
                        msg_copy = dict(msg)
                        target._enrich_msg_with_cache_id(msg_copy)
                        enriched.append(msg_copy)
                    with target._msg_lock:
                        total_count = len(target._messages)
                    self._send_json({
                        'messages': enriched,
                        'total': total_count,
                        'found': len(history_msgs),
                        'user_id': user_id or '',
                        'limit': limit
                    })
                except Exception as e:
                    self._send_json({
                        'messages': [],
                        'total': 0,
                        'found': 0,
                        'user_id': '',
                        'limit': 200
                    })
            
            def _serve_captcha(self):
                parsed_qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
                session_token = parsed_qs.get('token', [None])[0]
                print(f"[DEBUG] serve-captcha: token from URL query={session_token!r}")
                if not session_token:
                    session_token = self.headers.get('X-Session-Token')
                    print(f"[DEBUG] serve-captcha: token from X-Session-Token header={session_token!r}")
                if not session_token:
                    cookie_header = self.headers.get('Cookie', '')
                    if cookie_header:
                        for part in cookie_header.split(';'):
                            part = part.strip()
                            if part.startswith('session_token='):
                                t = part.split('=', 1)[1]
                                if t:
                                    session_token = t
                                    break
                    print(f"[DEBUG] serve-captcha: token from Cookie={session_token!r}")
                if not session_token or not bot._verify_session_token(session_token):
                    if session_token:
                        bot._session_tokens[session_token] = time.time() + 3600
                    else:
                        session_token = bot._generate_session_token()
                print(f"[DEBUG] serve-captcha: final session_token={session_token!r}, captcha_store_keys={list(bot._captcha_store.keys())}")
                try:
                    content_type, img_data = bot._generate_captcha(session_token)
                    print(f"[DEBUG] serve-captcha: generated captcha for token={session_token!r}, answer={bot._captcha_store.get(session_token, {}).get('answer')!r}, store_keys_after={list(bot._captcha_store.keys())}")
                    self.send_response(200)
                    self.send_header('Content-type', content_type)
                    self.send_header('Cache-Control', 'no-cache, no-store, must-revalidate')
                    self.send_header('Set-Cookie', f'session_token={session_token}; Path=/; SameSite=Lax; HttpOnly')
                    self.end_headers()
                    self.wfile.write(img_data)
                except Exception:
                    self.send_error(500)

            def _serve_lock_page(self):
                session_token = None
                cookie_header = self.headers.get('Cookie', '')
                if cookie_header:
                    for part in cookie_header.split(';'):
                        part = part.strip()
                        if part.startswith('session_token='):
                            t = part.split('=', 1)[1]
                            if t and bot._verify_session_token(t):
                                session_token = t
                                break
                if not session_token:
                    session_token = bot._generate_session_token()
                html = '''<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0, maximum-scale=1.0, user-scalable=no, viewport-fit=cover">
<title>Zyn iLink ChatBox</title>
<style>
* { margin: 0; padding: 0; box-sizing: border-box; -webkit-tap-highlight-color: transparent; }
html, body { font-family: "SF Pro Display", "SF Pro Text", "PingFang SC", "Hiragino Sans GB", "Microsoft YaHei", system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; font-weight: 450; -webkit-font-smoothing: antialiased; -moz-osx-font-smoothing: grayscale; text-rendering: optimizeLegibility; height: 100%; background: var(--bg-primary); color: var(--text-primary); }
:root { --bg-primary: #FFFFFF; --bg-secondary: #F2F2F6; --accent: #0A84FF; --accent-hover: #0973E0; --accent-light: rgba(10,132,255,0.08); --text-primary: #1C1C1E; --text-secondary: #8E8E93; --text-hint: #C6C6C8; --divider: #E5E5EA; --chat-bg: #F2F2F6; }
[data-theme="dark"] { --bg-primary: #2C2C2E; --bg-secondary: #1C1C1E; --accent: #0A84FF; --accent-hover: #0973E0; --accent-light: rgba(10,132,255,0.15); --text-primary: #F5F5F7; --text-secondary: #8E8E93; --text-hint: #636366; --divider: #38383A; --chat-bg: #1C1C1E; }
.lock-screen { position: fixed; top: 0; left: 0; right: 0; bottom: 0; background: var(--chat-bg); z-index: 100000; display: flex; flex-direction: column; align-items: center; justify-content: center; padding: 20px; overflow-y: auto; }
.lock-screen-logo { font-size: 48px; margin-bottom: 16px; }
.lock-screen-title { font-size: 24px; font-weight: 700; color: var(--text-primary); margin-bottom: 8px; }
.lock-screen-subtitle { font-size: 14px; color: var(--text-secondary); margin-bottom: 32px; }
.lock-screen-form { width: 100%; max-width: 320px; }
.lock-screen-input { width: 100%; height: 48px; border: 1px solid var(--divider); border-radius: var(--card-round); padding: 0 20px; font-size: 16px; outline: none; background: var(--bg-secondary); color: var(--text-primary); margin-bottom: 12px; text-align: center; letter-spacing: 4px; box-sizing: border-box; }
.lock-screen-input:focus { border-color: var(--accent); box-shadow: 0 4px 20px rgba(10,132,255,0.15), inset 0 0 0 1.5px var(--accent); }
.lock-screen-btn { width: 100%; height: 48px; border-radius: var(--card-round); border: none; background: var(--accent); color: #fff; font-size: 16px; font-weight: 600; cursor: pointer; }
.lock-screen-btn:active { transform: scale(0.96); }
.lock-screen-btn:disabled { opacity: 0.5; cursor: not-allowed; }
.lock-screen-forgot { margin-top: 16px; text-align: center; }
.lock-screen-forgot-btn { background: none; border: none; color: var(--accent); font-size: 14px; cursor: pointer; padding: 8px 16px; }
.lock-screen-forgot-btn:active { opacity: 0.7; }
.lock-screen-error { color: #FF3B30; font-size: 13px; text-align: center; margin-bottom: 8px; min-height: 18px; }
.lock-screen-reset { width: 100%; max-width: 320px; }
.lock-screen-reset-input { width: 100%; height: 44px; border: none; border-radius: 12px; padding: 0 16px; font-size: 15px; outline: none; background: var(--bg-secondary); color: var(--text-primary); margin-bottom: 10px; text-align: center; }
.lock-screen-reset-input:focus { background: var(--bg-primary); box-shadow: 0 0 0 1.5px var(--accent); }
.lock-screen-reset-btn { width: 100%; height: 44px; border-radius: 12px; border: none; background: var(--accent); color: #fff; font-size: 15px; font-weight: 600; cursor: pointer; transition: transform 0.2s; margin-bottom: 8px; }
.lock-screen-reset-btn:active { transform: scale(0.96); }
.lock-screen-reset-btn:disabled { opacity: 0.5; cursor: not-allowed; }
.lock-screen-back-btn { background: none; border: none; color: var(--text-secondary); font-size: 14px; cursor: pointer; padding: 8px 16px; margin-top: 4px; }
.lock-screen-back-btn:active { opacity: 0.7; }
.captcha-row { display: flex; align-items: center; gap: 6px; width: 100%; margin-bottom: 10px; }
.captcha-input { flex: 1; min-width: 0; height: 44px; border: 1px solid var(--divider); border-radius: var(--card-round); padding: 0 12px; font-size: 15px; outline: none; background: var(--bg-secondary); color: var(--text-primary); text-align: center; letter-spacing: 2px; box-sizing: border-box; }
.captcha-input:focus { background: var(--bg-primary); border-color: var(--accent); box-shadow: 0 0 0 1.5px var(--accent); }
.captcha-img { height: 44px; max-width: 100px; border-radius: 8px; cursor: pointer; flex-shrink: 0; }
.captcha-refresh { width: 32px; height: 44px; border: none; border-radius: 10px; background: var(--bg-primary); color: var(--text-secondary); font-size: 18px; cursor: pointer; display: flex; align-items: center; justify-content: center; flex-shrink: 0; box-shadow: 0 2px 8px rgba(0,0,0,0.04), inset 0 0 0 0.5px var(--divider); }
.captcha-refresh:active { transform: scale(0.96); }
@media (max-width: 480px) { .lock-screen { padding: 16px; } .lock-screen-logo { font-size: 36px; margin-bottom: 12px; } .lock-screen-title { font-size: 20px; } .lock-screen-subtitle { font-size: 12px; margin-bottom: 24px; } .lock-screen-input { height: 44px; font-size: 15px; } .lock-screen-btn { height: 44px; font-size: 15px; } .lock-screen-reset-input { height: 40px; font-size: 14px; } .lock-screen-reset-btn { height: 40px; font-size: 14px; } .captcha-row { gap: 4px; } .captcha-input { padding: 0 8px; font-size: 14px; } .captcha-img { height: 40px; max-width: 80px; } .captcha-refresh { width: 28px; height: 40px; font-size: 16px; } }
</style>
</head>
<body>
<div id="lock-screen" class="lock-screen">
    <div class="lock-screen-logo">ZynSync</div>
    <div class="lock-screen-title">Zyn ChatBox</div>
    <div class="lock-screen-subtitle">多账户机器人管理平台</div>
    <div id="lock-screen-error" class="lock-screen-error"></div>
    <div id="bot-login-form" class="lock-screen-form">
        <input type="text" id="bot-login-username" class="lock-screen-input" placeholder="用户名" autocomplete="off" style="letter-spacing:normal;" />
        <input type="password" id="bot-login-password" class="lock-screen-input" placeholder="密码" autocomplete="off" />
        <button id="bot-login-submit" class="lock-screen-btn">登录</button>
        <div class="lock-screen-forgot">
            <button id="bot-show-register" class="lock-screen-forgot-btn">没有账户？注册</button>
            <button id="bot-show-forgot" class="lock-screen-forgot-btn">忘记密码</button>
        </div>
    </div>
    <div id="bot-register-form" class="lock-screen-form" style="display:none;">
        <input type="text" id="bot-register-username" class="lock-screen-input" placeholder="用户名 (2-32位)" autocomplete="off" style="letter-spacing:normal;" />
        <input type="email" id="bot-register-email" class="lock-screen-input" placeholder="邮箱" autocomplete="off" style="display:none;" />
        <div id="bot-register-email-code-row" class="captcha-row" style="display:none;">
            <input type="text" id="bot-register-email-code" class="captcha-input" placeholder="邮箱验证码" autocomplete="off" />
            <button id="bot-register-send-email-code" class="captcha-refresh" style="width:auto;padding:0 12px;font-size:13px;white-space:nowrap;">发送</button>
        </div>
        <input type="password" id="bot-register-password" class="lock-screen-input" placeholder="密码 (至少6位)" autocomplete="off" />
        <input type="password" id="bot-register-password2" class="lock-screen-input" placeholder="确认密码" autocomplete="off" />
        <button id="bot-register-submit" class="lock-screen-btn">注册</button>
        <div style="margin-top:8px;padding:10px 14px;background:var(--accent-light);border:0.5px solid var(--accent);border-radius:var(--card-round);font-size:13px;color:var(--accent);text-align:center;line-height:1.5;">若邮箱没有收到邮件，请前往垃圾邮件查看（目前邮箱只适配QQ邮箱）</div>
        <div class="lock-screen-forgot">
            <button id="bot-show-login" class="lock-screen-forgot-btn">已有账户？登录</button>
        </div>
    </div>
    <div id="bot-forgot-form" class="lock-screen-form" style="display:none;">
        <input type="text" id="bot-forgot-username" class="lock-screen-input" placeholder="用户名" autocomplete="off" style="letter-spacing:normal;" />
        <div class="captcha-row">
            <input type="text" id="bot-forgot-captcha" class="captcha-input" placeholder="图形验证码" autocomplete="off" />
            <img id="bot-forgot-captcha-img" class="captcha-img" src="" alt="验证码" />
            <button id="bot-forgot-captcha-refresh" class="captcha-refresh">↻</button>
        </div>
        <button id="bot-forgot-send-code" class="lock-screen-reset-btn">发送验证码</button>
        <input type="text" id="bot-forgot-code" class="lock-screen-reset-input" placeholder="输入验证码" autocomplete="off" />
        <input type="password" id="bot-forgot-new-password" class="lock-screen-reset-input" placeholder="新密码 (至少6位)" autocomplete="off" />
        <input type="password" id="bot-forgot-confirm-password" class="lock-screen-reset-input" placeholder="确认新密码" autocomplete="off" />
        <button id="bot-forgot-submit" class="lock-screen-reset-btn">重置密码</button>
        <div class="lock-screen-forgot">
            <button id="bot-forgot-back-login" class="lock-screen-back-btn">返回登录</button>
        </div>
    </div>
</div>
<script>
var _token = "__SESSION_TOKEN__";
(function() {
    var lockScreen = document.getElementById("lock-screen");
    var loginForm = document.getElementById("bot-login-form");
    var registerForm = document.getElementById("bot-register-form");
    var showRegisterBtn = document.getElementById("bot-show-register");
    var showLoginBtn = document.getElementById("bot-show-login");
    var loginSubmitBtn = document.getElementById("bot-login-submit");
    var registerSubmitBtn = document.getElementById("bot-register-submit");
    var errorEl = document.getElementById("lock-screen-error");
    var _emailRegisterEnabled = false;

    function _generateFingerprint() {
        try {
            var canvas = document.createElement('canvas');
            canvas.width = 200; canvas.height = 50;
            var ctx = canvas.getContext('2d');
            ctx.textBaseline = 'top';
            ctx.font = '14px Arial';
            ctx.fillStyle = '#f60';
            ctx.fillRect(125, 1, 62, 20);
            ctx.fillStyle = '#069';
            ctx.fillText('ZynChatBox', 2, 15);
            ctx.fillStyle = 'rgba(102, 204, 0, 0.7)';
            ctx.fillText('ZynChatBox', 4, 17);
            var dataUrl = canvas.toDataURL();
            var fp = navigator.userAgent + '|' + screen.width + 'x' + screen.height + '|' + screen.colorDepth + '|' + new Date().getTimezoneOffset() + '|' + dataUrl;
            var hash = 0;
            for (var i = 0; i < fp.length; i++) {
                var char = fp.charCodeAt(i);
                hash = ((hash << 5) - hash) + char;
                hash = hash & hash;
            }
            return Math.abs(hash).toString(36) + fp.length.toString(36);
        } catch(e) {
            return navigator.userAgent.replace(/[^a-zA-Z0-9]/g, '').substring(0, 32);
        }
    }

    if (showRegisterBtn) showRegisterBtn.addEventListener("click", function() {
        if (loginForm) loginForm.style.display = "none";
        if (registerForm) registerForm.style.display = "";
        if (errorEl) errorEl.textContent = "";
    });
    if (showLoginBtn) showLoginBtn.addEventListener("click", function() {
        if (registerForm) registerForm.style.display = "none";
        if (loginForm) loginForm.style.display = "";
        if (errorEl) errorEl.textContent = "";
    });

    var loginUsernameInput = document.getElementById("bot-login-username");
    var loginPasswordInput = document.getElementById("bot-login-password");
    if (loginPasswordInput) loginPasswordInput.addEventListener("keypress", function(e) { if (e.key === "Enter") _doBotLogin(); });
    if (loginUsernameInput) loginUsernameInput.addEventListener("keypress", function(e) { if (e.key === "Enter") _doBotLogin(); });
    if (loginSubmitBtn) loginSubmitBtn.addEventListener("click", _doBotLogin);

    var regUsernameInput = document.getElementById("bot-register-username");
    var regPasswordInput = document.getElementById("bot-register-password");
    var regPassword2Input = document.getElementById("bot-register-password2");
    if (regPassword2Input) regPassword2Input.addEventListener("keypress", function(e) { if (e.key === "Enter") _doBotRegister(); });
    if (regPasswordInput) regPasswordInput.addEventListener("keypress", function(e) { if (e.key === "Enter") _doBotRegister(); });
    if (registerSubmitBtn) registerSubmitBtn.addEventListener("click", _doBotRegister);

    var sendEmailCodeBtn = document.getElementById("bot-register-send-email-code");
    if (sendEmailCodeBtn) sendEmailCodeBtn.addEventListener("click", async function() {
        var emailInput = document.getElementById("bot-register-email");
        var email = emailInput ? emailInput.value.trim() : "";
        if (!email) { errorEl.textContent = "请输入邮箱"; return; }
        sendEmailCodeBtn.disabled = true;
        errorEl.textContent = "";
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/send-email-code", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                    } else { resolve({success: false, error: "请求失败"}); }
                };
                o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                o.send(JSON.stringify({email: email}));
            });
            if (resp && resp.success) {
                errorEl.textContent = "";
                var _toastEl = document.getElementById("lock-screen-error");
                if (_toastEl) _toastEl.textContent = "验证码已发送";
                var _countdown = 60;
                var _origText = sendEmailCodeBtn.textContent;
                sendEmailCodeBtn.textContent = _countdown + "秒";
                var _timer = setInterval(function() {
                    _countdown--;
                    if (_countdown <= 0) {
                        clearInterval(_timer);
                        sendEmailCodeBtn.disabled = false;
                        sendEmailCodeBtn.textContent = _origText;
                    } else {
                        sendEmailCodeBtn.textContent = _countdown + "秒";
                    }
                }, 1000);
            } else {
                errorEl.textContent = (resp && resp.error) || "发送失败";
                sendEmailCodeBtn.disabled = false;
            }
        } catch(e) {
            errorEl.textContent = "发送失败";
            sendEmailCodeBtn.disabled = false;
        }
    });

    (async function() {
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("GET", "/api/wasm/email-register-status", true);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                    } else { resolve({enabled: false}); }
                };
                o.onerror = function() { resolve({enabled: false}); };
                o.send();
            });
            if (resp && resp.enabled) {
                var emailInput = document.getElementById("bot-register-email");
                var emailCodeRow = document.getElementById("bot-register-email-code-row");
                if (emailInput) emailInput.style.display = "";
                if (emailCodeRow) emailCodeRow.style.display = "";
                _emailRegisterEnabled = true;
            }
        } catch(e) {}
    })();

    var forgotForm = document.getElementById("bot-forgot-form");
    var showForgotBtn = document.getElementById("bot-show-forgot");
    var forgotBackLoginBtn = document.getElementById("bot-forgot-back-login");
    var forgotSendCodeBtn = document.getElementById("bot-forgot-send-code");
    var forgotSubmitBtn = document.getElementById("bot-forgot-submit");
    var forgotCaptchaRefreshBtn = document.getElementById("bot-forgot-captcha-refresh");
    var forgotCaptchaImg = document.getElementById("bot-forgot-captcha-img");

    function _showForgotForm() {
        if (loginForm) loginForm.style.display = "none";
        if (registerForm) registerForm.style.display = "none";
        if (forgotForm) forgotForm.style.display = "";
        if (errorEl) errorEl.textContent = "";
        _loadForgotCaptcha();
    }

    function _loadForgotCaptcha() {
        if (forgotCaptchaImg) {
            forgotCaptchaImg.src = "/api/wasm/captcha?token=" + encodeURIComponent(_token) + "&t=" + Date.now();
        }
    }

    if (showForgotBtn) showForgotBtn.addEventListener("click", _showForgotForm);
    if (forgotBackLoginBtn) forgotBackLoginBtn.addEventListener("click", function() {
        if (forgotForm) forgotForm.style.display = "none";
        if (loginForm) loginForm.style.display = "";
        if (errorEl) errorEl.textContent = "";
    });
    if (forgotCaptchaRefreshBtn) forgotCaptchaRefreshBtn.addEventListener("click", _loadForgotCaptcha);
    if (forgotCaptchaImg) forgotCaptchaImg.addEventListener("click", _loadForgotCaptcha);

    if (forgotSendCodeBtn) forgotSendCodeBtn.addEventListener("click", async function() {
        var username = document.getElementById("bot-forgot-username");
        var captchaInput = document.getElementById("bot-forgot-captcha");
        var u = username ? username.value.trim() : "";
        var c = captchaInput ? captchaInput.value.trim() : "";
        if (!u) { errorEl.textContent = "请输入用户名"; return; }
        if (!c) { errorEl.textContent = "请输入图形验证码"; return; }
        forgotSendCodeBtn.disabled = true;
        errorEl.textContent = "";
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/bot-send-reset-code", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                    } else { resolve({success: false, error: "请求失败"}); }
                };
                o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                o.send(JSON.stringify({username: u, captcha: c}));
            });
            if (resp && resp.success) {
                errorEl.textContent = "";
                _toast("验证码已发送");
                var _countdown = 60;
                var _origText = forgotSendCodeBtn.textContent;
                forgotSendCodeBtn.textContent = _countdown + "秒后重发";
                var _timer = setInterval(function() {
                    _countdown--;
                    if (_countdown <= 0) {
                        clearInterval(_timer);
                        forgotSendCodeBtn.disabled = false;
                        forgotSendCodeBtn.textContent = _origText;
                    } else {
                        forgotSendCodeBtn.textContent = _countdown + "秒后重发";
                    }
                }, 1000);
            } else {
                errorEl.textContent = (resp && resp.error) || "发送失败";
                _loadForgotCaptcha();
                forgotSendCodeBtn.disabled = false;
            }
        } catch(e) {
            errorEl.textContent = "发送失败";
            _loadForgotCaptcha();
            forgotSendCodeBtn.disabled = false;
        }
    });

    if (forgotSubmitBtn) forgotSubmitBtn.addEventListener("click", async function() {
        var username = document.getElementById("bot-forgot-username");
        var code = document.getElementById("bot-forgot-code");
        var newPwd = document.getElementById("bot-forgot-new-password");
        var confirmPwd = document.getElementById("bot-forgot-confirm-password");
        var u = username ? username.value.trim() : "";
        var cd = code ? code.value.trim() : "";
        var np = newPwd ? newPwd.value : "";
        var cp = confirmPwd ? confirmPwd.value : "";
        if (!u) { errorEl.textContent = "请输入用户名"; return; }
        if (!cd) { errorEl.textContent = "请输入验证码"; return; }
        if (!np || np.length < 6) { errorEl.textContent = "新密码长度不能少于6位"; return; }
        if (np !== cp) { errorEl.textContent = "两次密码不一致"; return; }
        forgotSubmitBtn.disabled = true;
        errorEl.textContent = "";
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/bot-reset-password", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                    } else { resolve({success: false, error: "请求失败"}); }
                };
                o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                o.send(JSON.stringify({username: u, code: cd, new_password: np}));
            });
            if (resp && resp.success) {
                errorEl.textContent = "密码重置成功，请重新登录";
                if (forgotForm) forgotForm.style.display = "none";
                if (loginForm) loginForm.style.display = "";
            } else {
                errorEl.textContent = (resp && resp.error) || "重置失败";
            }
        } catch(e) {
            errorEl.textContent = "重置失败";
        }
        forgotSubmitBtn.disabled = false;
    });

    async function _doBotLogin() {
        var username = loginUsernameInput ? loginUsernameInput.value.trim() : "";
        var password = loginPasswordInput ? loginPasswordInput.value : "";
        if (!username) { errorEl.textContent = "请输入用户名"; return; }
        if (!password) { errorEl.textContent = "请输入密码"; return; }
        if (loginSubmitBtn) loginSubmitBtn.disabled = true;
        errorEl.textContent = "";
        var fingerprint = _generateFingerprint();
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/bot-login", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                    } else { resolve({success: false, error: "请求失败"}); }
                };
                o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                o.send(JSON.stringify({username: username, password: password, fingerprint: fingerprint, client_ip: window._clientIP || ""}));
            });
            if (resp && resp.success && resp.session_token) {
                _token = resp.session_token;
                window.location.href = '/';
            } else {
                errorEl.textContent = (resp && resp.error) || "登录失败";
                if (loginSubmitBtn) loginSubmitBtn.disabled = false;
            }
        } catch(e) {
            errorEl.textContent = "登录失败";
            if (loginSubmitBtn) loginSubmitBtn.disabled = false;
        }
    }

    async function _doBotRegister() {
        var username = regUsernameInput ? regUsernameInput.value.trim() : "";
        var password = regPasswordInput ? regPasswordInput.value : "";
        var password2 = regPassword2Input ? regPassword2Input.value : "";
        var emailInput = document.getElementById("bot-register-email");
        var emailCodeInput = document.getElementById("bot-register-email-code");
        var email = emailInput ? emailInput.value.trim() : "";
        var emailCode = emailCodeInput ? emailCodeInput.value.trim() : "";
        if (!username) { errorEl.textContent = "请输入用户名"; return; }
        if (_emailRegisterEnabled) {
            if (!email) { errorEl.textContent = "请输入邮箱"; return; }
            if (!emailCode) { errorEl.textContent = "请输入邮箱验证码"; return; }
        }
        if (!password) { errorEl.textContent = "请输入密码"; return; }
        if (password !== password2) { errorEl.textContent = "两次密码不一致"; return; }
        if (registerSubmitBtn) registerSubmitBtn.disabled = true;
        errorEl.textContent = "";
        var reqData = {username: username, password: password, client_ip: window._clientIP || ""};
        if (_emailRegisterEnabled) {
            reqData.email = email;
            reqData.email_code = emailCode;
        }
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/bot-register", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({}); }
                    } else { resolve({success: false, error: "请求失败"}); }
                };
                o.onerror = function() { resolve({success: false, error: "网络错误"}); };
                o.send(JSON.stringify(reqData));
            });
            if (resp && resp.success && resp.session_token) {
                _token = resp.session_token;
                window.location.href = '/';
            } else {
                errorEl.textContent = (resp && resp.error) || "注册失败";
                if (registerSubmitBtn) registerSubmitBtn.disabled = false;
            }
        } catch(e) {
            errorEl.textContent = "注册失败";
            if (registerSubmitBtn) registerSubmitBtn.disabled = false;
        }
    }

    (async function() {
        var fingerprint = _generateFingerprint();
        try {
            var resp = await new Promise(function(resolve, reject) {
                var o = new XMLHttpRequest();
                o.open("POST", "/api/wasm/bot-fingerprint-login", true);
                o.setRequestHeader("Content-Type", "application/json");
                o.setRequestHeader("X-Session-Token", _token);
                o.onload = function() {
                    if (o.status >= 200 && o.status < 300) {
                        try { resolve(JSON.parse(o.responseText)); } catch(e) { resolve({success: false}); }
                    } else { resolve({success: false}); }
                };
                o.onerror = function() { resolve({success: false}); };
                o.send(JSON.stringify({fingerprint: fingerprint, client_ip: window._clientIP || ""}));
            });
            if (resp && resp.success && resp.session_token) {
                _token = resp.session_token;
                window.location.href = '/';
                return;
            }
        } catch(e) {}
    })();
})();
</script>
</body>
</html>'''
                try:
                    final_html = html.replace('__SESSION_TOKEN__', session_token)
                    encoded = final_html.encode('utf-8')
                    self.send_response(200)
                    self.send_header('Content-type', 'text/html; charset=utf-8')
                    self.send_header('Content-Length', str(len(encoded)))
                    self.send_header('Set-Cookie', f'session_token={session_token}; Path=/; SameSite=Lax; HttpOnly')
                    self.send_header('X-Content-Type-Options', 'nosniff')
                    self.send_header('X-Frame-Options', 'DENY')
                    self.send_header('X-XSS-Protection', '1; mode=block')
                    self.end_headers()
                    self.wfile.write(encoded)
                except BrokenPipeError:
                    pass
                except Exception:
                    pass

            def _send_json(self, data, status=200, cookies=None):
                try:
                    self.send_response(status)
                    self.send_header('Content-type', 'application/json; charset=utf-8')
                    if cookies:
                        for cookie_str in cookies:
                            self.send_header('Set-Cookie', cookie_str)
                    origin = self.headers.get('Origin', '')
                    if origin and origin.startswith(('http://localhost', 'http://127.0.0.1')):
                        self.send_header('Access-Control-Allow-Origin', origin)
                    self.end_headers()
                    self.wfile.write(json.dumps(data, ensure_ascii=False).encode('utf-8'))
                except BrokenPipeError:
                    pass
                except Exception:
                    pass
        
        return WebHandler
    
    def _print_ascii_qrcode(self, qrcode_url: str):
        pass
    
    def login_with_qrcode(self) -> bool:
        while True:
            print("正在获取连接二维码...")
            try:
                url = f"{self.ILINK_BASE_URL}/ilink/bot/get_bot_qrcode?bot_type=3"
                req = urllib.request.Request(url)
                with urllib.request.urlopen(req, timeout=self._timeout) as resp:
                    data = json.loads(resp.read().decode('utf-8'))
            except Exception as e:
                print(f"获取二维码失败: {e}，3秒后重试...")
                time.sleep(3)
                continue
            
            self._qrcode_key = data.get("qrcode")
            qrcode_url = data.get("qrcode_img_content")
            
            if not self._qrcode_key:
                print("获取二维码失败，3秒后重试...")
                time.sleep(3)
                continue
            
            self._qrcode_matrix = self._get_qrcode_matrix(qrcode_url)
            self._print_ascii_qrcode(qrcode_url)
            print("请使用微信扫码并确认连接...")
            print("Zyn")
            
            while not self._login_done:
                if sys.stdin.isatty():
                    if sys.platform == "win32":
                        try:
                            import msvcrt
                            if msvcrt.kbhit():
                                cmd = sys.stdin.readline().strip()
                                if cmd.lower() in ["/http", "/web"]:
                                    self._open_browser()
                                    continue
                        except (ImportError, AttributeError):
                            pass
                    elif is_termux():
                        try:
                            try:
                                import select as sel_module
                                try:
                                    rlist, _, _ = sel_module.select([sys.stdin], [], [], 0.1)
                                    if rlist:
                                        cmd = sys.stdin.readline().strip()
                                        if cmd.lower() in ["/http", "/web"]:
                                            self._open_browser()
                                            continue
                                except (OSError, ValueError, ImportError):
                                    pass
                            except Exception:
                                pass
                        except Exception:
                            pass
                    else:
                        try:
                            rlist, _, _ = select.select([sys.stdin], [], [], 0.1)
                            if rlist:
                                cmd = sys.stdin.readline().strip()
                                if cmd.lower() in ["/http", "/web"]:
                                    self._open_browser()
                                    continue
                        except (OSError, ValueError):
                            pass
                
                try:
                    status_url = f"{self.ILINK_BASE_URL}/ilink/bot/get_qrcode_status?qrcode={self._qrcode_key}"
                    status_req = urllib.request.Request(status_url, headers={"iLink-App-ClientVersion": "1"})
                    with urllib.request.urlopen(status_req, timeout=5) as status_resp:
                        status = json.loads(status_resp.read().decode('utf-8'))
                except Exception as e:
                    time.sleep(1)
                    continue
                
                if status.get("status") == "scaned":
                    print("已扫码，请在手机上确认...")
                elif status.get("status") == "confirmed":
                    self.token = status.get("bot_token")
                    self.bot_id = status.get("ilink_bot_id")
                    self.user_id = status.get("ilink_user_id")
                    print(f"连接成功!")
                    print(f"   bot_id: {self.bot_id}")
                    print(f"   user_id: {self.user_id}")
                    
                    self._bot_accounts[self.token] = {
                        "bot_id": self.bot_id or "",
                        "user_id": self.user_id or "",
                        "cursor": self._cursor,
                        "context_tokens": {}
                    }
                    
                    print("正在拉取历史消息，恢复会话...")
                    self._fetch_and_restore_conversations()
                    
                    self._save_config()
                    self._login_done = True
                    print(f"[WEB] 连接成功！网页端应该会自动跳转到聊天界面")
                    print(f"[WEB] 如果没有跳转，请刷新浏览器页面: http://localhost:{self._web_port}")
                    return True
                elif status.get("status") == "expired":
                    print("二维码已过期，正在重新获取...")
                    break
                time.sleep(2)
    
    def _fetch_and_restore_conversations(self):
        for _ in range(5):
            body = {"get_updates_buf": self._cursor}
            result = self._post("getupdates", body, timeout=5)
            if result.get("get_updates_buf"):
                self._cursor = result["get_updates_buf"]
            messages = result.get("msgs", [])
            for msg in messages:
                from_user = msg.get("from_user_id")
                ctx_token = msg.get("context_token")
                if from_user and ctx_token:
                    is_new = from_user not in self._context_tokens
                    self._register_user_to_account(from_user, ctx_token, self.token)
                    if is_new:
                        print(f"恢复会话: {from_user}")
                    
                    text = ""
                    for item in msg.get("item_list", []):
                        if item.get("type") == 1:
                            text = item.get("text_item", {}).get("text", "")
                    if text:
                        new_msg = {
                            'from': from_user,
                            'to': 'me',
                            'text': text,
                            'time': datetime.now().strftime('%H:%M:%S'),
                            'type': 'in'
                        }
                        self._add_message_to_history(new_msg)
            if not messages:
                break
        if self._context_tokens:
            print(f"已恢复 {len(self._context_tokens)} 个会话，{len(self._messages)} 条本地消息")
            print(f"当前会话用户: {self._current_user}")
            for user_id in self._context_tokens.keys():
                self._on_new_user(user_id)
        else:
            print("没有找到历史会话")
    
    def _build_headers(self, token: str = None) -> dict:
        random_uin = random.randint(0, 0xFFFFFFFF)
        wechat_uin = base64.b64encode(str(random_uin).encode()).decode()
        use_token = token or self.token
        return {
            "Content-Type": "application/json",
            "AuthorizationType": "ilink_bot_token",
            "Authorization": f"Bearer {use_token}",
            "X-WECHAT-UIN": wechat_uin,
        }
    
    def _post(self, endpoint: str, body: dict, timeout: int = 30, token: str = None) -> dict:
        if is_termux():
            timeout = max(timeout, 30)
            if "getupdates" in endpoint:
                timeout = 30
        
        body = dict(body)
        body["base_info"] = {"channel_version": "1.0.3"}
        headers = self._build_headers(token=token)
        url = f"{self.ILINK_BASE_URL}/ilink/bot/{endpoint}"
        
        data = json.dumps(body).encode('utf-8')
        req = urllib.request.Request(url, data=data, headers=headers, method='POST')
        
        max_retries = 2 if is_termux() else 0
        
        for attempt in range(max_retries + 1):
            try:
                with urllib.request.urlopen(req, timeout=timeout) as response:
                    result = response.read().decode('utf-8')
                    if result.strip() == "{}":
                        return {"ret": 0}
                    return json.loads(result)
            except (urllib.error.URLError, Exception) as e:
                is_timeout = (
                    isinstance(e, urllib.error.URLError) and isinstance(e.reason, TimeoutError)
                ) or "timeout" in str(e).lower() or "timed out" in str(e).lower()
                
                if is_timeout:
                    if attempt < max_retries:
                        print(f"[TERMUX] 网络超时，重试 ({attempt + 1}/{max_retries})...")
                        time.sleep(2)
                        continue
                    return {"ret": -1, "errmsg": "timeout"}
                
                if attempt < max_retries:
                    print(f"[TERMUX] 请求失败: {e}，重试 ({attempt + 1}/{max_retries})...")
                    time.sleep(3)
                    continue
                    
                return {"ret": -1, "errmsg": str(e)}
        
        return {"ret": -1, "errmsg": "max retries exceeded"}
    
    _MEDIA_ITEM_KEYS = ["image_item", "video_item", "file_item", "voice_item"]
    
    def _extract_cdn_media(self, item: dict) -> Optional[dict]:
        for ik in self._MEDIA_ITEM_KEYS:
            mi = item.get(ik)
            if mi and isinstance(mi, dict) and mi.get("media"):
                cdn_media = dict(mi["media"])
                if not cdn_media.get("aes_key") and mi.get("aeskey"):
                    cdn_media["aes_key"] = base64.b64encode(mi["aeskey"].encode('utf-8')).decode('utf-8')
                return cdn_media
        return None
    
    def _process_message_items(self, item_list: list) -> tuple:
        text = ""
        media_info = None
        
        for item in item_list:
            if item.get("text_item"):
                text_item = item["text_item"]
                if isinstance(text_item, dict):
                    text = text_item.get("text", "")
                    
            if item.get("image_item"):
                img_item = item["image_item"]
                if isinstance(img_item, dict):
                    media_info = {
                        "type": "image",
                        "filename": img_item.get("filename", "image.jpg"),
                        "item": item
                    }
                    
            elif item.get("video_item"):
                video_item = item["video_item"]
                if isinstance(video_item, dict):
                    media_info = {
                        "type": "video",
                        "filename": video_item.get("filename", "video.mp4"),
                        "duration": video_item.get("duration", 0),
                        "item": item
                    }
                    
            elif item.get("file_item"):
                file_item = item["file_item"]
                if isinstance(file_item, dict):
                    media_info = {
                        "type": "file",
                        "filename": file_item.get("filename", "file.bin"),
                        "description": file_item.get("description", ""),
                        "item": item
                    }
                    
            elif item.get("voice_item"):
                voice_item = item["voice_item"]
                if isinstance(voice_item, dict):
                    media_info = {
                        "type": "voice",
                        "filename": voice_item.get("filename", "voice.silk"),
                        "duration": voice_item.get("duration", 0),
                        "item": item
                    }
        
        return text, media_info
    
    def start_add_user_qrcode(self) -> str:
        def _gen_qrcode():
            try:
                url = f"{self.ILINK_BASE_URL}/ilink/bot/get_bot_qrcode?bot_type=3"
                req = urllib.request.Request(url)
                with urllib.request.urlopen(req, timeout=self._timeout) as resp:
                    data = json.loads(resp.read().decode('utf-8'))
                
                qrcode_url = data.get("qrcode_img_content")
                qrcode_key = data.get("qrcode")
                
                with self._add_user_lock:
                    self._pending_qrcode = {
                        "key": qrcode_key,
                        "matrix": self._get_qrcode_matrix(qrcode_url) if qrcode_url else None,
                        "status": "waiting",
                        "started_at": time.time()
                    }
                
                if qrcode_url:
                    self._print_ascii_qrcode(qrcode_url)
                    print(f"[添加用户] 二维码已生成，请扫描")
                
                start_ts = time.time()
                while time.time() - start_ts < 120:
                    if not self._running:
                        break
                    try:
                        status_url = f"{self.ILINK_BASE_URL}/ilink/bot/get_qrcode_status?qrcode={qrcode_key}"
                        status_req = urllib.request.Request(status_url, headers={"iLink-App-ClientVersion": "1"})
                        with urllib.request.urlopen(status_req, timeout=5) as status_resp:
                            status = json.loads(status_resp.read().decode('utf-8'))
                    except Exception:
                        time.sleep(1)
                        continue
                    
                    st = status.get("status", "")
                    
                    with self._add_user_lock:
                        if st == "scaned":
                            if self._pending_qrcode:
                                self._pending_qrcode["status"] = "scaned"
                            print("[添加用户] 已扫码，请在手机上确认...")
                        elif st == "confirmed":
                            new_token = status.get("bot_token")
                            new_bot_id = status.get("ilink_bot_id")
                            new_user_id = status.get("ilink_user_id")
                            
                            if not new_token:
                                print("[添加用户] 错误：未获取到 bot_token")
                                if self._pending_qrcode:
                                    self._pending_qrcode["status"] = "error"
                                break
                            
                            new_account = {
                                "bot_id": new_bot_id or "",
                                "user_id": new_user_id or "",
                                "cursor": "",
                                "context_tokens": {}
                            }
                            self._bot_accounts[new_token] = new_account
                            
       
                            if not self.token:
                                self.token = new_token
                                self.bot_id = new_bot_id
                                self.user_id = new_user_id
                                self._login_done = True
                            
                            print(f"[添加用户] 新 bot 账号已创建: {new_token[:8]}... (bot_id: {new_bot_id})")
                            
                            self._fetch_and_restore_for_account(new_token, new_account)
                            self._save_config()
                            
                            self._start_account_poll(new_token, new_account)
                            
                            with self._add_user_lock:
                                if self._pending_qrcode:
                                    self._pending_qrcode["status"] = "done"
                                    self._pending_qrcode["users"] = list(self._context_tokens.keys())
                            break
                        elif st == "expired":
                            with self._add_user_lock:
                                if self._pending_qrcode:
                                    self._pending_qrcode["status"] = "expired"
                            print("[添加用户] 二维码已过期")
                            break
                    
                    time.sleep(1.5)
                else:
                    with self._add_user_lock:
                        if self._pending_qrcode and self._pending_qrcode.get("status") == "waiting":
                            self._pending_qrcode["status"] = "timeout"
                    print("[添加用户] 二维码等待超时")
                    
            except Exception as e:
                print(f"[添加用户] 获取二维码失败: {e}")
                with self._add_user_lock:
                    self._pending_qrcode = {"key": "", "matrix": None, "status": "error", "error": str(e)}
        
        qrcode_key = uuid.uuid4().hex[:12]
        with self._add_user_lock:
            self._pending_qrcode = {"key": qrcode_key, "matrix": None, "status": "generating"}
        
        thread = threading.Thread(target=_gen_qrcode, daemon=True)
        thread.start()
        return qrcode_key
    
    def _fetch_and_restore_for_account(self, bot_token: str, account: dict):
        for _ in range(5):
            body = {"get_updates_buf": account.get("cursor", "")}
            result = self._post("getupdates", body, timeout=5, token=bot_token)
            if result.get("get_updates_buf"):
                account["cursor"] = result["get_updates_buf"]
            messages = result.get("msgs", [])
            for msg in messages:
                from_user = msg.get("from_user_id")
                ctx_token = msg.get("context_token")
                if from_user and ctx_token:
                    is_new = from_user not in self._context_tokens
                    self._register_user_to_account(from_user, ctx_token, bot_token)
                    
                    text = ""
                    for item in msg.get("item_list", []):
                        if item.get("type") == 1:
                            text = item.get("text_item", {}).get("text", "")
                    if text:
                        new_msg = {
                            'from': from_user,
                            'to': 'me',
                            'text': text,
                            'time': datetime.now().strftime('%H:%M:%S'),
                            'type': 'in'
                        }
                        self._add_message_to_history(new_msg)
                    
                    if is_new:
                        self._on_new_user(from_user)
            if not messages:
                break
        
        user_count = len(account.get("context_tokens", {}))
        print(f"[账号 {bot_token[:8]}...] 已恢复 {user_count} 个会话")
    
    def get_add_user_status(self) -> dict:
        with self._add_user_lock:
            if not self._pending_qrcode:
                return {"status": "none", "message": "没有进行中的添加操作"}
            return dict(self._pending_qrcode)
    
    def get_add_user_status_for_account(self, account) -> dict:
        with account._add_user_lock:
            if not account._pending_qrcode:
                return {"status": "none", "message": "没有进行中的添加操作"}
            return dict(account._pending_qrcode)
    
    def start_add_user_qrcode_for_account(self, account) -> str:
        account._pending_qrcode = {"status": "generating", "key": None, "matrix": None}
        def _do_start():
            try:
                url = f"{self.ILINK_BASE_URL}/ilink/bot/get_bot_qrcode?bot_type=3"
                req = urllib.request.Request(url)
                with urllib.request.urlopen(req, timeout=self._timeout) as resp:
                    data = json.loads(resp.read().decode('utf-8'))
                
                qrcode_key = data.get("qrcode")
                qrcode_url = data.get("qrcode_img_content")
                
                if not qrcode_key:
                    print(f"[ADD-USER] 获取二维码失败, data={data}")
                    with account._add_user_lock:
                        account._pending_qrcode = {"status": "error", "message": "获取二维码失败", "key": None, "matrix": None}
                    return
                
                matrix = self._get_qrcode_matrix(qrcode_url) if qrcode_url else None
                print(f"[ADD-USER] 二维码已生成, key={qrcode_key[:8] if qrcode_key else 'None'}, matrix={'有' if matrix else '无'}")
                
                with account._add_user_lock:
                    account._pending_qrcode = {"status": "waiting", "key": qrcode_key, "matrix": matrix}
                
                start_ts = time.time()
                while time.time() - start_ts < 120:
                    if not self._running:
                        break
                    with account._add_user_lock:
                        if not account._pending_qrcode or account._pending_qrcode.get("status") in ("done", "error", "expired"):
                            return
                    try:
                        status_url = f"{self.ILINK_BASE_URL}/ilink/bot/get_qrcode_status?qrcode={qrcode_key}"
                        status_req = urllib.request.Request(status_url, headers={"iLink-App-ClientVersion": "1"})
                        with urllib.request.urlopen(status_req, timeout=5) as status_resp:
                            status = json.loads(status_resp.read().decode('utf-8'))
                    except Exception:
                        time.sleep(1)
                        continue
                    
                    st = status.get("status", "")
                    
                    if st == "scaned":
                        with account._add_user_lock:
                            if account._pending_qrcode:
                                account._pending_qrcode["status"] = "scaned"
                        print("[ADD-USER] 已扫码，请在手机上确认...")
                    elif st == "confirmed":
                        new_token = status.get("bot_token")
                        new_bot_id = status.get("ilink_bot_id")
                        new_user_id = status.get("ilink_user_id")
                        
                        if not new_token:
                            print("[ADD-USER] 错误：未获取到 bot_token")
                            with account._add_user_lock:
                                account._pending_qrcode = {"status": "error", "message": "未获取到bot_token", "key": qrcode_key, "matrix": None}
                            return
                        
                        new_account_data = {
                            "bot_id": new_bot_id or "",
                            "user_id": new_user_id or "",
                            "cursor": "",
                            "context_tokens": {}
                        }
                        account._bot_accounts[new_token] = new_account_data
                        
                        if not account.token:
                            account.token = new_token
                            account.bot_id = new_bot_id
                            account.user_id = new_user_id
                            account._login_done = True
                        
                        print(f"[ADD-USER] 新 bot 账号已创建: {new_token[:8]}... (bot_id: {new_bot_id})")
                        
                        self._fetch_and_restore_for_account(new_token, new_account_data)
                        
                        for uid, ctx in new_account_data.get("context_tokens", {}).items():
                            if uid not in account._context_tokens:
                                account._context_tokens[uid] = ctx
                                account._user_token_map[uid] = new_token
                                if not account._current_user:
                                    account._current_user = uid
                        
                        account._save_config()
                        
                        self._start_account_poll_thread(account, new_token, new_account_data)
                        
                        with account._add_user_lock:
                            account._pending_qrcode = {"status": "done", "key": qrcode_key, "matrix": None, "users": list(account._context_tokens.keys())}
                        return
                    elif st == "expired":
                        with account._add_user_lock:
                            account._pending_qrcode = {"status": "expired", "message": "二维码已过期", "key": qrcode_key, "matrix": None}
                        print("[ADD-USER] 二维码已过期")
                        return
                    
                    time.sleep(1.5)
                else:
                    with account._add_user_lock:
                        if account._pending_qrcode and account._pending_qrcode.get("status") == "waiting":
                            account._pending_qrcode["status"] = "expired"
                    print("[ADD-USER] 二维码等待超时")
            except Exception as e:
                print(f"[ADD-USER] 线程异常: {e}")
                with account._add_user_lock:
                    account._pending_qrcode = {"status": "error", "message": str(e), "key": None, "matrix": None}
        threading.Thread(target=_do_start, daemon=True).start()
        return "pending"
    
    def _call_vision_api_for_account(self, account, image_base64: str, system_prompt: str, history: list, original_text: str = "", media_memory_text: str = "") -> Optional[str]:
        if not account.ai_config.get("vision_api_url"):
            return None
        try:
            vision_url = account.ai_config.get("vision_api_url")
            vision_key = account.ai_config.get("vision_api_key")
            vision_model = account.ai_config.get("vision_model", "gpt-4o")
            messages = [{"role": "system", "content": system_prompt}]
            for msg in history[-20:]:
                if msg.get("type") == "in":
                    messages.append({"role": "user", "content": msg.get("text", "")})
                elif msg.get("type") == "out":
                    messages.append({"role": "assistant", "content": msg.get("text", "")})
            user_content = []
            if original_text:
                user_content.append({"type": "text", "text": original_text})
            else:
                user_content.append({"type": "text", "text": "请描述这张图片"})
            user_content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"}})
            messages.append({"role": "user", "content": user_content})
            payload = {"model": vision_model, "messages": messages, "max_tokens": 500}
            headers = {"Content-Type": "application/json"}
            if vision_key:
                headers["Authorization"] = f"Bearer {vision_key}"
            req = urllib.request.Request(vision_url, data=json.dumps(payload).encode("utf-8"), headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                return result.get("choices", [{}])[0].get("message", {}).get("content", "")
        except Exception as e:
            print(f"[VISION API] 异常: {e}")
            return None
    
    def _call_image_gen_api_for_account(self, target, prompt: str) -> Optional[bytes]:
        if not target.ai_config.get("image_gen_api_url"):
            return None
        try:
            api_url = target.ai_config.get("image_gen_api_url")
            api_key = target.ai_config.get("image_gen_api_key")
            model = target.ai_config.get("image_gen_model", "dall-e-3")
            payload = {"model": model, "prompt": prompt, "n": 1, "size": "1024x1024"}
            headers = {"Content-Type": "application/json"}
            if api_key:
                headers["Authorization"] = f"Bearer {api_key}"
            req = urllib.request.Request(api_url, data=json.dumps(payload).encode("utf-8"), headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=120) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                image_url = result.get("data", [{}])[0].get("url", "")
                if image_url:
                    img_req = urllib.request.Request(image_url)
                    with urllib.request.urlopen(img_req, timeout=60) as img_resp:
                        return img_resp.read()
                b64 = result.get("data", [{}])[0].get("b64_json", "")
                if b64:
                    return base64.b64decode(b64)
            return None
        except Exception as e:
            print(f"[IMAGE GEN API] 异常: {e}")
            return None
    
    def _call_file_recognize_api_for_account(self, account, file_text: str, filename: str, system_prompt: str, history: list, original_text: str = "", media_memory_text: str = "") -> Optional[str]:
        if not account.ai_config.get("file_recognize_api_url"):
            return None
        try:
            api_url = account.ai_config.get("file_recognize_api_url")
            api_key = account.ai_config.get("file_recognize_api_key")
            model = account.ai_config.get("file_recognize_model", "gpt-4o")
            messages = [{"role": "system", "content": system_prompt}]
            for msg in history[-20:]:
                if msg.get("type") == "in":
                    messages.append({"role": "user", "content": msg.get("text", "")})
                elif msg.get("type") == "out":
                    messages.append({"role": "assistant", "content": msg.get("text", "")})
            max_size = account.ai_config.get("file_recognize_max_size", 512) * 1024
            truncated = file_text[:max_size]
            user_msg = f"[用户发送了文件: {filename}]\n文件内容如下:\n{truncated}"
            if original_text:
                user_msg += f"\n额外要求: {original_text}"
            messages.append({"role": "user", "content": user_msg})
            payload = {"model": model, "messages": messages, "max_tokens": 500}
            headers = {"Content-Type": "application/json"}
            if api_key:
                headers["Authorization"] = f"Bearer {api_key}"
            req = urllib.request.Request(api_url, data=json.dumps(payload).encode("utf-8"), headers=headers, method="POST")
            with urllib.request.urlopen(req, timeout=60) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                return result.get("choices", [{}])[0].get("message", {}).get("content", "")
        except Exception as e:
            print(f"[FILE RECOGNIZE API] 异常: {e}")
            return None
    
    def _send_ai_reply_in_segments_for_account(self, account, to_user_id: str, response_text: str):
        if not response_text or not to_user_id:
            return
        max_len = 4096
        if len(response_text) <= max_len:
            self.send_text_for_account(account, to_user_id, response_text)
            return
        segments = []
        current = ""
        for char in response_text:
            current += char
            if len(current) >= max_len:
                segments.append(current)
                current = ""
        if current:
            segments.append(current)
        for seg in segments:
            self.send_text_for_account(account, to_user_id, seg)
            time.sleep(0.5)
    
    def start_polling(self):
        if self.token and self.token not in self._bot_accounts:
            self._bot_accounts[self.token] = {
                "bot_id": self.bot_id or "",
                "user_id": self.user_id or "",
                "cursor": self._cursor,
                "context_tokens": dict(self._context_tokens)
            }
        
        seen_tokens = set()
        for user_id, bot_token in self._user_token_map.items():
            if bot_token and bot_token not in seen_tokens:
                seen_tokens.add(bot_token)
                account = self._bot_accounts.get(bot_token)
                if account:
                    self._start_account_poll(bot_token, account)
    
    def _start_account_poll(self, bot_token: str, account: dict):
        for t in self._poll_threads:
            if t.is_alive() and getattr(t, '_bot_token', None) == bot_token:
                return
        def poll():
            cursor = account.get("cursor", "")
            while self._running:
                try:
                    body = {"get_updates_buf": cursor}
                    result = self._post("getupdates", body, timeout=25, token=bot_token)
                    
                    if result.get("get_updates_buf"):
                        cursor = result["get_updates_buf"]
                        account["cursor"] = cursor
                        self._save_config()
                    
                    messages = result.get("msgs", [])
                    for msg in messages:
                        from_user = msg.get("from_user_id")
                        ctx_token = msg.get("context_token")
                        
                        text, media_info = self._process_message_items(msg.get("item_list", []))
                        
                        msg_text = text
                        msg_type = 'in'
                        msg_metadata = {}
                        
                        if media_info:
                            media_type_int = self.MEDIA_TYPE_MAP.get(media_info["type"], 0)
                            media_prefix = self.MEDIA_TYPE_PREFIXES.get(media_info["type"], f"[{media_info['type']}]")
                            
                            if text:
                                msg_text = f"{media_prefix} {text}"
                            else:
                                msg_text = f"{media_prefix} {media_info.get('filename', '')}"
                            
                            msg_metadata = {
                                'media_type': media_type_int,
                                'media_filename': media_info.get('filename', ''),
                                'media_duration': media_info.get('duration', 0),
                                'has_media': True
                            }
                            
                            media_item = media_info.get("item", {})
                            cdn_media = self._extract_cdn_media(media_item)
                            if cdn_media:
                                msg_metadata['media_cdn'] = json.dumps(cdn_media)
                                _prefetch_fn = media_info.get('filename', '')
                                threading.Thread(target=self._prefetch_media, args=(cdn_media, _prefetch_fn, from_user), daemon=True).start()
                            
                            print(f"\n[收到{media_info['type']}] {from_user}: {media_info.get('filename', '')}")
                        elif text:
                            print(f"\n[收到消息] {from_user}: {text}")
                        
                        if msg_text:
                            new_msg = {
                                'from': from_user,
                                'to': 'me',
                                'text': msg_text,
                                'time': datetime.now().strftime('%H:%M:%S'),
                                'type': msg_type,
                                **msg_metadata
                            }
                            
                            self._add_message_to_history(new_msg)
                            
                            if self._message_callback:
                                self._message_callback(new_msg)
                            
                            if media_info and media_info.get("type") == "image" and cdn_media:
                                def _vision_reply(fn_from_user=from_user, fn_cdn=cdn_media, fn_text=text, fn_filename=media_info.get('filename', '')):
                                    try:
                                        downloaded = self.download_media(fn_cdn, filename=fn_filename, user_id=fn_from_user)
                                        if downloaded:
                                            img_b64 = base64.b64encode(downloaded).decode('utf-8')
                                            self._auto_ai_reply_with_vision(fn_from_user, img_b64, original_text=fn_text, cdn_info=fn_cdn)
                                        else:
                                            if fn_text and self.ai_config.get("auto_reply"):
                                                self._auto_ai_reply(fn_from_user, fn_text)
                                    except Exception as e:
                                        print(f"[VISION] 识图处理异常: {e}")
                                        if fn_text and self.ai_config.get("auto_reply"):
                                            self._auto_ai_reply(fn_from_user, fn_text)
                                threading.Thread(target=_vision_reply, daemon=True).start()
                            elif media_info and media_info.get("type") == "file" and cdn_media:
                                def _file_reply(fn_from_user=from_user, fn_cdn=cdn_media, fn_text=text, fn_filename=media_info.get('filename', '')):
                                    try:
                                        downloaded = self.download_media(fn_cdn, filename=fn_filename, user_id=fn_from_user)
                                        if downloaded:
                                            self._auto_ai_reply_with_file(fn_from_user, downloaded, fn_filename, original_text=fn_text)
                                        else:
                                            if fn_text and self.ai_config.get("auto_reply"):
                                                self._auto_ai_reply(fn_from_user, fn_text)
                                    except Exception as e:
                                        print(f"[FILE_RECOGNIZE] 文件识别处理异常: {e}")
                                        if fn_text and self.ai_config.get("auto_reply"):
                                            self._auto_ai_reply(fn_from_user, fn_text)
                                threading.Thread(target=_file_reply, daemon=True).start()
                            elif text:
                                threading.Thread(target=self._auto_ai_reply, args=(from_user, text), daemon=True).start()
                        
                        if from_user and ctx_token:
                            is_new = from_user not in self._context_tokens
                            self._register_user_to_account(from_user, ctx_token, bot_token)
                            self._save_config()
                            if is_new:
                                self._on_new_user(from_user)
                                print(f"[USER] 新用户 {from_user} (账号 {bot_token[:8]}...)")
                except Exception as e:
                    time.sleep(0.5)
        
        thread = threading.Thread(target=poll, daemon=True)
        thread._bot_token = bot_token
        thread.start()
        self._poll_threads.append(thread)
        token_short = bot_token[:8] if bot_token else "?"
        print(f"[POLL] 已启动轮询线程: {token_short}...")
    
    def send_typing(self, to_user_id: str) -> bool:
        context_token = self._context_tokens.get(to_user_id)
        if not context_token:
            return False
        use_token = self._get_token_for_user(to_user_id)
        try:
            config_body = {
                "ilink_user_id": to_user_id,
                "context_token": context_token
            }
            config_result = self._post("getconfig", config_body, token=use_token)
            typing_ticket = config_result.get("typing_ticket")
            if not typing_ticket:
                return False
            typing_body = {
                "ilink_user_id": to_user_id,
                "typing_ticket": typing_ticket,
                "status": 1
            }
            self._post("sendtyping", typing_body, token=use_token)
            return True
        except Exception:
            return False
    
    def send_text(self, to_user_id: str, text: str) -> bool:
        context_token = self._context_tokens.get(to_user_id)
        if not context_token:
            print(f"[发送失败] 没有 {to_user_id} 的会话，让对方先发一条消息")
            return False
        
        use_token = self._get_token_for_user(to_user_id)
        
        client_id = f"msg-{uuid.uuid4().hex[:16]}"
        body = {
            "msg": {
                "from_user_id": "",
                "to_user_id": to_user_id,
                "client_id": client_id,
                "message_type": 2,
                "message_state": 2,
                "context_token": context_token,
                "item_list": [{"type": 1, "text_item": {"text": text}}]
            }
        }
        result = self._post("sendmessage", body, token=use_token)
        
        errcode = result.get("errcode")
        ret = result.get("ret")
        
        success = (ret is None or ret == 0) and (errcode is None or errcode == 0)
        
        if success:
            print(f"[发送成功] 给 {to_user_id}: {text[:50]}...")
            out_msg = {
                'from': 'me',
                'to': to_user_id,
                'text': text,
                'time': datetime.now().strftime('%H:%M:%S'),
                'type': 'out'
            }
            self._add_message_to_history(out_msg)
            return True
        
        if ret == -1:
            print(f"[发送失败] {result.get('errmsg', '未知错误')}")
            return False
        
        if errcode in self.EXPIRED_CODES:
            self._context_tokens.pop(to_user_id, None)
            bot_token = self._user_token_map.pop(to_user_id, None)
            if bot_token and bot_token in self._bot_accounts:
                self._bot_accounts[bot_token].get("context_tokens", {}).pop(to_user_id, None)
            self._save_config()
        print(f"[发送失败] ret={ret}, errcode={errcode}, errmsg={result.get('errmsg', '')}")
        return False

    CDN_BASE = "https://novac2c.cdn.weixin.qq.com/c2c"

    def _random_hex(self, num_bytes: int) -> str:
        raw = os.urandom(num_bytes)
        return raw.hex()

    def _md5_hex(self, data: bytes) -> str:
        return hashlib.md5(data).hexdigest()

    _AES_SBOX = [
        0x63,0x7c,0x77,0x7b,0xf2,0x6b,0x6f,0xc5,0x30,0x01,0x67,0x2b,0xfe,0xd7,0xab,0x76,
        0xca,0x82,0xc9,0x7d,0xfa,0x59,0x47,0xf0,0xad,0xd4,0xa2,0xaf,0x9c,0xa4,0x72,0xc0,
        0xb7,0xfd,0x93,0x26,0x36,0x3f,0xf7,0xcc,0x34,0xa5,0xe5,0xf1,0x71,0xd8,0x31,0x15,
        0x04,0xc7,0x23,0xc3,0x18,0x96,0x05,0x9a,0x07,0x12,0x80,0xe2,0xeb,0x27,0xb2,0x75,
        0x09,0x83,0x2c,0x1a,0x1b,0x6e,0x5a,0xa0,0x52,0x3b,0xd6,0xb3,0x29,0xe3,0x2f,0x84,
        0x53,0xd1,0x00,0xed,0x20,0xfc,0xb1,0x5b,0x6a,0xcb,0xbe,0x39,0x4a,0x4c,0x58,0xcf,
        0xd0,0xef,0xaa,0xfb,0x43,0x4d,0x33,0x85,0x45,0xf9,0x02,0x7f,0x50,0x3c,0x9f,0xa8,
        0x51,0xa3,0x40,0x8f,0x92,0x9d,0x38,0xf5,0xbc,0xb6,0xda,0x21,0x10,0xff,0xf3,0xd2,
        0xcd,0x0c,0x13,0xec,0x5f,0x97,0x44,0x17,0xc4,0xa7,0x7e,0x3d,0x64,0x5d,0x19,0x73,
        0x60,0x81,0x4f,0xdc,0x22,0x2a,0x90,0x88,0x46,0xee,0xb8,0x14,0xde,0x5e,0x0b,0xdb,
        0xe0,0x32,0x3a,0x0a,0x49,0x06,0x24,0x5c,0xc2,0xd3,0xac,0x62,0x91,0x95,0xe4,0x79,
        0xe7,0xc8,0x37,0x6d,0x8d,0xd5,0x4e,0xa9,0x6c,0x56,0xf4,0xea,0x65,0x7a,0xae,0x08,
        0xba,0x78,0x25,0x2e,0x1c,0xa6,0xb4,0xc6,0xe8,0xdd,0x74,0x1f,0x4b,0xbd,0x8b,0x8a,
        0x70,0x3e,0xb5,0x66,0x48,0x03,0xf6,0x0e,0x61,0x35,0x57,0xb9,0x86,0xc1,0x1d,0x9e,
        0xe1,0xf8,0x98,0x11,0x69,0xd9,0x8e,0x94,0x9b,0x1e,0x87,0xe9,0xce,0x55,0x28,0xdf,
        0x8c,0xa1,0x89,0x0d,0xbf,0xe6,0x42,0x68,0x41,0x99,0x2d,0x0f,0xb0,0x54,0xbb,0x16,
    ]
    _AES_INV_SBOX = [
        0x52,0x09,0x6a,0xd5,0x30,0x36,0xa5,0x38,0xbf,0x40,0xa3,0x9e,0x81,0xf3,0xd7,0xfb,
        0x7c,0xe3,0x39,0x82,0x9b,0x2f,0xff,0x87,0x34,0x8e,0x43,0x44,0xc4,0xde,0xe9,0xcb,
        0x54,0x7b,0x94,0x32,0xa6,0xc2,0x23,0x3d,0xee,0x4c,0x95,0x0b,0x42,0xfa,0xc3,0x4e,
        0x08,0x2e,0xa1,0x66,0x28,0xd9,0x24,0xb2,0x76,0x5b,0xa2,0x49,0x6d,0x8b,0xd1,0x25,
        0x72,0xf8,0xf6,0x64,0x86,0x68,0x98,0x16,0xd4,0xa4,0x5c,0xcc,0x5d,0x65,0xb6,0x92,
        0x6c,0x70,0x48,0x50,0xfd,0xed,0xb9,0xda,0x5e,0x15,0x46,0x57,0xa7,0x8d,0x9d,0x84,
        0x90,0xd8,0xab,0x00,0x8c,0xbc,0xd3,0x0a,0xf7,0xe4,0x58,0x05,0xb8,0xb3,0x45,0x06,
        0xd0,0x2c,0x1e,0x8f,0xca,0x3f,0x0f,0x02,0xc1,0xaf,0xbd,0x03,0x01,0x13,0x8a,0x6b,
        0x3a,0x91,0x11,0x41,0x4f,0x67,0xdc,0xea,0x97,0xf2,0xcf,0xce,0xf0,0xb4,0xe6,0x73,
        0x96,0xac,0x74,0x22,0xe7,0xad,0x35,0x85,0xe2,0xf9,0x37,0xe8,0x1c,0x75,0xdf,0x6e,
        0x47,0xf1,0x1a,0x71,0x1d,0x29,0xc5,0x89,0x6f,0xb7,0x62,0x0e,0xaa,0x18,0xbe,0x1b,
        0xfc,0x56,0x3e,0x4b,0xc6,0xd2,0x79,0x20,0x9a,0xdb,0xc0,0xfe,0x78,0xcd,0x5a,0xf4,
        0x1f,0xdd,0xa8,0x33,0x88,0x07,0xc7,0x31,0xb1,0x12,0x10,0x59,0x27,0x80,0xec,0x5f,
        0x60,0x51,0x7f,0xa9,0x19,0xb5,0x4a,0x0d,0x2d,0xe5,0x7a,0x9f,0x93,0xc9,0x9c,0xef,
        0xa0,0xe0,0x3b,0x4d,0xae,0x2a,0xf5,0xb0,0xc8,0xeb,0xbb,0x3c,0x83,0x53,0x99,0x61,
        0x17,0x2b,0x04,0x7e,0xba,0x77,0xd6,0x26,0xe1,0x69,0x14,0x63,0x55,0x21,0x0c,0x7d,
    ]
    _AES_RCON = [0x01,0x02,0x04,0x08,0x10,0x20,0x40,0x80,0x1b,0x36]

    @staticmethod
    def _xtime(a):
        return ((a << 1) ^ 0x1b) & 0xff if a & 0x80 else (a << 1) & 0xff

    @staticmethod
    def _gmul(a, b):
        p = 0
        for _ in range(8):
            if b & 1:
                p ^= a
            hi = a & 0x80
            a = (a << 1) & 0xff
            if hi:
                a ^= 0x1b
            b >>= 1
        return p

    @classmethod
    def _aes_key_expansion(cls, key: bytes) -> list:
        Nk = len(key) // 4
        Nr = Nk + 6
        W = []
        for i in range(Nk):
            W.append(list(key[4*i:4*i+4]))
        for i in range(Nk, 4*(Nr+1)):
            t = list(W[i-1])
            if i % Nk == 0:
                t = t[1:] + t[:1]
                t = [cls._AES_SBOX[b] for b in t]
                t[0] ^= cls._AES_RCON[i//Nk - 1]
            elif Nk > 6 and i % Nk == 4:
                t = [cls._AES_SBOX[b] for b in t]
            W.append([W[i-Nk][j] ^ t[j] for j in range(4)])
        return W

    @classmethod
    def _aes_encrypt_block(cls, block: bytes, round_keys: list) -> bytes:
        Nr = len(round_keys) // 4 - 1
        s = [[0]*4 for _ in range(4)]
        for i in range(16):
            s[i%4][i//4] = block[i]
        for c in range(4):
            for r in range(4):
                s[r][c] ^= round_keys[c][r]
        for rnd in range(1, Nr):
            s = [[cls._AES_SBOX[s[r][c]] for c in range(4)] for r in range(4)]
            for r in range(1, 4):
                s[r] = s[r][r:] + s[r][:r]
            for c in range(4):
                a = [s[r][c] for r in range(4)]
                s[0][c] = cls._xtime(a[0]) ^ cls._xtime(a[1]) ^ a[1] ^ a[2] ^ a[3]
                s[1][c] = a[0] ^ cls._xtime(a[1]) ^ cls._xtime(a[2]) ^ a[2] ^ a[3]
                s[2][c] = a[0] ^ a[1] ^ cls._xtime(a[2]) ^ cls._xtime(a[3]) ^ a[3]
                s[3][c] = cls._xtime(a[0]) ^ a[0] ^ a[1] ^ a[2] ^ cls._xtime(a[3])
            for c in range(4):
                for r in range(4):
                    s[r][c] ^= round_keys[rnd*4+c][r]
        s = [[cls._AES_SBOX[s[r][c]] for c in range(4)] for r in range(4)]
        for r in range(1, 4):
            s[r] = s[r][r:] + s[r][:r]
        for c in range(4):
            for r in range(4):
                s[r][c] ^= round_keys[Nr*4+c][r]
        out = []
        for i in range(16):
            out.append(s[i%4][i//4])
        return bytes(out)

    @classmethod
    def _aes_decrypt_block(cls, block: bytes, round_keys: list) -> bytes:
        Nr = len(round_keys) // 4 - 1
        s = [[0]*4 for _ in range(4)]
        for i in range(16):
            s[i%4][i//4] = block[i]
        for c in range(4):
            for r in range(4):
                s[r][c] ^= round_keys[Nr*4+c][r]
        for rnd in range(Nr-1, 0, -1):
            for r in range(1, 4):
                s[r] = s[r][-r:] + s[r][:-r]
            s = [[cls._AES_INV_SBOX[s[r][c]] for c in range(4)] for r in range(4)]
            for c in range(4):
                for r in range(4):
                    s[r][c] ^= round_keys[rnd*4+c][r]
            for c in range(4):
                a = [s[r][c] for r in range(4)]
                s[0][c] = cls._gmul(a[0],14) ^ cls._gmul(a[1],11) ^ cls._gmul(a[2],13) ^ cls._gmul(a[3],9)
                s[1][c] = cls._gmul(a[0],9) ^ cls._gmul(a[1],14) ^ cls._gmul(a[2],11) ^ cls._gmul(a[3],13)
                s[2][c] = cls._gmul(a[0],13) ^ cls._gmul(a[1],9) ^ cls._gmul(a[2],14) ^ cls._gmul(a[3],11)
                s[3][c] = cls._gmul(a[0],11) ^ cls._gmul(a[1],13) ^ cls._gmul(a[2],9) ^ cls._gmul(a[3],14)
        for r in range(1, 4):
            s[r] = s[r][-r:] + s[r][:-r]
        s = [[cls._AES_INV_SBOX[s[r][c]] for c in range(4)] for r in range(4)]
        for c in range(4):
            for r in range(4):
                s[r][c] ^= round_keys[c][r]
        out = []
        for i in range(16):
            out.append(s[i%4][i//4])
        return bytes(out)

    @staticmethod
    def _pkcs7_pad(data: bytes, block_size: int = 16) -> bytes:
        pad_len = block_size - (len(data) % block_size)
        return data + bytes([pad_len] * pad_len)

    @staticmethod
    def _pkcs7_unpad(data: bytes) -> bytes:
        if not data:
            raise ValueError("Empty data")
        pad_len = data[-1]
        if pad_len < 1 or pad_len > 16:
            raise ValueError(f"Invalid padding: {pad_len}")
        if data[-pad_len:] != bytes([pad_len] * pad_len):
            raise ValueError("Invalid PKCS7 padding")
        return data[:-pad_len]

    def _aes_ecb_encrypt(self, plain: bytes, key: bytes) -> bytes:
        if _HAS_PYCRYPTODOME:
            cipher = _CryptoAES.new(key, _CryptoAES.MODE_ECB)
            padded = self._pkcs7_pad(plain)
            return cipher.encrypt(padded)
        round_keys = self._aes_key_expansion(key)
        padded = self._pkcs7_pad(plain)
        out = bytearray()
        for i in range(0, len(padded), 16):
            out.extend(self._aes_encrypt_block(padded[i:i+16], round_keys))
        return bytes(out)

    def _aes_ecb_decrypt(self, encrypted: bytes, key: bytes) -> bytes:
        if _HAS_PYCRYPTODOME:
            cipher = _CryptoAES.new(key, _CryptoAES.MODE_ECB)
            decrypted = cipher.decrypt(encrypted)
            return self._pkcs7_unpad(decrypted)
        round_keys = self._aes_key_expansion(key)
        if len(encrypted) % 16 != 0:
            raise ValueError("Encrypted data length must be multiple of 16")
        out = bytearray()
        for i in range(0, len(encrypted), 16):
            out.extend(self._aes_decrypt_block(encrypted[i:i+16], round_keys))
        return self._pkcs7_unpad(bytes(out))

    def _upload_media(self, file_bytes: bytes, filename: str, media_type: int, to_user_id: str) -> Optional[dict]:
        try:
            print(f"[媒体上传] 正在上传 {filename}, 类型={media_type}, 大小={len(file_bytes)} bytes")

            use_token = self._get_token_for_user(to_user_id)

            aes_key_hex = self._random_hex(16)
            aes_key_bytes = bytes.fromhex(aes_key_hex)

            encrypted = self._aes_ecb_encrypt(file_bytes, aes_key_bytes)

            filekey = self._random_hex(16)
            raw_md5 = self._md5_hex(file_bytes)

            body = {
                "filekey": filekey,
                "media_type": media_type,
                "to_user_id": to_user_id,
                "rawsize": len(file_bytes),
                "rawfilemd5": raw_md5,
                "filesize": len(encrypted),
                "no_need_thumb": True,
                "aeskey": aes_key_hex
            }

            result = self._post("getuploadurl", body, token=use_token)

            ret = result.get("ret")
            errcode = result.get("errcode")

            if ret is not None and ret != 0:
                print(f"[媒体上传失败] getuploadurl 失败: ret={ret}, errcode={errcode}, errmsg={result.get('errmsg', '')}")
                return None
            if errcode is not None and errcode != 0:
                print(f"[媒体上传失败] getuploadurl 失败: ret={ret}, errcode={errcode}, errmsg={result.get('errmsg', '')}")
                return None

            upload_param = result.get("upload_param")
            if not upload_param:
                print(f"[媒体上传失败] 未获取到 upload_param, 返回数据: {json.dumps(result, ensure_ascii=False)[:300]}")
                return None

            cdn_url = self.CDN_BASE + "/upload?encrypted_query_param=" + urllib.parse.quote(upload_param, safe='') + "&filekey=" + urllib.parse.quote(filekey, safe='')

            print(f"[媒体上传] 获取到上传参数，正在上传到 CDN...")

            req = urllib.request.Request(
                cdn_url,
                data=encrypted,
                method='POST',
                headers={'Content-Type': 'application/octet-stream'}
            )

            with urllib.request.urlopen(req, timeout=120) as resp:
                encrypted_param = resp.headers.get('x-encrypted-param', '')
                if not encrypted_param:
                    resp_body = resp.read()
                    print(f"[媒体上传失败] CDN 响应缺少 x-encrypted-param 头, status={resp.status}, body={resp_body[:200]}")
                    return None

                aes_key_b64 = base64.b64encode(aes_key_hex.encode('utf-8')).decode('utf-8')

                cdn_media = {
                    "encrypt_query_param": encrypted_param,
                    "aes_key": aes_key_b64,
                    "encrypt_type": 1
                }

                uploaded = {
                    "filekey": filekey,
                    "media": cdn_media,
                    "aes_key_hex": aes_key_hex,
                    "raw_size": len(file_bytes),
                    "encrypted_size": len(encrypted),
                    "md5": raw_md5,
                    "filename": filename
                }

                print(f"[媒体上传成功] filekey={filekey}, enc_size={len(encrypted)}")
                return uploaded

        except Exception as e:
            print(f"[媒体上传异常] {e}")
            traceback.print_exc()
            return None

    def _send_media_message(self, to_user_id: str, media_item: dict,
                            description: str = "", media_data: str = "",
                            media_filename: str = "", media_duration: int = 0) -> bool:
        context_token = self._context_tokens.get(to_user_id)
        if not context_token:
            print(f"[发送失败] 没有 {to_user_id} 的会话，让对方先发一条消息")
            return False

        use_token = self._get_token_for_user(to_user_id)
        client_id = f"ilink-sdk:{int(time.time() * 1000)}-{uuid.uuid4().hex[:8]}"

        if description:
            text_item = {"type": 1, "text_item": {"text": description}}
            media_item_list = [media_item, text_item]
        else:
            media_item_list = [media_item]

        body = {
            "msg": {
                "from_user_id": "",
                "to_user_id": to_user_id,
                "client_id": client_id,
                "message_type": 2,
                "message_state": 2,
                "context_token": context_token,
                "item_list": media_item_list
            }
        }

        result = self._post("sendmessage", body, token=use_token)

        errcode = result.get("errcode")
        ret = result.get("ret")

        success = (ret is None or ret == 0) and (errcode is None or errcode == 0)

        if success:
            type_name = self.MEDIA_TYPE_NAMES.get(media_item.get("type", 0), "媒体")
            print(f"[发送成功] {type_name} 给 {to_user_id}")
            out_msg = {
                'from': 'me',
                'to': to_user_id,
                'text': f"[{type_name}]" + (f" {description}" if description else ""),
                'time': datetime.now().strftime('%H:%M:%S'),
                'type': 'out',
                'media_type': media_item.get("type"),
                'media_data': media_data,
                'media_filename': media_filename or description,
                'media_duration': media_duration
            }
            cdn_media = self._extract_cdn_media(media_item)
            if cdn_media:
                out_msg['media_cdn'] = json.dumps(cdn_media)
            self._add_message_to_history(out_msg)
            return True

        print(f"[发送失败] ret={ret}, errcode={errcode}, errmsg={result.get('errmsg', '')}")
        return False

    def send_image(self, to_user_id: str, image_bytes: bytes,
                   filename: str = "image.jpg", description: str = "",
                   media_data: str = "") -> bool:
        print(f"[发送图片] 准备发送图片给 {to_user_id}: {filename} ({len(image_bytes)} bytes)")

        uploaded = self._upload_media(image_bytes, filename, media_type=1, to_user_id=to_user_id)
        if not uploaded:
            print("[发送图片失败] 上传失败")
            return False

        image_item = {
            "media": uploaded["media"],
            "aeskey": uploaded["aes_key_hex"],
            "mid_size": uploaded["encrypted_size"]
        }

        media_item = {
            "type": 2,
            "image_item": image_item
        }

        return self._send_media_message(to_user_id, media_item, description,
                                        media_data=media_data, media_filename=filename)

    def send_file(self, to_user_id: str, file_bytes: bytes,
                  filename: str = "file.bin", description: str = "",
                  media_data: str = "") -> bool:
        print(f"[发送文件] 准备发送文件给 {to_user_id}: {filename} ({len(file_bytes)} bytes)")

        uploaded = self._upload_media(file_bytes, filename, media_type=3, to_user_id=to_user_id)
        if not uploaded:
            print("[发送文件失败] 上传失败")
            return False

        file_item = {
            "media": uploaded["media"],
            "file_name": filename,
            "md5": uploaded["md5"],
            "len": str(uploaded["raw_size"])
        }

        media_item = {
            "type": 4,
            "file_item": file_item
        }

        return self._send_media_message(to_user_id, media_item, description,
                                        media_filename=filename)

    def send_voice(self, to_user_id: str, voice_bytes: bytes,
                   filename: str = "voice.silk", duration_ms: int = 1000,
                   sample_rate: int = 16000) -> bool:
        print(f"[发送语音] 准备发送语音给 {to_user_id}: {filename} ({len(voice_bytes)} bytes, {duration_ms}ms)")

        uploaded = self._upload_media(voice_bytes, filename, media_type=4, to_user_id=to_user_id)
        if not uploaded:
            print("[发送语音失败] 上传失败")
            return False

        voice_item = {
            "media": uploaded["media"],
            "encode_type": 6,
            "bits_per_sample": 16,
            "playtime": duration_ms,
            "sample_rate": sample_rate
        }

        media_item = {
            "type": 3,
            "voice_item": voice_item
        }

        return self._send_media_message(to_user_id, media_item,
                                        media_filename=filename, media_duration=duration_ms)

    def send_video(self, to_user_id: str, video_bytes: bytes,
                   filename: str = "video.mp4", duration_ms: int = 5000,
                   description: str = "", media_data: str = "") -> bool:
        print(f"[发送视频] 准备发送视频给 {to_user_id}: {filename} ({len(video_bytes)} bytes, {duration_ms}ms)")

        uploaded = self._upload_media(video_bytes, filename, media_type=2, to_user_id=to_user_id)
        if not uploaded:
            print("[发送视频失败] 上传失败")
            return False

        video_item = {
            "media": uploaded["media"],
            "video_size": uploaded["encrypted_size"],
            "play_length": duration_ms,
            "video_md5": uploaded["md5"]
        }

        media_item = {
            "type": 5,
            "video_item": video_item
        }

        return self._send_media_message(to_user_id, media_item, description,
                                        media_data=media_data, media_filename=filename,
                                        media_duration=duration_ms)

    def _media_cache_key(self, cdn_media_info: dict) -> str:
        eqp = cdn_media_info.get("encrypt_query_param") or cdn_media_info.get("encrypted_query_param") or ""
        return hashlib.md5(eqp.encode('utf-8')).hexdigest()

    def send_video_for_account(self, account, to_user_id: str, video_bytes: bytes,
                   filename: str = "video.mp4", duration_ms: int = 5000,
                   description: str = "", media_data: str = "") -> bool:
        uploaded = self._upload_media_for_account(account, video_bytes, filename, media_type=2, to_user_id=to_user_id)
        if not uploaded:
            return False
        video_item = {
            "media": uploaded["media"],
            "video_size": uploaded["encrypted_size"],
            "play_length": duration_ms,
            "video_md5": uploaded["md5"]
        }
        media_item = {"type": 5, "video_item": video_item}
        return self._send_media_message_for_account(account, to_user_id, media_item, description,
                                        media_data=media_data, media_filename=filename,
                                        media_duration=duration_ms)

    def send_file_for_account(self, account, to_user_id: str, file_bytes: bytes,
                  filename: str = "file.bin", description: str = "",
                  media_data: str = "") -> bool:
        uploaded = self._upload_media_for_account(account, file_bytes, filename, media_type=3, to_user_id=to_user_id)
        if not uploaded:
            return False
        file_item = {
            "media": uploaded["media"],
            "file_name": filename,
            "md5": uploaded["md5"],
            "len": str(uploaded["raw_size"])
        }
        media_item = {"type": 4, "file_item": file_item}
        return self._send_media_message_for_account(account, to_user_id, media_item, description,
                                        media_filename=filename)

    def _enrich_msg_with_cache_id(self, msg: dict) -> dict:
        if msg.get('media_cdn') and msg.get('media_type'):
            try:
                cdn_info = json.loads(msg['media_cdn']) if isinstance(msg['media_cdn'], str) else msg['media_cdn']
                cache_key = self._media_cache_key(cdn_info)
                
                user_id = msg.get('from') if msg.get('type') == 'in' else msg.get('to')
                
                cached = None
                if user_id:
                    cached = self._get_user_cached_media(user_id, cache_key)
                if not cached:
                    cached = self._get_cached_media(cache_key)
                if cached:
                    msg['media_cache_id'] = cache_key
                    msg['media_cache_user'] = user_id
            except Exception:
                pass
        return msg

    def _media_cache_path(self, cache_key: str) -> Path:
        return self._media_cache_dir / cache_key

    def _media_meta_path(self, cache_key: str) -> Path:
        return self._media_cache_dir / (cache_key + ".meta")

    def _get_cached_media(self, cache_key: str) -> Optional[tuple]:
        data_path = self._media_cache_path(cache_key)
        meta_path = self._media_meta_path(cache_key)
        if data_path.exists() and meta_path.exists():
            try:
                media_data = data_path.read_bytes()
                meta = json.loads(meta_path.read_text('utf-8'))
                return (media_data, meta.get('mime', 'application/octet-stream'), meta.get('filename', ''))
            except Exception:
                return None
        return None

    def _save_media_cache(self, cache_key: str, media_data: bytes, mime: str, filename: str = ""):
        try:
            self._media_cache_path(cache_key).write_bytes(media_data)
            meta = {'mime': mime, 'filename': filename, 'size': len(media_data)}
            self._media_meta_path(cache_key).write_text(json.dumps(meta, ensure_ascii=False), 'utf-8')
        except Exception as e:
            print(f"[媒体缓存] 保存失败: {e}")

    def _prefetch_media(self, cdn_media_info: dict, filename: str = "", user_id: str = ""):
        try:
            cache_key = self._media_cache_key(cdn_media_info)
            
            if user_id and self._get_user_cached_media(user_id, cache_key):
                return
            if self._get_cached_media(cache_key):
                return
            
            print(f"[媒体预取] 开始下载: {cache_key[:12]}...")
            result = self.download_media(cdn_media_info, filename=filename, user_id=user_id)
            if result:
                print(f"[媒体预取] 完成: {cache_key[:12]}..., {len(result)} bytes")
            else:
                print(f"[媒体预取] 失败: {cache_key[:12]}...")
        except Exception as e:
            print(f"[媒体预取] 异常: {e}")

    def _detect_mime(self, data: bytes) -> str:
        if data[:8] == b'\x89PNG\r\n\x1a\n':
            return 'image/png'
        if data[:4] == b'GIF8':
            return 'image/gif'
        if data[:4] == b'RIFF' and len(data) > 12 and data[8:12] == b'WEBP':
            return 'image/webp'
        if data[:2] == b'\xff\xd8':
            return 'image/jpeg'
        if data[:4] == b'RIFF' and len(data) > 12 and data[8:12] == b'WAVE':
            return 'audio/wav'
        if len(data) > 3 and (data[:3] == b'ID3' or data[:2] in (b'\xff\xfb', b'\xff\xf3', b'\xff\xf2')):
            return 'audio/mpeg'
        if data[:4] == b'fLaC':
            return 'audio/flac'
        if data[:4] == b'OggS':
            return 'audio/ogg'
        if len(data) > 9 and data[:9] == b'#!SILK_V3':
            return 'audio/silk'
        if len(data) > 10 and data[:1] == b'\x02' and data[1:10] == b'#!SILK_V3':
            return 'audio/silk'
        if len(data) > 5 and data[:5] == b'#!AMR':
            return 'audio/amr'
        if data[:4] == b'MThd':
            return 'audio/midi'
        if len(data) > 8 and data[:4] == b'\x00\x00\x00':
            box_type = data[4:8]
            if box_type == b'ftyp':
                return 'video/mp4'
        if data[:4] == b'\x1a\x45\xdf\xa3':
            return 'video/webm'
        return 'application/octet-stream'

    def _silk_to_wav(self, silk_data: bytes) -> Optional[bytes]:
        try:
            import pilk
        except ImportError:
            print("[SILK转WAV] pilk 未安装，尝试 ffmpeg")
            return self._ffmpeg_to_wav(silk_data)
        if silk_data[:1] == b'\x02' and len(silk_data) > 10 and silk_data[1:10] == b'#!SILK_V3':
            silk_data = silk_data[1:]
        if silk_data[:9] != b'#!SILK_V3':
            print("[SILK转WAV] 非 SILK V3 格式")
            return self._ffmpeg_to_wav(silk_data)
        try:
            tmp_in = self._media_cache_dir / ('_silk_tmp_in_' + uuid.uuid4().hex[:12] + '.silk')
            tmp_out = self._media_cache_dir / ('_silk_tmp_out_' + uuid.uuid4().hex[:12] + '.pcm')
            tmp_in.write_bytes(silk_data)
            pilk.decode(str(tmp_in), str(tmp_out), pcm_rate=24000)
            if not tmp_out.exists() or tmp_out.stat().st_size == 0:
                print("[SILK转WAV] pilk 解码无输出，尝试 ffmpeg")
                return self._ffmpeg_to_wav(silk_data)
            pcm_data = tmp_out.read_bytes()
            sample_rate = 24000
            num_channels = 1
            bits_per_sample = 16
            byte_rate = sample_rate * num_channels * bits_per_sample // 8
            block_align = num_channels * bits_per_sample // 8
            data_size = len(pcm_data)
            wav_buf = io.BytesIO()
            wav_buf.write(b'RIFF')
            wav_buf.write(struct.pack('<I', 36 + data_size))
            wav_buf.write(b'WAVE')
            wav_buf.write(b'fmt ')
            wav_buf.write(struct.pack('<I', 16))
            wav_buf.write(struct.pack('<H', 1))
            wav_buf.write(struct.pack('<H', num_channels))
            wav_buf.write(struct.pack('<I', sample_rate))
            wav_buf.write(struct.pack('<I', byte_rate))
            wav_buf.write(struct.pack('<H', block_align))
            wav_buf.write(struct.pack('<H', bits_per_sample))
            wav_buf.write(b'data')
            wav_buf.write(struct.pack('<I', data_size))
            wav_buf.write(pcm_data)
            print(f"[SILK转WAV] pilk 转换成功: {len(silk_data)} bytes SILK -> {wav_buf.tell()} bytes WAV")
            return wav_buf.getvalue()
        except Exception as e:
            print(f"[SILK转WAV] pilk 转换失败: {e}，尝试 ffmpeg")
            return self._ffmpeg_to_wav(silk_data)
        finally:
            try:
                tmp_in.unlink(missing_ok=True)
            except Exception:
                pass
            try:
                tmp_out.unlink(missing_ok=True)
            except Exception:
                pass

    def _ffmpeg_to_wav(self, audio_data: bytes) -> Optional[bytes]:
        tmp_in = None
        tmp_out = None
        try:
            tmp_in = self._media_cache_dir / ('_ffmpeg_tmp_in_' + uuid.uuid4().hex[:12])
            tmp_out = self._media_cache_dir / ('_ffmpeg_tmp_out_' + uuid.uuid4().hex[:12] + '.wav')
            tmp_in.write_bytes(audio_data)
            result = subprocess.run(
                ['ffmpeg', '-y', '-i', str(tmp_in), '-f', 'wav', '-ar', '24000', '-ac', '1', str(tmp_out)],
                capture_output=True, timeout=30
            )
            if tmp_out.exists() and tmp_out.stat().st_size > 44:
                wav_data = tmp_out.read_bytes()
                print(f"[ffmpeg转WAV] 转换成功: {len(audio_data)} bytes -> {len(wav_data)} bytes")
                return wav_data
            print(f"[ffmpeg转WAV] 转换失败: {result.stderr.decode('utf-8', errors='replace')[:200]}")
            return None
        except Exception as e:
            print(f"[ffmpeg转WAV] 异常: {e}")
            return None
        finally:
            for tmp in (tmp_in, tmp_out):
                if tmp:
                    try:
                        if tmp.exists(): tmp.unlink()
                    except Exception:
                        pass

    def download_media(self, cdn_media_info: dict, filename: str = "", user_id: str = "") -> Optional[bytes]:
        cache_key = self._media_cache_key(cdn_media_info)

        if user_id:
            cached = self._get_user_cached_media(user_id, cache_key)
            if cached:
                return cached[0]
        
        cached = self._get_cached_media(cache_key)
        if cached:
            return cached[0]

        with self._media_download_lock:
            if cache_key in self._media_downloading:
                wait_event = self._media_downloading[cache_key]
            else:
                wait_event = None

        if wait_event:
            wait_event.wait(timeout=60)
            if user_id:
                cached = self._get_user_cached_media(user_id, cache_key)
                if cached:
                    return cached[0]
            cached = self._get_cached_media(cache_key)
            if cached:
                return cached[0]
            return None

        event = threading.Event()
        with self._media_download_lock:
            self._media_downloading[cache_key] = event

        try:
            encrypt_query_param = cdn_media_info.get("encrypt_query_param")
            aes_key_b64 = cdn_media_info.get("aes_key")
            
            if not encrypt_query_param:
                encrypt_query_param = cdn_media_info.get("encrypted_query_param")
            if not encrypt_query_param:
                return None
            
            if not aes_key_b64:
                aes_key_hex = cdn_media_info.get("aeskey") or cdn_media_info.get("aes_key_hex")
                if aes_key_hex:
                    aes_key_b64 = base64.b64encode(aes_key_hex.encode('utf-8')).decode('utf-8')
            
            if not aes_key_b64:
                return None

            download_url = self.CDN_BASE + "/download?encrypted_query_param=" + urllib.parse.quote(encrypt_query_param, safe='')

            print(f"[媒体下载] 正在从 CDN 下载...")
            req = urllib.request.Request(download_url)

            with urllib.request.urlopen(req, timeout=60) as resp:
                data = resp.read()

                decoded_key = base64.b64decode(aes_key_b64)
                if len(decoded_key) == 16:
                    aes_key_bytes = decoded_key
                else:
                    aes_key_hex = decoded_key.decode('utf-8')
                    aes_key_bytes = bytes.fromhex(aes_key_hex)

                decrypted = self._aes_ecb_decrypt(data, aes_key_bytes)
                print(f"[媒体下载成功] 解密后大小: {len(decrypted)} bytes")

                mime = self._detect_mime(decrypted)
                if mime == 'audio/silk':
                    wav_data = self._silk_to_wav(decrypted)
                    if wav_data:
                        decrypted = wav_data
                        mime = 'audio/wav'
                        filename = filename.replace('.silk', '.wav') if filename else 'voice.wav'
                elif mime == 'audio/amr':
                    wav_data = self._ffmpeg_to_wav(decrypted)
                    if wav_data:
                        decrypted = wav_data
                        mime = 'audio/wav'
                        filename = filename.replace('.amr', '.wav') if filename else 'voice.wav'

                self._save_media_cache(cache_key, decrypted, mime, filename)
                
                if user_id:
                    self._save_user_media_cache(user_id, cache_key, decrypted, mime, filename)

                return decrypted

        except Exception as e:
            print(f"[媒体下载异常] {e}")
            return None
        finally:
            with self._media_download_lock:
                self._media_downloading.pop(cache_key, None)
            event.set()

    def download_media_from_message_item(self, message_item: dict) -> Optional[bytes]:
        cdn_media_info = self._extract_cdn_media(message_item)

        if cdn_media_info and cdn_media_info.get("encrypt_query_param"):
            return self.download_media(cdn_media_info)

        print("[下载失败] 消息项中未找到有效的媒体信息")
        return None
    
    def list_users(self) -> list:
        return list(self._context_tokens.keys())
    
    def get_current_user(self):
        return self._current_user
    
    def set_current_user(self, user_id: str):
        if user_id in self._context_tokens:
            self._current_user = user_id
            self._save_config()
            print(f"已切换到: {user_id}")
    
    def remove_user(self, user_id: str):
        if not user_id or user_id not in self._context_tokens:
            return False
        
        self._context_tokens.pop(user_id, None)
        
        bot_token = self._user_token_map.pop(user_id, None)
        
        if bot_token and bot_token in self._bot_accounts:
            self._bot_accounts[bot_token].get("context_tokens", {}).pop(user_id, None)
        
        if user_id in self._active_timers:
            timer = self._active_timers.pop(user_id)
            if timer:
                timer.cancel()
        
        self._messages = [m for m in self._messages 
                         if m.get('from') != user_id and m.get('to') != user_id]
        
        try:
            user_dir = self._get_user_dir_path(user_id)
            if user_dir.exists():
                shutil.rmtree(str(user_dir))
        except Exception as e:
            print(f"[USER] 删除用户目录失败 ({user_id}): {e}")
        
        if self._current_user == user_id:
            remaining = list(self._context_tokens.keys())
            self._current_user = remaining[0] if remaining else None
        
        self._save_config()
        print(f"[USER] 已删除用户: {user_id}")
        return True
    
    def stop(self):
        self._running = False
        for timer in self._active_timers.values():
            timer.cancel()
        self._active_timers.clear()
        if self._http_server:
            self._http_server.shutdown()

def main():
    print("╔" + "═" * 58 + "╗")
    print("║" + " " * 12 + "Zynsync iLink ChatBox" + " " * 18 + "║")
    print("║" + " " * 18 + "-无忧传递-" + " " * 22 + "║")
    print("║" + " " * 20 + "v3.1.9" + " " * 26 + "║")
    print("╚" + "═" * 58 + "╝")
    print('注意，一定要在当前控制台为一个用户设置为管理员账号（命令 3 ），否则管理员面板功能无法使用!-管理员账号打开网页后即可在设置里看到管理员面板选项这一项!')
    
    if is_termux():
        print("\n[TERMUX] 运行环境: Android/Termux")
        print("[TERMUX] 网络模式: 可能需要代理或 VPN 访问微信服务器")
        print("[TERMUX] 提示: 如果网络不稳定，程序会自动重试")
        print()
    
    bot = WeChatiLinkBot()
    
    for username, account in bot._accounts.items():
        if account.token:
            print(f"[多账户] 启动账户 {username} 的轮询...")
            bot._start_account_polling(account)
    
    bot.start_web_interface()
    
    if is_termux():
        print(f"\n[Zyn] 📱 网页地址: http://127.0.0.1:{bot._web_port}")
        print("[TERMUX] 💡 使用方法:")
        print("   方法1: 在手机浏览器访问上述地址")
        print(f"   方法2: 在电脑浏览器访问 http://<你的IP>:{bot._web_port}")
        print("   (需确保手机和电脑在同一网络)")
        print("[TERMUX] 输入 /web 可再次显示地址\n")
    
    if bot.load_config():
        print("[zyn]已获取到连接缓存")
    else:
        print(f"[zyn]首次运行，请在网页端扫码连接: http://127.0.0.1:{bot._web_port}")
        print("[zyn]点击网页上的加号按钮扫码即可，控制台可正常输入指令")
    
    bot.start_polling()
    
    print("[zyn]后台监听已启动，等待消息...")
    print(f"[zyn]网页地址: http://127.0.0.1:{bot._web_port}")
    
    users = bot.list_users()
    if users:
        print(f"\n已保存 {len(users)} 个会话")
        for uid in users:
            marker = "[zyn]" if uid == bot.get_current_user() else "   "
            print(f"{marker}{uid}")
    else:
        print("\n暂未有任何会话")
        print("[zyn]对方扫完二维码后必须先发送一条消息才能建立联系!")
    
    print(f"\n[多账户] 已加载 {len(bot._accounts)} 个账户")
    for username, account in bot._accounts.items():
        user_count = len(account._context_tokens)
        status = "已连接" if account.token else "未连接"
        print(f"  - {username}: {status}, {user_count} 个会话")
    
    print("\n" + "┌" + "─" * 58 + "┐")
    print("│ 1.查看所有用户" + " " * 30 + "│")
    print("│ 2.注销选定用户账号" + " " * 24 + "│")
    print("│ 3.为选定用户设置管理员" + " " * 20 + "│")
    print("│ /quit 退出" + " " * 36 + "│")
    print("└" + "─" * 58 + "┘" + "\n")
    
    try:
        while True:
            user_input = input("命令:").strip()
            if not user_input:
                continue
            if user_input == "/quit":
                break
            elif user_input == "1":
                accounts = bot._accounts
                if accounts:
                    print("[zyn]用户列表:")
                    for i, (uname, acc) in enumerate(accounts.items(), 1):
                        admin_mark = " [管理员]" if acc.is_admin else ""
                        status = "已连接" if acc.token else "未连接"
                        print(f"  {i}. {uname} ({status}){admin_mark}")
                else:
                    print("[zyn]暂无用户")
                continue
            elif user_input == "2":
                accounts = bot._accounts
                if not accounts:
                    print("[zyn]暂无用户")
                    continue
                print("[zyn]选择要注销的用户:")
                for i, (uname, acc) in enumerate(accounts.items(), 1):
                    admin_mark = " [管理员]" if acc.is_admin else ""
                    print(f"  {i}. {uname}{admin_mark}")
                try:
                    choice = input("[zyn]请输入序号: ").strip()
                    idx = int(choice) - 1
                    uname_list = list(accounts.keys())
                    if 0 <= idx < len(uname_list):
                        target = uname_list[idx]
                        confirm = input(f"[zyn]确认注销用户 {target} 并删除其目录? (y/n): ").strip().lower()
                        if confirm == 'y':
                            acc = accounts[target]
                            old_sessions = [tok for tok, un in bot._account_sessions.items() if un == target]
                            for tok in old_sessions:
                                bot._account_sessions.pop(tok, None)
                                bot._verified_sessions.pop(tok, None)
                            for fp, un in list(bot._fingerprint_sessions.items()):
                                if un == target:
                                    del bot._fingerprint_sessions[fp]
                            del bot._accounts[target]
                            bot._save_accounts()
                            try:
                                acc_dir = Path(f"accounts/{target}")
                                if acc_dir.exists():
                                    shutil.rmtree(str(acc_dir))
                            except Exception as e:
                                print(f"[zyn]删除用户目录失败: {e}")
                            print(f"[zyn]已注销用户: {target}")
                        else:
                            print("[zyn]已取消")
                    else:
                        print("[zyn]无效序号")
                except ValueError:
                    print("[zyn]请输入数字")
                continue
            elif user_input == "3":
                accounts = bot._accounts
                if not accounts:
                    print("[zyn]暂无用户")
                    continue
                print("[zyn]选择要设为管理员的用户:")
                for i, (uname, acc) in enumerate(accounts.items(), 1):
                    admin_mark = " [当前管理员]" if acc.is_admin else ""
                    print(f"  {i}. {uname}{admin_mark}")
                try:
                    choice = input("[zyn]请输入序号 (输入0取消所有管理员): ").strip()
                    idx = int(choice)
                    uname_list = list(accounts.keys())
                    if idx == 0:
                        for acc in accounts.values():
                            acc.is_admin = False
                        bot._save_accounts()
                        print("[zyn]已取消所有管理员设置")
                    elif 1 <= idx <= len(uname_list):
                        target = uname_list[idx - 1]
                        for acc in accounts.values():
                            acc.is_admin = False
                        accounts[target].is_admin = True
                        bot._save_accounts()
                        print(f"[zyn]已将 {target} 设为管理员")
                    else:
                        print("[zyn]无效序号")
                except ValueError:
                    print("[zyn]请输入数字")
                continue
    except KeyboardInterrupt:
        print()
    finally:
        bot.stop()

if __name__ == "__main__":
    main()
#我是屎山代码的底座